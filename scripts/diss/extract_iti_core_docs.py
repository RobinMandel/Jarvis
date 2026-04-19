from pathlib import Path
from docx import Document
from pptx import Presentation
import json

base = Path(r'E:\OneDrive\Dokumente\Studium\!!Medizin\03 Doktorarbeit\ITI')
files = [
    base / 'Dissertation' / 'Dissertation Robin Mandel neu.docx',
    base / 'Dissertation' / 'Dissertation Robin Mandel.docx',
    base / 'Dissertation' / 'Abbildungen.docx',
    base / 'Präsis' / '!Text Abschlussbericht.docx',
    base / 'Präsis' / '!Abschlussbericht Robin Mandel.pptx',
]
outdir = base / 'Jarvis Workspace' / '06_Data-Extraction'
outdir.mkdir(parents=True, exist_ok=True)
manifest = []
for f in files:
    if not f.exists():
        manifest.append({'file': str(f), 'status': 'missing'})
        continue
    text_parts = []
    status = 'ok'
    try:
        if f.suffix.lower() == '.docx':
            doc = Document(str(f))
            for p in doc.paragraphs:
                t = (p.text or '').strip()
                if t:
                    text_parts.append(t)
            for ti, table in enumerate(doc.tables, 1):
                text_parts.append(f'\n[TABLE {ti}]')
                for row in table.rows:
                    vals = [c.text.strip().replace('\n',' ') for c in row.cells]
                    if any(vals):
                        text_parts.append(' | '.join(vals))
        elif f.suffix.lower() == '.pptx':
            prs = Presentation(str(f))
            for si, slide in enumerate(prs.slides, 1):
                text_parts.append(f'\n[SLIDE {si}]')
                for shape in slide.shapes:
                    txt = ''
                    if hasattr(shape, 'text'):
                        txt = (shape.text or '').strip()
                    elif getattr(shape, 'has_text_frame', False):
                        txt = shape.text_frame.text.strip()
                    if txt:
                        text_parts.append(txt)
    except Exception as e:
        status = f'error: {e}'
    out = outdir / (f.name + '.txt')
    out.write_text('\n'.join(text_parts), encoding='utf-8')
    manifest.append({'file': str(f), 'status': status, 'out': str(out), 'chars': sum(len(x) for x in text_parts)})

(outdir / 'text-extraction-manifest.json').write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding='utf-8')
print(json.dumps(manifest, ensure_ascii=False, indent=2))