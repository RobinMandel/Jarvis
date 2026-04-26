"""Mission Control V3  -  Backend Server
aiohttp-based server bridging Claude Code CLI to a web dashboard.
"""

import asyncio
import json
import time
import uuid
import os
import sys
import faulthandler
from pathlib import Path
from datetime import datetime

# faulthandler ASAP so native crashes (Access Violation, segfault, abort)
# from C-extensions like faiss / whisper / torch leave a Python traceback
# in stderr instead of just rc=4294967295. This was the missing piece for
# the 2026-04-25 hard-crash loop where mc_stderr.log had nothing useful.
# When launched by watchdog.py, sys.stderr is the per-run err logfile,
# so the dump lands next to the crash automatically.
try:
    faulthandler.enable(all_threads=True)
    _boot_marker = (
        f"[Boot] faulthandler enabled, pid={os.getpid()}, "
        f"ts={datetime.now().isoformat(timespec='seconds')}\n"
    )
    sys.stderr.write(_boot_marker)
    sys.stderr.flush()
except Exception as _e:
    # Non-fatal — server still starts, we just lose native-crash diagnostics
    sys.stderr.write(f"[Boot] faulthandler.enable failed: {_e}\n")

import aiohttp
from aiohttp import web

from notify import send_push

BASE_DIR = Path(__file__).parent
_JARVIS_ROOT = BASE_DIR.parent
# Mission-Control/ must come before Jarvis root so that Mission-Control/rag/
# takes priority over any rag/ package that may exist in the Jarvis root.
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))
if str(_JARVIS_ROOT) not in sys.path:
    sys.path.append(str(_JARVIS_ROOT))

# Suppress console windows when spawning subprocesses on Windows
_SUBPROCESS_FLAGS = {}
if sys.platform == 'win32':
    import subprocess as _sp
    _si = _sp.STARTUPINFO()
    _si.dwFlags |= _sp.STARTF_USESHOWWINDOW
    _si.wShowWindow = 0  # SW_HIDE
    _SUBPROCESS_FLAGS = {
        'creationflags': 0x08000000,  # CREATE_NO_WINDOW
        'startupinfo': _si,
    }
    # Force UTF-8 stdout/stderr so Unicode chars like → don't crash on cp1252
    try:
        sys.stdout.reconfigure(encoding='utf-8', errors='replace')
        sys.stderr.reconfigure(encoding='utf-8', errors='replace')
    except AttributeError:
        pass
DATA_DIR = BASE_DIR / "data"
UPLOADS_DIR = DATA_DIR / "uploads"
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
IMAGE_JOBS_DIR = DATA_DIR / "image_jobs"
IMAGE_JOBS_DIR.mkdir(parents=True, exist_ok=True)
IMAGES_LIBRARY_DIR = DATA_DIR / "images-library"
IMAGES_LIBRARY_DIR.mkdir(parents=True, exist_ok=True)
(IMAGES_LIBRARY_DIR / "_inbox").mkdir(parents=True, exist_ok=True)
TRANSCRIPTS_DIR = DATA_DIR / "transcripts"
TRANSCRIPTS_DIR.mkdir(parents=True, exist_ok=True)


def atomic_write_text(path, content, *, encoding="utf-8", keep_backup=True):
    """Crash-safe write: tmp-file + rename, keeps a .bak of the previous good copy.

    Prevents half-written sessions.json / tasks.json when the server is killed
    mid-write (Ctrl+C, crash, power loss). If the write succeeds we rotate the
    previous file to .bak so the on-load fallback has something to recover from.
    """
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(path.suffix + ".tmp")
    tmp.write_text(content, encoding=encoding)
    if keep_backup and path.exists():
        try:
            bak = path.with_suffix(path.suffix + ".bak")
            if bak.exists():
                bak.unlink()
            path.replace(bak)
        except Exception:
            pass  # Backup is best-effort; don't block the actual write
    tmp.replace(path)


def safe_read_json(path, default=None):
    """Read JSON with automatic fallback to .bak if the main file is corrupt."""
    path = Path(path)
    for candidate in (path, path.with_suffix(path.suffix + ".bak")):
        if not candidate.exists():
            continue
        try:
            return json.loads(candidate.read_text(encoding="utf-8"))
        except Exception as e:
            print(f"[Persist] {candidate.name} corrupt ({e}); trying fallback", flush=True)
    return default

# Redirect stdout/stderr to log file so all print() calls are captured
import logging as _logging
_log_path = DATA_DIR / "server.log"
_log_fh = open(_log_path, "a", encoding="utf-8", buffering=1)

class _Tee:
    """Write to both original stream and log file."""
    def __init__(self, orig, fh):
        self._orig = orig
        self._fh = fh
    def write(self, s):
        try:
            self._orig.write(s)
        except UnicodeEncodeError:
            self._orig.write(s.encode('utf-8', errors='replace').decode(
                getattr(self._orig, 'encoding', 'utf-8') or 'utf-8', errors='replace'))
        self._fh.write(s)
    def flush(self):
        self._orig.flush()
        self._fh.flush()
    def __getattr__(self, name):
        return getattr(self._orig, name)

sys.stdout = _Tee(sys.stdout, _log_fh)
sys.stderr = _Tee(sys.stderr, _log_fh)
print(f"[Boot] Mission Control started at {datetime.now().isoformat()}", flush=True)
STATIC_DIR = BASE_DIR / "static"
with open(Path(__file__).parent / "config.json") as _cf:
    _config_memory = json.load(_cf).get("memory_dir", "data/memory")
MEMORY_DIR = Path(_config_memory) if Path(_config_memory).is_absolute() else BASE_DIR / _config_memory
AVATARS_DIR = BASE_DIR / "avatars"

with open(BASE_DIR / "config.json") as f:
    CONFIG = json.load(f)

# Config-Keys in os.environ propagieren damit Provider-Clients (image_client.py,
# Anthropic-SDK, OpenAI-SDK) sie ohne Extra-Wiring finden. Wir überschreiben
# NICHT wenn die Env-Var schon gesetzt ist (Env gewinnt vor Config).
for _cfg_key, _env_var in [
    ("gemini_api_key", "GEMINI_API_KEY"),
    ("openai_api_key", "OPENAI_API_KEY"),
    ("anthropic_api_key", "ANTHROPIC_API_KEY"),
    ("openrouter_api_key", "OPENROUTER_API_KEY"),
]:
    _v = (CONFIG.get(_cfg_key) or "").strip()
    if _v and not os.environ.get(_env_var):
        os.environ[_env_var] = _v

import re


# ---------------------------------------------------------------------------
# Session Management
# ---------------------------------------------------------------------------

class _SafeLock(asyncio.Lock):
    """asyncio.Lock whose context-manager exit tolerates a prior force_unlock().

    Without this, `force_unlock()` releasing the lock mid-run causes the `async
    with`'s __aexit__ to raise RuntimeError('Lock is not acquired.'), which
    kills the `_run_claude` Task silently — the frontend then sees no chat.done
    event and hangs forever in "crunching...".
    """

    async def __aexit__(self, exc_type, exc, tb):
        try:
            self.release()
        except RuntimeError:
            pass


class ClaudeSession:
    """Represents a single chat session. Despite the legacy name, this is now
    the generic provider-agnostic session — see `provider`/`model`/`provider_thread_id`.

    Provider-Switching (hart verworfenes Resume bei Wechsel, siehe Codex-Review 2026-04-23):
      - provider und model sind getrennt persistiert, kombiniert nur in UI
      - provider_thread_id ist provider-spezifisch (claude-cli: CLI-Session-UUID,
        codex-cli: Codex thread_id). Bei Provider-Wechsel wird es hart auf None
        gesetzt; eine System-Message in history markiert den Kontextbruch.
    """

    def __init__(self, session_id=None):
        self.session_id = session_id or str(uuid.uuid4())
        # Provider-Konfiguration (neu, 2026-04-23)
        # Default: claude-cli — schnell, zuverlässig, kein API-Key-Management.
        # OpenClaw bleibt als EXPLIZITE Picker-Wahl für agentische Tasks (sonst
        # triggert jeder Quick-Chat unnötig einen vollen Agent-Turn mit Skills+Memory).
        self.provider = "claude-cli"
        self.model = ""               # leer = Provider nimmt Account-Default
        self.provider_thread_id = None  # generisches Resume-Token (ersetzt cli_session_id semantisch)
        # Global-Context-Bundles für diese Session (Geminis Vorschlag #2, 2026-04-23).
        # Leere Liste = nur RAG-Semantic-Hits. Liste von Bundle-Namen aus
        # VAULT_CONTEXT_BUNDLES = alle .md dieser Bereiche werden als Volltext
        # an den System-Prompt gehängt.
        self.context_bundles: list[str] = []
        # Shared Scratchpad (Stufe 3, 2026-04-23) — key→value dict, beide
        # Provider können darüber Zwischen-Ergebnisse teilen.
        self.shared_notes: dict[str, str] = {}
        # Legacy-Feld: bleibt erhalten für Rückwärts-Kompatibilität mit dem
        # Persistenz-Format. Wert ist synchronisiert mit provider_thread_id
        # WENN provider == "claude-cli", sonst None.
        self.cli_session_id = None
        self.process = None
        self.locked = False
        self.lock = _SafeLock()
        self.created_at = datetime.now().isoformat()
        self.message_count = 0
        self.history = []
        self.custom_title = None
        self.auto_title = "Neue Session"
        self.topic = None
        self.topic_color = None
        self.last_message = None
        self.last_heartbeat = None  # updated on stream-start and message-done; shows real liveness
        self._streaming_partial = None  # partial response text saved during streaming
        self._lock_acquired_at = None  # timestamp when lock was acquired
        self.needs_resume = False  # set on load if stream was interrupted by server restart
        self._stream_active = False  # True while Claude subprocess is running; persisted

    async def cancel(self):
        """Kill the running Claude subprocess if any."""
        if self.process and self.process.returncode is None:
            self.process.kill()
            await self.process.wait()
            self.process = None
        self._stream_active = False
        self._streaming_partial = None

    def force_unlock(self):
        """Force-release lock if stuck (e.g. after WS disconnect or crash)."""
        if self.lock.locked():
            try:
                self.lock.release()
                print(f"[Session] Force-unlocked session {self.session_id[:8]}")
            except RuntimeError:
                pass  # already unlocked

    def to_dict(self):
        """Serialize session for JSON persistence and API responses."""
        import re as _re
        last_preview = ''
        for msg in reversed(self.history):
            if msg.get('role') == 'assistant':
                raw = msg.get('content', '')
                clean = _re.sub(r'```[\s\S]*?```', '', raw)
                clean = _re.sub(r'[#*_`~>\[\]()!|]', '', clean)
                clean = ' '.join(clean.split())
                last_preview = clean[:110]
                break
        return {
            "session_id": self.session_id,
            "provider": self.provider,
            "model": self.model,
            "provider_thread_id": self.provider_thread_id,
            "context_bundles": list(self.context_bundles),
            "cli_session_id": self.cli_session_id,  # Legacy; gleich provider_thread_id wenn claude-cli
            "created_at": self.created_at,
            "message_count": self.message_count,
            "history": self.history,
            "custom_title": self.custom_title,
            "auto_title": self.auto_title,
            "title": self.custom_title or self.auto_title,
            "topic": self.topic,
            "topic_color": self.topic_color,
            "last_message": self.last_message,
            "last_heartbeat": self.last_heartbeat,
            "last_preview": last_preview,
        }


class SessionManager:
    """Manages all chat sessions, persisted to data/sessions.json."""

    def __init__(self):
        self.sessions: dict[str, ClaudeSession] = {}
        self._file = DATA_DIR / "sessions.json"
        self._topic_file = DATA_DIR / "topic-colors.json"
        self._load()

    def _load(self):
        """Load sessions from disk (with .bak fallback on corruption)."""
        raw = safe_read_json(self._file, default=None)
        if raw is not None:
            try:
                _skipped = 0
                for sid, data in raw.items():
                    if not isinstance(data, dict) or "history" not in data:
                        _skipped += 1
                        print(f"[Sessions] Skip dead session {sid[:8]} — invalid data (keys: {list(data.keys()) if isinstance(data, dict) else type(data).__name__})", flush=True)
                        continue
                    s = ClaudeSession(session_id=sid)
                    # Don't restore provider_thread_id/cli_session_id — old
                    # Provider-Threads sind nach Server-Restart unzuverlässig
                    # und führen zu Zombie-Prozessen bei --resume.
                    s.provider_thread_id = None
                    s.cli_session_id = None
                    # Provider-Config aus JSON (neue Felder). Fallback:
                    # alte Session ohne provider → "claude-cli" (war der Default).
                    s.provider = data.get("provider") or "claude-cli"
                    s.model = data.get("model") or ""
                    bundles_raw = data.get("context_bundles") or []
                    if isinstance(bundles_raw, list):
                        s.context_bundles = [str(b) for b in bundles_raw if b]
                    notes_raw = data.get("shared_notes") or {}
                    if isinstance(notes_raw, dict):
                        s.shared_notes = {str(k): str(v) for k, v in notes_raw.items()}
                    s.created_at = data.get("created_at", s.created_at)
                    s.message_count = data.get("message_count", 0)
                    s.history = data.get("history", [])
                    s.custom_title = data.get("custom_title")
                    s.auto_title = data.get("auto_title")
                    s.topic = data.get("topic")
                    s.topic_color = data.get("topic_color")
                    s.last_message = data.get("last_message")
                    s.last_heartbeat = data.get("last_heartbeat", s.last_message)
                    # Recover partial streaming response from crash. If the
                    # partial exists, append it as truncated assistant reply.
                    partial = data.get("_streaming_partial")
                    interrupted_stream = bool(partial) and (
                        not s.history or s.history[-1].get("role") != "assistant"
                    )
                    if interrupted_stream:
                        s.history.append({"role": "assistant", "content": partial + "\n\n_(unvollständig — Server wurde neugestartet)_"})
                        print(f"[Sessions] Recovered partial response for {sid[:8]} ({len(partial)} chars)")
                    s._streaming_partial = None  # clear after recovery
                    # Auto-resume trigger: _stream_active was True at last persist
                    # means the subprocess was alive when the server died — even
                    # if no text tokens had streamed yet (tool-phase kill).
                    if data.get("_stream_active"):
                        s.needs_resume = True
                        print(f"[Sessions] Session {sid[:8]} was mid-stream at shutdown → queued for resume")
                    s._stream_active = False  # reset after load
                    self.sessions[sid] = s
            except Exception as e:
                _skipped = 0
                print(f"[Sessions] Failed to parse sessions.json payload: {e}")
        else:
            _skipped = 0
        skip_note = f", {_skipped} tot übersprungen" if _skipped else ""
        print(f"[Sessions] Loaded {len(self.sessions)} session(s) from disk{skip_note}", flush=True)

    def persist(self):
        """Save all sessions to disk."""
        data = {}
        for sid, s in self.sessions.items():
            # Während des Chat-Refactors schreibt der existierende _run_claude auf
            # cli_session_id, der neue _run_chat auf provider_thread_id. Beide
                # Pfade schreiben in den gleichen Persistenz-Slot:
            effective_thread = s.provider_thread_id or s.cli_session_id
            data[sid] = {
                "provider": s.provider,
                "model": s.model,
                "provider_thread_id": effective_thread,
                "context_bundles": list(s.context_bundles),
                "shared_notes": dict(getattr(s, "shared_notes", {}) or {}),
                # Legacy-Feld für Backward-Compat; gleicher Wert wenn claude-cli
                "cli_session_id": effective_thread if s.provider == "claude-cli" else None,
                "created_at": s.created_at,
                "message_count": s.message_count,
                "history": s.history,
                "custom_title": s.custom_title,
                "auto_title": s.auto_title,
                "topic": s.topic,
                "topic_color": s.topic_color,
                "last_message": s.last_message,
                "last_heartbeat": s.last_heartbeat,
                "_streaming_partial": s._streaming_partial,
                "_stream_active": s._stream_active,
            }
        atomic_write_text(self._file, json.dumps(data, ensure_ascii=False, indent=2))

    def list_sessions(self):
        """Return dict of all sessions (serialized for API/WS responses)."""
        result = {}
        for sid, s in self.sessions.items():
            result[sid] = s.to_dict()
        return result

    def create_session(self, session_id=None):
        """Create a new session and persist."""
        s = ClaudeSession(session_id=session_id)
        self.sessions[s.session_id] = s
        self.persist()
        return s

    def get_or_create(self, session_id=None):
        """Get existing session by ID or create a new one."""
        if session_id and session_id in self.sessions:
            return self.sessions[session_id]
        return self.create_session(session_id=session_id)

    def delete_session(self, sid):
        """Delete a session by ID and persist."""
        # jarvis-main = interactive Jarvis chat, jarvis-heartbeat = background
        # tick log. Beide always-on, nie loeschbar.
        if sid in ("jarvis-main", "jarvis-heartbeat"):
            return
        if sid in self.sessions:
            del self.sessions[sid]
            self.persist()

    def rename_session(self, sid, title):
        """Set a custom title for a session."""
        if sid in self.sessions:
            self.sessions[sid].custom_title = title
            self.persist()

    def set_session_topic(self, session_id, topic, color):
        """Assign a topic and color to a session."""
        if session_id in self.sessions:
            self.sessions[session_id].topic = topic
            self.sessions[session_id].topic_color = color
            self.persist()

    def get_topics(self):
        """Return dict of topic name -> color from all sessions."""
        topics = {}
        for s in self.sessions.values():
            if s.topic and s.topic_color:
                topics[s.topic] = s.topic_color
        return topics

    def _topic_colors(self):
        """Load topic color mapping from topic-colors.json."""
        if self._topic_file.exists():
            try:
                return json.loads(self._topic_file.read_text(encoding="utf-8"))
            except (json.JSONDecodeError, Exception):
                pass
        return {}

    def _save_topic_color(self, topic, color):
        """Save a topic->color mapping to topic-colors.json."""
        colors = self._topic_colors()
        colors[topic] = color
        atomic_write_text(self._topic_file, json.dumps(colors, ensure_ascii=False, indent=2))

    @staticmethod
    def _default_color(topic):
        """Generate a deterministic color from a topic name."""
        palette = [
            "#f59e0b", "#06b6d4", "#6366f1", "#e11d48",
            "#10b981", "#8b5cf6", "#ec4899", "#14b8a6",
            "#f97316", "#3b82f6", "#a855f7", "#ef4444",
        ]
        h = sum(ord(c) for c in topic)
        return palette[h % len(palette)]


sessions = SessionManager()


# ---------------------------------------------------------------------------
# Window Hider  -  hide any console windows spawned by Claude CLI / MCP
# ---------------------------------------------------------------------------
if sys.platform == 'win32':
    import ctypes
    import ctypes.wintypes as _wt

    _user32 = ctypes.windll.user32
    _WNDENUMPROC = ctypes.WINFUNCTYPE(ctypes.c_bool, _wt.HWND, _wt.LPARAM)
    _SW_HIDE = 0
    _GW_OWNER = 4

    def _hide_console_windows():
        """Find and hide any visible console windows from claude/cmd/node MCP."""
        targets = {b'claude', b'cmd', b'node', b'npx'}

        @_WNDENUMPROC
        def _cb(hwnd, _):
            if not _user32.IsWindowVisible(hwnd):
                return True
            # Skip owned windows (dialogs etc.)
            if _user32.GetWindow(hwnd, _GW_OWNER):
                return True
            buf = ctypes.create_string_buffer(512)
            _user32.GetWindowTextA(hwnd, buf, 512)
            title = buf.value.lower()
            if any(t in title for t in targets):
                pid = _wt.DWORD()
                _user32.GetWindowThreadProcessId(hwnd, ctypes.byref(pid))
                # Don't hide our own terminal or user-opened windows
                # Only hide windows whose PID matches a known Claude subprocess
                # DISABLED: _user32.ShowWindow(hwnd, _SW_HIDE)
                print(f"[WinHider] Hidden window: '{buf.value.decode(errors='replace')}' PID={pid.value}")
            return True

        _user32.EnumWindows(_cb, 0)

    async def _window_hider_loop():
        return  # DISABLED — was hiding all windows with "claude" in title
else:
    async def _window_hider_loop():
        return  # no-op on non-Windows


# ---------------------------------------------------------------------------
# System Prompt Builder
# ---------------------------------------------------------------------------

_TOOLS_AND_MEMORY_BRIEF = """\
--- WERKZEUGE & GEDÄCHTNIS ---

MCP-Tools (über den `jarvis`-MCP-Server):
  • vault_search(query, n)        → semantische Suche in Robins Obsidian-Vault + Diss-PDFs (bge-m3, 9412 chunks)
  • vault_read_bundle(name)       → ganzes Themen-Bundle als Volltext (dissertation, medizin, trading, openclaw, famulatur)
  • vault_list_bundles()          → verfügbare Bundles
  • memory_list() / memory_read(f)→ MCs interne memory-Files (SOUL.md, projects.md, decisions.md, todos.md, robin.md, …)
  • skills_list() / skill_get(s)  → 129 Jarvis-Skills (academic-deep-research, algorithmic-art, …)
  • skill_activate(slug)          → Skill für Claude-CLI dauerhaft aktivieren (Wrapper in ~/.claude/skills/)
  • anki_list_decks()             → Decks aus laufendem Anki
  • generate_image(prompt)        → Nano-Banana oder GPT-Image — kann Bilder direkt in den Chat posten
  • ask_peer(provider, query)     → einen anderen LLM-Provider einmalig befragen (cross-provider second-opinion)
  • mc_health()                   → MC-Status-Check

WANN welches Tool:
  • Faktenfrage über Robin/Projekte  → vault_search
  • Ganzes Thema "atmen"             → vault_read_bundle
  • Spezialisierter Workflow nötig   → skills_list → skill_get → adoptieren
  • Medizin-Karte erzeugen           → skill_get("anki") + AnkiConnect
  • Visualisierung/Diagramm          → generate_image
  • Zweitmeinung eines anderen LLMs  → ask_peer

MEMORY-PFLEGE (wichtig für Kontinuität):
  Robins Jarvis-Gedächtnis lebt unter `data/memory/`. Struktur:
    data/memory/SOUL.md           — Identity, was Jarvis ist
    data/memory/robin.md          — Fakten über Robin (Person, Rolle, Prioritäten)
    data/memory/projects.md       — laufende Projekte
    data/memory/decisions.md      — getroffene Entscheidungen mit Warum
    data/memory/todos.md          — offene Todos
    data/memory/YYYY-MM-DD.md     — Tages-Log (was heute geschah, neue Erkenntnisse)

  Wenn du in dieser Session etwas Neues lernst — eine relevante Entscheidung, eine
  Projekt-Erkenntnis, etwas das Jarvis in 3 Monaten noch wissen sollte — biete
  am Ende deiner Antwort einen **kurzen Memory-Patch** an im Format:

      MEMORY-UPDATE (zur Bestätigung):
      Datei: data/memory/{zielfile}.md
      Anhang: <1-3 Zeilen, knapp, faktenbasiert>

  Robin entscheidet dann ob er das übernimmt. Niemals ungefragt Files schreiben
  außer der User bittet explizit darum oder du bist im Exec-Mode mit klarem Auftrag.
"""


def _build_system_prompt():
    """Build full system prompt — identity + memory + tools + context."""
    parts = []
    # SOUL.md — core identity
    soul = MEMORY_DIR / "SOUL.md"
    if soul.exists():
        parts.append(soul.read_text(encoding="utf-8", errors="replace"))
    # Identity + key memory files
    for name in ("identity.md", "robin.md", "decisions.md", "projects.md", "todos.md"):
        p = MEMORY_DIR / name
        if p.exists():
            content = p.read_text(encoding="utf-8", errors="replace")
            parts.append(f"--- {name} ---\n{content}")
    # Claude Code memory files (feedback, project state) — check all relevant project dirs
    for proj_dir in ("C--WINDOWS-system32", "C--Users-Robin-Jarvis-Mission-Control"):
        claude_memory = Path.home() / ".claude" / "projects" / proj_dir / "memory"
        if not claude_memory.exists():
            continue
        for f in claude_memory.glob("*.md"):
            try:
                content = f.read_text(encoding="utf-8", errors="replace")
                if content.strip():
                    parts.append(f"--- Claude Memory: {f.name} ---\n{content}")
            except Exception:
                pass
    # Werkzeug-Übersicht + Memory-Pflege-Anleitung (gleicher Inhalt für alle Bots)
    parts.append(_TOOLS_AND_MEMORY_BRIEF)
    return "\n\n".join(parts)


_VAULT_ROOT = Path("E:/OneDrive/AI/Obsidian-Vault/Jarvis-Brain")


#: Kuratierte "Themen-Bundles" — jedes Bundle ist eine Liste von Vault-Pfaden
#: (relativ zu _VAULT_ROOT), die on-demand als Volltext-Context mitgeschickt
#: werden. Ergänzt das RAG (semantisch getriggert, Chunks): Global-Context lädt
#: den kompletten Inhalt der Bundle-Pfade — nützlich wenn Robin über längere
#: Zeit in einem Thema arbeitet und das ganze "Kapitel" greifbar sein soll.
VAULT_CONTEXT_BUNDLES: dict[str, list[str]] = {
    "dissertation": [
        "Dissertation-ITI.md",
        "Jarvis-Knowledge/Dissertation-ITI",
    ],
    "medizin": [
        "Medizin.md",
        "OSCE.md",
        "Anki.md",
        "Jarvis-Knowledge/Medizin",
    ],
    "trading": [
        "Trading.md",
        "Jarvis-Knowledge/Trading",
    ],
    "openclaw": [
        "OpenClaw-Tech.md",
        "Jarvis-Knowledge/OpenClaw",
    ],
    "famulatur": [
        "Famulatur-Tanzania.md",
    ],
}


def _load_vault_context(
    bundle_names: list[str],
    max_total_chars: int = 40000,
    max_file_chars: int = 12000,
) -> tuple[str, int, list[str]]:
    """Lädt alle .md-Dateien aus den angegebenen Bundles konkateniert.

    Returns: (context_text, num_files, loaded_paths). Nicht-existierende Pfade
    werden still übersprungen.
    """
    if not bundle_names:
        return "", 0, []

    raw_paths: list[str] = []
    for bn in bundle_names:
        bundle = VAULT_CONTEXT_BUNDLES.get(bn.lower())
        if not bundle:
            continue
        raw_paths.extend(bundle)
    if not raw_paths:
        return "", 0, []

    loaded: list[tuple[str, str]] = []
    total = 0
    seen: set[Path] = set()

    for rel in raw_paths:
        abs_path = _VAULT_ROOT / rel
        if not abs_path.exists():
            continue
        if abs_path.is_file() and abs_path.suffix.lower() == ".md":
            if abs_path in seen:
                continue
            seen.add(abs_path)
            try:
                content = abs_path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
            if len(content) > max_file_chars:
                content = content[:max_file_chars].rsplit("\n", 1)[0] + "\n…[gekürzt]"
            total += len(content)
            loaded.append((rel, content))
            if total >= max_total_chars:
                break
        elif abs_path.is_dir():
            for md_file in sorted(abs_path.rglob("*.md")):
                if md_file in seen:
                    continue
                seen.add(md_file)
                try:
                    content = md_file.read_text(encoding="utf-8", errors="replace")
                except Exception:
                    continue
                if len(content) > max_file_chars:
                    content = content[:max_file_chars].rsplit("\n", 1)[0] + "\n…[gekürzt]"
                rel_to_vault = str(md_file.relative_to(_VAULT_ROOT)).replace("\\", "/")
                total += len(content)
                loaded.append((rel_to_vault, content))
                if total >= max_total_chars:
                    break
        if total >= max_total_chars:
            break

    if not loaded:
        return "", 0, []

    blocks = [f"### Vault/{p}\n{c}" for p, c in loaded]
    context = (
        "\n\n--- GLOBAL VAULT CONTEXT (aktive Themen-Bundles) ---\n"
        "Robin hat für diese Session folgende Vault-Bereiche als dauerhaften "
        "Kontext markiert. Nutze sie als Faktenbasis.\n\n"
        + "\n\n".join(blocks)
    )
    return context, len(loaded), [p for p, _ in loaded]


async def _build_rag_context(
    user_message: str,
    n_results: int = 4,
    min_score: float = 0.35,
    max_chars_per_chunk: int = 800,
) -> tuple[str, int]:
    """Query the RAG index with the user's message and return top-N chunks as a
    system-prompt-appendable context block. Works for ALL providers — no tool
    definitions needed, because it's injected into the system prompt.

    Gemini's "Stufe 1"-Vorschlag (2026-04-23): macht das bestehende ChromaDB +
    bge-m3 Setup für Claude-CLI/API, Codex-CLI, Gemini-CLI und OpenAI nutzbar
    ohne dass jeder Provider separat Tool-wiring braucht.

    Returns: (context_block, num_chunks) — context_block ist "" wenn keine
    Hits über Schwelle. num_chunks zum Loggen/UI.
    """
    if not user_message or len(user_message.strip()) < 10:
        return "", 0

    def _search():
        try:
            from rag.search import query as rag_query
            return rag_query(user_message, n_results)
        except Exception as e:
            print(f"[RAG] query failed: {e}")
            return []

    try:
        hits = await asyncio.to_thread(_search)
    except Exception as e:
        print(f"[RAG] to_thread failed: {e}")
        return "", 0

    # bge-m3 scores: 0-1 cosine similarity. 0.35 = "topically related",
    # 0.5+ = "clearly relevant". Threshold bewusst niedrig damit loose matches
    # auch reinkommen (der Provider entscheidet dann selbst was er nutzt).
    good = [h for h in (hits or []) if float(h.get("score") or 0) >= min_score]
    if not good:
        return "", 0

    blocks: list[str] = []
    for h in good:
        src_raw = str(h.get("source") or "unknown")
        if "Obsidian-Vault" in src_raw:
            src_short = "Vault/" + src_raw.split("Obsidian-Vault", 1)[-1].lstrip("\\/")
        else:
            src_short = src_raw.rsplit("\\", 1)[-1].rsplit("/", 1)[-1]
        text = (h.get("text") or "").strip()
        if len(text) > max_chars_per_chunk:
            text = text[:max_chars_per_chunk].rsplit(" ", 1)[0] + "…"
        score = float(h.get("score") or 0)
        blocks.append(f"### {src_short}  (relevance {score:.2f})\n{text}")

    context = (
        "\n\n--- JARVIS KNOWLEDGE (via RAG) ---\n"
        "Folgende Ausschnitte aus Robins Obsidian-Vault und PDF-Dissertation sind "
        "semantisch ähnlich zur aktuellen Frage. Nutze sie als zusätzlichen Kontext "
        "wenn sie relevant sind — ignoriere sonst.\n\n"
        + "\n\n".join(blocks)
    )
    return context, len(good)


# ---------------------------------------------------------------------------
# Memory Files API (reconstructed after accidental deletion)
# ---------------------------------------------------------------------------

async def api_memory_list(request):
    """List all memory/vault files for the Memory tab."""
    daily_claude = []
    daily_openclaw = []
    meta_files = []
    other_files = []
    CLAUDE_CODE_START = "2026-04-10"

    if not MEMORY_DIR.exists():
        return web.json_response([])

    for f in MEMORY_DIR.rglob("*.md"):
        stat = f.stat()
        entry = {
            "name": f.name,
            "path": str(f.relative_to(MEMORY_DIR)),
            "size": stat.st_size,
            "modified": stat.st_mtime,
        }

        if f.name[:4].isdigit() and len(f.name) == 13:
            if f.name[:10] >= CLAUDE_CODE_START:
                entry["era"] = "claude-code"
                daily_claude.append(entry)
            else:
                entry["era"] = "openclaw"
                daily_openclaw.append(entry)
        elif f.name in ("SOUL.md", "identity.md", "robin.md", "decisions.md",
                        "projects.md", "todos.md", "inbox.md", "failsafe.md",
                        "INDEX.md", "README.md"):
            entry["type"] = "meta"
            meta_files.append(entry)
        else:
            entry["type"] = "other"
            other_files.append(entry)

    # Sort: meta first, then claude daily (newest first), then openclaw daily (newest first), then other
    daily_claude.sort(key=lambda x: x["name"], reverse=True)
    daily_openclaw.sort(key=lambda x: x["name"], reverse=True)
    meta_files.sort(key=lambda x: x["name"])
    other_files.sort(key=lambda x: x["name"])

    all_files = meta_files + daily_claude + daily_openclaw + other_files
    return web.json_response(all_files)


# ---------------------------------------------------------------------------
# Tasks / Timeline APIs (reconstructed after accidental deletion)
# ---------------------------------------------------------------------------

async def api_tasks(request):
    """CRUD for simple task list stored in data/tasks.json."""
    tasks_file = DATA_DIR / "tasks.json"
    if request.method == "GET":
        if tasks_file.exists():
            return web.json_response(json.loads(tasks_file.read_text(encoding="utf-8")))
        return web.json_response([])
    # PUT or POST — save tasks
    body = await request.json()
    atomic_write_text(tasks_file, json.dumps(body, ensure_ascii=False, indent=2))
    return web.json_response({"ok": True})


async def api_timeline(request):
    """Return timeline events from data/timeline.json."""
    tl_file = DATA_DIR / "timeline.json"
    if tl_file.exists():
        return web.json_response(json.loads(tl_file.read_text(encoding="utf-8")))
    return web.json_response([])


async def api_memory_file(request):
    name = request.match_info["filename"]
    if not name.endswith(".md"):
        raise web.HTTPNotFound()
    # Search in Memory, Knowledge, then Vault root + ChatGPT-Import hubs
    VAULT_BASE = Path("E:/OneDrive/AI/Obsidian-Vault/Jarvis-Brain")
    candidates = [
        MEMORY_DIR / name,
        VAULT_BASE / "Jarvis-Knowledge" / name,
        VAULT_BASE / name,
        VAULT_BASE / "ChatGPT-Import" / "_Hubs" / name,
    ]
    for path in candidates:
        if path.exists():
            text = path.read_text(encoding="utf-8")
            return web.Response(text=text, content_type="text/plain", charset="utf-8")
    raise web.HTTPNotFound()


async def api_sessions(request):
    return web.json_response(sessions.list_sessions())


# ---------------------------------------------------------------------------
# New Panel APIs: Email, Calendar, System
# ---------------------------------------------------------------------------

SCRIPTS_DIR = Path("C:/Users/Robin/Jarvis/scripts")
JARVIS_DATA = Path("C:/Users/Robin/Jarvis/data")


async def _run_script(cmd, timeout=20):
    """Run a python script and return stdout."""
    env = {**os.environ, "PYTHONIOENCODING": "utf-8",
           "MS_CLIENT_ID": "59386692-536e-4afe-9ae2-0f728d617727",
           "MS_TENANT_ID": "common"}
    proc = await asyncio.create_subprocess_exec(
        *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
        env=env, **_SUBPROCESS_FLAGS,
    )
    try:
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
        return stdout.decode("utf-8", errors="replace").strip()
    except asyncio.TimeoutError:
        proc.kill()
        return ""


async def api_email(request):
    """Get unread emails from Outlook + Uni."""
    outlook_raw = await _run_script([
        sys.executable, str(SCRIPTS_DIR / "outlook_graph.py"), "unread", "--top", "15"
    ])
    # Parse outlook output
    outlook_mails = []
    for line in outlook_raw.split("\n"):
        if line.startswith("- "):
            parts = line[2:].split(" | ")
            if len(parts) >= 3:
                outlook_mails.append({
                    "date": parts[0].strip(),
                    "from": parts[1].strip(),
                    "subject": parts[2].strip(),
                    "account": "Outlook"
                })

    # Uni mail via IMAP
    uni_mails = []
    try:
        import imaplib
        import ssl
        cred_file = Path("C:/Users/Robin/Jarvis/secrets/uni-mail-cred.json")
        if cred_file.exists():
            creds = json.loads(cred_file.read_text(encoding="utf-8"))
            ctx = ssl.create_default_context()
            m = imaplib.IMAP4_SSL("imap.uni-ulm.de", 993, ssl_context=ctx)
            m.login("robin.mandel@uni-ulm.de", creds["password"])
            m.select("INBOX")
            _, data = m.search(None, "UNSEEN")
            ids = data[0].split()
            for mid in ids[-10:]:
                _, msg_data = m.fetch(mid, "(BODY[HEADER.FIELDS (FROM SUBJECT DATE)])")
                raw = msg_data[0][1].decode("utf-8", errors="replace")
                mail = {"account": "Uni", "from": "", "subject": "", "date": ""}
                for l in raw.strip().split("\n"):
                    l = l.strip()
                    if l.lower().startswith("from:"): mail["from"] = l[5:].strip()
                    elif l.lower().startswith("subject:"): mail["subject"] = l[8:].strip()
                    elif l.lower().startswith("date:"): mail["date"] = l[5:].strip()
                if mail["subject"] or mail["from"]:
                    uni_mails.append(mail)
            m.logout()
    except Exception as e:
        print(f"[Email] Uni error: {e}")

    return web.json_response({
        "outlook": outlook_mails,
        "uni": uni_mails,
        "total_unread": len(outlook_mails) + len(uni_mails)
    })


async def api_calendar(request):
    """Get upcoming calendar events."""
    days = request.query.get("days", "7")
    raw = await _run_script([
        sys.executable, str(SCRIPTS_DIR / "icloud-calendar.py"),
        "list-events", "--days", days, "--json"
    ], timeout=30)
    try:
        events = json.loads(raw) if raw else []
    except json.JSONDecodeError:
        events = []
    return web.json_response({"events": events, "days": int(days)})


async def api_quick_action(request):
    """Unified quick-action endpoint: Anki card / Email / Calendar event.
    POST body: {"type": "anki|email|calendar", ...params}
    """
    try:
        data = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "invalid_json"}, status=400)

    action_type = (data.get("type") or "").lower()

    if action_type == "anki":
        deck = (data.get("deck") or "Default").strip()
        front = (data.get("front") or "").strip()
        back = (data.get("back") or "").strip()
        tags = data.get("tags") or []
        if not front or not back:
            return web.json_response({"ok": False, "error": "front/back required"}, status=400)
        payload = {
            "action": "addNote",
            "version": 6,
            "params": {"note": {
                "deckName": deck,
                "modelName": "Basic",
                "fields": {"Front": front, "Back": back},
                "tags": tags if isinstance(tags, list) else [str(tags)],
                "options": {"allowDuplicate": False},
            }},
        }
        try:
            import urllib.request
            req = urllib.request.Request(
                "http://localhost:8765",
                data=json.dumps(payload).encode("utf-8"),
                headers={"Content-Type": "application/json"},
            )
            with urllib.request.urlopen(req, timeout=get_action_timeout("anki.quick")) as resp:
                result = json.loads(resp.read().decode("utf-8"))
            if result.get("error"):
                return web.json_response({"ok": False, "error": result["error"]}, status=500)
            return web.json_response({"ok": True, "note_id": result.get("result"), "deck": deck})
        except Exception as e:
            return web.json_response({"ok": False, "error": f"anki_unreachable: {e}"}, status=500)

    elif action_type == "email":
        to = (data.get("to") or "").strip()
        subject = (data.get("subject") or "").strip()
        body = (data.get("body") or "").strip()
        if not to:
            return web.json_response({"ok": False, "error": "to required"}, status=400)
        # write body to tmp file to avoid shell escaping issues
        import tempfile
        with tempfile.NamedTemporaryFile("w", encoding="utf-8", delete=False, suffix=".txt") as f:
            f.write(body)
            body_path = f.name
        try:
            out = await _run_script([
                sys.executable, str(SCRIPTS_DIR / "outlook_graph.py"), "send",
                "--to", to, "--subject", subject, "--body-file", body_path,
            ], timeout=get_action_timeout("email.send"))
        finally:
            try: os.unlink(body_path)
            except Exception: pass
        ok = "OK" in (out or "")
        return web.json_response({"ok": ok, "output": out})

    elif action_type == "calendar":
        title = (data.get("title") or "").strip()
        start = (data.get("start") or "").strip()  # ISO string
        duration = int(data.get("duration") or 60)
        location = (data.get("location") or "").strip()
        description = (data.get("description") or "").strip()
        if not title or not start:
            return web.json_response({"ok": False, "error": "title/start required"}, status=400)
        cmd = [
            sys.executable, str(SCRIPTS_DIR / "icloud-calendar.py"), "create-event",
            "--title", title, "--start", start, "--duration", str(duration),
        ]
        if location: cmd += ["--location", location]
        if description: cmd += ["--description", description]
        out = await _run_script(cmd, timeout=get_action_timeout("calendar.create"))
        ok = "error" not in (out or "").lower() and "traceback" not in (out or "").lower()
        return web.json_response({"ok": ok, "output": out})

    return web.json_response({"ok": False, "error": f"unknown type: {action_type}"}, status=400)


async def api_system(request):
    """Get system status: scheduled tasks, bridge, heartbeat."""
    result = {"tasks": [], "bridge": {}, "heartbeat": {}, "telegram": {}}

    # Heartbeat state
    hb = JARVIS_DATA / "heartbeat-state.json"
    if hb.exists():
        try:
            result["heartbeat"] = json.loads(hb.read_text(encoding="utf-8"))
        except Exception:
            pass

    # Telegram bridge state
    lock = JARVIS_DATA / "telegram-bridge.lock"
    result["bridge"]["running"] = lock.exists()
    tg_state = JARVIS_DATA / "telegram-state.json"
    if tg_state.exists():
        try:
            result["telegram"] = json.loads(tg_state.read_text(encoding="utf-8"))
        except Exception:
            pass

    # Memory sync state
    sync = JARVIS_DATA / "memory-sync-state.json"
    if sync.exists():
        try:
            result["memory_sync"] = json.loads(sync.read_text(encoding="utf-8"))
        except Exception:
            pass

    return web.json_response(result)



# ---------------------------------------------------------------------------
# RAG — Semantic search over Obsidian Vault + Dissertation
# ---------------------------------------------------------------------------

def _rag_available() -> bool:
    try:
        from rag.search import stats
        return True
    except Exception:
        return False


async def api_rag_search(request):
    """POST /api/rag/search  {query, n_results?, type?}"""
    try:
        body = await request.json()
        q = (body.get("query") or "").strip()
        if not q:
            return web.json_response({"error": "query required"}, status=400)
        n = int(body.get("n_results", 5))
        src_type = body.get("type")  # "markdown" | "pdf" | None

        import asyncio
        from rag.search import query as rag_query
        loop = asyncio.get_event_loop()
        hits = await loop.run_in_executor(None, lambda: rag_query(q, n, src_type))
        return web.json_response({"results": hits})
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_rag_status(request):
    """GET /api/rag/status"""
    try:
        from rag.search import stats
        import asyncio
        loop = asyncio.get_event_loop()
        s = await loop.run_in_executor(None, stats)
        return web.json_response({"ok": True, **s})
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)})


async def api_rag_reindex(request):
    """POST /api/rag/reindex  {force?}"""
    try:
        body = await request.json() if request.content_length else {}
        force = bool(body.get("force", False))
        import asyncio
        from rag.ingest import build_index
        loop = asyncio.get_event_loop()
        result = await loop.run_in_executor(None, lambda: build_index(force=force))
        return web.json_response({"ok": True, **result})
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)


async def api_rag_reload(request):
    """POST /api/rag/reload  -  Invalidate in-process Chroma singleton and reload stats."""
    try:
        import asyncio
        import sys
        # Ensure rag package is loaded first
        from rag import store as store_mod
        from rag.search import stats
        # Reset singletons so next call opens a fresh PersistentClient
        with store_mod._lock:
            store_mod._client = None
            store_mod._collection = None
        loop = asyncio.get_event_loop()
        s = await loop.run_in_executor(None, stats)
        print(f"[RAG] Singleton reloaded: {s['total_chunks']} chunks")
        return web.json_response({"ok": True, "reloaded": True, **s})
    except Exception as e:
        import traceback
        return web.json_response({"ok": False, "error": str(e), "trace": traceback.format_exc()}, status=500)


async def api_system_restart(request):
    """POST /api/system/restart  -  Sauberer Self-Shutdown, watchdog.py bringt
    den Server in ~15s wieder hoch. Trigger-File bleibt aus Diagnose-Gruenden,
    falls jemand spaeter doch noch einen separaten restart_watcher startet.
    """
    try:
        trigger_file = DATA_DIR / "restart.v2.trigger"
        trigger_file.write_text(f"restart at {datetime.now().isoformat()}")
        print(f"[System] Restart triggered via UI — exiting in 1s, watchdog respawn fast")

        async def _shutdown_soon():
            await asyncio.sleep(1.0)
            print("[System] Self-exit for restart")
            os._exit(0)

        asyncio.create_task(_shutdown_soon())

        return web.json_response({
            "ok": True,
            "message": "Restart ausgeloest — Server beendet sich in 1s, watchdog startet neu (~15s).",
            "trigger_file": str(trigger_file),
        })
    except Exception as e:
        print(f"[System] Restart trigger failed: {e}")
        return web.json_response({
            "ok": False,
            "error": str(e),
        }, status=500)


async def api_bots(request):
    """Get all cron jobs (Windows Scheduled Tasks) + active bot/session status."""
    import subprocess

    result = {
        "cron_jobs": [],
        "services": [],
        "mc_sessions": [],
    }

    # ? Windows Scheduled Tasks ?
    try:
        ps = (
            "$tasks = Get-ScheduledTask | Where-Object {$_.TaskName -match 'Jarvis|Mission Control|Telegram|Claude|Token|Dream|Trading|Bridge' -and $_.TaskPath -notmatch 'Microsoft'};"
            "$tasks | ForEach-Object {"
            "  $info = $_ | Get-ScheduledTaskInfo;"
            "  [PSCustomObject]@{"
            "    TaskName    = $_.TaskName;"
            "    State       = [int]$_.State;"
            "    Description = $_.Description;"
            "    LastRun     = if($info.LastRunTime -gt [DateTime]'2000-01-01'){$info.LastRunTime.ToString('o')}else{$null};"
            "    NextRun     = if($info.NextRunTime -gt [DateTime]'2000-01-01'){$info.NextRunTime.ToString('o')}else{$null};"
            "  }"
            "} | ConvertTo-Json -Compress"
        )
        proc = subprocess.run(
            ["powershell", "-NonInteractive", "-Command", ps],
            capture_output=True, text=True, timeout=10,
            **_SUBPROCESS_FLAGS,
        )
        if proc.returncode == 0 and proc.stdout.strip():
            raw = json.loads(proc.stdout.strip())
            if isinstance(raw, dict):
                raw = [raw]
            STATE_MAP = {0: "unknown", 1: "disabled", 2: "queued", 3: "ready", 4: "running"}
            for t in raw:
                state_num = t.get("State", 0)
                state = STATE_MAP.get(state_num, str(state_num))
                # parse timestamps
                def _fmt(ts):
                    if not ts:
                        return None
                    try:
                        from datetime import timezone
                        dt = datetime.fromisoformat(ts.replace("Z", "+00:00"))
                        if dt.year < 2000:
                            return None
                        return dt.strftime("%Y-%m-%d %H:%M")
                    except Exception:
                        return None
                result["cron_jobs"].append({
                    "name": t.get("TaskName", ""),
                    "state": state,
                    "last_run": _fmt(t.get("LastRun")),
                    "next_run": _fmt(t.get("NextRun")),
                    "description": t.get("Description") or "",
                })
    except Exception as e:
        result["cron_jobs_error"] = str(e)

    # ? Running Services / Bots ?
    services = []

    # Mission Control (this server)
    services.append({"name": "Mission Control", "type": "server", "status": "running",
                     "detail": "Port 8090", "icon": "✅"})

    # Telegram Bridge
    lock = JARVIS_DATA / "telegram-bridge.lock"
    tg_running = lock.exists()
    tg_state = {}
    tg_state_file = JARVIS_DATA / "telegram-state.json"
    if tg_state_file.exists():
        try:
            tg_state = json.loads(tg_state_file.read_text(encoding="utf-8"))
        except Exception:
            pass
    services.append({
        "name": "Telegram Bridge",
        "type": "bot",
        "status": "running" if tg_running else "stopped",
        "detail": f"@RobinMandels_Jarvis_bot" + (f" · {tg_state.get('last_msg_time','')}" if tg_state.get('last_msg_time') else ""),
        "icon": "⚙",
    })

    # Trading Bot (Port 8080)
    try:
        import socket
        s = socket.socket()
        s.settimeout(0.5)
        s.connect(("127.0.0.1", 8080))
        s.close()
        trading_running = True
    except Exception:
        trading_running = False
    services.append({
        "name": "Trading Bot Engine",
        "type": "bot",
        "status": "running" if trading_running else "stopped",
        "detail": "Port 8080",
        "icon": "📈",
    })

    # Night Research cron last result
    night_dir = JARVIS_DATA / "night-results"
    night_last = None
    if night_dir.exists():
        files = sorted(night_dir.glob("*.md"), key=lambda f: f.stat().st_mtime, reverse=True)
        if files:
            night_last = files[0].name
    services.append({
        "name": "Nacht-Recherche",
        "type": "cron",
        "status": "scheduled",
        "detail": f"Letztes Ergebnis: {night_last}" if night_last else "Noch kein Ergebnis",
        "icon": "🔬",
    })

    # Memory Sync
    sync_state = {}
    sync_file = JARVIS_DATA / "memory-sync-state.json"
    if sync_file.exists():
        try:
            sync_state = json.loads(sync_file.read_text(encoding="utf-8"))
        except Exception:
            pass
    services.append({
        "name": "Memory Sync",
        "type": "cron",
        "status": "scheduled",
        "detail": f"Letzter Sync: {sync_state.get('last_sync', '?')}",
        "icon": "🧠",
    })

    # RAG Index
    rag_manifest = JARVIS_DATA.parent / "chroma" / "files.json"
    rag_detail = "Nicht indexiert"
    rag_status = "stopped"
    if rag_manifest.exists():
        try:
            m = json.loads(rag_manifest.read_text(encoding="utf-8"))
            n_files = len(m)
            mtime = datetime.fromtimestamp(rag_manifest.stat().st_mtime).strftime("%d.%m. %H:%M")
            rag_detail = f"{n_files} Dateien · Stand {mtime}"
            rag_status = "running"
        except Exception:
            pass
    services.append({
        "name": "RAG Index",
        "type": "cron",
        "status": rag_status,
        "detail": rag_detail,
        "icon": "🔍",
    })

    result["services"] = services

    # ? MC Chat Sessions ?
    result["mc_sessions"] = sessions.list_sessions()

    return web.json_response(result)


async def api_bots_diss_research(request):
    """Return all night-research results for the Diss-Recherche bot detail view."""
    night_dir = JARVIS_DATA / "night-results"
    results = []
    if night_dir.exists():
        for f in sorted(night_dir.glob("*-diss-recherche.json"), reverse=True):
            try:
                data = json.loads(f.read_text(encoding="utf-8"))
                results.append(data)
            except Exception:
                pass
    # Aggregate: all unique papers across all runs, deduplicated by doi/title
    seen = set()
    all_papers = []
    for run in results:
        for item in run.get("items", []):
            key = item.get("doi") or item.get("title", "")
            if key and key not in seen:
                seen.add(key)
                all_papers.append({**item, "run_date": run.get("date")})
    # Sort by stars desc
    all_papers.sort(key=lambda p: p.get("stars", 0), reverse=True)
    return web.json_response({
        "runs": results,
        "all_papers": all_papers,
        "total_papers": len(all_papers),
        "total_runs": len(results),
    })


async def api_knowledge_graph(request):
    """Scan Jarvis-Memory and Jarvis-Knowledge for nodes and wikilinks."""
    import re
    OBSIDIAN_BASE = Path("E:/OneDrive/AI/Obsidian-Vault/Jarvis-Brain")
    scan_dirs = [
        ("daily",     OBSIDIAN_BASE / "Jarvis-Memory"),
        ("knowledge", OBSIDIAN_BASE / "Jarvis-Knowledge"),
        ("knowledge", OBSIDIAN_BASE),                          # root-level .md files
        ("knowledge", OBSIDIAN_BASE / "ChatGPT-Import" / "_Hubs"),  # ChatGPT hub pages
    ]
    CLAUDE_CODE_START = "2026-04-10"
    META_NAMES = {
        "SOUL.md", "identity.md", "robin.md", "decisions.md",
        "projects.md", "todos.md", "inbox.md", "failsafe.md",
        "INDEX.md", "README.md",
    }
    wikilink_re = re.compile(r'\[\[([^\]|#]+?)(?:\|[^\]]*)?\]\]')

    nodes = []
    links = []
    seen_nodes = set()

    for default_type, folder in scan_dirs:
        if not folder.exists():
            continue
        # Root vault dir: only direct files (not recursive, to avoid double-scanning subdirs)
        use_glob = (folder == OBSIDIAN_BASE)
        for f in (folder.glob("*.md") if use_glob else folder.rglob("*.md")):
            name = f.stem  # node id = filename without .md
            size = f.stat().st_size

            # Determine type
            if f.name in META_NAMES:
                ntype = "meta"
            elif f.name[:4].isdigit() and len(f.name) == 13:
                ntype = "daily"
            else:
                ntype = default_type

            if name not in seen_nodes:
                seen_nodes.add(name)
                nodes.append({"id": name, "type": ntype, "size": size})

            # Extract wikilinks
            try:
                text = f.read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
            for target in wikilink_re.findall(text):
                target = target.strip()
                if target and target != name:
                    links.append({"source": name, "target": target})
                    # Ensure target node exists (may be in another folder)
                    if target not in seen_nodes:
                        seen_nodes.add(target)
                        nodes.append({"id": target, "type": "knowledge", "size": 0})

    return web.json_response({"nodes": nodes, "links": links})


# ---------------------------------------------------------------------------
# Smart Tasks  -  Claude CLI Worker Queue
# ---------------------------------------------------------------------------

SMART_TASKS_FILE = Path("C:/Users/Robin/Jarvis/data/smart-tasks.json")
CLAUDE_EXE = r"C:\Users\Robin\.local\bin\claude.exe"

# Timeouts (seconds) pro Action-Typ fuer Bot-Orchestrierung.
# Zentrale Quelle — jede neue Action hier eintragen statt Magic-Numbers im Code.
# Greift automatisch via get_action_timeout() und api_quick_action.
ACTION_TIMEOUTS: dict = {
    # Quick-Actions (api_quick_action)
    "anki.quick":        15,   # AnkiConnect HTTP call
    "email.send":        45,   # outlook_graph.py send
    "calendar.create":   45,   # icloud-calendar.py create-event
    # LLM-basierte Generatoren
    "anki.generate":     240,  # anki-from-text.py (Claude-subprocess drin)
    "lernplan.generate": 240,  # lernplan-generator.py
    # Bot-Orchestrierung
    "smart_task.run":    900,  # voller claude -p run (hartes Limit)
    "arena.turn":        180,  # ein Bot-Turn in einer Arena-Diskussion
    "claude.bg":         600,  # generischer Claude-Background-Run
    # Fallback
    "_default":          30,
}


def get_action_timeout(name: str, default: int | None = None) -> int:
    """Liefert Timeout (s) fuer eine Action. Fallback: ACTION_TIMEOUTS['_default']."""
    if name in ACTION_TIMEOUTS:
        return int(ACTION_TIMEOUTS[name])
    if default is not None:
        return int(default)
    return int(ACTION_TIMEOUTS.get("_default", 30))

# Dict of running subprocesses: task_id -> asyncio.subprocess.Process
_running_tasks: dict = {}


_smart_tasks_lock = asyncio.Lock() if hasattr(asyncio, 'Lock') else None

def _load_smart_tasks() -> list:
    return safe_read_json(SMART_TASKS_FILE, default=[]) or []


def _save_smart_tasks(tasks: list):
    atomic_write_text(SMART_TASKS_FILE, json.dumps(tasks, ensure_ascii=False, indent=2))


def _update_task_field(task_id: str, **fields):
    """Thread-safe update of a single task's fields."""
    tasks = _load_smart_tasks()
    for t in tasks:
        if t["id"] == task_id:
            t.update(fields)
            break
    _save_smart_tasks(tasks)


def _find_task(tasks: list, task_id: str):
    for t in tasks:
        if t["id"] == task_id:
            return t
    return None


# ── Task Sizer ────────────────────────────────────────────────────────────────

_SIZE_LARGE_KEYWORDS = {
    "refactor", "refaktoriere", "implementiere", "implementier", "baue", "build",
    "erstelle", "create", "system", "architektur", "migration", "migriere",
    "redesign", "komplett", "vollständig", "umstrukturiere", "pipeline",
    "framework", "integration", "scraper", "crawler", "dashboard", "analyse",
    "analyze", "research", "recherche", "alle", "multiple", "several",
}
_SIZE_SMALL_KEYWORDS = {
    "zeig", "show", "liste", "list", "check", "prüf", "was", "what",
    "wann", "when", "wie viel", "how many", "status", "get", "find",
    "such", "öffne", "open", "schreib kurz", "kurze", "quick", "fix",
    "fixe", "korrigiere", "rename", "umbenennen", "delete", "lösch",
}

def classify_task_size(title: str, description: str = "") -> dict:
    """Heuristic classifier: returns size (small/medium/large) with confidence and reason."""
    combined = f"{title} {description}".lower().strip()
    word_count = len(combined.split())

    score = 0  # negative = small, positive = large

    # Length signals
    if word_count <= 8:
        score -= 2
    elif word_count <= 20:
        score -= 1
    elif word_count >= 40:
        score += 2
    elif word_count >= 25:
        score += 1

    # Keyword signals
    for kw in _SIZE_LARGE_KEYWORDS:
        if kw in combined:
            score += 2
            break
    for kw in _SIZE_SMALL_KEYWORDS:
        if kw in combined:
            score -= 2
            break

    # Multi-step signals
    step_markers = combined.count("\n") + combined.count("- ") + combined.count(". ")
    multi_keywords = ["außerdem", "zusätzlich", "dann", "anschließend", "danach", "furthermore"]
    step_markers += sum(1 for mk in multi_keywords if mk in combined)
    if step_markers >= 4:
        score += 2
    elif step_markers >= 2:
        score += 1

    # File/path mentions suggest real work
    if "/" in combined or "\\" in combined or ".py" in combined or ".js" in combined:
        score += 1

    if score <= -2:
        size, reason = "small", "Kurze/einfache Anfrage, einzelne Aktion"
    elif score <= 1:
        size, reason = "medium", "Mehrstufige Anfrage mit klarem Scope"
    else:
        size, reason = "large", "Komplexe Aufgabe, viele Schritte oder System-Level"

    return {"size": size, "score": score, "reason": reason, "word_count": word_count}


async def api_smart_tasks_size(request):
    """POST /api/smart-tasks/size  -  Klassifiziert einen Task als small/medium/large."""
    body = await request.json()
    title = body.get("title", "")
    description = body.get("description", "")
    result = classify_task_size(title, description)
    return web.json_response(result)


async def api_smart_tasks_list(request):
    """GET /api/smart-tasks  -  Liste aller Tasks."""
    tasks = _load_smart_tasks()
    return web.json_response(tasks)


async def api_smart_tasks_create(request):
    """POST /api/smart-tasks  -  Neuen Task erstellen."""
    body = await request.json()
    task = {
        "id": str(uuid.uuid4()),
        "title": body.get("title", "Unbenannter Task"),
        "description": body.get("description", ""),
        "original_request": body.get("original_request", ""),
        "chat_context": body.get("chat_context", ""),
        "system_prompt": body.get("system_prompt", ""),
        "priority": body.get("priority", "medium"),
        "status": "pending",
        "created": datetime.now().isoformat(),
        "started": None,
        "completed": None,
        "output": "",
        "error": "",
    }
    sizing = classify_task_size(task["title"], task["description"])
    task["size"] = sizing["size"]
    task["size_reason"] = sizing["reason"]
    tasks = _load_smart_tasks()
    tasks.append(task)
    _save_smart_tasks(tasks)
    return web.json_response(task, status=201)


async def _run_smart_task(task_id: str):
    """Startet claude -p als Subprocess mit vollem Jarvis-Kontext."""
    tasks = _load_smart_tasks()
    task = _find_task(tasks, task_id)
    if not task:
        return

    task["status"] = "running"
    task["started"] = datetime.now().isoformat()
    task["output"] = ""
    task["error"] = ""
    _save_smart_tasks(tasks)

    desc = task.get('description', '') or ''
    original = task.get('original_request', '') or ''
    chat_context = task.get('chat_context', '') or ''
    system_prompt = task.get('system_prompt', '') or ''

    # Build rich prompt with full context
    context_parts = [f"# Aufgabe: {task['title']}"]
    if system_prompt:
        context_parts.append(f"\n## Kontext\n{system_prompt}")
    if original:
        context_parts.append(f"\n## Original-Anfrage\n{original}")
    if desc:
        context_parts.append(f"\n## Beschreibung\n{desc}")
    if chat_context:
        context_parts.append(f"\n## Chat-Kontext (letzte Nachrichten)\n{chat_context}")
    context_parts.append(
        "\n## Anweisungen"
        "\n- Fuehre diese Aufgabe Schritt fuer Schritt aus."
        "\n- Lies zuerst CLAUDE.md fuer Projektkontext und Regeln."
        "\n- Wenn die Aufgabe Code/Scripts beinhaltet: schreibe die Dateien direkt."
        "\n- Pruefe dein Ergebnis bevor du es als fertig meldest."
        "\n- Antworte am Ende mit einer kurzen Zusammenfassung was du gemacht hast."
        "\n- Sprache: Deutsch."
    )
    prompt = "\n".join(context_parts)

    cmd = [
        CLAUDE_EXE,
        "-p", prompt,
        "--output-format", "text",
        "--max-turns", "25",
    ]

    # (encoding-fixed) print(f"[SmartTask] Start: {task_id[:8]} €" {task['title'][:50]}")
    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            stdin=asyncio.subprocess.DEVNULL,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        _running_tasks[task_id] = proc

        # Use communicate()  -  text format buffers everything until done
        _st_timeout = get_action_timeout("smart_task.run")
        try:
            stdout_data, stderr_data = await asyncio.wait_for(
                proc.communicate(), timeout=_st_timeout
            )
        except asyncio.TimeoutError:
            proc.kill()
            await proc.wait()
            _update_task_field(task_id,
                status="error", error=f"Timeout nach {_st_timeout}s",
                output="", completed=datetime.now().isoformat())
            _running_tasks.pop(task_id, None)
            return

        full_output = stdout_data.decode("utf-8", errors="replace").strip()
        stderr_text = stderr_data.decode("utf-8", errors="replace").strip()

        if proc.returncode == 0:
            _update_task_field(task_id,
                status="done", output=full_output or "(Kein Output)",
                error="", completed=datetime.now().isoformat())
        else:
            _update_task_field(task_id,
                status="error", output=full_output,
                error=stderr_text or f"Exit code {proc.returncode}",
                completed=datetime.now().isoformat())

        # (encoding-fixed) print(f"[SmartTask] Done: {task_id[:8]} €" exit={proc.returncode}")

    except Exception as e:
        # (encoding-fixed) print(f"[SmartTask] Error: {task_id[:8]} €" {e}")
        _update_task_field(task_id,
            status="error", error=str(e),
            completed=datetime.now().isoformat())
    finally:
        _running_tasks.pop(task_id, None)


async def api_smart_tasks_start(request):
    """POST /api/smart-tasks/{id}/start  -  Task starten."""
    task_id = request.match_info["id"]
    tasks = _load_smart_tasks()
    task = _find_task(tasks, task_id)
    if not task:
        raise web.HTTPNotFound(reason="Task nicht gefunden")
    if task["status"] == "running":
        return web.json_response({"ok": False, "error": "Task laeuft bereits"}, status=409)

    # Reset für Neustart
    task["status"] = "pending"
    _save_smart_tasks(tasks)

    asyncio.create_task(_run_smart_task(task_id))
    return web.json_response({"ok": True, "task_id": task_id})


async def api_smart_tasks_get(request):
    """GET /api/smart-tasks/{id}  -  Task-Status + Output."""
    task_id = request.match_info["id"]
    tasks = _load_smart_tasks()
    task = _find_task(tasks, task_id)
    if not task:
        raise web.HTTPNotFound(reason="Task nicht gefunden")
    return web.json_response(task)


async def api_smart_tasks_delete(request):
    """DELETE /api/smart-tasks/{id}  -  Task löschen."""
    task_id = request.match_info["id"]
    # Laufenden Prozess abbrechen falls vorhanden
    proc = _running_tasks.get(task_id)
    if proc and proc.returncode is None:
        proc.terminate()
        try:
            await asyncio.wait_for(proc.wait(), timeout=3)
        except asyncio.TimeoutError:
            proc.kill()
        _running_tasks.pop(task_id, None)

    tasks = _load_smart_tasks()
    new_tasks = [t for t in tasks if t["id"] != task_id]
    if len(new_tasks) == len(tasks):
        raise web.HTTPNotFound(reason="Task nicht gefunden")
    _save_smart_tasks(new_tasks)
    return web.json_response({"ok": True})


# ---------------------------------------------------------------------------
# Roadmap / projects.md Endpoints
# ---------------------------------------------------------------------------

PROJECTS_MD_PATH = Path("E:/OneDrive/AI/Obsidian-Vault/Jarvis-Brain/Jarvis-Memory/projects.md")


def _item_id(text: str) -> str:
    """Stabiler Hash aus Item-Text (8 Hex-Zeichen)."""
    import hashlib, unicodedata
    normalized = unicodedata.normalize("NFC", text.strip())
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()[:8]


def _parse_projects_md() -> dict:
    """Parst projects.md und gibt strukturiertes JSON zurück."""
    import re
    if not PROJECTS_MD_PATH.exists():
        return {"categories": [], "stats": {"total": 0, "done": 0, "open": 0, "percent": 0}}

    text = PROJECTS_MD_PATH.read_text(encoding="utf-8")
    lines = text.splitlines()

    categories = []
    current_cat = None

    for line in lines:
        # Kategorie-Überschrift: ## 1. Titel
        m = re.match(r'^##\s+(.+)', line)
        if m:
            current_cat = {"name": m.group(1).strip(), "items": []}
            categories.append(current_cat)
            continue

        if current_cat is None:
            continue

        # Item mit ✅ oder ⬜
        m_done = re.match(r'^\s*-\s+✅\s+(.+)', line)
        m_open = re.match(r'^\s*-\s+⬜\s+(.+)', line)

        if m_done:
            txt = m_done.group(1).strip()
            current_cat["items"].append({
                "id": _item_id(txt),
                "text": txt,
                "done": True,
            })
        elif m_open:
            txt = m_open.group(1).strip()
            current_cat["items"].append({
                "id": _item_id(txt),
                "text": txt,
                "done": False,
            })

    total = sum(len(c["items"]) for c in categories)
    done  = sum(sum(1 for i in c["items"] if i["done"]) for c in categories)
    open_ = total - done
    pct   = round(done / total * 100) if total else 0

    return {
        "categories": categories,
        "stats": {"total": total, "done": done, "open": open_, "percent": pct},
    }


def _find_roadmap_item(item_id: str):
    """Gibt (kategorie_name, item_dict) zurück oder (None, None)."""
    data = _parse_projects_md()
    for cat in data["categories"]:
        for item in cat["items"]:
            if item["id"] == item_id:
                return cat["name"], item
    return None, None


async def api_roadmap_get(request):
    """GET /api/roadmap  -  Parst projects.md, gibt JSON zurück."""
    return web.json_response(_parse_projects_md())


async def api_roadmap_toggle(request):
    """POST /api/roadmap/{id}/toggle  -  Togglet ✅/â¬œ in projects.md."""
    item_id = request.match_info["id"]
    _, item = _find_roadmap_item(item_id)
    if not item:
        raise web.HTTPNotFound(reason="Item nicht gefunden")

    text = PROJECTS_MD_PATH.read_text(encoding="utf-8")
    import re

    item_text_escaped = re.escape(item["text"])
    if item["done"]:
        # ✅ â†' â¬œ
        pattern = r'(- )✅( ' + item_text_escaped + r')'
        replacement = r'\1⬜\2'
    else:
        # â¬œ â†' ✅
        pattern = r'(- )⬜( ' + item_text_escaped + r')'
        replacement = r'\1✅\2'

    new_text = re.sub(pattern, replacement, text, count=1)
    if new_text == text:
        return web.json_response({"ok": False, "error": "Kein Match gefunden"}, status=400)

    PROJECTS_MD_PATH.write_text(new_text, encoding="utf-8")
    return web.json_response({"ok": True, "done": not item["done"]})


async def api_roadmap_build(request):
    """POST /api/roadmap/{id}/build  -  Erstellt Smart-Task für Item und startet ihn."""
    item_id = request.match_info["id"]
    cat_name, item = _find_roadmap_item(item_id)
    if not item:
        raise web.HTTPNotFound(reason="Item nicht gefunden")
    if item["done"]:
        return web.json_response({"ok": False, "error": "Item ist bereits erledigt"}, status=400)

    # Check if there's an original request stored for this item
    context_file = JARVIS_DATA / "roadmap-context.json"
    original_request = ""
    if context_file.exists():
        try:
            ctx = json.loads(context_file.read_text(encoding="utf-8"))
            original_request = ctx.get(item_id, "")
        except Exception:
            pass

    # Smart-Task erstellen
    task = {
        "id": str(uuid.uuid4()),
        "title": item["text"],
        "description": (
            f"Jarvis-Roadmap Item aus Kategorie '{cat_name}'.\n"
            f"Baue / implementiere: {item['text']}"
        ),
        "original_request": original_request,
        "priority": "high",
        "status": "pending",
        "created": datetime.now().isoformat(),
        "started": None,
        "completed": None,
        "output": "",
        "error": "",
        "roadmap_id": item_id,
    }
    tasks = _load_smart_tasks()
    tasks.append(task)
    _save_smart_tasks(tasks)

    # Direkt starten
    asyncio.create_task(_run_smart_task(task["id"]))

    return web.json_response({"ok": True, "task_id": task["id"]}, status=201)


async def api_roadmap_suggest(request):
    """POST /api/roadmap/suggest  -  Idee -> Claude -> Items in projects.md einfügen."""
    import re
    import subprocess

    body = await request.json()
    message = body.get("message", "").strip()
    if not message:
        raise web.HTTPBadRequest(reason="message darf nicht leer sein")

    prompt = (
        f'Du bist Jarvis. Robin hat eine Idee fuer ein neues Feature:\n'
        f'"{message}"\n\n'
        f'Wandle das in 1-3 konkrete, umsetzbare Todo-Items um.\n'
        f'Bestimme fuer jedes Item die passende Kategorie aus dieser Liste:\n'
        f'1. Grundsystem & Identitaet\n'
        f'2. Gedaechtnis & Selbstverbesserung\n'
        f'3. Heartbeats & Cronjobs\n'
        f'4. Modelle & Routing\n'
        f'5. Skills\n'
        f'6. Agents\n'
        f'7. Kommunikation\n'
        f'8. Email & Kalender\n'
        f'9. Studium & Forschung\n'
        f'10. Mission Control\n'
        f'11. Trading\n'
        f'12. Infrastruktur\n'
        f'13. Medien-KI\n\n'
        f'Antworte NUR als JSON Array:\n'
        f'[{{"category": "3. Heartbeats & Cronjobs", "text": "Morning Briefing per Telegram"}}]'
    )

    print(f"[RoadmapSuggest] Calling Claude for: {message[:60]}")
    try:
        result = subprocess.run(
            [CLAUDE_EXE, "-p", prompt, "--output-format", "json",
             "--max-turns", "1",
             "--system-prompt", "Antworte NUR als JSON Array. Keine Tools nutzen.",
             "--disallowed-tools", "Bash", "Read", "Edit", "Write", "Glob", "Grep", "WebFetch", "WebSearch", "Agent"],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            stdin=subprocess.DEVNULL,
            timeout=90,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        raw_output = result.stdout.strip()
        print(f"[RoadmapSuggest] Raw output ({len(raw_output)} bytes)")
    except subprocess.TimeoutExpired:
        raise web.HTTPInternalServerError(reason="Claude Timeout (90s)")
    except Exception as e:
        raise web.HTTPInternalServerError(reason=f"Claude Fehler: {e}")

    # Parse Claude JSON envelope
    try:
        envelope = json.loads(raw_output)
        inner = envelope.get("result", raw_output)
    except json.JSONDecodeError:
        inner = raw_output

    # Extract JSON array from response (might be wrapped in markdown or extra text)
    # Try multiple strategies
    items = None
    # Strategy 1: direct parse
    try:
        parsed = json.loads(inner)
        if isinstance(parsed, list):
            items = parsed
    except (json.JSONDecodeError, TypeError):
        pass

    # Strategy 2: find array with bracket matching
    if items is None:
        start = inner.find('[')
        if start >= 0:
            depth = 0
            end = start
            for i, ch in enumerate(inner[start:], start):
                if ch == '[': depth += 1
                elif ch == ']': depth -= 1
                if depth == 0:
                    end = i + 1
                    break
            try:
                items = json.loads(inner[start:end])
            except json.JSONDecodeError:
                pass

    if not items or not isinstance(items, list):
        raise web.HTTPInternalServerError(reason=f"Kein JSON-Array in Antwort: {inner[:200]}")

    if not PROJECTS_MD_PATH.exists():
        raise web.HTTPInternalServerError(reason="projects.md nicht gefunden")

    text = PROJECTS_MD_PATH.read_text(encoding="utf-8")
    lines = text.splitlines(keepends=True)
    inserted = []

    for item in items:
        cat_name = item.get("category", "").strip()
        item_text = item.get("text", "").strip()
        if not item_text or not cat_name:
            continue

        # Find the category header line
        # Match: ## <number>. <rest> or ## <cat_name> exactly
        cat_number = cat_name.split(".")[0].strip() if "." in cat_name else ""
        cat_rest   = cat_name.split(".", 1)[1].strip() if "." in cat_name else cat_name

        insert_idx = None
        for i, line in enumerate(lines):
            stripped = line.strip()
            if stripped.startswith("##"):
                header = stripped.lstrip("#").strip()
                # Match by full name or by number prefix
                if header == cat_name or header.startswith(cat_number + "."):
                    # Find end of this category's block (next ## or EOF)
                    j = i + 1
                    while j < len(lines):
                        if lines[j].strip().startswith("##"):
                            break
                        j += 1
                    # Insert before next category or at end of file
                    insert_idx = j
                    break

        if insert_idx is None:
            # Category not found  -  append at end
            insert_idx = len(lines)

        new_line = f"- ⬜ {item_text}\n"
        lines.insert(insert_idx, new_line)
        inserted.append({
            "category": cat_name,
            "text": item_text,
            "id": _item_id(item_text),
            "original_request": message,  # Preserve Robin's original input
        })
        print(f"[RoadmapSuggest] Inserted '{item_text}' into '{cat_name}' at line {insert_idx}")

    PROJECTS_MD_PATH.write_text("".join(lines), encoding="utf-8")

    # Save original request context for each item (so Build knows what Robin wanted)
    context_file = JARVIS_DATA / "roadmap-context.json"
    ctx = {}
    if context_file.exists():
        try:
            ctx = json.loads(context_file.read_text(encoding="utf-8"))
        except Exception:
            pass
    for item in inserted:
        ctx[item["id"]] = message
    context_file.write_text(json.dumps(ctx, ensure_ascii=False, indent=2), encoding="utf-8")

    return web.json_response({"ok": True, "inserted": inserted})


DISS_RESEARCH_SESSIONS_FILE = JARVIS_DATA / "diss_research_sessions.json"


async def api_diss_research(request):
    """GET — list research sessions  |  POST — save a session entry."""
    if request.method == "POST":
        try:
            entry = await request.json()
        except Exception:
            return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
        sessions = safe_read_json(DISS_RESEARCH_SESSIONS_FILE, default=[]) or []
        sessions.insert(0, entry)
        atomic_write_text(DISS_RESEARCH_SESSIONS_FILE, json.dumps(sessions[:100], ensure_ascii=False, indent=2))
        return web.json_response({"ok": True})

    # GET — merge persisted sessions + live night-results files
    sessions = safe_read_json(DISS_RESEARCH_SESSIONS_FILE, default=[]) or []
    night_dir = JARVIS_DATA / "night-results"
    if night_dir.exists():
        for f in sorted(night_dir.glob("*diss*"), reverse=True)[:5]:
            try:
                data = json.loads(f.read_text(encoding="utf-8"))
                if isinstance(data, dict):
                    sessions.append({
                        "date": f.stem.split("-diss")[0],
                        "takeaway": data.get("takeaway", ""),
                        "papers": data.get("papers", []),
                    })
                else:
                    sessions.append({"date": f.stem, "takeaway": str(data)[:200]})
            except Exception:
                pass
    return web.json_response({"results": sessions})


DISS_WIKI_DIR = Path("E:/OneDrive/Dokumente/Studium/!!Medizin/03 Doktorarbeit/ITI/Jarvis Workspace/wiki")
DISS_SESSION_TOPIC_FILE = JARVIS_DATA / "diss-session-topic.json"


async def api_diss_session_topic(request):
    """GET/POST current dissertation session topic."""
    if request.method == "POST":
        data = await request.json()
        DISS_SESSION_TOPIC_FILE.write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
        return web.json_response({"ok": True})
    if DISS_SESSION_TOPIC_FILE.exists():
        return web.json_response(json.loads(DISS_SESSION_TOPIC_FILE.read_text(encoding="utf-8")))
    return web.json_response({"topic": "", "date": ""})


async def api_diss_wiki_log(request):
    """GET recent wiki log entries."""
    log_file = DISS_WIKI_DIR / "log.md"
    if not log_file.exists():
        return web.json_response({"entries": []})
    content = log_file.read_text(encoding="utf-8")
    entries = []
    current = None
    for line in content.split("\n"):
        if line.startswith("## ["):
            if current:
                entries.append(current)
            current = {"header": line[3:].strip(), "body": ""}
        elif current and line.strip():
            current["body"] += line.strip() + " "
    if current:
        entries.append(current)
    return web.json_response({"entries": entries[:8]})


MEDIZIN_LOG_FILE = Path("C:/Users/Robin/Jarvis/data/medizin-log.json")

async def api_medizin_log(request):
    """GET or POST Medizin knowledge log entries."""
    if request.method == "POST":
        try:
            data = await request.json()
            entries = []
            if MEDIZIN_LOG_FILE.exists():
                entries = json.loads(MEDIZIN_LOG_FILE.read_text(encoding="utf-8"))
            entries.insert(0, {
                "title":   data.get("title", "Unbenannt"),
                "type":    data.get("type", "sonstiges"),
                "preview": data.get("preview", ""),
                "ts":      data.get("ts", ""),
            })
            # Keep last 200 entries
            entries = entries[:200]
            MEDIZIN_LOG_FILE.parent.mkdir(parents=True, exist_ok=True)
            MEDIZIN_LOG_FILE.write_text(json.dumps(entries, ensure_ascii=False, indent=2), encoding="utf-8")
            return web.json_response({"ok": True})
        except Exception as e:
            return web.json_response({"error": str(e)}, status=500)
    else:
        if not MEDIZIN_LOG_FILE.exists():
            return web.json_response({"entries": []})
        try:
            entries = json.loads(MEDIZIN_LOG_FILE.read_text(encoding="utf-8"))
            return web.json_response({"entries": entries})
        except Exception:
            return web.json_response({"entries": []})


async def api_anki_generate(request):
    """POST /api/anki/generate — Erzeugt Anki-Karten aus Vorlesungstext.
    Body: { text, deck?, title?, tags? }
    Ruft scripts/anki-from-text.py auf (Claude CLI + AnkiConnect)."""
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    text = (body.get("text") or "").strip()
    if not text:
        return web.json_response({"ok": False, "error": "text ist leer"}, status=400)

    deck = body.get("deck") or "Medizin::Sem8::Auto"
    title = body.get("title") or ""
    tags = body.get("tags") or []
    if isinstance(tags, str):
        tags = [t.strip() for t in tags.split(",") if t.strip()]

    script = Path("C:/Users/Robin/Jarvis/scripts/anki-from-text.py")
    if not script.exists():
        return web.json_response({"ok": False, "error": f"Script fehlt: {script}"}, status=500)

    # Mode A: Skript NUR um cards generieren zu lassen (--cards-only).
    # Image-Pass + Anki-Sync passieren danach hier im Server, mit der existierenden
    # Bulk-Sync-Pipeline — die kann Bilder via storeMediaFile pushen.
    try:
        proc = await asyncio.create_subprocess_exec(
            sys.executable, str(script),
            "--deck", deck,
            "--title", title,
            "--tags", ",".join(tags),
            "--cards-only",
            stdin=asyncio.subprocess.PIPE,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(
            proc.communicate(text.encode("utf-8")),
            timeout=get_action_timeout("anki.generate"),
        )
        out = stdout.decode("utf-8", errors="replace").strip()
        try:
            cards_result = json.loads(out)
        except Exception:
            return web.json_response({
                "ok": False,
                "error": "Konnte Script-Output nicht parsen",
                "stdout": out[:500],
                "stderr": stderr.decode("utf-8", errors="replace")[:500],
            })

        if not cards_result.get("ok"):
            return web.json_response(cards_result)

        cards = cards_result.get("cards") or []
        if not cards:
            return web.json_response({"ok": False, "error": "Keine Karten generiert"})

        # Bild-Pass: parallel Nano Banana 2 fuer Karten mit image_prompt
        await _enrich_anki_cards_with_images(cards)

        # Anki-Push via existierende Bulk-Sync-Pipeline (storeMediaFile + addNotes).
        # _bulk_sync_cards rendert per _card_to_anki_note die Bilder mit ein.
        try:
            sync_report = await _bulk_sync_cards(cards, deck=deck, subject=title or "Karten")
        except Exception as e:
            return web.json_response({
                "ok": False,
                "error": f"Anki-Sync fehlgeschlagen: {e}",
                "cards_generated": len(cards),
                "images_generated": sum(1 for c in cards if c.get("image")),
            })

        return web.json_response({
            "ok": True,
            "deck": deck,
            "cards_generated": len(cards),
            "images_generated": sum(1 for c in cards if c.get("image")),
            **sync_report,
        })
    except asyncio.TimeoutError:
        return web.json_response({"ok": False, "error": f"Timeout nach {get_action_timeout('anki.generate')}s"}, status=504)
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)


# Limits fuer Anki-Auto-Image-Generation
_ANKI_IMG_MAX_TOTAL = 8           # max Bilder pro Generate-Call (Cost-Cap)
_ANKI_IMG_PARALLEL = 4            # Concurrency
_ANKI_IMG_MODEL = "gemini-3.1-flash-image-preview"  # Nano Banana 2


async def _enrich_anki_cards_with_images(cards: list[dict]) -> None:
    """In-place: Fuer Karten mit image_prompt generiert Nano Banana 2 ein Bild
    und schreibt absoluten Library-Pfad in card['image']. Karten ohne
    image_prompt werden uebersprungen. Cap bei _ANKI_IMG_MAX_TOTAL."""
    from image_client import generate_image, ImageGenError
    import library as _lib

    candidates = [c for c in cards if c.get("image_prompt") and not c.get("image")]
    if not candidates:
        return
    candidates = candidates[:_ANKI_IMG_MAX_TOTAL]

    sem = asyncio.Semaphore(_ANKI_IMG_PARALLEL)

    async def _one(card: dict) -> None:
        prompt = (card.get("image_prompt") or "").strip()
        if not prompt:
            return
        async with sem:
            try:
                img = await generate_image(
                    provider="gemini",
                    prompt=prompt,
                    uploads_dir=UPLOADS_DIR,
                    model=_ANKI_IMG_MODEL,
                    size="1024x1024",
                )
                # Bild von uploads → library/_inbox verschieben (wie api_image_generate)
                inbox = IMAGES_LIBRARY_DIR / "_inbox"
                inbox.mkdir(parents=True, exist_ok=True)
                target = inbox / img.path.name
                if target.exists():
                    stem, suf = target.stem, target.suffix
                    i = 2
                    while (inbox / f"{stem}_{i}{suf}").exists():
                        i += 1
                    target = inbox / f"{stem}_{i}{suf}"
                img.path.rename(target)
                img.path = target
                # Sidecar-Meta damit Lightbox auch hier den Prompt zeigt
                try:
                    _lib.write_image_meta(target, {
                        "user_prompt": prompt,
                        "enhanced_prompt": None,
                        "final_prompt": prompt,
                        "provider": "gemini",
                        "model": _ANKI_IMG_MODEL,
                        "size": "1024x1024",
                        "source": "anki-from-text",
                        "created_at": datetime.now().isoformat(),
                    })
                except Exception:
                    pass
                # Absoluten Pfad in Card legen — _card_to_anki_note resolved abs paths
                card["image"] = str(target)
                print(f"[Anki/Image] generated for card: {prompt[:60]}")
            except ImageGenError as e:
                print(f"[Anki/Image] gen failed ({e.code}): {e.message[:120]}")
            except Exception as e:
                print(f"[Anki/Image] unexpected: {type(e).__name__}: {e}")

    await asyncio.gather(*(_one(c) for c in candidates), return_exceptions=True)


_ANKI_CONNECT_URL = "http://localhost:8765"


async def _anki_invoke(action: str, params: dict | None = None, timeout: float = 15.0) -> dict:
    """Zentraler Helper für AnkiConnect-Calls. Wirft RuntimeError bei Anki-seitigen
    Fehlern. Erwartet dass Anki läuft mit AnkiConnect-Addon auf localhost:8765.
    """
    import urllib.request
    payload = {"action": action, "version": 6, "params": params or {}}

    def _call() -> dict:
        req = urllib.request.Request(
            _ANKI_CONNECT_URL,
            data=json.dumps(payload).encode("utf-8"),
            headers={"Content-Type": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            return json.loads(resp.read().decode("utf-8"))

    try:
        result = await asyncio.to_thread(_call)
    except Exception as e:
        raise RuntimeError(f"anki_unreachable: {e}") from e
    if result.get("error"):
        raise RuntimeError(f"anki_error: {result['error']}")
    return result


async def api_anki_decks(request):
    """GET /api/anki/decks — Liste aller Deck-Namen aus Anki (via AnkiConnect).
    Rückgabe: { ok, decks: [str, ...] } oder { ok: false, error }.
    """
    try:
        result = await _anki_invoke("deckNames", timeout=5.0)
        decks = result.get("result") or []
        return web.json_response({"ok": True, "decks": sorted(decks)})
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e), "decks": []}, status=200)


def _card_to_anki_note(card: dict, deck: str) -> tuple[dict, dict[str, Path]]:
    """Wandelt eine Generator-Card in ein AnkiConnect `note`-Dict um.
    Gibt (note, images_used) zurück — images_used ist {basename: abs_path} für
    vorherigen storeMediaFile-Upload. `note` hat allowDuplicate:false mit
    duplicateScope:"deck" — Anki's eigener Dedup greift pro Deck.

    Logik analog zu api_learn_export_zip: Basic vs. Cloze-Entscheidung,
    Back-Extra mit Bild/FunFact/Merkhilfe/Footer, Occlusion-Mapping zu Cloze.
    """
    import re as _re_sync

    images_used: dict[str, Path] = {}

    def _collect(rel: str) -> str:
        rel = (rel or "").strip()
        if not rel:
            return ""
        # 1) Absoluter Pfad (z.B. von _enrich_anki_cards_with_images mit Library-Pfad)
        try:
            p_abs = Path(rel)
            if p_abs.is_absolute() and p_abs.is_file():
                images_used[p_abs.name] = p_abs
                return p_abs.name
        except Exception:
            pass
        # 2) URL-Prefix abstreifen (Library + Uploads)
        if rel.startswith("/library/"):
            sub = rel[len("/library/"):]
            abs_img = IMAGES_LIBRARY_DIR / sub
            if abs_img.is_file():
                images_used[abs_img.name] = abs_img
                return abs_img.name
        if rel.startswith("/api/uploads/"):
            rel = rel[len("/api/uploads/"):]
        abs_img = UPLOADS_DIR / rel
        if abs_img.exists():
            images_used[abs_img.name] = abs_img
            return abs_img.name
        return ""

    def _rewrite_imgs(html: str) -> str:
        def _replace(m):
            src = m.group(1)
            if src.startswith(("http://", "https://", "data:")):
                return m.group(0)
            basename = _collect(src)
            return f'<img src="{basename}">' if basename else m.group(0)
        return _re_sync.sub(r'<img\s+src="([^"]+)"', _replace, html)

    front = _rewrite_imgs(card.get("front") or "")
    back = _rewrite_imgs(card.get("back") or "")

    # Back-Extra (Bild + FunFact + Merkhilfe + Footer)
    parts: list[str] = []
    img_rel = (card.get("image") or "").strip()
    if img_rel:
        basename = _collect(img_rel)
        if basename:
            parts.append(f'<div style="margin:8px 0"><img src="{basename}" style="max-width:100%"></div>')
    if card.get("funfact"):
        parts.append(f'<div style="margin:6px 0;padding:6px 10px;background:#f59e0b15;border-left:3px solid #f59e0b;border-radius:4px"><b>💡 Fun Fact:</b> {card["funfact"]}</div>')
    if card.get("merkhilfe"):
        parts.append(f'<div style="margin:6px 0;padding:6px 10px;background:#a78bfa15;border-left:3px solid #a78bfa;border-radius:4px"><b>🧠 Merkhilfe:</b> {card["merkhilfe"]}</div>')
    footer_bits: list[str] = []
    if card.get("source"):
        footer_bits.append(str(card["source"]))
    if card.get("external_source"):
        footer_bits.append(f'Ergänzt aus {card["external_source"]}')
    if footer_bits:
        parts.append(f'<div style="margin-top:10px;padding-top:6px;border-top:1px solid #1e293b;font-size:10px;color:#64748b">{" · ".join(footer_bits)}</div>')
    back_extra = "".join(parts)

    # Tags: Generator liefert space- oder comma-separierten String
    raw_tags = card.get("tags") or ""
    if isinstance(raw_tags, list):
        tags = [str(t).strip() for t in raw_tags if str(t).strip()]
    else:
        tags = [t.strip() for t in str(raw_tags).replace(",", " ").split() if t.strip()]

    occs = card.get("occlusions") or []
    has_cloze_syntax = bool(_re_sync.search(r"\{\{c\d+::", front + " " + back))

    if occs and img_rel:
        # Occlusion → Cloze mit Bild + Labels als {{cN::…}}
        basename = _collect(img_rel)
        img_tag = f'<img src="{basename}" style="max-width:100%">' if basename else ""
        labels_cloze = "; ".join(
            f"{{{{c{i+1}::{(o.get('label') or '').strip()}}}}}"
            for i, o in enumerate(occs) if o.get("label")
        )
        text = f'{img_tag}<br><br><b>Strukturen im Bild:</b> {labels_cloze}'
        note = {
            "deckName": deck,
            "modelName": "Cloze",
            "fields": {"Text": text, "Back Extra": back_extra},
            "tags": tags,
        }
    elif has_cloze_syntax:
        combined = f"{front}<br><br>{back}" if front and back else (front or back)
        note = {
            "deckName": deck,
            "modelName": "Cloze",
            "fields": {"Text": combined, "Back Extra": back_extra},
            "tags": tags,
        }
    else:
        full_back = back + (("<br><br>" + back_extra) if back_extra else "")
        note = {
            "deckName": deck,
            "modelName": "Basic",
            "fields": {"Front": front, "Back": full_back},
            "tags": tags,
        }

    note["options"] = {
        "allowDuplicate": False,
        "duplicateScope": "deck",
        "duplicateScopeOptions": {"deckName": deck, "checkChildren": False},
    }
    return note, images_used


async def _bulk_sync_cards(cards: list[dict], *, deck: str, subject: str = "Karten") -> dict:
    """Reusable Bulk-Sync: deck anlegen, Bilder hochladen (storeMediaFile), notes
    batch-adden. Returnt das gleiche Report-Dict wie api_anki_bulk_sync (ohne ok-Key —
    Caller wrappt). Wirft RuntimeError bei AnkiConnect-Errors.
    """
    import base64 as _b64

    await _anki_invoke("createDeck", {"deck": deck}, timeout=5.0)

    notes: list[dict] = []
    all_images: dict[str, Path] = {}
    for c in cards:
        note, imgs = _card_to_anki_note(c, deck)
        notes.append(note)
        all_images.update(imgs)

    media_ok = 0
    media_errors: list[str] = []
    for basename, abs_path in all_images.items():
        try:
            data_b64 = _b64.b64encode(abs_path.read_bytes()).decode("ascii")
            await _anki_invoke("storeMediaFile", {"filename": basename, "data": data_b64}, timeout=30.0)
            media_ok += 1
        except Exception as e:
            media_errors.append(f"{basename}: {e}")
            print(f"[Anki/BulkSync] storeMediaFile failed {basename}: {e}")

    result = await _anki_invoke("addNotes", {"notes": notes}, timeout=60.0)
    ids = result.get("result") or []
    try:
        can_add = (await _anki_invoke("canAddNotes", {"notes": notes}, timeout=15.0)).get("result") or []
    except Exception:
        can_add = [None] * len(notes)

    added = duplicates = errors = 0
    card_report = []
    for i, note_id in enumerate(ids):
        if note_id is not None:
            added += 1
            card_report.append({"index": i, "ok": True, "note_id": note_id})
        else:
            is_dup = can_add[i] is False if i < len(can_add) else False
            if is_dup:
                duplicates += 1
                card_report.append({"index": i, "ok": False, "is_duplicate": True})
            else:
                errors += 1
                card_report.append({"index": i, "ok": False, "error": "addNote returned null"})

    print(f"[Anki/BulkSync] deck={deck}: {added} neu · {duplicates} Dup · {errors} Fehler · {media_ok} Bilder")
    return {
        "deck": deck,
        "added": added,
        "duplicates": duplicates,
        "errors": errors,
        "media": media_ok,
        "media_errors": media_errors,
        "cards": card_report,
    }


async def api_anki_bulk_sync(request):
    """POST /api/anki/bulk-sync — Pusht generierte Karten + Bilder direkt in Anki.
    Body: { subject, deck?, cards: [...] }

    Ablauf:
      1. Sicherstellen dass Deck existiert (createDeck — idempotent)
      2. Für jede Card: note-Dict bauen, Bilder einsammeln
      3. Alle Bilder per storeMediaFile hochladen (base64)
      4. Batch addNotes → AnkiConnect gibt pro Karte note_id oder null (Dup/Fehler)
      5. Report pro Karte zurückgeben

    Rückgabe: {
      ok, deck,
      added: int, duplicates: int, errors: int, media: int,
      cards: [{index, ok, note_id?, error?, is_duplicate?}]
    }
    """
    import base64 as _b64

    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    cards_in = body.get("cards") or []
    if not cards_in:
        return web.json_response({"ok": False, "error": "Keine Karten übergeben"}, status=400)

    subject = (body.get("subject") or "Karten").strip()
    deck = (body.get("deck") or subject or "Default").strip() or "Default"

    # 1. Deck anlegen (falls nicht vorhanden) — idempotent
    try:
        await _anki_invoke("createDeck", {"deck": deck}, timeout=5.0)
    except Exception as e:
        return web.json_response({"ok": False, "error": f"Deck anlegen fehlgeschlagen: {e}"}, status=502)

    # 2. Notes bauen + Bilder sammeln
    notes: list[dict] = []
    all_images: dict[str, Path] = {}
    for c in cards_in:
        note, imgs = _card_to_anki_note(c, deck)
        notes.append(note)
        all_images.update(imgs)

    # 3. Bilder hochladen (storeMediaFile akzeptiert base64)
    media_ok = 0
    media_errors: list[str] = []
    for basename, abs_path in all_images.items():
        try:
            data_b64 = _b64.b64encode(abs_path.read_bytes()).decode("ascii")
            await _anki_invoke("storeMediaFile", {"filename": basename, "data": data_b64}, timeout=30.0)
            media_ok += 1
        except Exception as e:
            media_errors.append(f"{basename}: {e}")
            print(f"[Anki/BulkSync] storeMediaFile failed {basename}: {e}")

    # 4. Batch addNotes — Return ist Liste mit note_id oder null (bei Dup/Fehler)
    try:
        result = await _anki_invoke("addNotes", {"notes": notes}, timeout=60.0)
    except Exception as e:
        return web.json_response({
            "ok": False, "error": f"addNotes fehlgeschlagen: {e}",
            "deck": deck, "media": media_ok,
        }, status=502)

    ids = result.get("result") or []
    # Bei canAddNotes-Check fragen wir zusätzlich ab um Dup vs. echter Fehler zu trennen
    try:
        can_add = (await _anki_invoke("canAddNotes", {"notes": notes}, timeout=15.0)).get("result") or []
    except Exception:
        can_add = [None] * len(notes)

    # 5. Report
    added = duplicates = errors = 0
    card_report = []
    for i, note_id in enumerate(ids):
        if note_id is not None:
            added += 1
            card_report.append({"index": i, "ok": True, "note_id": note_id})
        else:
            # null → entweder Dup oder Fehler; canAddNotes false = Dup
            is_dup = can_add[i] is False if i < len(can_add) else False
            if is_dup:
                duplicates += 1
                card_report.append({"index": i, "ok": False, "is_duplicate": True})
            else:
                errors += 1
                card_report.append({"index": i, "ok": False, "error": "addNote returned null"})

    print(f"[Anki/BulkSync] deck={deck}: {added} neu · {duplicates} Dup · {errors} Fehler · {media_ok} Bilder")

    return web.json_response({
        "ok": True,
        "deck": deck,
        "added": added,
        "duplicates": duplicates,
        "errors": errors,
        "media": media_ok,
        "media_errors": media_errors,
        "cards": card_report,
    })


async def api_diss_wiki_index(request):
    """GET wiki index.md content."""
    index_file = DISS_WIKI_DIR / "index.md"
    if not index_file.exists():
        return web.json_response({"content": "Wiki-Index nicht gefunden."})
    return web.json_response({"content": index_file.read_text(encoding="utf-8")})


async def api_diss_word(request):
    """GET status / POST upload content to OneDrive Dissertation.docx."""
    import subprocess
    script = SCRIPTS_DIR / "onedrive_diss.py"
    if not script.exists():
        return web.json_response({"ok": False, "error": "onedrive_diss.py nicht gefunden"}, status=500)

    if request.method == "GET":
        proc = await asyncio.create_subprocess_exec(
            sys.executable, str(script), "status",
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        out, err = await proc.communicate()
        try:
            return web.json_response(json.loads(out.decode("utf-8", errors="replace")))
        except Exception:
            return web.json_response({"exists": False, "error": err.decode("utf-8", errors="replace")[:200]})

    # POST — upload content
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
    content = (body.get("content") or "").strip()
    if not content:
        return web.json_response({"ok": False, "error": "content required"}, status=400)

    import tempfile, os as _os
    with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as tf:
        tf.write(content)
        tmppath = tf.name
    try:
        proc = await asyncio.create_subprocess_exec(
            sys.executable, str(script), "upload", f"--content={content}",
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        out, err = await proc.communicate()
    finally:
        _os.unlink(tmppath)
    try:
        return web.json_response(json.loads(out.decode("utf-8", errors="replace")))
    except Exception:
        return web.json_response({"ok": False, "error": err.decode("utf-8", errors="replace")[:400]})


async def api_diss_word_url(request):
    """GET/POST — persist the Word Online share/embed URL."""
    store = DATA_DIR / "diss_word_url.json"
    if request.method == "GET":
        if store.exists():
            try:
                return web.json_response(json.loads(store.read_text("utf-8")))
            except Exception:
                pass
        return web.json_response({"url": "", "shareUrl": ""})
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
    payload = {"url": body.get("url", ""), "shareUrl": body.get("shareUrl", "")}
    store.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    return web.json_response({"ok": True})


async def api_diss_prism_results(request):
    """GET — list saved Prism analyses  |  POST — save one analysis result."""
    store = DATA_DIR / "diss_prism_results.json"

    if request.method == "GET":
        if store.exists():
            try:
                return web.json_response(json.loads(store.read_text("utf-8")))
            except Exception:
                pass
        return web.json_response([])

    # POST — save result atomically
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    required = ("context", "mode", "output")
    if not all(body.get(k) for k in required):
        return web.json_response({"ok": False, "error": "context, mode, output required"}, status=400)

    results = []
    if store.exists():
        try:
            results = json.loads(store.read_text("utf-8"))
        except Exception:
            results = []

    entry = {
        "id": f"{int(__import__('time').time()*1000)}",
        "ts": __import__('datetime').datetime.utcnow().isoformat() + "Z",
        "context": body["context"],
        "mode": body["mode"],
        "output": body["output"],
    }
    results.insert(0, entry)
    results = results[:50]  # keep last 50

    tmp = store.with_suffix(".tmp")
    tmp.write_text(json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(store)

    return web.json_response({"ok": True, "id": entry["id"]})


async def api_diss_prism_analyse(request):
    """POST — analyse Prism CSV/table data with Claude and persist result."""
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    csv_text = (body.get("csv") or "").strip()
    context  = (body.get("context") or "Prism-Daten aus meiner Dissertation (Kardiochirurgie/ITI)").strip()
    if not csv_text:
        return web.json_response({"ok": False, "error": "csv required"}, status=400)

    prompt = (
        f"Du analysierst Prism-Exportdaten für eine medizinische Dissertation (Kardiochirurgie/ITI).\n\n"
        f"Kontext: {context}\n\n"
        f"Prism-Daten:\n{csv_text}\n\n"
        "Erstelle eine präzise, wissenschaftliche Zusammenfassung:\n"
        "1. Beschreibe die Daten (Gruppen, n, Zeitpunkte falls vorhanden)\n"
        "2. Zentrale statistische Ergebnisse (Mittelwert±SD/SEM, p-Werte)\n"
        "3. Klinische Interpretation (1-2 Sätze)\n"
        "4. Formulierungsvorschlag für den Diss-Text (Methods+Results, deutsch, akademisch)\n\n"
        "Antworte strukturiert mit Markdown-Überschriften."
    )

    raw = ""
    try:
        import anthropic as _ant
        _aclient = _ant.AsyncAnthropic()
        _model = CONFIG.get("claude_model_force") or CONFIG.get("claude_model") or "claude-sonnet-4-6"
        _resp = await _aclient.messages.create(
            model=_model, max_tokens=3000,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = _resp.content[0].text.strip()
    except Exception:
        try:
            proc = await asyncio.create_subprocess_exec(
                "claude", "-p",
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            stdout, _ = await asyncio.wait_for(proc.communicate(prompt.encode("utf-8")), timeout=180)
            raw = stdout.decode("utf-8", errors="replace").strip()
        except Exception as e:
            return web.json_response({"ok": False, "error": str(e)}, status=500)

    if not raw:
        return web.json_response({"ok": False, "error": "Leere Antwort von Claude"}, status=500)

    # Persist to prism_results store
    store = DATA_DIR / "diss_prism_results.json"
    results = []
    if store.exists():
        try:
            results = json.loads(store.read_text("utf-8"))
        except Exception:
            pass
    import time as _time
    entry = {
        "id":      str(int(_time.time() * 1000)),
        "ts":      __import__('datetime').datetime.utcnow().isoformat() + "Z",
        "context": context,
        "mode":    "prism-analyse",
        "output":  raw,
    }
    results.insert(0, entry)
    results = results[:50]
    tmp = store.with_suffix(".tmp")
    tmp.write_text(json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(store)

    return web.json_response({"ok": True, "output": raw, "id": entry["id"]})


# ── Diss Delta-Log ─────────────────────────────────────────────────────────────
DISS_DELTA_FILE = DATA_DIR / "diss_delta_log.json"

async def api_diss_delta(request):
    """GET — list delta entries  |  POST — append  |  DELETE ?id=... — remove."""
    if request.method == "GET":
        if DISS_DELTA_FILE.exists():
            try:
                return web.json_response(json.loads(DISS_DELTA_FILE.read_text("utf-8")))
            except Exception:
                pass
        return web.json_response([])

    if request.method == "DELETE":
        entry_id = request.rel_url.query.get("id", "")
        if not entry_id:
            return web.json_response({"ok": False, "error": "id required"}, status=400)
        entries = []
        if DISS_DELTA_FILE.exists():
            try:
                entries = json.loads(DISS_DELTA_FILE.read_text("utf-8"))
            except Exception:
                pass
        entries = [e for e in entries if e.get("id") != entry_id]
        tmp = DISS_DELTA_FILE.with_suffix(".tmp")
        tmp.write_text(json.dumps(entries, ensure_ascii=False, indent=2), encoding="utf-8")
        tmp.replace(DISS_DELTA_FILE)
        return web.json_response({"ok": True})

    # POST — append
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    text = (body.get("text") or "").strip()
    if not text:
        return web.json_response({"ok": False, "error": "text required"}, status=400)

    entries = []
    if DISS_DELTA_FILE.exists():
        try:
            entries = json.loads(DISS_DELTA_FILE.read_text("utf-8"))
        except Exception:
            pass

    import time as _time
    entry = {
        "id":       str(int(_time.time() * 1000)),
        "ts":       __import__('datetime').datetime.utcnow().isoformat() + "Z",
        "category": body.get("category", "finding"),
        "text":     text,
        "source":   (body.get("source") or "").strip(),
    }
    entries.insert(0, entry)
    entries = entries[:200]

    tmp = DISS_DELTA_FILE.with_suffix(".tmp")
    tmp.write_text(json.dumps(entries, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(DISS_DELTA_FILE)
    return web.json_response({"ok": True, "id": entry["id"]})


# ── Diss Structure ─────────────────────────────────────────────────────────────
DISS_STRUCTURE_FILE = DATA_DIR / "diss_structure.json"

async def api_diss_structure(request):
    """GET — load chapter structure  |  POST — save chapter structure."""
    if request.method == "GET":
        if DISS_STRUCTURE_FILE.exists():
            try:
                return web.json_response(json.loads(DISS_STRUCTURE_FILE.read_text("utf-8")))
            except Exception:
                pass
        return web.json_response({"chapters": []})
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
    chapters = body.get("chapters", [])
    tmp = DISS_STRUCTURE_FILE.with_suffix(".tmp")
    tmp.write_text(json.dumps({"chapters": chapters}, ensure_ascii=False, indent=2), encoding="utf-8")
    tmp.replace(DISS_STRUCTURE_FILE)
    return web.json_response({"ok": True})


async def api_medizin_lernplan(request):
    """POST Lernplan erzeugen + optional in iCloud schreiben.
    Body: { topic, exam_date, start_date?, time?, duration?, subtopics?, calendar?, dry_run? }
    Ruft scripts/lernplan-generator.py auf und gibt JSON zurueck.
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    topic = (body.get("topic") or "").strip()
    exam_date = (body.get("exam_date") or "").strip()
    if not topic or not exam_date:
        return web.json_response({"ok": False, "error": "topic + exam_date required"}, status=400)

    script = SCRIPTS_DIR / "lernplan-generator.py"
    if not script.exists():
        return web.json_response({"ok": False, "error": f"Script fehlt: {script}"}, status=500)

    cmd = [
        sys.executable, str(script),
        "--topic", topic,
        "--exam-date", exam_date,
        "--duration", str(int(body.get("duration") or 45)),
        "--time", (body.get("time") or "18:00"),
        "--calendar", (body.get("calendar") or "Kalender"),
        "--json",
    ]
    if body.get("start_date"):
        cmd += ["--start-date", body["start_date"]]
    subtopics = body.get("subtopics") or ""
    if isinstance(subtopics, list):
        subtopics = ",".join(subtopics)
    if subtopics:
        cmd += ["--subtopics", subtopics]
    if body.get("dry_run"):
        cmd.append("--dry-run")

    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=get_action_timeout("lernplan.generate"))
        out = stdout.decode("utf-8", errors="replace").strip()
        try:
            result = json.loads(out)
        except Exception:
            result = {
                "ok": False,
                "error": "Konnte Script-Output nicht parsen",
                "stdout": out[:500],
                "stderr": stderr.decode("utf-8", errors="replace")[:500],
            }
        return web.json_response(result)
    except asyncio.TimeoutError:
        return web.json_response({"ok": False, "error": f"Timeout nach {get_action_timeout('lernplan.generate')}s"}, status=504)
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)


async def api_trading_account(request):
    """Get Alpaca account info."""
    try:
        cred_file = Path("C:/Users/Robin/Jarvis/secrets/alpaca-cred.json")
        if not cred_file.exists():
            return web.json_response({"error": "No Alpaca credentials"})
        creds = json.loads(cred_file.read_text(encoding="utf-8"))
        import requests as req
        headers = {
            "APCA-API-KEY-ID": creds.get("key_id", creds.get("api_key", "")),
            "APCA-API-SECRET-KEY": creds.get("secret_key", creds.get("api_secret", "")),
        }
        base = creds.get("base_url", "https://paper-api.alpaca.markets")
        r = req.get(f"{base}/v2/account", headers=headers, timeout=10)
        if r.status_code != 200:
            return web.json_response({"error": f"HTTP {r.status_code}"})
        acc = r.json()
        return web.json_response({
            "equity": acc.get("equity"),
            "cash": acc.get("cash"),
            "buying_power": acc.get("buying_power"),
            "pnl": float(acc.get("equity", 0)) - float(acc.get("last_equity", 0)),
            "position_count": acc.get("position_count", 0),
        })
    except Exception as e:
        return web.json_response({"error": str(e)})


async def api_trading_positions(request):
    """Get Alpaca positions."""
    try:
        cred_file = Path("C:/Users/Robin/Jarvis/secrets/alpaca-cred.json")
        if not cred_file.exists():
            return web.json_response([])
        creds = json.loads(cred_file.read_text(encoding="utf-8"))
        import requests as req
        headers = {
            "APCA-API-KEY-ID": creds.get("key_id", creds.get("api_key", "")),
            "APCA-API-SECRET-KEY": creds.get("secret_key", creds.get("api_secret", "")),
        }
        base = creds.get("base_url", "https://paper-api.alpaca.markets")
        r = req.get(f"{base}/v2/positions", headers=headers, timeout=10)
        if r.status_code != 200:
            return web.json_response([])
        return web.json_response(r.json())
    except Exception:
        return web.json_response([])


async def api_trading_summary(request):
    """Get latest trading summary from night results."""
    night_dir = JARVIS_DATA / "night-results"
    if night_dir.exists():
        files = sorted(night_dir.glob("*trading*"), reverse=True)
        if files:
            try:
                data = json.loads(files[0].read_text(encoding="utf-8"))
                return web.json_response(data)
            except Exception:
                pass
    return web.json_response({})


def _alpaca_headers():
    """Read Alpaca credentials and return (headers_dict, base_url)."""
    cred_path = Path("C:/Users/Robin/Jarvis/secrets/alpaca-cred.json")
    creds = json.loads(cred_path.read_text(encoding="utf-8"))
    headers = {
        "APCA-API-KEY-ID": creds.get("key_id", creds.get("api_key", "")),
        "APCA-API-SECRET-KEY": creds.get("secret_key", creds.get("api_secret", "")),
    }
    base = creds.get("base_url", "https://paper-api.alpaca.markets")
    return headers, base


async def api_trading_orders(request):
    """GET /api/trading/orders  -  Recent orders from Alpaca."""
    try:
        headers, base = _alpaca_headers()
        async with aiohttp.ClientSession(headers=headers) as session:
            async with session.get(f"{base}/v2/orders?status=all&limit=20&direction=desc") as resp:
                data = await resp.json()
                return web.json_response(data)
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_trading_history(request):
    """GET /api/trading/history  -  Portfolio equity history from Alpaca."""
    period = request.query.get("period", "1D")
    period_map = {
        "1D": ("1D", "5Min"),
        "1W": ("1W", "1H"),
        "1M": ("1M", "1D"),
        "3M": ("3M", "1D"),
    }
    p, tf = period_map.get(period, ("1D", "15Min"))

    try:
        headers, base = _alpaca_headers()
        async with aiohttp.ClientSession(headers=headers) as session:
            url = f"{base}/v2/account/portfolio/history?period={p}&timeframe={tf}"
            async with session.get(url) as resp:
                data = await resp.json()
                return web.json_response(data)
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_trading_order(request):
    """POST /api/trading/order  -  Place a market order."""
    body = await request.json()
    symbol = body.get("symbol", "").upper().strip()
    qty = int(body.get("qty", 0))
    side = body.get("side", "buy").lower()

    if not symbol or qty < 1 or side not in ("buy", "sell"):
        return web.json_response({"error": "Invalid order params"}, status=400)

    try:
        headers, base = _alpaca_headers()
        headers["Content-Type"] = "application/json"
        order_data = {
            "symbol": symbol,
            "qty": str(qty),
            "side": side,
            "type": "market",
            "time_in_force": "day",
        }
        async with aiohttp.ClientSession(headers=headers) as session:
            async with session.post(f"{base}/v2/orders", json=order_data) as resp:
                data = await resp.json()
                if resp.status >= 400:
                    return web.json_response({"error": data.get("message", str(data))}, status=resp.status)
                return web.json_response(data)
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


# ---------------------------------------------------------------------------
# Trading Bots  -  Start / Stop / Restart Grundgeruest
# ---------------------------------------------------------------------------
import trading_bot_manager as _tbm  # noqa: E402


async def api_trading_bots_list(request):
    """GET /api/trading/bots  -  Liste aller registrierten Bots (mit Live-Status)."""
    try:
        return web.json_response({"bots": _tbm.list_bots()})
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_trading_bots_create(request):
    """POST /api/trading/bots  -  Neuen Bot registrieren.

    Body: {name, strategy?, params?, interval_sec?}
    """
    try:
        body = await request.json() if request.content_length else {}
    except Exception:
        body = {}
    name = (body.get("name") or "").strip()
    if not name:
        return web.json_response({"error": "name required"}, status=400)
    try:
        bot = _tbm.create_bot(
            name=name,
            strategy=body.get("strategy") or "demo",
            params=body.get("params") or {},
            interval_sec=int(body.get("interval_sec") or 60),
        )
        return web.json_response(bot, status=201)
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_trading_bots_get(request):
    bot_id = request.match_info["id"]
    bot = _tbm.get_bot(bot_id)
    if not bot:
        raise web.HTTPNotFound(reason="Bot not found")
    return web.json_response(bot)


async def api_trading_bots_update(request):
    bot_id = request.match_info["id"]
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"error": "invalid JSON"}, status=400)
    updated = _tbm.update_bot(bot_id, body)
    if not updated:
        raise web.HTTPNotFound(reason="Bot not found")
    return web.json_response(updated)


async def api_trading_bots_delete(request):
    bot_id = request.match_info["id"]
    ok = _tbm.delete_bot(bot_id)
    if not ok:
        raise web.HTTPNotFound(reason="Bot not found")
    return web.json_response({"ok": True})


async def api_trading_bots_start(request):
    bot_id = request.match_info["id"]
    try:
        result = _tbm.start_bot(bot_id)
        return web.json_response(result)
    except KeyError:
        raise web.HTTPNotFound(reason="Bot not found")
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_trading_bots_stop(request):
    bot_id = request.match_info["id"]
    try:
        result = _tbm.stop_bot(bot_id)
        return web.json_response(result)
    except KeyError:
        raise web.HTTPNotFound(reason="Bot not found")
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_trading_bots_restart(request):
    bot_id = request.match_info["id"]
    try:
        result = _tbm.restart_bot(bot_id)
        return web.json_response(result)
    except KeyError:
        raise web.HTTPNotFound(reason="Bot not found")
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_trading_bots_rank(request):
    """GET /api/trading/bots/rank  -  Bots sortiert nach Sortino Ratio."""
    try:
        return web.json_response({"bots": _tbm.rank_bots_by_sortino()})
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


async def api_trading_bots_kick_mutate(request):
    """POST /api/trading/bots/kick-mutate  -  Kickt Verlierer, spawnt 2 Mutanten.

    Body: {loser_id, winner_id}
    """
    try:
        body = await request.json()
        loser_id = body["loser_id"]
        winner_id = body["winner_id"]
        result = _tbm.kick_and_mutate(loser_id, winner_id)
        return web.json_response(result)
    except (KeyError, TypeError) as e:
        return web.json_response({"error": f"bad request: {e}"}, status=400)
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


# ---------------------------------------------------------------------------
# World Monitor  -  Intelligence Dashboard (kostenlose APIs, kein Key nötig)
# ---------------------------------------------------------------------------

import csv
import io
from collections import Counter

# Cache: (data_dict, timestamp_float) oder None
_world_monitor_cache: tuple | None = None
_WORLD_MONITOR_TTL = 900  # 15 Minuten


async def _fetch_json(session, url: str, timeout: int = 10) -> dict | list | None:
    """GET-Request mit Timeout, gibt None bei Fehler zurück."""
    try:
        async with session.get(url, timeout=aiohttp.ClientTimeout(total=timeout),
                               headers={"User-Agent": "Jarvis-MissionControl/4.0"}) as resp:
            if resp.status != 200:
                # (encoding-fixed) print(f"[WorldMonitor] HTTP {resp.status} €" {url}")
                return None
            return await resp.json(content_type=None)
    except Exception as e:
        print(f"[WorldMonitor] Error {url}: {e}")
        return None


async def _fetch_text(session, url: str, timeout: int = 10) -> str | None:
    """GET-Request, gibt Plaintext zurück."""
    try:
        async with session.get(url, timeout=aiohttp.ClientTimeout(total=timeout),
                               headers={"User-Agent": "Jarvis-MissionControl/4.0"}) as resp:
            if resp.status != 200:
                # (encoding-fixed) print(f"[WorldMonitor] HTTP {resp.status} €" {url}")
                return None
            return await resp.text()
    except Exception as e:
        print(f"[WorldMonitor] Error {url}: {e}")
        return None


async def _fetch_news(session) -> list:
    """Reddit r/worldnews  -  Top 10 Weltnachrichten (kein API Key noetig)."""
    url = "https://www.reddit.com/r/worldnews/.json?limit=10"
    data = await _fetch_json(session, url)
    if not data:
        return []
    posts = data.get("data", {}).get("children", [])
    return [
        {
            "title": p["data"].get("title", ""),
            "url": p["data"].get("url", ""),
            "source": p["data"].get("domain", "reddit"),
            "date": datetime.utcfromtimestamp(p["data"].get("created_utc", 0)).isoformat() + "Z",
            "score": p["data"].get("score", 0),
        }
        for p in posts[:10] if p.get("data")
    ]


async def _fetch_earthquakes(session) -> dict:
    """USGS  -  Erdbeben weltweit M4.5+ (letzter Tag, immer verfuegbar)."""
    url = "https://earthquake.usgs.gov/earthquakes/feed/v1.0/summary/4.5_day.geojson"
    data = await _fetch_json(session, url)
    if not data:
        return {"count": 0, "quakes": []}
    features = data.get("features", [])
    quakes = []
    for f in features[:10]:
        props = f.get("properties", {})
        coords = f.get("geometry", {}).get("coordinates", [0, 0, 0])
        quakes.append({
            "place": props.get("place", ""),
            "mag": props.get("mag", 0),
            "time": datetime.utcfromtimestamp(props.get("time", 0) / 1000).isoformat() + "Z",
            "lat": coords[1] if len(coords) > 1 else 0,
            "lon": coords[0] if len(coords) > 0 else 0,
        })
    return {"count": len(features), "quakes": quakes}


async def _fetch_crypto(session) -> list:
    """CoinGecko  -  BTC, ETH, SOL Preise (kein API Key noetig)."""
    url = "https://api.coingecko.com/api/v3/simple/price?ids=bitcoin,ethereum,solana&vs_currencies=usd&include_24hr_change=true"
    data = await _fetch_json(session, url)
    if not data:
        return []
    result = []
    names = {"bitcoin": "BTC", "ethereum": "ETH", "solana": "SOL"}
    for coin_id, label in names.items():
        if coin_id in data:
            result.append({
                "symbol": label,
                "price": data[coin_id].get("usd", 0),
                "change_24h": round(data[coin_id].get("usd_24h_change", 0), 2),
            })
    return result


async def _fetch_weather_alerts(session) -> list:
    """NOAA Weather.gov  -  Aktive US Wetterwarnungen."""
    url = "https://api.weather.gov/alerts/active?status=actual&limit=10"
    data = await _fetch_json(session, url)
    if not data:
        return []
    features = data.get("features") or []
    alerts = []
    for f in features[:10]:
        props = f.get("properties", {})
        alerts.append({
            "headline": props.get("headline", ""),
            "severity": props.get("severity", ""),
            "event": props.get("event", ""),
            "area": props.get("areaDesc", ""),
        })
    return alerts


async def _fetch_flights(session) -> dict:
    """OpenSky Network  -  Flüge über DACH Region (lamin=45..55, lomin=5..15)."""
    url = "https://opensky-network.org/api/states/all?lamin=45&lamax=55&lomin=5&lomax=15"
    data = await _fetch_json(session, url, timeout=15)
    if not data:
        return {"count": 0, "region": "DACH", "error": "no_data"}
    states = data.get("states") or []
    return {"count": len(states), "region": "DACH"}


async def _fetch_conflicts(session) -> list:
    """GDELT GeoJSON  -  Konfliktereignisse der letzten 24h."""
    url = "https://api.gdeltproject.org/api/v2/geo/geo?query=conflict+OR+war+OR+military&mode=PointData&format=GeoJSON&timespan=24h&maxpoints=30"
    data = await _fetch_json(session, url, timeout=15)
    if not data:
        return []
    features = data.get("features", [])
    return [
        {
            "lat": f["geometry"]["coordinates"][1],
            "lon": f["geometry"]["coordinates"][0],
            "name": f.get("properties", {}).get("name", ""),
            "url": f.get("properties", {}).get("url", ""),
            "type": f.get("properties", {}).get("eventtype", ""),
        }
        for f in features[:30] if f.get("geometry", {}).get("coordinates")
    ]


async def _fetch_market_indices(session) -> list:
    """Yahoo Finance  -  VIX, Gold, S&P 500."""
    indices = []
    for symbol, name in [("^VIX", "VIX"), ("GC=F", "Gold"), ("^GSPC", "S&P 500"), ("^GDAXI", "DAX")]:
        url = f"https://query1.finance.yahoo.com/v8/finance/chart/{symbol}?range=1d&interval=1d"
        data = await _fetch_json(session, url, timeout=8)
        if data:
            try:
                meta = data["chart"]["result"][0]["meta"]
                price = meta.get("regularMarketPrice", 0)
                prev = meta.get("chartPreviousClose", 0)
                change = ((price - prev) / prev * 100) if prev else 0
                indices.append({"symbol": name, "price": round(price, 2), "change": round(change, 2)})
            except (KeyError, IndexError):
                pass
    return indices


async def api_world_monitor(request):
    """GET /api/world-monitor  -  Aggregiertes Intelligence Dashboard."""
    import time

    global _world_monitor_cache

    now = time.time()
    if _world_monitor_cache is not None:
        cached_data, cached_ts = _world_monitor_cache
        if now - cached_ts < _WORLD_MONITOR_TTL:
            print(f"[WorldMonitor] Cache hit ({int(now - cached_ts)}s old)")
            return web.json_response(cached_data)

    print("[WorldMonitor] Fetching fresh data from all sources…")

    async with aiohttp.ClientSession() as session:
        results = await asyncio.gather(
            _fetch_news(session),
            _fetch_earthquakes(session),
            _fetch_crypto(session),
            _fetch_weather_alerts(session),
            _fetch_flights(session),
            _fetch_conflicts(session),
            _fetch_market_indices(session),
            return_exceptions=True,
        )

    def safe(val, fallback):
        return fallback if isinstance(val, Exception) else val

    data = {
        "timestamp": datetime.utcnow().isoformat() + "Z",
        "news":           safe(results[0], []),
        "earthquakes":    safe(results[1], {"count": 0, "quakes": []}),
        "crypto":         safe(results[2], []),
        "weather_alerts": safe(results[3], []),
        "flights":        safe(results[4], {"count": 0, "region": "DACH", "error": "fetch_failed"}),
        "conflicts":      safe(results[5], []),
        "market_indices": safe(results[6], []),
    }

    _world_monitor_cache = (data, now)
    print(f"[WorldMonitor] Done - news:{len(data['news'])} quakes:{data['earthquakes'].get('count',0)} "
          f"crypto:{len(data['crypto'])} alerts:{len(data['weather_alerts'])} "
          f"flights:{data['flights'].get('count',0)} conflicts:{len(data['conflicts'])} "
          f"markets:{len(data['market_indices'])}")

    return web.json_response(data)


# ---------------------------------------------------------------------------
# Skills endpoint
# ---------------------------------------------------------------------------

SKILLS_DIR = Path("C:/Users/Robin/Jarvis/skills")

async def api_skills(request):
    """GET /api/skills  -  Alle installierten Skills mit Metadaten."""
    import re as _re
    skills = []
    if SKILLS_DIR.exists():
        for skill_dir in sorted(SKILLS_DIR.iterdir()):
            if not skill_dir.is_dir():
                continue
            slug = skill_dir.name
            name = slug
            description = ""
            version = None
            tags = []
            triggers = []
            capabilities = []
            # Try SKILL.md frontmatter first
            for fname in ["SKILL.md", "README.md"]:
                fpath = skill_dir / fname
                if fpath.exists():
                    try:
                        content = fpath.read_text(encoding="utf-8", errors="ignore")
                        body_text = content
                        if content.startswith("---"):
                            end = content.find("---", 3)
                            if end != -1:
                                fm = content[3:end].strip()
                                body_text = content[end+3:].strip()
                                in_triggers = False
                                for line in fm.splitlines():
                                    stripped = line.strip()
                                    if line.startswith("name:"):
                                        name = line[5:].strip().strip('"')
                                        in_triggers = False
                                    elif line.startswith("description:"):
                                        description = line[12:].strip().strip('"')
                                        in_triggers = False
                                    elif line.startswith("version:"):
                                        version = line[8:].strip().strip('"')
                                        in_triggers = False
                                    elif line.startswith("tags:"):
                                        raw = line[5:].strip()
                                        tags = [t.strip().strip('"') for t in raw.strip("[]").split(",") if t.strip()]
                                        in_triggers = False
                                    elif line.startswith("triggers:") or line.startswith("trigger:"):
                                        raw = line.split(":", 1)[1].strip()
                                        if raw.startswith("["):
                                            triggers = [t.strip().strip('"\'') for t in raw.strip("[]").split(",") if t.strip()]
                                        elif raw:
                                            triggers = [raw.strip('"\'')]
                                        in_triggers = True
                                    elif in_triggers and stripped.startswith("- "):
                                        triggers.append(stripped[2:].strip().strip('"\''))
                                    elif stripped and not stripped.startswith("-"):
                                        in_triggers = False
                        # If no frontmatter description, take first non-empty line after frontmatter
                        if not description:
                            for bline in body_text.splitlines():
                                stripped = bline.strip().lstrip("#").strip()
                                if stripped and not stripped.startswith("---"):
                                    description = stripped[:200]
                                    break
                        # Extract capabilities: first 3 bullet points from body
                        for m in _re.finditer(r'^[-*]\s+(.+)', body_text, _re.MULTILINE):
                            cap = m.group(1).strip()
                            if len(cap) > 8 and not cap.lower().startswith("trigger"):
                                capabilities.append(cap[:120])
                            if len(capabilities) >= 3:
                                break
                    except Exception:
                        pass
                    break
            # Try _meta.json for version
            meta_path = skill_dir / "_meta.json"
            if meta_path.exists():
                try:
                    import json as _json
                    meta = _json.loads(meta_path.read_text(encoding="utf-8"))
                    if not version:
                        version = meta.get("version")
                except Exception:
                    pass
            cc_wrapper = Path.home() / ".claude" / "skills" / f"{slug}.md"
            skills.append({
                "slug": slug,
                "name": name,
                "description": description,
                "version": version,
                "tags": tags,
                "triggers": triggers[:3],
                "capabilities": capabilities[:3],
                "hasSkillMd": (skill_dir / "SKILL.md").exists(),
                "hasReadme": (skill_dir / "README.md").exists(),
                "hasWrapper": cc_wrapper.exists(),
            })
    return web.json_response({"skills": skills, "total": len(skills)})


# ---------------------------------------------------------------------------
# Skill detail + usage stats
# ---------------------------------------------------------------------------

_TRANSCRIPTS_DIR = Path.home() / ".claude" / "projects"
_usage_cache = {}      # slug -> count
_usage_cache_ts = 0.0  # timestamp of last scan

def _scan_skill_usage():
    """Scan session transcripts for skill usage. Returns {slug: session_count}.
    Fast approach: full-text search per session file, skip very short slugs."""
    global _usage_cache, _usage_cache_ts
    import time
    now = time.time()
    if _usage_cache and (now - _usage_cache_ts) < 600:  # cache 10 min
        return _usage_cache
    counts = {}
    if not SKILLS_DIR.exists():
        return counts
    # Only count slugs >= 5 chars and skip generic English words that cause false positives
    _NOISY_SLUGS = {"adapt", "audit", "shape", "polish", "connect", "layout", "clarify",
                    "animate", "optimize", "distill", "delight", "critique", "harden",
                    "bolder", "quieter", "kaizen", "overdrive", "marketing", "evaluation",
                    "playwright", "remotion", "typeset", "colorize", "imagen"}
    slugs = [d.name for d in SKILLS_DIR.iterdir()
             if d.is_dir() and len(d.name) >= 5 and d.name not in _NOISY_SLUGS]
    # Only scan Jarvis project transcripts (not all projects)
    jarvis_dir = _TRANSCRIPTS_DIR / "C--Users-Robin-Jarvis"
    if not jarvis_dir.exists():
        _usage_cache = counts
        _usage_cache_ts = now
        return counts
    for jsonl in jarvis_dir.glob("*.jsonl"):
        try:
            text = jsonl.read_text(encoding="utf-8", errors="ignore")
            # Quick check: only count if slug appears in user content
            # Filter out the skill_listing block which lists all skill names
            for slug in slugs:
                # Check for the slug in contexts that indicate actual usage:
                # - user typed it or mentioned it
                # - skill was invoked
                # Simple heuristic: count if slug appears AND the file has user content
                if slug in text:
                    # Verify it's not ONLY in system prompts by checking multiple occurrences
                    # or presence in non-system content
                    idx = text.find(slug)
                    # Check a window around the match — if it's in "skill_listing" skip
                    window_start = max(0, idx - 200)
                    window = text[window_start:idx + len(slug) + 50]
                    if 'skill_listing' not in window and 'system-reminder' not in window:
                        counts[slug] = counts.get(slug, 0) + 1
        except Exception:
            continue
    _usage_cache = counts
    _usage_cache_ts = now
    return counts


def _skill_file_tree(skill_dir, prefix=""):
    """Return list of files in skill dir as [{name, size, type}]."""
    entries = []
    try:
        for item in sorted(skill_dir.iterdir()):
            rel = item.name
            if item.is_dir():
                children = _skill_file_tree(item, prefix + rel + "/")
                entries.append({"name": rel, "type": "dir", "children": children})
            else:
                try:
                    size = item.stat().st_size
                except Exception:
                    size = 0
                entries.append({"name": rel, "type": "file", "size": size})
    except Exception:
        pass
    return entries


def _parse_skill_md(content):
    """Parse SKILL.md frontmatter + body into structured data."""
    meta = {}
    body = content
    if content.startswith("---"):
        end = content.find("---", 3)
        if end != -1:
            fm = content[3:end].strip()
            body = content[end+3:].strip()
            current_key = None
            current_list = []
            for line in fm.splitlines():
                stripped = line.strip()
                if not stripped:
                    continue
                if stripped.startswith("- ") and current_key:
                    current_list.append(stripped[2:].strip().strip('"\''))
                    meta[current_key] = current_list
                    continue
                if ":" in stripped:
                    if current_key and isinstance(meta.get(current_key), list):
                        pass  # already saved
                    key, val = stripped.split(":", 1)
                    key = key.strip()
                    val = val.strip().strip('"\'')
                    current_key = key
                    current_list = []
                    if val.startswith("["):
                        # inline list
                        items = [t.strip().strip('"\'') for t in val.strip("[]").split(",") if t.strip()]
                        meta[key] = items
                        current_list = items
                    elif val:
                        meta[key] = val
                    else:
                        meta[key] = []
                        current_list = meta[key]
    return meta, body


async def api_skill_detail(request):
    """GET /api/skills/{slug} - Full detail for one skill."""
    slug = request.match_info["slug"]
    skill_dir = SKILLS_DIR / slug
    if not skill_dir.exists() or not skill_dir.is_dir():
        return web.json_response({"error": "Skill nicht gefunden"}, status=404)

    # Basic metadata
    name = slug
    description = ""
    version = None
    tags = []
    skill_md_content = ""
    readme_content = ""
    meta = {}
    body = ""

    # Read SKILL.md
    skill_md_path = skill_dir / "SKILL.md"
    if skill_md_path.exists():
        try:
            skill_md_content = skill_md_path.read_text(encoding="utf-8", errors="ignore")
            meta, body = _parse_skill_md(skill_md_content)
            name = meta.get("name", name)
            description = meta.get("description", "")
            version = meta.get("version")
            if isinstance(meta.get("tags"), list):
                tags = meta["tags"]
        except Exception:
            pass

    # Read README.md as fallback
    readme_path = skill_dir / "README.md"
    if readme_path.exists():
        try:
            readme_content = readme_path.read_text(encoding="utf-8", errors="ignore")[:5000]
            if not name or name == slug:
                rm_meta, rm_body = _parse_skill_md(readme_content)
                name = rm_meta.get("name", name)
                if not description:
                    description = rm_meta.get("description", "")
        except Exception:
            pass

    # _meta.json
    meta_json = {}
    meta_path = skill_dir / "_meta.json"
    if meta_path.exists():
        try:
            import json as _json
            meta_json = _json.loads(meta_path.read_text(encoding="utf-8"))
            if not version:
                version = meta_json.get("version")
        except Exception:
            pass

    # File tree
    file_tree = _skill_file_tree(skill_dir)
    total_files = sum(1 for _ in skill_dir.rglob("*") if _.is_file())
    total_size = sum(f.stat().st_size for f in skill_dir.rglob("*") if f.is_file())

    # Check for Claude Code wrapper
    cc_wrapper = Path.home() / ".claude" / "skills" / f"{slug}.md"
    has_wrapper = cc_wrapper.exists()
    wrapper_content = ""
    if has_wrapper:
        try:
            wrapper_content = cc_wrapper.read_text(encoding="utf-8", errors="ignore")[:2000]
        except Exception:
            pass

    # Triggers
    triggers = meta.get("triggers", meta.get("trigger", []))
    if isinstance(triggers, str):
        triggers = [triggers]

    # Usage stats
    usage = _scan_skill_usage()
    usage_count = usage.get(slug, 0)

    # Extract examples from body
    examples = []
    if body:
        import re
        for m in re.finditer(r'(?:example|beispiel)[^\n]*\n(.*?)(?=\n#|\n---|\Z)', body, re.IGNORECASE | re.DOTALL):
            ex = m.group(1).strip()[:500]
            if ex:
                examples.append(ex)

    return web.json_response({
        "slug": slug,
        "name": name,
        "description": description,
        "version": version,
        "tags": tags,
        "triggers": triggers,
        "hasSkillMd": skill_md_path.exists(),
        "hasReadme": readme_path.exists(),
        "hasWrapper": has_wrapper,
        "skillMdContent": skill_md_content[:8000],
        "readmeContent": readme_content[:5000],
        "wrapperContent": wrapper_content,
        "fileTree": file_tree,
        "totalFiles": total_files,
        "totalSize": total_size,
        "usageCount": usage_count,
        "meta": {k: v for k, v in meta.items() if k not in ("name", "description", "version", "tags")},
        "metaJson": meta_json,
        "examples": examples[:3],
        "body": body[:4000],
    })


async def api_skills_usage(request):
    """GET /api/skills/usage - Usage stats for all skills."""
    usage = _scan_skill_usage()
    return web.json_response({"usage": usage, "total_sessions_scanned": sum(1 for p in _TRANSCRIPTS_DIR.rglob("*.jsonl"))})


async def api_skills_install(request):
    """POST /api/skills  -  Neuen Skill anlegen. Akzeptiert:
    - {"content": "---\\nname: ...\\n..."} â†' parsed Slug aus name: Frontmatter
    - {"content": "https://github.com/..."} â†' fetched SKILL.md von GitHub
    """
    import re as _re
    import aiohttp as _aiohttp
    try:
        data = await request.json()
        raw = data.get("content", "").strip()
        if not raw:
            return web.json_response({"error": "Inhalt erforderlich"}, status=400)

        content = raw
        # If it looks like a URL, fetch SKILL.md from GitHub
        if raw.startswith("http://") or raw.startswith("https://"):
            url = raw.rstrip("/")
            # Try common SKILL.md locations
            candidates = []
            if "github.com" in url:
                parts = url.replace("https://github.com/", "").replace("http://github.com/", "").split("/")
                if len(parts) >= 2:
                    user, repo = parts[0], parts[1]
                    sub_path = "/".join(parts[4:]) if len(parts) > 4 else ""
                    for branch in ["main", "master"]:
                        if sub_path:
                            candidates.append(f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/{sub_path}/SKILL.md")
                        candidates.append(f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/SKILL.md")
            if not candidates and "github.com" not in url:
                return web.json_response({"error": "URL nicht erkannt  -  nur GitHub URLs oder direkter SKILL.md Inhalt"}, status=400)

            # Helper to fetch a single SKILL.md and install it
            installed = []
            async with _aiohttp.ClientSession() as sess:
                async def _fetch(u):
                    try:
                        async with sess.get(u, timeout=_aiohttp.ClientTimeout(total=10)) as r:
                            return (await r.text()) if r.status == 200 else None
                    except Exception:
                        return None

                # 1) Try direct candidates
                content = None
                for cand_url in candidates:
                    content = await _fetch(cand_url)
                    if content:
                        break

                # 2) If not found, scan repo for skill directories
                if not content and "github.com" in url:
                    parts = url.replace("https://github.com/", "").replace("http://github.com/", "").split("/")
                    user, repo = parts[0], parts[1]
                    # Search common skill folder locations
                    skill_dirs_found = []
                    for search_path in ["skills", ".claude/skills", ""]:
                        api_url = f"https://api.github.com/repos/{user}/{repo}/contents/{search_path}".rstrip("/")
                        try:
                            async with sess.get(api_url, timeout=_aiohttp.ClientTimeout(total=10)) as resp:
                                if resp.status == 200:
                                    dirs = await resp.json()
                                    for d in dirs:
                                        if d.get("type") == "dir":
                                            name = d["name"]
                                            if search_path == "" and (name.startswith(".") or name in {"docs", "tests", "test", "node_modules", "__pycache__", ".github", "assets", "images", "img"}):
                                                continue
                                            prefix = f"{search_path}/" if search_path else ""
                                            skill_dirs_found.append(f"{prefix}{name}")
                                    if skill_dirs_found:
                                        break
                        except Exception:
                            continue

                    # Install ALL found skills
                    if skill_dirs_found:
                        for skill_path in skill_dirs_found:
                            skill_content = None
                            for branch in ["main", "master"]:
                                skill_url = f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/{skill_path}/SKILL.md"
                                skill_content = await _fetch(skill_url)
                                if skill_content:
                                    break
                            if not skill_content:
                                continue
                            # Derive slug
                            s_slug = skill_path.split("/")[-1].lower().replace(" ", "-")
                            s_slug = _re.sub(r"[^a-z0-9\-_]", "", s_slug)
                            if not s_slug:
                                continue
                            s_dir = SKILLS_DIR / s_slug
                            if s_dir.exists():
                                installed.append({"slug": s_slug, "skipped": True})
                                continue
                            s_dir.mkdir(parents=True, exist_ok=True)
                            (s_dir / "SKILL.md").write_text(skill_content, encoding="utf-8")
                            installed.append({"slug": s_slug, "ok": True})

                        if installed:
                            new_count = sum(1 for i in installed if i.get("ok"))
                            skip_count = sum(1 for i in installed if i.get("skipped"))
                            return web.json_response({
                                "ok": True,
                                "slug": f"{new_count} neu, {skip_count} übersprungen",
                                "installed": installed,
                                "batch": True
                            })

            if not content and not installed:
                # Fallback: try README.md + package.json
                if "github.com" in url:
                    parts = url.replace("https://github.com/", "").replace("http://github.com/", "").split("/")
                    if len(parts) >= 2:
                        user, repo = parts[0], parts[1]
                        async with _aiohttp.ClientSession() as sess:
                            async def _qf(u):
                                try:
                                    async with sess.get(u, timeout=_aiohttp.ClientTimeout(total=10)) as r:
                                        return (await r.text()) if r.status == 200 else None
                                except Exception:
                                    return None
                            readme_text = None
                            pkg_data = {}
                            for branch in ["main", "master"]:
                                if not readme_text:
                                    readme_text = await _qf(f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/README.md")
                                pkg_raw = await _qf(f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/package.json")
                                if pkg_raw:
                                    try:
                                        pkg_data = json.loads(pkg_raw)
                                    except Exception:
                                        pass
                                if readme_text:
                                    break
                            if readme_text or pkg_data:
                                gen_name = pkg_data.get("name", repo).replace("@", "").replace("/", "-")
                                gen_desc = pkg_data.get("description", "")
                                gen_version = pkg_data.get("version", "")
                                if not gen_desc and readme_text:
                                    for line in readme_text.splitlines():
                                        stripped = line.strip()
                                        if stripped and not stripped.startswith("#") and not stripped.startswith("!") and not stripped.startswith("<") and len(stripped) > 20:
                                            gen_desc = stripped[:300]
                                            break
                                content = f"---\nname: {gen_name}\ndescription: {gen_desc}\nversion: {gen_version}\n---\n\n"
                                if readme_text:
                                    content += readme_text[:3000]
                if not content and not installed:
                    return web.json_response({"error": "Kein Skill gefunden — weder SKILL.md, README.md noch package.json im Repo."}, status=404)

        # Extract slug from frontmatter name: field
        slug = None
        if content.startswith("---"):
            end = content.find("---", 3)
            if end != -1:
                fm = content[3:end]
                for line in fm.splitlines():
                    if line.strip().startswith("name:"):
                        name_val = line.split(":", 1)[1].strip().strip('"').strip("'")
                        slug = name_val.lower().replace(" ", "-")
                        break

        # Fallback: use first 40 chars of content as slug
        if not slug:
            slug = _re.sub(r"[^a-z0-9]+", "-", raw[:40].lower()).strip("-") or "unnamed-skill"

        slug = _re.sub(r"[^a-z0-9\-_]", "", slug)
        if not slug:
            return web.json_response({"error": "Konnte keinen Slug ableiten"}, status=400)

        skill_dir = SKILLS_DIR / slug
        if skill_dir.exists():
            return web.json_response({"error": f"Skill '{slug}' existiert bereits"}, status=409)
        skill_dir.mkdir(parents=True, exist_ok=False)
        (skill_dir / "SKILL.md").write_text(content, encoding="utf-8")
        return web.json_response({"ok": True, "slug": slug})
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


def _skill_recommendation(skill_names, is_batch=False, repo_name="", name="", description="", verdict=""):
    """Generate a Jarvis recommendation based on skill content + Robin's profile."""
    text = " ".join(skill_names).lower() + " " + name.lower() + " " + description.lower()

    # Robin's interests: medical, research, trading, anki, code, data, writing
    relevant = []
    irrelevant = []
    INTEREST_KEYWORDS = {
        "medizin": ["medical", "health", "clinical", "pubmed", "osce", "anki", "flashcard"],
        "research": ["research", "paper", "academic", "science", "zotero", "deep-research"],
        "trading": ["trading", "stock", "market", "finance", "alpaca", "crypto"],
        "code/dev": ["code", "debug", "git", "dev", "python", "refactor", "test", "api", "mcp"],
        "data": ["data", "csv", "excel", "analytics", "extract"],
        "writing": ["write", "doc", "essay", "content", "article"],
        "design": ["design", "ui", "ux", "css", "canvas", "image"],
        "audio": ["voice", "tts", "speech", "transcrib", "mikro", "audio", "whisper", "elevenlabs"],
    }
    for category, keywords in INTEREST_KEYWORDS.items():
        matches = [k for k in keywords if k in text]
        if matches:
            relevant.append(category)

    # Detect clearly irrelevant skills
    IRRELEVANT_KEYWORDS = ["invoice", "lead-research-assistant", "internal-comms", "competitive-ads",
                           "developer-growth", "langsmith", "composio", "connect-apps"]
    irrelevant_hits = [k for k in IRRELEVANT_KEYWORDS if k in text]

    if is_batch:
        useful = [n for n in skill_names if any(k in n.lower() for cat_kws in INTEREST_KEYWORDS.values() for k in cat_kws)]
        noise = len(skill_names) - len(useful)
        if useful:
            rec_text = f"Davon für dich relevant: {', '.join(useful[:5])}{'…' if len(useful) > 5 else ''}."
            if noise > len(useful):
                rec_text += f" Aber {noise} von {len(skill_names)} sind eher Business/Marketing-Skills die du nicht brauchst."
                rec_action = "selektiv"
                rec_reason = "Einzelne Skills rauspicken statt alles installieren — die meisten sind für SaaS/Business, nicht für Medizin oder Dev."
            else:
                rec_action = "ja"
                rec_reason = "Gutes Verhältnis von nützlichen zu irrelevanten Skills."
        else:
            rec_text = "Keiner der Skills passt direkt zu deinem Setup (Medizin, Research, Trading, Dev)."
            rec_action = "nein"
            rec_reason = "Eher Business/Marketing-Fokus, nicht dein Bereich."
    else:
        if verdict == "schwach":
            rec_action = "nein"
            rec_text = "Qualität zu niedrig — keine Beschreibung oder Triggers."
            rec_reason = "Ein Skill ohne klare Beschreibung und Trigger ist Ballast."
        elif relevant:
            rec_action = "ja"
            rec_text = f"Passt zu: {', '.join(relevant)}."
            rec_reason = f"Direkt relevant für dein {relevant[0]}-Setup."
        elif irrelevant_hits:
            rec_action = "nein"
            rec_text = "Kein Match mit deinem Profil."
            rec_reason = "Eher für Business/SaaS, nicht für Medizinstudium oder Dev."
        else:
            rec_action = "vielleicht"
            rec_text = "Nicht klar einzuordnen."
            rec_reason = "Kein direkter Match, aber auch nicht offensichtlich irrelevant."

    return {
        "action": rec_action,  # ja, nein, vielleicht, selektiv
        "text": rec_text,
        "reason": rec_reason,
        "relevant_categories": relevant,
    }


async def api_skills_preview(request):
    """POST /api/skills/preview  -  Fetch + parse skill, return review without installing."""
    import re as _re
    import aiohttp as _aiohttp
    try:
        data = await request.json()
        raw = data.get("content", "").strip()
        if not raw:
            return web.json_response({"error": "Inhalt erforderlich"}, status=400)

        content = raw
        source_url = None

        # Fetch from URL if needed
        if raw.startswith("http://") or raw.startswith("https://"):
            source_url = raw.rstrip("/")
            candidates = []
            if "github.com" in source_url:
                parts = source_url.replace("https://github.com/", "").replace("http://github.com/", "").split("/")
                if len(parts) >= 2:
                    user, repo = parts[0], parts[1]
                    sub_path = "/".join(parts[4:]) if len(parts) > 4 else ""
                    for branch in ["main", "master"]:
                        if sub_path:
                            candidates.append(f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/{sub_path}/SKILL.md")
                        candidates.append(f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/SKILL.md")
            if not candidates:
                return web.json_response({"error": "Nur GitHub-URLs oder direkter SKILL.md Inhalt"}, status=400)

            content = None
            async with _aiohttp.ClientSession() as sess:
                for cand_url in candidates:
                    try:
                        async with sess.get(cand_url, timeout=_aiohttp.ClientTimeout(total=10)) as r:
                            if r.status == 200:
                                content = await r.text()
                                break
                    except Exception:
                        continue

            if not content:
                if "github.com" in source_url:
                    parts = source_url.replace("https://github.com/", "").replace("http://github.com/", "").split("/")
                    user, repo = parts[0], parts[1]

                    # Strategy 1: Scan for skill dirs (batch repo)
                    skill_names = []
                    async with _aiohttp.ClientSession() as sess:
                        async def _qfetch(u):
                            try:
                                async with sess.get(u, timeout=_aiohttp.ClientTimeout(total=10)) as r:
                                    return (await r.text()) if r.status == 200 else None
                            except Exception:
                                return None

                        for search_path in ["skills", ".claude/skills", ""]:
                            api_url = f"https://api.github.com/repos/{user}/{repo}/contents/{search_path}".rstrip("/")
                            try:
                                async with sess.get(api_url, timeout=_aiohttp.ClientTimeout(total=8)) as resp:
                                    if resp.status == 200:
                                        dirs = await resp.json()
                                        dir_names = [d["name"] for d in dirs if d.get("type") == "dir"]
                                        if search_path == "":
                                            skip = {".github", "docs", "tests", "test", "node_modules", "__pycache__", ".vscode", "assets", "images", "img"}
                                            dir_names = [n for n in dir_names if n not in skip and not n.startswith(".")]
                                            if len(dir_names) >= 3:
                                                skill_names = dir_names
                                                break
                                        else:
                                            skill_names = dir_names
                                            break
                            except Exception:
                                continue

                        if skill_names:
                            already = [d.name for d in SKILLS_DIR.iterdir() if d.is_dir()] if SKILLS_DIR.exists() else []
                            new_skills = [s for s in skill_names if s not in already]
                            rec = _skill_recommendation(skill_names, is_batch=True, repo_name=f"{user}/{repo}")
                            return web.json_response({
                                "ok": True, "batch": True,
                                "skill_count": len(skill_names), "new_count": len(new_skills),
                                "skill_names": skill_names[:20],
                                "name": f"Skill-Paket: {user}/{repo}",
                                "description": f"{len(skill_names)} Skills gefunden, davon {len(new_skills)} neu.",
                                "verdict": "batch", "content": raw, "recommendation": rec,
                            })

                        # Strategy 2: Fallback — README.md + package.json → auto-generate SKILL.md
                        readme_text = None
                        pkg_data = {}
                        for branch in ["main", "master"]:
                            if not readme_text:
                                readme_text = await _qfetch(f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/README.md")
                            pkg_raw = await _qfetch(f"https://raw.githubusercontent.com/{user}/{repo}/{branch}/package.json")
                            if pkg_raw:
                                try:
                                    pkg_data = json.loads(pkg_raw)
                                except Exception:
                                    pass
                            if readme_text:
                                break

                        if readme_text or pkg_data:
                            # Build a synthetic SKILL.md from what we found
                            gen_name = pkg_data.get("name", repo).replace("@", "").replace("/", "-")
                            gen_desc = pkg_data.get("description", "")
                            gen_version = pkg_data.get("version", "")
                            if not gen_desc and readme_text:
                                # Grab first non-heading, non-empty line from README
                                for line in readme_text.splitlines():
                                    stripped = line.strip()
                                    if stripped and not stripped.startswith("#") and not stripped.startswith("!") and not stripped.startswith("<") and len(stripped) > 20:
                                        gen_desc = stripped[:300]
                                        break

                            # Extract capabilities from README headers/lists
                            caps = []
                            if readme_text:
                                import re as _re2
                                for m in _re2.finditer(r'^(?:[-*]|\d+\.)\s+[`*]*([^`*\n]{10,80})', readme_text, _re2.MULTILINE):
                                    cap = m.group(1).strip().rstrip('.')
                                    if cap and not cap.startswith("http") and not cap.startswith("npm") and not cap.startswith("git"):
                                        caps.append(cap)
                                    if len(caps) >= 8:
                                        break

                            content = f"---\nname: {gen_name}\ndescription: {gen_desc}\nversion: {gen_version}\n---\n\n"
                            if readme_text:
                                content += readme_text[:3000]

                            # Check if already installed
                            gen_slug = _re.sub(r"[^a-z0-9\-_]", "", gen_name.lower().replace(" ", "-"))
                            already_installed = (SKILLS_DIR / gen_slug).exists() if gen_slug else False

                            rec = _skill_recommendation([], name=gen_name, description=gen_desc)
                            return web.json_response({
                                "ok": True, "batch": False,
                                "name": gen_name, "description": gen_desc, "version": gen_version,
                                "source": "readme", "source_url": source_url,
                                "capabilities": caps,
                                "already_installed": already_installed,
                                "verdict": "ok",
                                "score": max(2, min(5, len(caps))), "max_score": 7,
                                "pros": [p for p in [
                                    f"README vorhanden ({len(readme_text)} Zeichen)" if readme_text else None,
                                    f"package.json mit Version {gen_version}" if gen_version else None,
                                    f"{len(caps)} Features erkannt" if caps else None,
                                ] if p],
                                "cons": ["Kein SKILL.md — wurde aus README/package.json generiert"],
                                "triggers": [],
                                "content": content,
                                "recommendation": rec,
                            })

                return web.json_response({"error": "Kein Skill gefunden — weder SKILL.md, README.md noch package.json im Repo."}, status=404)

        # Parse frontmatter
        name, description, triggers, when_to_use, version = "", "", [], "", ""
        if content.startswith("---"):
            end = content.find("---", 3)
            if end != -1:
                fm = content[3:end]
                body = content[end+3:].strip()
                for line in fm.splitlines():
                    line = line.strip()
                    if line.startswith("name:"):
                        name = line.split(":", 1)[1].strip().strip('"\'')
                    elif line.startswith("description:"):
                        description = line.split(":", 1)[1].strip().strip('"\'')
                    elif line.startswith("version:"):
                        version = line.split(":", 1)[1].strip()
                    elif line.startswith("triggers:") or line.startswith("trigger:"):
                        val = line.split(":", 1)[1].strip()
                        if val:
                            triggers.append(val.strip('"\''))
                    elif line.startswith("-") and triggers is not None:
                        triggers.append(line.lstrip("- ").strip('"\''))
        else:
            body = content

        if not name:
            name = _re.sub(r"[^a-z0-9]+", "-", raw[:40].lower()).strip("-") or "Unbekannt"

        # Derive slug
        slug = _re.sub(r"[^a-z0-9\-_]", "", name.lower().replace(" ", "-"))

        # Heuristic quality score
        score = 0
        pros, cons = [], []
        if name and len(name) > 3:
            score += 1
        if description and len(description) > 20:
            score += 2
            pros.append("Klare Beschreibung vorhanden")
        else:
            cons.append("Beschreibung fehlt oder sehr kurz")
        if triggers:
            score += 2
            pros.append(f"{len(triggers)} Trigger-Bedingung(en) definiert")
        else:
            cons.append("Keine Trigger definiert — manueller Aufruf nötig")
        if len(content) > 300:
            score += 1
            pros.append("Detaillierte Anweisungen")
        if "SKILL.md" in content or "skill" in content.lower():
            score += 1
        if slug in [d.name for d in SKILLS_DIR.iterdir() if d.is_dir()] if SKILLS_DIR.exists() else []:
            return web.json_response({
                "ok": True,
                "already_installed": True,
                "slug": slug,
                "name": name,
                "description": description,
            })

        # Verdict
        if score >= 5:
            verdict = "empfohlen"
        elif score >= 3:
            verdict = "ok"
        else:
            verdict = "schwach"

        rec = _skill_recommendation([slug], is_batch=False, name=name, description=description, verdict=verdict)
        return web.json_response({
            "ok": True,
            "batch": False,
            "slug": slug,
            "name": name,
            "description": description,
            "version": version,
            "triggers": triggers[:5],
            "score": score,
            "max_score": 7,
            "pros": pros,
            "cons": cons,
            "verdict": verdict,
            "content_len": len(content),
            "already_installed": False,
            "recommendation": rec,
        })
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)


# ---------------------------------------------------------------------------
# ---------------------------------------------------------------------------
# File Upload endpoint
# ---------------------------------------------------------------------------

MAX_UPLOAD_BYTES = 500 * 1024 * 1024  # 500 MB — hochgesetzt für große GoodNotes/Skript-PDFs

async def api_upload(request):
    """POST /api/upload  -  multipart/form-data, stores file in data/uploads/."""
    try:
        reader = await request.multipart()
        field = await reader.next()
        if field is None or field.name != "file":
            return web.json_response({"error": "No file field"}, status=400)

        original_name = field.filename or "upload"
        # Sanitize filename
        safe_name = "".join(c for c in original_name if c.isalnum() or c in "._- ").strip()
        if not safe_name:
            safe_name = "upload"
        # Unique prefix to avoid collisions
        unique_name = f"{uuid.uuid4().hex[:8]}_{safe_name}"
        dest = UPLOADS_DIR / unique_name

        size = 0
        with dest.open("wb") as f:
            while True:
                chunk = await field.read_chunk(65536)
                if not chunk:
                    break
                size += len(chunk)
                if size > MAX_UPLOAD_BYTES:
                    f.close()
                    dest.unlink(missing_ok=True)
                    return web.json_response({"error": "File too large (max 500 MB)"}, status=413)
                f.write(chunk)

        print(f"[Upload] Saved {unique_name} ({size} bytes)")
        return web.json_response({"filename": unique_name, "size": size})
    except Exception as e:
        print(f"[Upload] Error: {e}")
        return web.json_response({"error": str(e)}, status=500)


# ---------------------------------------------------------------------------
# Knowledge Funnel  -  ingest any content into Obsidian Brain
# ---------------------------------------------------------------------------

KNOWLEDGE_DIR = Path("E:/OneDrive/AI/Obsidian-Vault/Jarvis-Brain/Jarvis-Knowledge")
FUNNEL_LOG    = DATA_DIR / "funnel-log.json"


def _funnel_load_log() -> list:
    """Load funnel log from disk."""
    if FUNNEL_LOG.exists():
        try:
            return json.loads(FUNNEL_LOG.read_text(encoding="utf-8"))
        except Exception:
            pass
    return []


def _funnel_save_log(items: list):
    """Save funnel log to disk."""
    FUNNEL_LOG.write_text(json.dumps(items, ensure_ascii=False, indent=2), encoding="utf-8")


async def api_funnel_log(request):
    """GET /api/funnel/log  -  Return log of all ingested items."""
    items = _funnel_load_log()
    return web.json_response({"items": items})


async def api_funnel_ingest(request):
    """POST /api/funnel/ingest  -  Ingest text/files into Obsidian Knowledge Base.

    Body: { type, tags, text, files: [{filename, original, size}] }
    Creates a smart-task that runs Claude to process the content.
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"error": "Invalid JSON"}, status=400)

    content_type  = body.get("type", "general")
    tags          = body.get("tags", [])
    text          = (body.get("text") or "").strip()
    files         = body.get("files", [])

    if not text and not files:
        return web.json_response({"error": "Kein Inhalt"}, status=400)

    type_labels = {
        "notes":   "Eigene Notizen",
        "paper":   "Paper/Abstract",
        "data":    "Daten/Ergebnisse",
        "quote":   "Zitat/Beobachtung",
        "general": "Allgemein",
        "file":    "Datei",
    }
    type_label = type_labels.get(content_type, content_type)
    tags_str   = ", ".join(tags) if tags else "keine"

    # Build prompt for Claude
    parts = [
        f"Du bist Jarvis. Ingestiere den folgenden Inhalt in die Wissensbasis.",
        f"",
        f"**Typ:** {type_label}",
        f"**Tags:** {tags_str}",
        f"**Ziel-Ordner:** {KNOWLEDGE_DIR}",
        f"",
        f"**Aufgabe:**",
        f"1. Entscheide in welchem Unterordner der Inhalt am besten passt (Forschung, Medizin, Technik, Kontext, Skills, Meta — oder neuer Ordner)",
        f"2. Erstelle eine strukturierte Markdown-Note mit passendem Dateinamen",
        f"3. Füge YAML-Frontmatter ein mit: tags, type, date, source (falls erkennbar)",
        f"4. Setze Obsidian-Wikilinks zu verwandten existierenden Notes",
        f"5. Aktualisiere `{KNOWLEDGE_DIR}/INDEX.md` (neue Zeile in der passenden Kategorie-Tabelle)",
        f"6. Gib am Ende aus: NOTE_CREATED: <relativer Pfad>",
    ]

    if text:
        parts += ["", "**INHALT:**", text]

    if files:
        parts += ["", "**DATEIEN:**"]
        for f in files:
            orig = f.get("original", f.get("filename", ""))
            fpath = UPLOADS_DIR / f.get("filename", "")
            parts.append(f"- {orig} (gespeichert als: {fpath})")

    prompt = "\n".join(parts)

    # Generate title preview
    preview = text[:80] if text else (files[0].get("original", "") if files else "")
    title = (text.split("\n")[0][:60] if text else (files[0].get("original", "Upload") if files else "Ingest"))

    # Log entry
    log_id  = uuid.uuid4().hex[:12]
    log_item = {
        "id":      log_id,
        "title":   title,
        "preview": preview,
        "type":    content_type,
        "tags":    tags,
        "date":    datetime.now().strftime("%d.%m.%Y %H:%M"),
        "files":   [f.get("original", "") for f in files],
        "status":  "processing",
    }
    items = _funnel_load_log()
    items.insert(0, log_item)
    _funnel_save_log(items)

    # Create and start smart-task (use shared smart-tasks storage)
    task_id  = str(uuid.uuid4())
    task_obj = {
        "id":               task_id,
        "title":            f"Brain-Ingest: {title[:50]}",
        "description":      f"Knowledge-Funnel Ingest ({type_label})",
        "prompt":           prompt,
        "original_request": prompt,
        "chat_context":     "",
        "priority":         "medium",
        "status":           "pending",
        "created":          datetime.now().isoformat(),
        "started":          None,
        "completed":        None,
        "output":           "",
        "error":            "",
        "funnel_log_id":    log_id,
    }
    tasks = _load_smart_tasks()
    tasks.append(task_obj)
    _save_smart_tasks(tasks)

    # Start task async (funnel-specific runner)
    asyncio.create_task(_run_funnel_task(task_id, log_id))

    print(f"[Funnel] Ingest started: {title[:40]} ({content_type}) -> task {task_id}")
    return web.json_response({"ok": True, "task_id": task_id, "log_id": log_id})


def _learn_extract_images(file_path: "Path") -> list:
    """Extract images from PPTX or PDF into UPLOADS_DIR/images/<stem>/. Returns relative paths."""
    import json as _jcap
    images = []
    captions = {}  # rel_path → slide text caption
    suffix = file_path.suffix.lower()
    img_dir = UPLOADS_DIR / "images" / file_path.stem
    try:
        if suffix == ".pptx":
            from pptx import Presentation
            img_dir.mkdir(parents=True, exist_ok=True)
            prs = Presentation(str(file_path))
            idx_ref = [0]

            def _slide_text(slide):
                texts = []
                for sh in slide.shapes:
                    try:
                        if sh.has_text_frame:
                            t = sh.text_frame.text.strip()
                            if t:
                                texts.append(t)
                    except Exception:
                        pass
                return " | ".join(texts[:4])[:200]  # max 200 chars, first 4 text boxes

            def _extract_shapes(shapes, slide_id, caption):
                for shape in shapes:
                    if shape.shape_type == 6:  # GROUP — recurse
                        try:
                            _extract_shapes(shape.shapes, slide_id, caption)
                        except Exception:
                            pass
                        continue
                    img_blob = None
                    ext = "png"
                    if shape.shape_type == 13:  # PICTURE
                        try:
                            img_blob = shape.image.blob
                            ext = shape.image.ext or "png"
                        except Exception:
                            pass
                    else:
                        # XML-level fallback for EMF/WMF embedded in non-PICTURE shapes
                        try:
                            from pptx.oxml.ns import qn
                            blip = shape.element.find(".//" + qn("a:blip"))
                            if blip is not None:
                                rId = blip.get(qn("r:embed"))
                                if rId:
                                    part = shape.part.related_parts.get(rId)
                                    if part and hasattr(part, "blob"):
                                        img_blob = part.blob
                                        ct = getattr(part, "content_type", "")
                                        ext = "wmf" if "wmf" in ct else ("emf" if "emf" in ct else "png")
                        except Exception:
                            pass
                    if img_blob:
                        # Convert EMF/WMF to PNG via wand so browsers can render
                        if ext in ("emf", "wmf"):
                            try:
                                from wand.image import Image as WandImage
                                with WandImage(blob=img_blob) as wimg:
                                    wimg.format = "png"
                                    img_blob = wimg.make_blob()
                                    ext = "png"
                            except ImportError:
                                print(f"[Learn] WARN: wand nicht installiert — EMF/WMF-Bild übersprungen (pip install wand). Datei: {file_path.name}")
                            except Exception as _wand_err:
                                print(f"[Learn] WARN: EMF/WMF-Konvertierung fehlgeschlagen ({_wand_err}) — Bild übersprungen. Datei: {file_path.name}")
                        img_name = f"s{slide_id}_{idx_ref[0]}.{ext}"
                        out = img_dir / img_name
                        out.write_bytes(img_blob)
                        rel = f"images/{file_path.stem}/{img_name}"
                        images.append(rel)
                        if caption:
                            captions[rel] = caption
                        idx_ref[0] += 1

            for slide in prs.slides:
                cap = _slide_text(slide)
                _extract_shapes(slide.shapes, slide.slide_id, cap)
        elif suffix == ".pdf":
            import fitz  # pymupdf
            img_dir.mkdir(parents=True, exist_ok=True)
            doc = fitz.open(str(file_path))
            idx = 0
            for pg_num, page in enumerate(doc):
                page_text = page.get_text("text").strip().replace("\n", " ")[:200]
                for img_info in page.get_images():
                    xref = img_info[0]
                    base = doc.extract_image(xref)
                    ext = base.get("ext", "png")
                    img_name = f"p{pg_num}_{idx}.{ext}"
                    out = img_dir / img_name
                    out.write_bytes(base["image"])
                    rel = f"images/{file_path.stem}/{img_name}"
                    images.append(rel)
                    if page_text:
                        captions[rel] = page_text
                    idx += 1
            doc.close()
    except ImportError as e:
        print(f"[Learn] Image extraction skipped (missing lib): {e}")
    except Exception as e:
        print(f"[Learn] Image extraction error for {file_path.name}: {e}")
    # Save captions alongside images for use during card generation
    # Always write captions.json — use filename as fallback if no text was extracted
    if images:
        try:
            cap_path = img_dir / "captions.json"
            img_dir.mkdir(parents=True, exist_ok=True)
            existing = {}
            if cap_path.exists():
                try:
                    existing = _jcap.loads(cap_path.read_text("utf-8"))
                except Exception:
                    pass
            existing.update(captions)
            # Fallback: any image without a caption gets its filename as caption
            for rel in images:
                if rel not in existing:
                    stem = Path(rel).stem.replace("_", " ").replace("-", " ")
                    existing[rel] = stem
            cap_path.write_text(_jcap.dumps(existing, ensure_ascii=False, indent=2), "utf-8")
        except Exception as e:
            print(f"[Learn] Caption save error: {e}")
    return images


async def api_learn_ingest(request):
    """POST /api/learn/ingest  -  multipart upload for Lernraum.
    Fields: files[] (multipart), subject (str), type (str).
    Saves to UPLOADS_DIR, triggers Funnel Ingest + RAG reindex.
    """
    try:
        reader = await request.multipart()
        saved_files = []
        subject = ""
        learn_type = "vorlesung"
        provider = "folie"

        async for field in reader:
            if field.name == "subject":
                subject = (await field.read()).decode("utf-8", errors="replace").strip()
            elif field.name == "type":
                learn_type = (await field.read()).decode("utf-8", errors="replace").strip()
            elif field.name == "provider":
                provider = (await field.read()).decode("utf-8", errors="replace").strip() or "folie"
            elif field.name in ("file", "files[]", "files"):
                original_name = field.filename or "upload"
                safe_name = "".join(c for c in original_name if c.isalnum() or c in "._- ").strip() or "upload"
                unique_name = f"{uuid.uuid4().hex[:8]}_{safe_name}"
                dest = UPLOADS_DIR / unique_name
                size = 0
                with dest.open("wb") as f:
                    while True:
                        chunk = await field.read_chunk(65536)
                        if not chunk:
                            break
                        size += len(chunk)
                        if size > MAX_UPLOAD_BYTES:
                            dest.unlink(missing_ok=True)
                            return web.json_response({"error": f"Datei zu groß (max 500 MB): {original_name}"}, status=413)
                        f.write(chunk)
                # Extract images from PPTX/PDF in thread pool (CPU-bound)
                loop = asyncio.get_event_loop()
                extracted = await loop.run_in_executor(None, _learn_extract_images, dest)
                saved_files.append({"filename": unique_name, "original": original_name, "size": size, "images": extracted})
                print(f"[Learn] Saved {unique_name} ({size} bytes), {len(extracted)} images extracted")
                # Write subject→stem mapping for reliable image lookup during generate
                # Always write meta — even if no images were extracted (EMF/WMF slides)
                meta_path = UPLOADS_DIR / "images" / "learn_meta.json"
                meta_path.parent.mkdir(parents=True, exist_ok=True)
                import json as _jmeta
                meta = {}
                if meta_path.exists():
                    try:
                        meta = _jmeta.loads(meta_path.read_text("utf-8"))
                    except Exception:
                        meta = {}
                subj_key = (subject or "medizin").lower()
                if subj_key not in meta:
                    meta[subj_key] = []
                if dest.stem not in meta[subj_key]:
                    meta[subj_key].append(dest.stem)
                # Cleanup: remove stale stems (image dir no longer exists)
                meta[subj_key] = [s for s in meta[subj_key] if (meta_path.parent / s).is_dir() or s == dest.stem]
                # Remove empty subject keys
                meta = {k: v for k, v in meta.items() if v}
                meta_path.write_text(_jmeta.dumps(meta, ensure_ascii=False, indent=2), "utf-8")
                print(f"[Learn] Meta updated: {subj_key} → {meta[subj_key]}")

        if not saved_files:
            return web.json_response({"error": "Keine Dateien erhalten"}, status=400)

        subj_label = subject or "Medizin"
        tags = ["lernraum", learn_type, provider]
        if subject:
            tags.append(subject.lower().replace(" ", "-"))

        # Build funnel-style request body and call internal logic directly
        fake_body = {
            "type": learn_type,
            "provider": provider,
            "tags": tags,
            "text": f"Lernmaterial für: {subj_label}\nTyp: {learn_type}\nQuelle: {provider}",
            "files": saved_files,
        }
        from aiohttp import web as _web
        import json as _json

        class _FakeRequest:
            async def json(self_inner):
                return fake_body
            content_length = 1

        resp = await api_funnel_ingest(_FakeRequest())
        result = _json.loads(resp.body)

        # Trigger RAG reindex in background (don't await — fire and forget)
        async def _reindex():
            try:
                from rag.ingest import build_index
                import asyncio as _aio
                loop = _aio.get_event_loop()
                await loop.run_in_executor(None, lambda: build_index(force=False))
                print("[Learn] RAG reindex done")
            except Exception as e:
                print(f"[Learn] RAG reindex error: {e}")
        asyncio.create_task(_reindex())

        total_images = sum(len(f.get("images", [])) for f in saved_files)
        return web.json_response({"ok": True, "files": len(saved_files), "images_extracted": total_images, **result})

    except Exception as e:
        print(f"[Learn] Ingest error: {e}")
        return web.json_response({"error": str(e)}, status=500)


async def api_learn_generate(request):
    """POST /api/learn/generate
    Streaming-Wrapper: schreibt alle paar Sekunden ein Whitespace-Byte als
    Heartbeat, damit Cloudflare's 100s-Proxy-Timeout nicht greift. JSON.parse
    ignoriert Leerzeichen am Anfang, daher braucht das Frontend keine Anpassung.
    """
    resp = web.StreamResponse(
        status=200,
        headers={
            "Content-Type": "application/json; charset=utf-8",
            "Cache-Control": "no-store",
            "X-Accel-Buffering": "no",
        },
    )
    await resp.prepare(request)

    async def _runner():
        try:
            return await _learn_generate_impl(request)
        except Exception as _e:
            import traceback as _tb
            _tb.print_exc()
            return {"ok": False, "error": f"{type(_e).__name__}: {_e}"}

    work_task = asyncio.create_task(_runner())
    HEARTBEAT_S = 8  # CF Free timeout = 100s, weit darunter bleiben
    try:
        while True:
            try:
                body = await asyncio.wait_for(asyncio.shield(work_task), timeout=HEARTBEAT_S)
                break
            except asyncio.TimeoutError:
                try:
                    await resp.write(b" ")
                except (ConnectionResetError, BrokenPipeError):
                    work_task.cancel()
                    return resp
    except asyncio.CancelledError:
        work_task.cancel()
        raise

    import json as _jstream
    await resp.write(_jstream.dumps(body, ensure_ascii=False).encode("utf-8"))
    await resp.write_eof()
    return resp


async def _learn_generate_impl(request):
    """Implementiert die eigentliche Karten-Generierung. Gibt ein Dict zurück
    (kein web.Response), damit der Streaming-Wrapper das letzte Stück als JSON
    schreiben kann.
    """
    try:
        body = await request.json()
    except Exception:
        return {"ok": False, "error": "Invalid JSON"}

    subject = (body.get("subject") or "").strip()
    topics = body.get("topics") or []
    if isinstance(topics, str):
        topics = [t.strip() for t in topics.split(",") if t.strip()]
    try:
        _raw_chunks = body.get("n_chunks", 12)
        n_chunks = int(_raw_chunks) if isinstance(_raw_chunks, (int, str, float)) else 12
    except (TypeError, ValueError):
        n_chunks = 12
    # n_cards: optionaler harter Wert. Wenn nicht gesetzt UND density gesetzt,
    # entscheidet Claude selbst basierend auf Stoffumfang.
    _raw_cards = body.get("n_cards")
    if _raw_cards in (None, "", "auto"):
        n_cards = None
    else:
        try:
            n_cards = max(3, min(60, int(_raw_cards)))
        except (TypeError, ValueError):
            n_cards = None
    density = (body.get("density") or "").strip().lower()
    if density not in ("low", "medium", "high"):
        density = ""
    # Wenn weder n_cards noch density: Default = medium (sinnvoller als fix 12)
    if n_cards is None and not density:
        density = "medium"
    feedback_new = (body.get("feedback") or "").strip()
    screenshots_b64 = body.get("screenshots") or []  # list of dataURL strings

    user_rules = body.get("rules") or []
    if isinstance(user_rules, str):
        user_rules = [r.strip() for r in user_rules.splitlines() if r.strip()]
    else:
        user_rules = [str(r).strip() for r in user_rules if str(r).strip()]

    # source_mode: bestimmt woraus die Karten gebaut werden.
    # - "lernraum" (default): RAG-Suche im Lernraum + Amboss/ViaMedici als Ergaenzung
    # - "amboss":  rein aus Amboss-Artikeln + Bildern (kein RAG, kein ViaMedici)
    # - "thieme":  rein aus ViaMedici/Thieme (kein RAG, kein Amboss)
    # - "mix":     Amboss + ViaMedici, kein RAG (wenn Robin keine eigenen Folien hat)
    source_mode = (body.get("source_mode") or "lernraum").strip().lower()
    if source_mode not in ("lernraum", "amboss", "thieme", "mix"):
        source_mode = "lernraum"
    use_rag = source_mode == "lernraum"
    use_amboss = source_mode in ("lernraum", "amboss", "mix")
    use_viamedici = source_mode in ("lernraum", "thieme", "mix")

    if not subject and not topics:
        return {"ok": False, "error": "subject oder topics erforderlich"}

    # Persistiertes Feedback pro Fach laden und mit neuem Feedback kombinieren
    import json as _jfb
    _fb_path = UPLOADS_DIR / "learn_feedback.json"
    try:
        _fb_store = _jfb.loads(_fb_path.read_text(encoding="utf-8")) if _fb_path.exists() else {}
    except Exception:
        _fb_store = {}
    _subj_key = subject.lower().strip()
    _persisted_fb = _fb_store.get(_subj_key, "")
    # Neues Feedback sofort persistieren (auch wenn es keine neuen Karten gibt)
    if feedback_new:
        existing = _persisted_fb.splitlines()
        if feedback_new not in existing:
            _persisted_fb = "\n".join([l for l in existing if l] + [feedback_new])
        _fb_store[_subj_key] = _persisted_fb
        try:
            _fb_path.write_text(_jfb.dumps(_fb_store, ensure_ascii=False, indent=2), encoding="utf-8")
            print(f"[Learn/Feedback] Gespeichert für '{subject}': {feedback_new[:80]}")
        except Exception as _fbe:
            print(f"[Learn/Feedback] Speichern fehlgeschlagen: {_fbe}")
    feedback = _persisted_fb

    # 0. Amboss-Screenshots aus Frontend speichern → werden als available_images eingebunden
    screenshot_paths = []
    if screenshots_b64:
        import base64 as _b64
        subj_dir = UPLOADS_DIR / "images" / f"amboss_{(subject or 'medizin').lower().replace(' ', '_')}"
        subj_dir.mkdir(parents=True, exist_ok=True)
        for i, data_url in enumerate(screenshots_b64):
            try:
                if data_url.startswith("data:image/svg"):
                    print(f"[Learn] Screenshot {i} übersprungen: SVG nicht unterstützt")
                    continue
                # Strip "data:image/...;base64," prefix
                if "," in data_url:
                    data_url = data_url.split(",", 1)[1]
                img_bytes = _b64.b64decode(data_url)
                # Auf 1024px verkleinern — verhindert stille Timeouts bei großen Screenshots
                try:
                    from PIL import Image as _SPIL
                    import io as _sio
                    _simg = _SPIL.open(_sio.BytesIO(img_bytes))
                    _simg.thumbnail((1024, 1024))
                    _sbuf = _sio.BytesIO()
                    _simg.save(_sbuf, format="PNG")
                    img_bytes = _sbuf.getvalue()
                except Exception:
                    pass  # Pillow nicht verfügbar — Original nutzen
                img_name = f"amboss_{uuid.uuid4().hex[:8]}.png"
                img_path = subj_dir / img_name
                img_path.write_bytes(img_bytes)
                screenshot_paths.append(str(img_path.relative_to(UPLOADS_DIR)).replace("\\", "/"))
                print(f"[Learn] Amboss-Screenshot gespeichert: {img_name}")
            except Exception as e:
                print(f"[Learn] Screenshot {i} konnte nicht gespeichert werden: {e}")

    # 0.5 Amboss: Artikel-Bilder + Text-Kontext via Session-Cookie (Cache > config.json)
    from amboss_client import _load_cached_cookie as _load_amboss_cookie
    _amboss_cookie = _load_amboss_cookie(UPLOADS_DIR) or CONFIG.get("amboss_cookie", "").strip()
    amboss_article_texts: list[dict] = []
    if not use_amboss:
        print(f"[Learn] Amboss uebersprungen (source_mode={source_mode})")
    elif _amboss_cookie:
        try:
            from amboss_client import fetch_amboss_context as _fetch_amboss_ctx
            _terms = topics if topics else [subject]
            _amboss_result = await _fetch_amboss_ctx(
                cookie=_amboss_cookie,
                terms=_terms,
                subject=subject,
                uploads_dir=UPLOADS_DIR,
            )
            _amboss_paths = _amboss_result.get("image_paths") or []
            amboss_article_texts = _amboss_result.get("article_texts") or []
            screenshot_paths.extend(_amboss_paths)
            print(f"[Learn] Amboss: {len(_amboss_paths)} Bilder + {len(amboss_article_texts)} Artikel-Texte")
        except Exception as _ae:
            print(f"[Learn] Amboss-Abruf fehlgeschlagen (ignoriert): {_ae}")
    else:
        print("[Learn] Amboss: kein Cookie — übersprungen (Login via Medizin-Tab oder config.json `amboss_cookie`)")

    # 0.6 ViaMedici (Thieme): Bilder + Artikel-Text via Keycloak Silent-Refresh
    viamedici_article_texts: list[dict] = []
    if not use_viamedici:
        print(f"[Learn] ViaMedici uebersprungen (source_mode={source_mode})")
    else:
        try:
            from viamedici_client import _load_kc_cookies as _load_vm_cookies, fetch_viamedici_context
            if _load_vm_cookies(UPLOADS_DIR) or CONFIG.get("viamedici_cookies", "").strip():
                _terms = topics if topics else [subject]
                _vm_result = await fetch_viamedici_context(
                    uploads_dir=UPLOADS_DIR, terms=_terms, subject=subject,
                )
                _vm_paths = _vm_result.get("image_paths") or []
                viamedici_article_texts = _vm_result.get("article_texts") or []
                screenshot_paths.extend(_vm_paths)
                print(f"[Learn] ViaMedici: {len(_vm_paths)} Bilder + {len(viamedici_article_texts)} Texte")
            else:
                print("[Learn] ViaMedici: keine Cookies — übersprungen")
        except Exception as _ve:
            print(f"[Learn] ViaMedici-Abruf fehlgeschlagen (ignoriert): {_ve}")

    # 1. RAG-Suche: nur wenn source_mode=lernraum (sonst kein eigenes Material)
    all_chunks = []
    if not use_rag:
        print(f"[Learn] RAG uebersprungen (source_mode={source_mode}) — Karten kommen aus externen Quellen")
    else:
        try:
            from rag.search import query as rag_query
            loop = asyncio.get_event_loop()

            search_terms = topics if topics else [subject]
            seen_ids = set()
            for term in search_terms:
                q = f"{subject} {term}".strip()
                hits = await loop.run_in_executor(None, lambda q=q: rag_query(q, n_chunks))
                for h in hits:
                    cid = h.get("id") or h.get("text", "")[:60]
                    if cid not in seen_ids:
                        seen_ids.add(cid)
                        all_chunks.append(h)
        except Exception as e:
            print(f"[Learn/Generate] RAG-Fehler (ignoriert): {e}")

    # 2. Bilder aus UPLOADS_DIR/images/ auflisten und nach Fach vorfiltern
    img_root = UPLOADS_DIR / "images"
    available_images = []
    if img_root.exists():
        import json as _jmeta2
        # Lese learn_meta.json für direkte Stem→Subject-Zuordnung
        meta_path = img_root / "learn_meta.json"
        meta = {}
        if meta_path.exists():
            try:
                meta = _jmeta2.loads(meta_path.read_text("utf-8"))
            except Exception:
                meta = {}

        subj_key = subject.lower()
        direct_stems = set(meta.get(subj_key, []))

        # Fuzzy-Fallback: subject-Wörter gegen gespeicherte Keys matchen
        if not direct_stems:
            subj_words = [w for w in subject.lower().split() if len(w) > 3]
            for key, stems in meta.items():
                if any(w in key or key in w for w in subj_words):
                    direct_stems.update(stems)

        # Keywords aus Subject + Topics für Dateiname-Matching (Score 1)
        match_keywords = set()
        for t in ([subject] + topics):
            for word in t.lower().split():
                if len(word) > 3:
                    match_keywords.add(word)

        all_imgs = []
        for img in img_root.rglob("*"):
            if img.suffix.lower() in (".jpg", ".jpeg", ".png", ".gif", ".webp") and img.name != "learn_meta.json":
                all_imgs.append(img)

        def img_score(img):
            # Score 2: Bild gehört zu einem direkt gemappten Ordner (via learn_meta.json)
            if img.parent.name in direct_stems:
                return 2
            name_lower = img.stem.lower()
            if any(kw in name_lower for kw in match_keywords):
                return 1
            return 0

        scored = sorted(all_imgs, key=img_score, reverse=True)
        relevant = [i for i in scored if img_score(i) > 0]
        if len(relevant) >= 3:
            available_images = [str(i.relative_to(UPLOADS_DIR)).replace("\\", "/") for i in relevant[:20]]
        else:
            available_images = [str(i.relative_to(UPLOADS_DIR)).replace("\\", "/") for i in scored[:20]]

    # Amboss-Screenshots haben höchste Priorität — vorne einreihen, Duplikate vermeiden
    if screenshot_paths:
        available_images = screenshot_paths + [p for p in available_images if p not in set(screenshot_paths)]

    # source_meta.json aus Thieme/Amboss-Upload-Ordnern einlesen (werden nicht durch img_score erfasst)
    # Für Amboss: Artikel-Bilder (article) vor Atlas-Bildern (atlas) einreihen
    import json as _jsmeta
    _subj_slug = (subject or "medizin").lower().replace(" ", "_")
    for _provider_prefix in ("thieme", "amboss"):
        _pmeta = img_root / f"{_provider_prefix}_{_subj_slug}" / "source_meta.json"
        if _pmeta.exists():
            try:
                _pdata = _jsmeta.loads(_pmeta.read_text("utf-8"))
                if _provider_prefix == "amboss":
                    # Artikel-Bilder haben Vorrang — Atlas nur als Bonus hinten
                    _art_paths = [e["path"] for e in _pdata if "path" in e and e.get("source") != "atlas"]
                    _atl_paths = [e["path"] for e in _pdata if "path" in e and e.get("source") == "atlas"]
                    _ppaths = _art_paths + _atl_paths
                    print(f"[Learn] Amboss source_meta: {len(_art_paths)} Artikel + {len(_atl_paths)} Atlas eingebunden")
                else:
                    _ppaths = [e["path"] for e in _pdata if "path" in e]
                    print(f"[Learn] {_provider_prefix} source_meta: {len(_ppaths)} Bilder eingebunden")
                available_images = _ppaths + [p for p in available_images if p not in set(_ppaths)]
            except Exception as _pme:
                print(f"[Learn] {_provider_prefix} source_meta Lesefehler: {_pme}")

    # 3. Prompt für Claude zusammenbauen
    chunks_text = "\n\n---\n\n".join(
        h.get("text", "") for h in all_chunks[:20]
    ) or "(Kein RAG-Material verfügbar — nutze dein Medizin-Wissen)"

    # Load captions from captions.json files in image directories
    all_captions = {}
    try:
        import json as _jcap2
        for cap_file in img_root.rglob("captions.json"):
            try:
                all_captions.update(_jcap2.loads(cap_file.read_text("utf-8")))
            except Exception:
                pass
    except Exception:
        pass

    # Captions + Provider-Tag aus source_meta.json mergen (Thieme/Amboss)
    import json as _jsmeta2
    for _provider_prefix in ("thieme", "amboss"):
        _pmeta2 = img_root / f"{_provider_prefix}_{_subj_slug}" / "source_meta.json"
        if _pmeta2.exists():
            try:
                for _e in _jsmeta2.loads(_pmeta2.read_text("utf-8")):
                    if "path" in _e:
                        _cap_parts = []
                        if _e.get("caption"):
                            _cap_parts.append(_e["caption"])
                        # Immer Provider-Tag hinzufügen — _provider_prefix als Fallback wenn "provider" fehlt
                        _src = (_e.get("provider") or _provider_prefix).capitalize()
                        _cap_parts.append(f"Quelle: {_src}")
                        all_captions[_e["path"]] = " — ".join(_cap_parts)
            except Exception:
                pass

    images_hint = ""
    if available_images:
        img_lines = []
        for p in available_images:
            cap = all_captions.get(p, "")
            if cap:
                img_lines.append(f"- {p}  → Thema: \"{cap}\"")
            else:
                img_lines.append(f"- {p}  → (kein Folientitel)")
        images_hint = (
            f"\n\nVerfügbare Bilder für **{subject}** (relative Pfade für ImagePath-Feld):\n"
            + "\n".join(img_lines)
            + "\n\nZuordnungsregel: Bevorzuge Amboss-Artikel-Bilder (Pfad enthält `amboss_art_`) — "
            "sie sind didaktisch aufbereitet und stehen immer zur Verfügung. "
            "Atlas-Bilder (Pfad enthält `amboss_atl_`) nur verwenden, wenn kein passendes Artikel-Bild existiert. "
            "Wähle das Bild dessen Thema-Label inhaltlich am stärksten zur Kartenfrage passt. "
            "Kopiere den Pfad wörtlich. Verteile alle Bilder möglichst gleichmäßig auf die Karten."
        )

    topics_line = ", ".join(topics) if topics else subject
    feedback_block = f"\n\nNUTZER-FEEDBACK zu vorherigen Karten (bitte umsetzen): {feedback}" if feedback else ""
    rules_block = (
        "\n\nZUSÄTZLICHE NUTZER-REGELN (überschreiben Defaults bei Konflikt):\n"
        + "\n".join(f"- {r}" for r in user_rules)
    ) if user_rules else ""

    # Source-Kontext für Karten: welche Quelldateien wurden genutzt
    source_files = list({
        p.split("/")[1] if "/" in p else p
        for p in available_images
    })[:5]
    source_hint = f"\n\nQuelldateien: {', '.join(source_files)}" if source_files else ""

    # Amboss-Artikel-Text als ergänzender Kontext (Claude soll vertiefende Karten generieren)
    amboss_context_block = ""
    if amboss_article_texts:
        _blocks = []
        for at in amboss_article_texts[:4]:  # max 4 Artikel → Prompt-Budget
            _blocks.append(f"── Amboss-Artikel: {at['title']} ──\n{at['text']}")
        amboss_context_block = (
            "\n\nZUSÄTZLICHER AMBOSS-KONTEXT (Amboss-Artikel-Auszüge zum Thema — "
            "nutze dies um das Vorlesungs-/Seminarmaterial zu **ergänzen und zu vertiefen**, "
            "z.B. mit Details/Leitlinien/Amboss-Perlen die im Lernraum-Material fehlen):\n\n"
            + "\n\n".join(_blocks)
        )

    # ViaMedici (Thieme)-Artikel-Text als ergänzender Kontext
    viamedici_context_block = ""
    if viamedici_article_texts:
        _vblocks = []
        for vt in viamedici_article_texts[:3]:
            _vblocks.append(f"── Thieme/ViaMedici-Lernmodul: {vt['title']} ──\n{vt['text']}")
        viamedici_context_block = (
            "\n\nZUSÄTZLICHER THIEME/VIAMEDICI-KONTEXT (Lernmodul-Auszüge — "
            "ergänze damit das Vorlesungsmaterial um Thieme-spezifische Erklärungen, "
            "IMPP-Hochleistungsfakten und didaktisch aufbereitete Zusammenhänge):\n\n"
            + "\n\n".join(_vblocks)
        )

    # Karten-Anzahl: entweder hart vorgegeben oder Density-Bereich für Claude-Selbsteinschätzung
    if n_cards is not None:
        count_directive = f"Erstelle genau **{n_cards}** hochwertige Anki-Karten"
    else:
        _ranges = {
            "low":    ("WENIG",  "5-10",  "Nur die absolut wichtigsten Konzepte. Highlights only — keine Details."),
            "medium": ("MITTEL", "12-20", "Solide Coverage der Hauptthemen + ein paar Vertiefungen / IMPP-relevante Details."),
            "high":   ("VIEL",   "22-40", "Tiefe Abdeckung inkl. Edge-Cases, IMPP-Fallen, Differenzialdiagnosen, Pharmako-Details."),
        }
        _label, _range, _desc = _ranges.get(density, _ranges["medium"])
        count_directive = (
            f"Wähle SELBST eine sinnvolle Karten-Anzahl im Bereich **{_range}** "
            f"(Modus: **{_label}** — {_desc})\n"
            f"Schau dir das Quellmaterial an und entscheide:\n"
            f"- Wenig/leichter Stoff → eher untere Bereichsgrenze\n"
            f"- Viel komplexer Stoff → eher obere Bereichsgrenze\n"
            f"- Lieber {_range.split('-')[0]} richtig gute Karten als {_range.split('-')[1]} dünne. "
            f"Qualität > Quantität, aber bleib im Bereich.\n"
            f"Erstelle die Anki-Karten dann"
        )

    prompt = f"""Du bist ein Medizin-Lernkarten-Experte. Erstelle Anki-Karten zum Thema: **{subject}**.
Fokus-Topics: {topics_line}

Quellmaterial aus dem Lernraum:
{chunks_text}{images_hint}{source_hint}{amboss_context_block}{viamedici_context_block}

{count_directive} im folgenden CSV-Format (TAB-separiert, kein Header):
Front[TAB]Back[TAB]Tags[TAB]ImagePath[TAB]FunFact[TAB]Merkhilfe[TAB]Source

KARTEN-TYPEN (wähle pro Karte) — **Mindestens 60% des Decks MUSS Cloze sein**:
- **Cloze/Lückentext (default!)**: Nutze Anki-Syntax IMMER mit DOPPELTEN geschweiften Klammern: `{{c1::Antwort}}`. **NIEMALS** `{{c1::Antwort}}` (single brace), das wird nicht als Cloze erkannt und landet als Roh-Text auf der Karte.
  - **Format-Konvention**: Front = kurzer Kontext-Header (z.B. "Glasgow Coma Scale - Augenoeffnung"), Back = der Lueckentext-Satz mit den `{{c1::}}`-Markern. Front darf KEINE Cloze-Marker haben — nur Back. Wenn die Karte nur ein einzelner Lueckentext-Satz ist, lass Front leer und schreib alles in Back.
  - Beispiel: Front = `Aortenstenose - Diagnostik`. Back = `Echokardiographie zeigt {{c1::AVA < 1,0 cm²}} und {{c2::dPmean > 40 mmHg}}; Goldstandard ist {{c3::TTE}}.`
  - Mehrere Cluecken pro Karte OK. Gut fuer: Definitionen, Zahlen/Werte, Medikamente+Dosen, Klassifikationen, Schemata. Werden rot+fett gerendert.
- **Basic (Front/Back klassisch)**: NUR fuer komplexe Erklaerungen, offene Fragen, Algorithmen — wenn Cloze inhaltlich gar nicht passt. Sollte die Minderheit sein (<= 40%).

EIGENE GRAFIK erzeugen — **MINDESTENS 25% der Karten sollen ein eigenes SVG/Mermaid-Element haben**, wenn der Inhalt visualisierbar ist:
- **SVG inline**: `<svg viewBox="0 0 320 140" xmlns="http://www.w3.org/2000/svg">…</svg>` — Pflicht-Trigger: Druckkurven (Aortenstenose), Achsen-Diagramme (BNP-Cutoffs), Schemata (Klappen, Hirnnerven, Reizleitung), EKG-Streifen, anatomische Abrisse. Max 4-6 Elemente, klar beschriftet, dunkler Hintergrund passt (stroke hell, z.B. `stroke="#60a5fa"` oder `#34d399`).
- **Mermaid-Diagramm**: `<div class="mermaid">graph TD; A[Dyspnoe] --> B{{BNP}}; B -->|hoch| C[Echokardiographie]; B -->|normal| D[DD prüfen]</div>` — Pflicht-Trigger: Algorithmen, Entscheidungsbaeume, Differenzialdiagnose-Trees, Therapie-Stufenschemata.
- **KI-Generierung (Nano Banana/Gemini)**: Wenn KEIN passendes Amboss-/Thieme-Bild existiert UND ein Foto-aehnliches Bild (klinisches Bild, makroskopische Anatomie) den Inhalt verstaerkt, schreibe in ImagePath: `[GENERATE: "Detaillierter Prompt fuer eine medizinische Illustration von ..."]`. Nur wenn echt noetig — bei abstrakten Konzepten lieber SVG/Mermaid.

**Bild-Quality-Filter (KRITISCH)**: Bevor du ein Bild aus der Liste oben zuordnest, frag dich: "Zeigt dieses Bild DIREKT was die Karte abfragt?" Wenn das Bild nur tangential passt oder nur dekorativ waere → **ImagePath leer lassen** und stattdessen ein eigenes SVG bauen oder ohne Bild auskommen. Schlechte Bild-Zuordnung ist schlimmer als kein Bild.

IMAGE OCCLUSION (optional, nur bei anatomischen Schema-Bildern mit klar erkennbaren Strukturen):
Wenn ImagePath ein Bild mit mehreren sichtbaren anatomischen/klinischen Strukturen zeigt (z.B. Herzklappen-Schema, Leber-Segmente, Hirnnerven-Austritt, EKG-Ableitungen), kannst du im Back Abdeckungen hinzufügen:
- Syntax pro Struktur: `[OCCLUDE x=25 y=30 w=20 h=15]Aortenklappe[/OCCLUDE]`
- Koordinaten in Prozent der Bildgröße (0-100). x,y = obere linke Ecke des Rechtecks.
- Nutze NUR wenn du das Bild wirklich verstehst (du siehst die Bilder oben per Vision) und Strukturen präzise lokalisieren kannst.
- Max 6 Abdeckungen pro Bild, sonst unleserlich.
- Bei Foto/Graph/Tabelle keine [OCCLUDE]-Tags, normale Karte.

FELDER:
- **Front**: bei Basic = praezise klinische Frage; bei Cloze = kurzer Kontext-Header (max 80 Zeichen), keine `{{c}}`-Marker hier.
- **Back**: bei Basic = strukturierte Antwort; bei Cloze = der Lueckentext-Satz mit `{{c1::}}`-Markern. HTML erlaubt (`<b>`, `<i>`, `<br>`, `<ul><li>`, `<span style="…">`), SVG/Mermaid inline OK; Box-Palette nur wenn echter Mehrwert.
- **Format-Pflicht**: Aufzaehlungen mit ≥3 Items IMMER untereinander (`<ul><li>` oder `<br>`-getrennt). NIEMALS komma-getrennt nebeneinander wenn es eine Liste ist. Beispiel falsch: "Symptome: Dyspnoe, Synkope, Angina." Beispiel richtig: `<ul><li>Dyspnoe</li><li>Synkope</li><li>Angina</li></ul>` oder mit `<br>` zwischen den Items.
- **Tags**: Komma-getrennt, z.B. "Kardiologie,EKG,sem8"
- **ImagePath**: EXAKT einen Pfad aus der Bilderliste oben kopieren (nicht erfinden). Bild-Quality-Filter (oben) anwenden — lieber leer als unpassend. Matche ueber das Thema-Label.
- **FunFact**: NUR wenn du eine wirklich ueberraschende klinische Pointe hast (z.B. Epidemiologie-Trivia, historischer Fakt, Eponym-Hintergrund). Keine Fueller. Sonst leer.
- **Merkhilfe**: NUR konkrete Eselsbruecken — Akronyme (z.B. "MONA fuer ACS: Morphin, O2, Nitro, ASS"), Reime, Visual-Bilder ("Aortenstenose-Trias = SAD: Synkope/Angina/Dyspnoe"). KEINE schwammigen "Denke an X"-Phrasen. Wenn du nichts Echtes hast, lass leer.
- **Source**: Format `Datei.pdf · S. N` (mit Seitenzahl wenn aus dem Bildpfad ableitbar, z.B. `slide_03.png` → `S. 3`). Wenn keine Seite ableitbar: nur Datei.pdf. Bei Amboss/Thieme: "Quelle: Amboss" / "Quelle: Thieme".

BOX-PALETTE (optional im Back, max 3 pro Karte, Syntax `[TAG]Inhalt[/TAG]`):
- `[THERAPIE]…[/THERAPIE]`    💊 grün   — Medikamente, Dosierungen, Therapie-Schemata
- `[NOTFALL]…[/NOTFALL]`      ⚠️ rot    — akute Situationen, was sofort tun
- `[DX]…[/DX]`                🩺 blau   — Diagnostik-Kriterien, Scores, Tests
- `[PATHO]…[/PATHO]`          🧬 lila   — Pathomechanismus
- `[LABOR]…[/LABOR]`          🔬 gelb   — Laborwerte, Cut-offs, Normbereiche
- `[IMPP]…[/IMPP]`            🎯 pink   — examensrelevante High-Yield-Fakten
- `[LEITLINIE]…[/LEITLINIE]`  📖 türkis — S3/S2k-Empfehlungen, Guideline-Zitate
- `[REDFLAG]…[/REDFLAG]`      🚨 dunkelrot — Alarm-/Warnsymptome
- `[DD]…[/DD]`                🎯 violett — Differentialdiagnosen
- `[AMBOSS]…[/AMBOSS]`        🟢 grün-Badge — wenn Inhalt klar aus Amboss-Kontext (wenn oben mitgegeben), explizit als Ergänzung markieren
- `[THIEME]…[/THIEME]`        🔵 blau-Badge — wenn Inhalt aus ViaMedici/Thieme

TEXT-FORMATIERUNG im Back:
- Schlüsselbegriffe fett: `<b>Aortenstenose</b>`
- Synonyme/Definitionen kursiv: `<i>(syn.: Valvula aortae)</i>`
- Farb-Codes direkt inline (nur wenn nicht in einer Box):
  • Therapie → `<span style="color:#34d399">Metoprolol 2×25mg</span>`
  • Notfall/Warning → `<span style="color:#ef4444;font-weight:700">Notfall!</span>`
  • Diagnostik → `<span style="color:#60a5fa">EKG, TTE</span>`
  • Laborwerte → `<span style="color:#fde047">BNP > 400 pg/ml</span>`
- Hervorgehobene Begriffe: `<span style="font-size:1.15em;font-weight:700">Kernaussage</span>`
- Emojis inline für Kategorien: 🩺 💊 ⚠️ 🔬 🧬 💉 📊 📖

INHALT-PRIORITÄT (WICHTIG):
1. **Primärquelle:** Robins Seminar-/VL-Material (oben unter „Quellmaterial aus dem Lernraum"). Der KERN jeder Karte kommt daher. Wenn das Material eine Aussage macht, ist DAS das Fundament.
2. **Amboss/Thieme ist ergänzend**, NICHT Ersatz. Pro Karte maximal EINE externe Box — entweder [AMBOSS] ODER [THIEME], nie beide gleichzeitig, nie mehrfach.
3. Nutze die externe Box nur wenn sie KONKRETEN Mehrwert hat, den das Seminar-Material nicht bringt (z.B. „Amboss-Perle: 80% rechtsventrikulär", „Thieme-Lernmodul ergänzt: ESC-Leitlinie 2021 sagt …", eine IMPP-Falle). Nicht einfach „auch Amboss bestätigt X".
4. Wenn das Seminar-Material die Frage vollständig abdeckt, LASS die externe Box weg. Nicht jede Karte braucht Extras.
5. Nicht jede Karte braucht Boxen überhaupt. Klare, präzise Antworten sind besser als voll-gestopfte Karten.

Gib NUR die TSV-Zeilen aus, keine Erklärungen, keinen Header.{feedback_block}{rules_block}"""

    # 4. Top-5-Bilder als base64 für Vision-API laden
    import base64 as _b64_vision
    _ext_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".gif": "image/gif", ".webp": "image/webp"}
    image_blocks = []
    for img_rel in available_images[:5]:
        img_file = UPLOADS_DIR / img_rel
        ext = img_file.suffix.lower()
        media = _ext_map.get(ext, "image/png")
        try:
            try:
                _fsize = img_file.stat().st_size
            except Exception:
                _fsize = 0
            if _fsize > 5 * 1024 * 1024:
                print(f"[Learn] Bild zu groß ({_fsize // (1024*1024)}MB), übersprungen: {img_rel}")
                continue
            raw_bytes = img_file.read_bytes()
            try:
                from PIL import Image as _PILImage
                import io as _io
                _pil_img = _PILImage.open(_io.BytesIO(raw_bytes))
                _pil_img.thumbnail((1024, 1024))
                _buf = _io.BytesIO()
                _pil_img.save(_buf, format="PNG")
                raw_bytes = _buf.getvalue()
                media = "image/png"
            except Exception:
                pass  # Pillow nicht verfügbar oder Fehler — Original nutzen
            img_data = _b64_vision.b64encode(raw_bytes).decode()
            image_blocks.append({
                "type": "image",
                "source": {"type": "base64", "media_type": media, "data": img_data}
            })
        except Exception as _img_err:
            print(f"[Learn] Bild nicht lesbar ({img_rel}): {_img_err}")

    # 5. Claude-Aufruf — SDK nur wenn ANTHROPIC_API_KEY gesetzt, sonst direkt CLI.
    # Robin's Setup nutzt OAuth via claude-cli, kein API-Key → SDK-Versuch ist
    # dort sinnlos (wirft "Could not resolve authentication method") und kostet
    # nur Latenz. CLI kann via Read-Tool die Bilder selbst lesen.
    raw = ""
    _claude_error = None
    _has_api_key = bool(os.environ.get("ANTHROPIC_API_KEY"))
    if _has_api_key and image_blocks:
        try:
            import anthropic as _anthropic
            _aclient = _anthropic.AsyncAnthropic()
            _model = CONFIG.get("claude_model_force") or CONFIG.get("claude_model") or "claude-sonnet-4-6"
            _content = image_blocks + [{"type": "text", "text": prompt}]
            _resp = await _aclient.messages.create(
                model=_model,
                max_tokens=4096,
                messages=[{"role": "user", "content": _content}]
            )
            raw = _resp.content[0].text.strip()
            print(f"[Learn] SDK Vision-Call: {len(image_blocks)} Bilder, Modell {_model}")
        except Exception as _sdk_err:
            _claude_error = f"SDK fail (CLI-Fallback aktiv): {_sdk_err}"
            print(f"[Learn] {_claude_error}")
            raw = ""

    if not raw:
        try:
            # CLI-Pfad: Bilder als absolute Pfade in den Prompt — Agent liest sie via Read-Tool.
            cli_prompt = prompt
            if available_images:
                paths = []
                for img_rel in available_images[:5]:
                    try:
                        abs_path = str((UPLOADS_DIR / img_rel).resolve()).replace("\\", "/")
                        paths.append(abs_path)
                    except Exception:
                        pass
                if paths:
                    img_section = (
                        "\n\n## Bilder zum Analysieren\n"
                        "Lies die folgenden Bilder mit deinem Read-Tool und beziehe sie "
                        "in die Karten ein. Schreibe in ImagePath **exakt** einen dieser Pfade — "
                        "kopiere den Pfad wörtlich, erfinde keinen eigenen Namen. "
                        "Wähle das Bild das inhaltlich am besten zur Karte passt. "
                        "Verteile alle verfügbaren Bilder auf die Karten:\n"
                        + "\n".join("- " + p for p in paths)
                    )
                    cli_prompt = prompt + img_section

            # claude-cli: --dangerously-skip-permissions damit Read-Tool ohne
            # Approval durchgeht. Sonst haengt der Subprocess auf der Bilder-Phase.
            _gen_model = (CONFIG.get("claude_model_generate") or "opus").strip()
            proc = await asyncio.create_subprocess_exec(
                "claude", "-p", "--model", _gen_model,
                "--dangerously-skip-permissions",
                "--max-turns", "30",
                stdin=asyncio.subprocess.PIPE,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            stdout, stderr = await asyncio.wait_for(
                proc.communicate(cli_prompt.encode("utf-8")), timeout=300
            )
            raw = stdout.decode("utf-8", errors="replace").strip()
            # Markdown-Codeblock-Wrap entfernen falls Modell ihn dazugepackt hat
            if raw.startswith("```"):
                lines = raw.splitlines()
                # erste + letzte ``` Zeile droppen
                if lines[0].startswith("```"):
                    lines = lines[1:]
                if lines and lines[-1].strip().startswith("```"):
                    lines = lines[:-1]
                raw = "\n".join(lines).strip()
            print(f"[Learn] CLI-Call: {len(raw)} chars output, {len(available_images[:5])} Bilder")
            if not raw:
                _claude_error = (_claude_error or "") + " | CLI lieferte leeren Output"
        except asyncio.TimeoutError:
            return {"ok": False, "error": "Claude Timeout (300s) — vermutlich zu viele Bilder"}
        except FileNotFoundError:
            return {"ok": False, "error": "claude CLI nicht gefunden"}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    # 5. TSV parsen → strukturierte Karten
    import re as _re
    # LLMs schreiben oft single-brace `{c1::...}` statt Anki's `{{c1::...}}` —
    # ohne Normalisierung schlaegt _is_cloze in apkg_export fehl und die Karte
    # landet als Basic mit Roh-Text "Augenoeffnung max {c1::4}" auf der Front.
    _SINGLE_BRACE_CLOZE = _re.compile(r'(?<!\{)\{c(\d+)::([^{}]+?)\}(?!\})')

    def _normalize_cloze(text: str) -> str:
        if not text:
            return text
        return _SINGLE_BRACE_CLOZE.sub(r'{{c\1::\2}}', text)

    # Page-Number-Extraktion aus Bildpfaden (z.B. slide_03.png, page-5.jpg, _p12_).
    _PAGE_RE = _re.compile(r'(?:slide|seite|page|_p|-p)[_\-]?0*(\d{1,4})', _re.IGNORECASE)

    def _extract_page(image_path: str) -> str:
        if not image_path:
            return ""
        m = _PAGE_RE.search(image_path)
        return f"S. {m.group(1)}" if m else ""

    cards = []
    fields_order = ["front", "back", "tags", "image", "funfact", "merkhilfe", "source"]
    for line in raw.splitlines():
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        parts = line.split("\t")
        if len(parts) < 2:
            continue
        # Pad to 7 fields
        while len(parts) < 7:
            parts.append("")
        card = dict(zip(fields_order, parts[:7]))
        # Nano-Banana-Tag aus image-Feld extrahieren: [GENERATE: "Prompt"]
        # → image_prompt gesetzt, image geleert. _enrich_anki_cards_with_images
        # generiert nach der Loop parallel via Gemini Flash Image.
        _gen_m = _re.match(r'\s*\[GENERATE:\s*"(.+?)"\s*\]\s*$', card.get("image") or "")
        if _gen_m:
            card["image_prompt"] = _gen_m.group(1).strip()
            card["image"] = ""
        # Cloze-Normalisierung (single-brace -> double-brace) BEVOR irgendwas anderes
        # die Felder anschaut. Sonst sieht _is_cloze ein "{c1::4}" als Plaintext.
        card["front"] = _normalize_cloze(card.get("front") or "")
        card["back"] = _normalize_cloze(card.get("back") or "")
        # External-Source vor Transform ableiten (nach Transform sind die [TAG]s weg)
        _raw_back = card.get("back") or ""
        if _re.search(r'\[AMBOSS\]', _raw_back, _re.IGNORECASE):
            card["external_source"] = "Amboss"
        elif _re.search(r'\[THIEME\]', _raw_back, _re.IGNORECASE):
            card["external_source"] = "Thieme/ViaMedici"
        else:
            card["external_source"] = None
        # Image-Occlusions aus Back extrahieren BEVOR Box-Transform läuft
        _raw_back, _occs = _extract_occlusions(_raw_back)
        card["occlusions"] = _occs
        card["back"] = _raw_back
        # Box-Tags [THERAPIE]…[/THERAPIE] etc. zu inline-HTML transformieren (Anki-kompatibel)
        if card.get("back"):
            card["back"] = _transform_card_boxes(card["back"])
        # Fallback: Source aus Bildpfad ableiten wenn Claude es leer gelassen hat
        if not card.get("source") and card.get("image"):
            img_p = card["image"].lower()
            if "amboss" in img_p:
                card["source"] = "Quelle: Amboss"
            elif "thieme" in img_p:
                card["source"] = "Quelle: Thieme"
        # Caption als Source-Override wenn Quelle bekannt
        elif card.get("image") and all_captions.get(card["image"]):
            cap = all_captions[card["image"]]
            if cap.startswith("Quelle:"):
                card["source"] = cap.split(" — ")[0]  # nur "Quelle: Amboss" ohne Caption-Text
        # Seitenzahl in Source ergaenzen wenn aus Bildpfad ableitbar und noch nicht drin
        _page = _extract_page(card.get("image") or "")
        if _page and card.get("source") and "S." not in card["source"] and "Seite" not in card["source"]:
            card["source"] = f'{card["source"]} · {_page}'
        cards.append(card)

    # 5b. Nano Banana 2: parallel Bilder generieren für Karten mit image_prompt
    # (extrahiert aus [GENERATE: "..."]-Tags des LLM). Best-effort — Fehler killen
    # den Generate-Flow nicht. Cap=8 / 4 parallel via _enrich_anki_cards_with_images.
    _img_prompts_n = sum(1 for c in cards if c.get("image_prompt"))
    if _img_prompts_n:
        print(f"[Learn/Image] Nano Banana: {_img_prompts_n} Karte(n) mit image_prompt → generiere Bilder")
        try:
            await _enrich_anki_cards_with_images(cards)
        except Exception as _enrich_e:
            print(f"[Learn/Image] enrich failed (ignored): {type(_enrich_e).__name__}: {_enrich_e}")


    # 6. CSV für Download zusammenbauen
    import io, csv as _csv
    buf = io.StringIO()
    writer = _csv.writer(buf, delimiter="\t")
    for c in cards:
        writer.writerow([c["front"], c["back"], c["tags"], c["image"], c["funfact"], c["merkhilfe"], c.get("source", "")])
    csv_content = buf.getvalue()

    # 7. In UPLOADS_DIR speichern für späteres Herunterladen
    out_filename = f"anki_{subject.lower().replace(' ', '_') or 'karten'}_{uuid.uuid4().hex[:6]}.tsv"
    out_path = UPLOADS_DIR / out_filename
    out_path.write_text(csv_content, encoding="utf-8")

    # 7b. Vollständigen Deck-Snapshot als JSON ablegen — erlaubt späteres
    # Re-Building von .apkg (mit aktuellem Design) und HTML-Preview.
    # Bei source_mode != lernraum kommt der Inhalt nicht aus dem gewaehlten
    # Lernraum (Subject), sondern aus Amboss/Thieme. Dann ist der Subject-Name
    # irrefuehrend — Topics + Source-Label sind ehrlicher.
    _source_label = {"amboss": "Amboss", "thieme": "Thieme", "mix": "Amboss+Thieme"}.get(source_mode)
    if _source_label and topics:
        _display_subject = f"{', '.join(topics)} ({_source_label})"
    else:
        _display_subject = subject
    import json as _jdeck
    deck_id = uuid.uuid4().hex[:12]
    decks_dir = UPLOADS_DIR / "decks"
    decks_dir.mkdir(parents=True, exist_ok=True)
    deck_payload = {
        "deck_id": deck_id,
        "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "subject": _display_subject,
        "topics": topics,
        "card_count": len(cards),
        "cards": cards,
    }
    (decks_dir / f"{deck_id}.json").write_text(
        _jdeck.dumps(deck_payload, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    # 8. History-Eintrag schreiben
    import json as _jh
    history_path = UPLOADS_DIR / "learn_history.json"
    try:
        history = _jh.loads(history_path.read_text(encoding="utf-8")) if history_path.exists() else []
    except Exception:
        history = []
    history.insert(0, {
        "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "subject": _display_subject,
        "topics": topics,
        "card_count": len(cards),
        "source_files": source_files,
        "screenshots_used": len(screenshot_paths),
        "filename": out_filename,
        "deck_id": deck_id,
    })
    history_path.write_text(_jh.dumps(history[:50], ensure_ascii=False, indent=2), encoding="utf-8")

    return {
        "ok": True,
        "subject": subject,
        "cards": cards,
        "count": len(cards),
        "csv": csv_content,
        "filename": out_filename,
        "deck_id": deck_id,
        "screenshots_used": len(screenshot_paths),
        "images_used": len(image_blocks),
        "claude_error": _claude_error,
    }


# Box-Palette: Claude schreibt [TAG]Inhalt[/TAG], wir transformieren zu inline-style HTML.
# Inline Styles → robust auch im Anki-Export ohne externes CSS.
_CARD_BOX_PALETTE = {
    "THERAPIE":    {"emoji": "💊", "label": "Therapie",          "bg": "#34d39912", "border": "#34d399", "fg": "#86efac"},
    "NOTFALL":     {"emoji": "⚠️", "label": "Notfall",           "bg": "#ef444418", "border": "#ef4444", "fg": "#fca5a5"},
    "DX":          {"emoji": "🩺", "label": "Diagnostik",        "bg": "#60a5fa15", "border": "#60a5fa", "fg": "#93c5fd"},
    "PATHO":       {"emoji": "🧬", "label": "Pathophysiologie",  "bg": "#c084fc15", "border": "#c084fc", "fg": "#d8b4fe"},
    "LABOR":       {"emoji": "🔬", "label": "Labor",             "bg": "#fde04718", "border": "#eab308", "fg": "#fde047"},
    "IMPP":        {"emoji": "🎯", "label": "IMPP-Perle",        "bg": "#f472b615", "border": "#f472b6", "fg": "#f9a8d4"},
    "LEITLINIE":   {"emoji": "📖", "label": "Leitlinie",         "bg": "#2dd4bf15", "border": "#2dd4bf", "fg": "#5eead4"},
    "REDFLAG":     {"emoji": "🚨", "label": "Red Flag",          "bg": "#dc262620", "border": "#dc2626", "fg": "#fca5a5"},
    "DD":          {"emoji": "🎯", "label": "DD",                "bg": "#a855f715", "border": "#a855f7", "fg": "#d8b4fe"},
    "AMBOSS":      {"emoji": "🟢", "label": "Ergänzung aus Amboss",  "bg": "#16a34a15", "border": "#16a34a", "fg": "#86efac"},
    "THIEME":      {"emoji": "🔵", "label": "Aus Thieme/ViaMedici",  "bg": "#0ea5e915", "border": "#0ea5e9", "fg": "#7dd3fc"},
}


def _transform_card_boxes(text: str) -> str:
    """Ersetzt [TAG]…[/TAG] durch inline-styled HTML-Boxen (Anki-kompatibel)."""
    if not text:
        return text
    import re as _re_cb
    for tag, style in _CARD_BOX_PALETTE.items():
        pattern = _re_cb.compile(r"\[" + tag + r"\](.*?)\[/" + tag + r"\]", _re_cb.DOTALL | _re_cb.IGNORECASE)
        replacement = (
            '<div style="background:' + style["bg"] + ';border-left:3px solid ' + style["border"] +
            ';padding:8px 12px;border-radius:6px;margin:8px 0;font-size:13px;line-height:1.5">'
            '<b style="color:' + style["fg"] + ';letter-spacing:.02em">' + style["emoji"] + ' ' + style["label"] + ':</b> '
            '\\1</div>'
        )
        text = pattern.sub(replacement, text)
    return text


# Image-Occlusion: [OCCLUDE x=25 y=30 w=20 h=15]Label[/OCCLUDE] → parsed aus Back,
# danach aus Back entfernt. Renderer baut SVG-Overlay aus `card.occlusions`.
_OCCLUDE_RE = __import__("re").compile(
    r"\[OCCLUDE\s+x=(\d+(?:\.\d+)?)\s*%?\s+y=(\d+(?:\.\d+)?)\s*%?\s+w=(\d+(?:\.\d+)?)\s*%?\s+h=(\d+(?:\.\d+)?)\s*%?\s*\](.*?)\[/OCCLUDE\]",
    __import__("re").DOTALL | __import__("re").IGNORECASE,
)


def _extract_occlusions(text: str) -> tuple[str, list[dict]]:
    """Zieht [OCCLUDE x=.. y=.. w=.. h=..]Label[/OCCLUDE] aus `text` raus.
    Gibt (text_ohne_occlude, [{label, x, y, w, h}]) zurück. Koordinaten in % (0-100).
    """
    if not text or "[OCCLUDE" not in text.upper():
        return text, []
    occlusions: list[dict] = []
    def _capture(m):
        try:
            x, y, w, h = float(m.group(1)), float(m.group(2)), float(m.group(3)), float(m.group(4))
        except ValueError:
            return ""
        # Clamp auf sinnvolle Bereiche
        x = max(0.0, min(99.0, x))
        y = max(0.0, min(99.0, y))
        w = max(1.0, min(100.0 - x, w))
        h = max(1.0, min(100.0 - y, h))
        occlusions.append({
            "label": (m.group(5) or "").strip(),
            "x": round(x, 2), "y": round(y, 2),
            "w": round(w, 2), "h": round(h, 2),
        })
        return ""  # Tag aus Back entfernen
    cleaned = _OCCLUDE_RE.sub(_capture, text)
    # Whitespace-Aufräumen (doppelte Leerzeichen/Umbrüche wo Tags waren)
    import re as _re_oc
    cleaned = _re_oc.sub(r"[ \t]+\n", "\n", cleaned)
    cleaned = _re_oc.sub(r"\n{3,}", "\n\n", cleaned).strip()
    return cleaned, occlusions


def _find_best_image_for_topic(uploads_dir: Path, topic: str, subject: str) -> dict | None:
    """Sucht in allen `source_meta.json`-Files unter `uploads/images/*/` nach dem
    besten Bild-Match zum Topic. Scoring: Anzahl Topic-Keywords in Caption (voll)
    + 0.5 pro Subject-Keyword. Gibt `{path, caption, provider, source}` oder None.
    Fallback-Strategie: wenn kein Caption-Match, bevorzuge Bilder aus dem Subject-Ordner.
    """
    import json as _jfi
    images_root = uploads_dir / "images"
    if not images_root.exists():
        return None

    topic_kws = {w.lower() for w in topic.split() if len(w) > 3}
    subj_kws  = {w.lower() for w in subject.split() if len(w) > 3}
    subj_slug = (subject or "").lower().replace(" ", "_")

    candidates: list[tuple[float, dict, str]] = []  # (score, entry, folder_name)
    for folder in images_root.iterdir():
        if not folder.is_dir():
            continue
        meta = folder / "source_meta.json"
        if not meta.exists():
            continue
        try:
            entries = _jfi.loads(meta.read_text("utf-8"))
        except Exception:
            continue
        for e in entries:
            if not isinstance(e, dict) or not e.get("path"):
                continue
            caption = (e.get("caption") or "").lower()
            score = sum(1 for kw in topic_kws if kw in caption)
            score += 0.5 * sum(1 for kw in subj_kws if kw in caption)
            # Kleine Bonus wenn der Ordner selbst zum Fach passt (z.B. amboss_kardiologie)
            if subj_slug and subj_slug in folder.name:
                score += 0.25
            candidates.append((score, e, folder.name))

    if not candidates:
        return None

    candidates.sort(key=lambda t: -t[0])
    best_score, best_entry, _ = candidates[0]
    # Nur zurückgeben wenn echter thematischer Match (Topic-Keyword in Caption).
    # Kein Fallback auf „irgendein Bild aus dem Subject-Ordner" mehr — das führt
    # zu irrelevanten Bildern („Aortenstenose" → generisches Herzklappen-Bild).
    # Lieber garkein Bild als ein unpassendes.
    if best_score >= 1.0:
        return best_entry
    return None


async def api_learn_preview_card(request):
    """POST /api/learn/preview-card
    Body: { topic: str, subject?: str }
    Returns: { ok, card: {front, back, tags, funfact, merkhilfe, image, subject, topic} }
    Leichtgewichtiger Preview ohne RAG — Haiku generiert 1 echte Karte.
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    topic = (body.get("topic") or "").strip()
    subject = (body.get("subject") or topic).strip()
    if not topic:
        return web.json_response({"ok": False, "error": "topic erforderlich"}, status=400)

    prompt = (
        f"Erstelle eine einzige hochwertige Anki-Karte auf Deutsch zum Thema: {topic} (Fach: {subject}).\n\n"
        "Format — genau eine Zeile, Felder durch Tab getrennt:\n"
        "Front\tBack\tTags\tFunFact\tMerkhilfe\n\n"
        "FELDER:\n"
        "- Front: präzise klinische Frage (1-2 Sätze). BEVORZUGT als Cloze/Lückentext mit {{c1::…}}, {{c2::…}} wenn das Thema Fakten/Zahlen/Definitionen sind.\n"
        "- Back: strukturierte Antwort (HTML erlaubt: <b>, <i>, <br>, <span style=…>); Box-Tags siehe unten. Cloze-Syntax {{c1::…}} auch hier erlaubt.\n"
        "- Tags: Fach::Unterthema Format\n"
        "- FunFact: **DEFAULT: LEER LASSEN.** Nur ausfüllen wenn du einen WIRKLICH überraschenden, nicht-offensichtlichen Fakt hast (z.B. 'Amiodaron enthält 37% Jod nach Gewicht'). Keine Füller wie 'Aortenstenose ist häufig bei Älteren'. 80% der Karten haben KEINEN guten FunFact — dann leer.\n"
        "- Merkhilfe: **DEFAULT: LEER LASSEN.** Nur ausfüllen wenn du eine echt einprägsame Eselsbrücke hast (z.B. 'SINKEN: Sinus-Na-Kalium-Ca'). Erzwungene Mnemonics wie 'A=Aortenklappe, S=Stenose' sind schlechter als keine. 70% der Karten haben keine gute — dann leer.\n\n"
        "BOX-PALETTE (optional im Back, max 3, Syntax [TAG]Inhalt[/TAG]):\n"
        "  [THERAPIE] 💊 Medikamente/Therapie   |  [NOTFALL] ⚠️ akute Situation\n"
        "  [DX] 🩺 Diagnostik/Scores            |  [PATHO] 🧬 Pathomechanismus\n"
        "  [LABOR] 🔬 Laborwerte/Cut-offs       |  [IMPP] 🎯 examensrelevante Perle\n"
        "  [LEITLINIE] 📖 S3/S2k-Guideline      |  [REDFLAG] 🚨 Alarm-Symptom\n"
        "  [DD] Differentialdiagnose            |  [AMBOSS]/[THIEME] Ergänzung aus externer Quelle\n\n"
        "FORMATIERUNG im Back: <b>wichtig</b>, <i>Synonym</i>, Farbcodes "
        "(Therapie=#34d399, Notfall=#ef4444, Diagnostik=#60a5fa, Labor=#fde047), "
        "Emojis inline (🩺💊⚠️🔬🧬).\n\n"
        "SVG/Mermaid inline im Back NUR wenn simples Diagramm den Inhalt klarer macht "
        "(Kurven, Algorithmen, Schemata). Reine Faktenkarten brauchen das nicht.\n\n"
        "Gib NUR die eine Tab-getrennte Zeile aus, kein Markdown, keine Überschrift. "
        "Leere Felder = einfach Tab, keine Platzhalter."
    )

    # Claude CLI Subprocess — MC nutzt Max-Subscription, kein API-Key.
    # Opus für beste Befolgung der Prompt-Regeln (Cloze-Nutzung, FunFact/Merkhilfe nur wenn echt sinnvoll,
    # Box-Auswahl). Dauer: ~20-30s statt ~8s mit Haiku, aber Qualität deutlich besser.
    _preview_model = (CONFIG.get("claude_model_preview") or "opus").strip()
    try:
        proc = await asyncio.create_subprocess_exec(
            "claude", "-p", "--model", _preview_model,
            stdin=asyncio.subprocess.PIPE,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(
            proc.communicate(prompt.encode("utf-8")), timeout=90
        )
        text = stdout.decode("utf-8", errors="replace").strip()
        if not text:
            err = stderr.decode("utf-8", errors="replace")[:300]
            return web.json_response({"ok": False, "error": f"Claude CLI leere Antwort: {err}"}, status=500)
        parts = text.split("\t")
        _back_raw = parts[1] if len(parts) > 1 else text
        import re as _re_es
        if _re_es.search(r'\[AMBOSS\]', _back_raw, _re_es.IGNORECASE):
            _ext_src = "Amboss"
        elif _re_es.search(r'\[THIEME\]', _back_raw, _re_es.IGNORECASE):
            _ext_src = "Thieme/ViaMedici"
        else:
            _ext_src = None
        card = {
            "front": parts[0] if len(parts) > 0 else topic,
            "back": _transform_card_boxes(_back_raw),
            "tags": parts[2] if len(parts) > 2 else subject,
            "funfact": parts[3] if len(parts) > 3 else "",
            "merkhilfe": parts[4] if len(parts) > 4 else "",
            "image": None,
            "external_source": _ext_src,
            "subject": subject,
            "topic": topic,
        }
        # Passendes Bild via source_meta.json (Caption-Match, nicht Filename-Keyword)
        # Amboss/ViaMedici-Bilder haben UUID-Filenames — Filename-Match fand nie was.
        # Bild-Feld bleibt relativer Pfad (wie beim Full-Generate); Renderer prefixt /api/uploads/.
        best = _find_best_image_for_topic(UPLOADS_DIR, topic, subject)
        if best:
            card["image"] = best["path"]
            card["image_caption"] = best.get("caption") or ""
            card["image_provider"] = best.get("provider") or ""
        # Cookie-Status für Frontend-Anzeige
        import time as _time
        from amboss_client import _load_cached_cookie as _lcc, _cookie_saved_at as _csa
        _ck = _lcc(UPLOADS_DIR)
        _saved = _csa(UPLOADS_DIR)
        _cookie_age_h = round(((_time.time() - _saved) / 3600), 1) if _saved else None
        return web.json_response({
            "ok": True,
            "card": card,
            "cookie_status": {
                "has_cookie": bool(_ck),
                "age_hours": _cookie_age_h,
                "warning": (_cookie_age_h or 0) > 48,
            }
        })
    except asyncio.TimeoutError:
        return web.json_response({"ok": False, "error": "Claude CLI Timeout (90s) — Opus ist lahmer, versuch's nochmal oder setze claude_model_preview=sonnet in config.json"}, status=504)
    except FileNotFoundError:
        return web.json_response({"ok": False, "error": "claude CLI nicht im PATH gefunden"}, status=500)
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)


async def api_learn_meta(request):
    """GET /api/learn/meta — Gibt learn_meta.json zurück: Fach → [Ordner/Stems]."""
    import json as _jlm
    meta_path = UPLOADS_DIR / "images" / "learn_meta.json"
    if not meta_path.exists():
        return web.json_response({"subjects": {}, "total_files": 0})
    try:
        meta = _jlm.loads(meta_path.read_text("utf-8"))
    except Exception as e:
        return web.json_response({"error": str(e)}, status=500)
    # Count original PDF/PPTX uploads per subject via stem name matching
    upload_exts = {".pdf", ".pptx", ".ppt"}
    uploaded_originals = list(UPLOADS_DIR.glob("*")) if UPLOADS_DIR.exists() else []
    uploaded_originals = [f for f in uploaded_originals if f.suffix.lower() in upload_exts]

    subjects_info = {}
    for subject, stems in meta.items():
        stem_list = stems if isinstance(stems, list) else [stems]
        image_count = 0
        for stem in stem_list:
            img_dir = UPLOADS_DIR / "images" / stem
            if img_dir.exists():
                image_count += sum(1 for _ in img_dir.iterdir())
        # Match uploaded files whose name contains any stem of this subject
        file_count = sum(
            1 for f in uploaded_originals
            if any(stem.lower() in f.name.lower() for stem in stem_list)
        )
        subjects_info[subject] = {
            "stems": stem_list,
            "image_count": image_count,
            "file_count": file_count,
        }
    # Total unique originals (may overlap if same file covers multiple subjects → deduplicate)
    total_originals = len(uploaded_originals)
    return web.json_response({
        "subjects": subjects_info,
        "total_files": sum(v["image_count"] for v in subjects_info.values()),
        "total_uploads": total_originals,
    })


async def api_learn_history(request):
    """GET /api/learn/history — Letzte 50 Generierungsläufe."""
    import json as _jh2
    history_path = UPLOADS_DIR / "learn_history.json"
    if not history_path.exists():
        return web.json_response({"history": []})
    try:
        history = _jh2.loads(history_path.read_text(encoding="utf-8"))
    except Exception:
        history = []
    return web.json_response({"history": history})


def _load_deck(deck_id: str) -> dict | None:
    import json as _jd
    safe = "".join(ch for ch in (deck_id or "") if ch.isalnum())[:32]
    if not safe:
        return None
    p = UPLOADS_DIR / "decks" / f"{safe}.json"
    if not p.exists():
        return None
    try:
        return _jd.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return None


async def api_learn_deck_apkg(request):
    """GET /api/learn/decks/{deck_id}/apkg — Aktueller Design-Stand der gespeicherten Karten."""
    deck = _load_deck(request.match_info.get("deck_id", ""))
    if not deck:
        return web.Response(status=404, text="Deck nicht gefunden")
    try:
        from apkg_export import build_apkg
        data, _stats = build_apkg(deck.get("cards") or [], deck.get("subject") or "Karten", UPLOADS_DIR)
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)
    slug = (deck.get("subject") or "karten").lower().replace(" ", "_")
    return web.Response(
        body=data,
        content_type="application/octet-stream",
        headers={"Content-Disposition": f'attachment; filename="{slug}_anki.apkg"'},
    )


async def api_learn_deck_preview(request):
    """GET /api/learn/decks/{deck_id}/preview — HTML-Vorschau aller Karten im MC-Design."""
    deck = _load_deck(request.match_info.get("deck_id", ""))
    if not deck:
        return web.Response(status=404, text="Deck nicht gefunden")
    try:
        from apkg_export import render_deck_preview_html
        html = render_deck_preview_html(deck.get("cards") or [], deck.get("subject") or "Karten",
                                        deck.get("topics") or [], deck.get("ts") or "")
    except Exception as e:
        import traceback as _tb
        _tb.print_exc()
        return web.Response(text=f"Preview-Fehler: {type(e).__name__}: {e}", status=500)
    return web.Response(text=html, content_type="text/html")


async def api_learn_feedback(request):
    """GET/DELETE /api/learn/feedback?subject=X — Persistiertes Feedback pro Fach lesen oder löschen."""
    import json as _jfb2
    _fb_path = UPLOADS_DIR / "learn_feedback.json"
    try:
        store = _jfb2.loads(_fb_path.read_text(encoding="utf-8")) if _fb_path.exists() else {}
    except Exception:
        store = {}
    subject = (request.rel_url.query.get("subject") or "").strip().lower()
    if request.method == "DELETE":
        if subject and subject in store:
            del store[subject]
            _fb_path.write_text(_jfb2.dumps(store, ensure_ascii=False, indent=2), encoding="utf-8")
        return web.json_response({"ok": True})
    if subject:
        return web.json_response({"subject": subject, "feedback": store.get(subject, "")})
    return web.json_response({"feedback": store})


async def api_learn_screenshot(request):
    """POST /api/learn/screenshot
    Body: { subject: str, images: [dataURL, ...] }
    Speichert Amboss-Screenshots für ein Fach, ohne Kartengenerierung.
    Returns: { ok, saved: int, paths: [...] }
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    subject = (body.get("subject") or "medizin").strip().lower().replace(" ", "_")
    images = body.get("images") or []
    if not images:
        return web.json_response({"ok": False, "error": "Keine Bilder übergeben"}, status=400)

    import base64 as _b64
    import re as _re
    subj_dir = UPLOADS_DIR / "images" / f"amboss_{subject}"
    subj_dir.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    saved_paths = []
    for idx, data_url in enumerate(images[:20]):
        try:
            m = _re.match(r"data:image/(\w+);base64,(.+)", data_url, _re.DOTALL)
            ext = m.group(1) if m else "png"
            b64_data = m.group(2) if m else data_url
            img_bytes = _b64.b64decode(b64_data)
            img_path = subj_dir / f"amboss_{ts}_{idx}.{ext}"
            img_path.write_bytes(img_bytes)
            saved_paths.append(str(img_path.relative_to(UPLOADS_DIR)).replace("\\", "/"))
        except Exception as e:
            print(f"[Learn/Screenshot] Fehler idx {idx}: {e}")

    print(f"[Learn/Screenshot] {len(saved_paths)} Bilder gespeichert → {subj_dir}")
    return web.json_response({"ok": True, "saved": len(saved_paths), "paths": saved_paths})


async def api_learn_upload(request):
    """POST /api/learn/upload
    Thieme-Zugang: Bilder und Quellen aus Thieme/Amboss speichern.
    Body: { subject: str, images: [dataURL, ...], source_url?: str, caption?: str, provider?: "thieme"|"amboss" }
    Speichert in UPLOADS_DIR/images/thieme_<subject>/ (oder amboss_ wenn provider=amboss).
    Optional: source_meta.json mit URL + Caption pro Bild für Quellenangabe.
    Returns: { ok, saved: int, paths: [...], provider: str }
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    subject = (body.get("subject") or "medizin").strip().lower().replace(" ", "_")
    images = body.get("images") or []
    source_url = (body.get("source_url") or "").strip()
    caption = (body.get("caption") or "").strip()
    provider = (body.get("provider") or "thieme").strip().lower()
    if provider not in ("thieme", "amboss"):
        provider = "thieme"

    if not images:
        return web.json_response({"ok": False, "error": "Keine Bilder übergeben"}, status=400)

    import base64 as _b64
    import re as _re
    import json as _jtu
    subj_dir = UPLOADS_DIR / "images" / f"{provider}_{subject}"
    subj_dir.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    saved_paths = []
    meta_entries = []

    for idx, data_url in enumerate(images[:30]):
        try:
            m = _re.match(r"data:image/(\w+);base64,(.+)", data_url, _re.DOTALL)
            ext = m.group(1) if m else "png"
            b64_data = m.group(2) if m else data_url
            img_bytes = _b64.b64decode(b64_data)
            img_name = f"{provider}_{ts}_{idx}.{ext}"
            img_path = subj_dir / img_name
            img_path.write_bytes(img_bytes)
            rel = str(img_path.relative_to(UPLOADS_DIR)).replace("\\", "/")
            saved_paths.append(rel)
            meta_entries.append({
                "path": rel,
                "source_url": source_url,
                "caption": caption or img_name,
                "provider": provider,
                "subject": subject,
                "saved_at": ts,
            })
        except Exception as e:
            print(f"[Learn/Upload] Fehler idx {idx}: {e}")

    # Quellenmetadaten persistieren für Kartengenerierung
    if meta_entries:
        meta_path = subj_dir / "source_meta.json"
        existing_meta = []
        if meta_path.exists():
            try:
                existing_meta = _jtu.loads(meta_path.read_text("utf-8"))
            except Exception:
                existing_meta = []
        existing_meta.extend(meta_entries)
        meta_path.write_text(_jtu.dumps(existing_meta, ensure_ascii=False, indent=2), "utf-8")

    print(f"[Learn/Upload] {len(saved_paths)} {provider}-Bilder gespeichert → {subj_dir}")
    return web.json_response({"ok": True, "saved": len(saved_paths), "paths": saved_paths, "provider": provider})


async def api_learn_export_zip(request):
    """POST /api/learn/export-zip
    Body: { subject: str, cards: [...] }
    Returns: .apkg-Paket mit Custom Note Types (MC-Basic + MC-Cloze) — bringt
    das Mockup-Design von card-mockup.html direkt in Anki, inkl. Image
    Occlusion (Cloze-Boxen über dem Bild), Cloze-Lücken und SVGs.

    Pfad heisst weiterhin /export-zip (Frontend-Kompat); Content ist aber .apkg.
    """
    try:
        body = await request.json()
    except Exception:
        return web.Response(status=400, text="Invalid JSON")

    subject_raw = (body.get("subject") or "karten").strip()
    subject_slug = subject_raw.lower().replace(" ", "_") or "karten"
    cards = body.get("cards") or []
    if not cards:
        return web.Response(status=400, text="Keine Karten übergeben")

    try:
        from apkg_export import build_apkg
        apkg_bytes, stats = build_apkg(cards, subject_raw, UPLOADS_DIR)
    except Exception as e:
        import traceback as _tb
        _tb.print_exc()
        return web.json_response({"ok": False, "error": f"apkg-Build fehlgeschlagen: {type(e).__name__}: {e}"}, status=500)

    print(
        f"[Learn/ExportApkg] {stats['basic']} basic + {stats['cloze']} cloze "
        f"({stats['occlusion']} occlusion) + {stats['images']} Bilder → {len(apkg_bytes)} bytes"
    )

    return web.Response(
        body=apkg_bytes,
        content_type="application/octet-stream",
        headers={
            "Content-Disposition": f'attachment; filename="{subject_slug}_anki.apkg"',
        },
    )


async def api_learn_amboss_login(request):
    """POST /api/learn/amboss-login
    Body: { cookie: str }  — der `next_auth_amboss_de`-Wert aus Browser-DevTools
    (oder der komplette Cookie-String). Wird in data/amboss_cookie_cache.json gecached.
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    cookie = (body.get("cookie") or "").strip()
    if not cookie:
        return web.json_response({"ok": False, "error": "cookie erforderlich"}, status=400)

    try:
        import aiohttp as _ahttp
        from amboss_client import _save_cookie, search_article_eids
        connector = _ahttp.TCPConnector(ssl=True)
        async with _ahttp.ClientSession(connector=connector) as sess:
            articles = await search_article_eids(sess, cookie, "Herz", limit=3)
        _save_cookie(UPLOADS_DIR, cookie)
        print(f"[Learn/AmbossLogin] Cookie akzeptiert, Testsuche 'Herz': {len(articles)} Artikel")
        return web.json_response({
            "ok": True,
            "message": f"Cookie gespeichert — Testsuche lieferte {len(articles)} Artikel",
        })
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)


async def api_learn_amboss_test(request):
    """GET /api/learn/amboss-test
    Testet Amboss-Verbindung: prüft Cookie-Cache, macht searchSuggestions-Testsuche,
    gibt Status + Artikel-Anzahl zurück. Fallback: config.json `amboss_cookie`.
    """
    from amboss_client import _load_cached_cookie, _cookie_saved_at, search_article_eids

    cookie = _load_cached_cookie(UPLOADS_DIR)
    saved_at = _cookie_saved_at(UPLOADS_DIR)
    source = "cache"
    if not cookie:
        _cfg_cookie = CONFIG.get("amboss_cookie", "").strip()
        if _cfg_cookie:
            cookie = _cfg_cookie
            source = "config"
        else:
            return web.json_response({
                "ok": False,
                "error": "Kein Cookie gecacht und keines in config.json. Bitte Cookie im Medizin-Tab einfügen.",
                "has_token": False,
                "source": None,
                "saved_at": None,
            })

    try:
        import aiohttp as _ahttp
        connector = _ahttp.TCPConnector(ssl=True)
        async with _ahttp.ClientSession(connector=connector) as sess:
            articles = await search_article_eids(sess, cookie, "Herz", limit=3)

        print(f"[Learn/AmbossTest] Testsuche 'Herz': {len(articles)} Artikel")
        return web.json_response({
            "ok": True,
            "has_token": True,
            "source": source,                  # "cache" | "config"
            "saved_at": saved_at,              # Unix-Timestamp, wann gecached
            "image_count": len(articles),
            "first_image": articles[0] if articles else None,
            "message": f"Verbindung OK — {len(articles)} Artikel für 'Herz' gefunden",
        })
    except Exception as e:
        err_s = str(e)
        if "abgelaufen" in err_s.lower() or "UNAUTH" in err_s.upper():
            return web.json_response({
                "ok": False, "has_token": False, "source": source, "saved_at": saved_at,
                "error": err_s,
            }, status=200)
        return web.json_response({
            "ok": False, "has_token": True, "source": source, "saved_at": saved_at,
            "error": err_s,
        }, status=500)


async def api_learn_viamedici_login(request):
    """POST /api/learn/viamedici-login
    Body: { cookie: str }  — Keycloak-Session-Cookies von authentication.thieme.de
    (z.B. `document.cookie` aus DevTools Console auf dieser Domain).
    Wird in data/viamedici_auth_cache.json gespeichert. Ein Silent-Refresh wird sofort
    versucht um zu prüfen ob die Cookies valid sind.
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    cookie = (body.get("cookie") or "").strip()
    if not cookie:
        return web.json_response({"ok": False, "error": "cookie erforderlich"}, status=400)

    try:
        from viamedici_client import save_kc_cookies, ensure_access_token, _normalize_kc_cookies, _jwt_exp
        normalized = _normalize_kc_cookies(cookie)
        save_kc_cookies(UPLOADS_DIR, normalized)
        # Silent-Refresh sofort triggern — liefert uns auch ein Access-Token + Refresh-Token
        try:
            access_token = await ensure_access_token(UPLOADS_DIR)
        except Exception as se:
            return web.json_response({
                "ok": False,
                "error": f"Cookies gespeichert, aber Silent-Refresh fehlgeschlagen: {se}",
            })
        exp = _jwt_exp(access_token) or 0
        remaining = max(0, exp - int(__import__("time").time()))
        print(f"[Learn/ViaMediciLogin] Silent-Refresh OK, Access-Token läuft in {remaining}s ab")
        return web.json_response({
            "ok": True,
            "message": f"Cookies OK — Access-Token frisch, gültig für {remaining // 60} Min {remaining % 60} s",
        })
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)


async def api_learn_viamedici_test(request):
    """GET /api/learn/viamedici-test
    Prüft ob Cookies da sind und macht einen Silent-Refresh um Access-Token zu refreshen.
    Gibt aktuellen Status zurück inkl. echter JWT-Expiry (KC-Identity + Refresh-Token).
    """
    from viamedici_client import (
        _load_cache, ensure_access_token, saved_at as _saved_at,
        _jwt_exp, get_session_info,
    )

    cache = _load_cache(UPLOADS_DIR)
    kc_cookies = cache.get("kc_cookies") or CONFIG.get("viamedici_cookies", "").strip()
    source = "cache" if cache.get("kc_cookies") else ("config" if kc_cookies else None)
    sess = get_session_info(UPLOADS_DIR)
    base = {
        "source": source,
        "saved_at": _saved_at(UPLOADS_DIR),
        "session_exp": sess["kc_identity_exp"],
        "refresh_exp": sess["refresh_token_exp"],
        "is_offline": sess["is_offline"],
    }
    if not kc_cookies:
        return web.json_response({
            **base,
            "ok": False,
            "has_cookies": False,
            "access_token_exp": None,
            "error": "Keine Keycloak-Cookies gecached. Bitte Cookies im Medizin-Tab einfügen.",
        })

    try:
        access_token = await ensure_access_token(UPLOADS_DIR)
        exp = _jwt_exp(access_token) or 0
        remaining = max(0, exp - int(__import__("time").time()))
        # Session-Info nach Refresh neu lesen — offline-Flag kann sich gerade erst gesetzt haben
        sess2 = get_session_info(UPLOADS_DIR)
        return web.json_response({
            **base,
            "ok": True,
            "has_cookies": True,
            "access_token_exp": exp,
            "session_exp": sess2["kc_identity_exp"],
            "refresh_exp": sess2["refresh_token_exp"],
            "is_offline": sess2["is_offline"],
            "message": f"Silent-Refresh OK — Access-Token gültig für {remaining // 60} Min {remaining % 60} s",
        })
    except Exception as e:
        return web.json_response({
            **base,
            "ok": False,
            "has_cookies": False,  # Cookies da, aber nicht mehr nutzbar
            "access_token_exp": None,
            "error": str(e),
        }, status=500)


async def _run_funnel_task(task_id: str, funnel_log_id: str = None):
    """Run a brain-ingest task via Claude CLI, no WS streaming."""
    _activity_set("funnel", active=True)
    # Mark running
    _update_task_field(task_id, status="running", started=datetime.now().isoformat())

    # Find prompt
    tasks   = _load_smart_tasks()
    task_obj = _find_task(tasks, task_id)
    if not task_obj:
        return

    prompt = task_obj.get("prompt", "") or task_obj.get("original_request", "")

    try:
        claude_exe = r"C:\Users\Robin\.local\bin\claude.exe"
        cwd        = r"C:\Users\Robin\Jarvis"
        cmd = [claude_exe, "-p", prompt,
               "--model", "sonnet",
               "--max-turns", "12",
               "--tools", "Bash,Read,Write,Edit,Glob,Grep,LS"]

        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=cwd,
            stdin=asyncio.subprocess.DEVNULL,
            **_SUBPROCESS_FLAGS,
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=300)
        output = (stdout or b"").decode("utf-8", errors="replace").strip()

        _update_task_field(task_id, status="done", completed=datetime.now().isoformat(),
                           output=output)

        # Update funnel log entry
        if funnel_log_id:
            items = _funnel_load_log()
            for it in items:
                if it.get("id") == funnel_log_id:
                    it["status"] = "done"
                    m = re.search(r"NOTE_CREATED:\s*(.+)", output)
                    if m: it["note_path"] = m.group(1).strip()
            _funnel_save_log(items)

    except asyncio.TimeoutError:
        _update_task_field(task_id, status="error", error="Timeout (5 min)")
        if funnel_log_id:
            items = _funnel_load_log()
            for it in items:
                if it.get("id") == funnel_log_id:
                    it["status"] = "error"
            _funnel_save_log(items)
    except Exception as e:
        _update_task_field(task_id, status="error", error=str(e))
        print(f"[Funnel] Task error: {e}")
    finally:
        _activity_set("funnel", active=False, pending_delta=1)


# ---------------------------------------------------------------------------
# Smart Model Routing  -  pick cheapest model that can handle the request
# ---------------------------------------------------------------------------

def _pick_model(message: str, session) -> str:
    """Pick the right model based on message complexity.

    - haiku: greetings, short questions, status checks, simple lookups
    - sonnet: code edits, explanations, medium tasks, most conversations
    - opus: deep analysis, architecture, multi-step reasoning, research

    Returns model string for --model flag.
    """
    import re
    msg = message.lower().strip()
    word_count = len(msg.split())

    # Config override: if user set a specific model, always use it
    forced = CONFIG.get("claude_model_force")
    if forced:
        return forced

    # === HAIKU: quick, simple, cheap ===
    haiku_patterns = [
        r'^(hi|hey|hallo|moin|guten (morgen|tag|abend)|wie geht)',
        r'^(ja|nein|ok|okay|danke|passt|perfekt|genau|klar|cool|nice|gut)[\s!?.]*$',
        r'^(was ist|wie heisst|wann ist|wo ist|wer ist)',
        r'^(status|zeit|uhrzeit|datum|wetter)',
        r'^(test|ping|check)[\s!?.]*$',
    ]
    if word_count <= 8 and any(re.search(p, msg) for p in haiku_patterns):
        print(f"[Router] haiku  -  short/simple ({word_count} words)")
        return "haiku"

    # === OPUS: complex, deep reasoning ===
    opus_patterns = [
        r'(analysier|analyse|architektur|refactor|design|strategie)',
        r'(erkl.r.*detail|tief.*eintauchen|ausf.hrlich)',
        r'(dissertation|doktorarbeit|paper|forschung|literatur)',
        r'(vergleich.*und.*und|trade.?off|vor.*und.*nachteil)',
        r'(debug.*komplex|multi.?file|mehrere dateien)',
        r'(plan.*erstell|roadmap|konzept.*entwickl)',
        r'(warum.*funktioniert.*nicht|root.?cause)',
    ]
    if any(re.search(p, msg) for p in opus_patterns):
        print(f"[Router] opus  -  complex task detected")
        return "opus"

    # Long messages (>100 words) or with code blocks â†' opus
    if word_count > 100 or '```' in message:
        print(f"[Router] opus  -  long/code message ({word_count} words)")
        return "opus"

    # === SONNET: default, good balance ===
    print(f"[Router] sonnet  -  default ({word_count} words)")
    return "sonnet"


# ---------------------------------------------------------------------------
# Bot Arena  -  Multi-Bot Discussion Rooms
# ---------------------------------------------------------------------------
ARENA_FILE = Path("C:/Users/Robin/Jarvis/data/arena-rooms.json")
ARENA_CRASH_DIR = Path("C:/Users/Robin/Jarvis/data/arena-crashes")
ARENA_FAIL_BACKLOG_FILE = Path("C:/Users/Robin/Jarvis/data/arena-fail-backlog.json")
ARENA_REVIEWS_LOG = Path("C:/Users/Robin/Jarvis/data/arena-reviews.jsonl")
ARENA_ARCHIVE_DIR = Path("C:/Users/Robin/Jarvis/data/arena-archive")

# Room renewal thresholds (trim oversized rooms to fight frontend lag)
ARENA_RENEW_MSG_COUNT = 120       # trim when message count exceeds this
ARENA_RENEW_PAYLOAD_KB = 300      # trim when messages payload exceeds this size
# Density trigger: renew when avg msg size is large even if total count/payload is below main thresholds
# Catches "long rooms" (few msgs, each very verbose) that still cause frontend lag
ARENA_RENEW_AVG_MSG_KB = 3.5     # avg KB per message
ARENA_RENEW_AVG_MSG_MIN_COUNT = 25  # only apply density check once room has this many messages


def _archive_and_renew_room(room: dict) -> dict:
    """Archive current messages to JSONL and reseed the room with a compact
    continuity seed (Fazit + offene ACTIONs). Returns info dict for logging.

    Mutates ``room`` in place: replaces ``messages`` with a compact seed,
    resets ``turn_count`` and keeps ``fazit_log`` / ``done_actions`` intact.
    """
    import re as _re
    messages = room.get("messages", []) or []
    room_id = room.get("id", "unknown")
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")

    # 1) Archive-Dump: jede Message eine Zeile + Meta-Header
    try:
        room_dir = ARENA_ARCHIVE_DIR / str(room_id)
        room_dir.mkdir(parents=True, exist_ok=True)
        archive_path = room_dir / f"{ts}.jsonl"
        with archive_path.open("w", encoding="utf-8") as fh:
            meta = {
                "_meta": True,
                "room_id": room_id,
                "title": room.get("title"),
                "topic": room.get("topic"),
                "archived_at": datetime.now().isoformat(),
                "msg_count": len(messages),
                "turn_count": room.get("turn_count", 0),
            }
            fh.write(json.dumps(meta, ensure_ascii=False) + "\n")
            for m in messages:
                try:
                    fh.write(json.dumps(m, ensure_ascii=False) + "\n")
                except Exception:
                    fh.write(json.dumps({"_err": "unserializable", "role": m.get("role")}) + "\n")
    except Exception as e:
        print(f"[Arena-Renew] Archive-Write failed: {e}")

    # 2) Offene [ACTION:...] Items sammeln (aus den letzten 40 Messages)
    done_keys = set(k.strip().lower() for k in room.get("done_actions", []))
    open_actions: list = []
    seen_keys: set = set()
    action_re = _re.compile(r"\[ACTION:([^\]]+)\]\s*([^\n]+)")
    for m in messages[-40:]:
        content = m.get("content", "") or ""
        for bot, body in action_re.findall(content):
            txt = body.strip().rstrip(".")
            key = f"{bot.strip()}{txt}".lower()
            if key in seen_keys or txt.lower() in done_keys:
                continue
            seen_keys.add(key)
            open_actions.append(f"[ACTION:{bot.strip()}] {txt}")
        if len(open_actions) >= 12:
            break

    # 3) Continuity-Seed (Gammas Vorschlag) als Orchestrator-System-Turn
    fazit = (room.get("orchestrator_summary") or "").strip()
    if not fazit:
        fazit_log = room.get("fazit_log") or []
        if fazit_log:
            fazit = (fazit_log[-1].get("content") or "").strip()
    fazit = fazit[:1800] if fazit else "(kein Fazit gespeichert)"

    # Konsens-Historie: letzte 3 Fazits aus renewals[] (sofern vorhanden)
    prior_renewals = [r for r in room.get("renewals", []) if r.get("fazit")][-3:]
    if prior_renewals:
        history_lines = []
        for i, r in enumerate(prior_renewals, 1):
            history_lines.append(f"[Runde -{len(prior_renewals)-i+1}] {r['fazit'][:600]}")
        history_block = "\n\n".join(history_lines)
        history_section = f"\nKONSENS-HISTORIE (letzte {len(prior_renewals)} Runden):\n{history_block}\n"
    else:
        history_section = ""

    open_block = "\n".join(open_actions) if open_actions else "(keine offenen Aktionen)"

    # [FAIL-CONTEXT] Block: letzte failed/timeout Actions als Lernmaterial für Bots
    fail_items = [
        a for a in room.get("action_items", [])
        if a.get("status") in ("failed", "timeout") and a.get("result")
    ][-3:]
    if fail_items:
        fail_lines = []
        for fa in fail_items:
            bot_tag = f"[ACTION:{fa.get('bot', '?')}]" if fa.get("bot") else "[ACTION]"
            txt = (fa.get("text") or "")[:120]
            reason = (fa.get("result") or "")[:180]
            fail_lines.append(f"{bot_tag} {txt}\n  → Fehler: {reason}")
        fail_block = "\n\n".join(fail_lines)
        fail_section = f"\nFEHLGESCHLAGENE AKTIONEN (analysiere Ursache, schlage kleinere Variante vor):\n{fail_block}\n"
    else:
        fail_section = ""

    seed_content = (
        f"RAUM WURDE AUTOMATISCH KOMPAKTIERT (Lag-Schutz)\n\n"
        f"KONTEXT AUS VORRUNDE:\n{fazit}\n"
        f"{history_section}"
        f"{fail_section}\n"
        f"OFFENE AKTIONEN:\n{open_block}"
    )

    # Beta's Token-Diff: alle [ACTION:...] aus dem gesamten Log gegen Seed prüfen.
    # Fehlende deterministisch anhängen — rettet Actions, die das LLM-Fazit verschluckt.
    all_action_keys: list = []
    seen_all: set = set()
    for m in messages:
        content = m.get("content", "") or ""
        for bot, body in action_re.findall(content):
            txt = body.strip().rstrip(".")
            key = f"{bot.strip()}{txt}".lower()
            if key in seen_all or txt.lower() in done_keys:
                continue
            seen_all.add(key)
            all_action_keys.append((bot.strip(), txt))
    seed_lower = seed_content.lower()
    missed: list = []
    for bot, txt in all_action_keys:
        if txt.lower() not in seed_lower:
            missed.append(f"[ACTION:{bot}] {txt}")
    if missed:
        seed_content += "\n\nNACHGETRAGEN (im Fazit vergessen):\n" + "\n".join(missed[:10])

    # Tag action_items: missed → recovered, surviving open → stale
    _now_iso = datetime.now().isoformat()
    missed_texts_lower = set()
    for _m in missed:
        _part = _m.split("] ", 1)[-1].strip().lower() if "] " in _m else _m.lower()
        missed_texts_lower.add(_part)
    for _item in room.get("action_items", []):
        if _item.get("status") == "offen":
            if _item["text"].lower() in missed_texts_lower:
                _item["recovered"] = True
                _item["recovered_at"] = _now_iso
                _item.pop("stale", None)
            else:
                _item["stale"] = True
                _item.pop("recovered", None)

    seed_msg = {
        "role": "orchestrator",
        "name": "Orchestrator",
        "content": seed_content,
        "timestamp": datetime.now().isoformat(),
        "color": "#a855f7",
        "turn": 0,
        "is_renewal_seed": True,
    }

    old_count = len(messages)
    room["messages"] = [seed_msg]
    room.setdefault("renewals", []).append({
        "at": datetime.now().isoformat(),
        "ts": ts,
        "old_msg_count": old_count,
        "open_actions": len(open_actions),
        "recovered_actions": len(missed),
        "archive": str(archive_path) if 'archive_path' in locals() else None,
        "fazit": fazit[:2000],
    })

    return {
        "archived_msgs": old_count,
        "open_actions": len(open_actions),
        "recovered": len(missed),
        "ts": ts,
    }

_arena_rooms: dict = {}  # room_id -> ArenaRoom state


def _extract_and_store_action_items(room: dict, fazit_text: str):
    """Extract [ACTION:Bot] tags from fazit and upsert into room['action_items'] with offen/erledigt."""
    action_re = re.compile(r"\[ACTION:([^\]]+)\]\s*([^\n]+)")
    existing = {item["key"]: item for item in room.get("action_items", [])}
    done_keys = set(k.strip().lower() for k in room.get("done_actions", []))
    now = datetime.now().isoformat()
    for bot, body in action_re.findall(fazit_text):
        bot = bot.strip()
        text = body.strip().rstrip(".")
        key = f"{bot}:{text}".lower()
        status = "erledigt" if text.lower() in done_keys else "offen"
        if key not in existing:
            existing[key] = {"key": key, "bot": bot, "text": text, "status": status, "created_at": now}
        else:
            existing[key]["status"] = status
    room["action_items"] = list(existing.values())


def _log_arena_crash(room_id: str, bot_name: str, turn: int, kind: str,
                     cmd: list, returncode, stderr: str, stdout_tail: str,
                     timeout_s: int, exception: str = "") -> dict:
    """Persist crash details to data/arena-crashes/ and return a short summary dict.

    kind: one of 'timeout', 'exception', 'empty_response', 'budget_exhausted'.
    Returns dict with 'file' (path str), 'summary' (one-line str) for room messaging.
    """
    try:
        ARENA_CRASH_DIR.mkdir(parents=True, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_bot = "".join(c if c.isalnum() else "_" for c in bot_name)[:30]
        fname = f"{ts}_{room_id[:8]}_{safe_bot}_t{turn}_{kind}.json"
        fpath = ARENA_CRASH_DIR / fname
        payload = {
            "timestamp": datetime.now().isoformat(),
            "room_id": room_id,
            "bot_name": bot_name,
            "turn": turn,
            "kind": kind,
            "cmd": cmd,
            "returncode": returncode,
            "timeout_s": timeout_s,
            "exception": exception,
            "stderr": stderr or "",
            "stdout_tail": stdout_tail or "",
        }
        fpath.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        # Short summary for the arena room
        err_snip = (stderr or "").strip().splitlines()[-1][:200] if (stderr or "").strip() else ""
        if kind == "timeout":
            summary = f"⏱️ Timeout nach {timeout_s}s"
        elif kind == "exception":
            summary = f"💥 Exception: {exception[:200]}"
        elif kind == "budget_exhausted":
            summary = f"📋 Schritte-Budget aufgebraucht (max_turns erreicht, exit={returncode})"
        elif kind == "thinking_budget_exhausted":
            summary = f"🧠 Denk-Budget erschöpft (Thinking-Tokens aufgebraucht, exit={returncode})"
        else:
            summary = f"🕳️ Echte Leer-Antwort (stdout leer, kein bekannter Fehler, exit={returncode})"
        if err_snip:
            summary += f" — stderr: {err_snip}"
        return {"file": str(fpath), "summary": summary, "payload": payload}
    except Exception as log_err:
        print(f"[Arena] CRASH-LOG FAILED: {log_err}")
        return {"file": "", "summary": f"Crash ({kind}) — Log fehlgeschlagen: {log_err}", "payload": {}}


def _load_arena_rooms() -> list:
    if ARENA_FILE.exists():
        try:
            return json.loads(ARENA_FILE.read_text(encoding="utf-8"))
        except Exception:
            return []
    return []


def _save_arena_rooms(rooms: list):
    ARENA_FILE.parent.mkdir(parents=True, exist_ok=True)
    ARENA_FILE.write_text(
        json.dumps(rooms, ensure_ascii=False, indent=2), encoding="utf-8"
    )


def _append_fail_backlog(entries: list):
    """Append failed/rejected actions to the persistent fail backlog.

    Dedup: Wenn (approval_id, action_num) schon existiert, wird kein neuer Eintrag
    angefuegt — stattdessen requeue_count auf dem bestehenden Eintrag hochgezaehlt.
    Fehlender approval_id (bei reinen fail/reject-Eintraegen) umgeht das Matching.
    """
    if not entries:
        return
    try:
        ARENA_FAIL_BACKLOG_FILE.parent.mkdir(parents=True, exist_ok=True)
        existing = []
        if ARENA_FAIL_BACKLOG_FILE.exists():
            existing = json.loads(ARENA_FAIL_BACKLOG_FILE.read_text(encoding="utf-8"))

        for entry in entries:
            aid = entry.get("approval_id")
            anum = entry.get("action_num")
            match = None
            if aid and anum is not None:
                for e in existing:
                    if e.get("approval_id") == aid and e.get("action_num") == anum:
                        match = e
                        break
            if match is not None:
                match["requeue_count"] = int(match.get("requeue_count", 1)) + 1
                match["last_requeue_ts"] = entry.get("timestamp")
                match["reason"] = entry.get("reason", match.get("reason"))
            else:
                entry.setdefault("requeue_count", 1)
                existing.append(entry)

        ARENA_FAIL_BACKLOG_FILE.write_text(
            json.dumps(existing, ensure_ascii=False, indent=2), encoding="utf-8"
        )
    except Exception as e:
        print(f"[Arena] fail-backlog write error: {e}")


def _get_requeue_count(approval_id: str, action_num: int) -> int:
    """Liest aktuellen requeue_count fuer (approval_id, action_num) aus dem Backlog."""
    try:
        if not ARENA_FAIL_BACKLOG_FILE.exists():
            return 0
        existing = json.loads(ARENA_FAIL_BACKLOG_FILE.read_text(encoding="utf-8"))
        for e in existing:
            if e.get("approval_id") == approval_id and e.get("action_num") == action_num:
                return int(e.get("requeue_count", 1))
    except Exception:
        pass
    return 0


async def api_arena_fail_backlog(request):
    """GET /api/arena/fail-backlog  — returns persisted failed/rejected actions."""
    try:
        if ARENA_FAIL_BACKLOG_FILE.exists():
            data = json.loads(ARENA_FAIL_BACKLOG_FILE.read_text(encoding="utf-8"))
        else:
            data = []
        return web.json_response({"ok": True, "items": data})
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)


async def api_arena_fail_backlog_clear(request):
    """DELETE /api/arena/fail-backlog  — clears the backlog after user acknowledges."""
    try:
        if ARENA_FAIL_BACKLOG_FILE.exists():
            ARENA_FAIL_BACKLOG_FILE.write_text("[]", encoding="utf-8")
        return web.json_response({"ok": True})
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)}, status=500)


def _requeue_nachbessern_actions(room: dict, reviewer: str, room_id: str, turn: int, action_nums: list = None):
    """Put NACHBESSERN-flagged actions back into the fail-backlog so the Orchestrator can retry them.

    If action_nums is given (1-based indices from 'URTEIL: NACHBESSERN #N'), only those
    specific actions are requeued. Otherwise all actions from the last approval are requeued.
    """
    last_approval = None
    for m in reversed(room.get("messages", [])):
        if m.get("role") == "approval" and m.get("actions"):
            last_approval = m
            break
    if not last_approval:
        return

    topic = room.get("topic", "")
    approval_id = last_approval.get("approval_id", "")
    all_actions = last_approval.get("actions", [])
    if action_nums:
        selected = [(n, all_actions[n - 1]) for n in action_nums if 1 <= n <= len(all_actions)]
    else:
        selected = [(i + 1, a) for i, a in enumerate(all_actions)]

    entries = []
    capped = []
    for num, a in selected:
        text = (a.get("text") or "").strip()
        if not text:
            continue
        prior = _get_requeue_count(approval_id, num) if approval_id else 0
        if prior >= 2:
            capped.append(num)
            print(f"[Arena] NACHBESSERN CAP: Action #{num} (approval={approval_id[:8] if approval_id else '-'}) bereits {prior}x requeued → FAIL, nicht neu angelegt")
            continue
        entries.append({
            "kind": "nachbessern",
            "action": text,
            "action_num": num,
            "reason": f"URTEIL: NACHBESSERN #{num} von {reviewer} (Turn {turn})",
            "room_id": room_id,
            "topic": topic,
            "approval_id": approval_id,
            "timestamp": datetime.now().isoformat(),
        })
    if entries:
        _append_fail_backlog(entries)
        print(f"[Arena] NACHBESSERN: {len(entries)} Aktion(en) zurück in Backlog (reviewer={reviewer}, room={room_id[:8]}, nums={action_nums or 'all'})")
    if capped:
        print(f"[Arena] NACHBESSERN: {len(capped)} Aktion(en) durch Cap blockiert (nums={capped})")
        for cap_num in capped:
            room["messages"].append({
                "role": "system",
                "name": "System",
                "content": f"🛑 Action #{cap_num} nach 2 Nachbesserungsrunden gestoppt — kein weiterer Retry.",
                "timestamp": datetime.now().isoformat(),
                "color": "#ef4444",
            })


def _get_arena_room(rooms, room_id):
    for r in rooms:
        if r["id"] == room_id:
            return r
    return None


ARENA_STALE_THRESHOLD_S = 300  # 5 min Stille → kein Blind-Revival


def _room_is_stale(room: dict) -> bool:
    """True wenn letzte Aktivität im Raum älter als 5 Minuten ist."""
    msgs = room.get("messages", [])
    last_ts_str = msgs[-1].get("timestamp", "") if msgs else ""
    if not last_ts_str:
        last_ts_str = room.get("last_active") or room.get("created", "")
    if not last_ts_str:
        return False
    try:
        last_ts = datetime.fromisoformat(last_ts_str)
        return (datetime.now() - last_ts).total_seconds() > ARENA_STALE_THRESHOLD_S
    except Exception:
        return False


async def api_arena_rooms(request):
    """GET /api/arena/rooms  -  List all arena rooms."""
    rooms = _load_arena_rooms()
    # Add live status
    for r in rooms:
        rid = r["id"]
        r["running"] = rid in _arena_rooms and _arena_rooms[rid].get("running", False)
    return web.json_response(rooms)


async def api_arena_create(request):
    """POST /api/arena/rooms  -  Create a new arena room."""
    body = await request.json()
    room = {
        "id": str(uuid.uuid4()),
        "title": body.get("title", "Neuer Raum"),
        "topic": body.get("topic", ""),
        "category": body.get("category", "general"),
        "bots": body.get("bots", [
            {"name": "Alpha", "model": "opus", "color": "#6366f1", "persona": "Du bist ein analytischer Denker. Argumentiere logisch und praezise. Erkenne gute Punkte der anderen an."},
            {"name": "Beta", "model": "opus", "color": "#22c55e", "persona": "Du bist ein kreativer Querdenker. Bringe unerwartete Perspektiven ein, aber erkenne auch gute Ideen an."},
            {"name": "Gamma", "model": "sonnet", "color": "#f59e0b", "persona": "Du bist ein pragmatischer Umsetzer. Fokussiere auf Machbarkeit und konkrete naechste Schritte."},
            {"name": "Delta", "model": "sonnet", "color": "#ec4899", "persona": "Du bist ein Vermittler. Finde Gemeinsamkeiten und schlage Kompromisse vor."},
        ]),
        "messages": [],
        "max_turns": body.get("max_turns", 20),
        "max_rounds": max(1, min(25, int(body.get("max_rounds", 8)))),
        "summarize_every": body.get("summarize_every", 10),
        "allow_read": body.get("allow_read", True),
        "created": datetime.now().isoformat(),
        "last_active": datetime.now().isoformat(),
        "status": "idle",
    }
    rooms = _load_arena_rooms()
    rooms.append(room)
    _save_arena_rooms(rooms)
    return web.json_response(room, status=201)


async def api_arena_delete(request):
    """DELETE /api/arena/rooms/{id}  -  Delete a room."""
    room_id = request.match_info["id"]
    # Stop if running
    if room_id in _arena_rooms:
        _arena_rooms[room_id]["running"] = False
    rooms = _load_arena_rooms()
    rooms = [r for r in rooms if r["id"] != room_id]
    _save_arena_rooms(rooms)
    return web.json_response({"ok": True})


async def api_arena_room(request):
    """GET /api/arena/rooms/{id}  -  Get room details (includes snapshots index)."""
    room_id = request.match_info["id"]
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")
    room["running"] = room_id in _arena_rooms and _arena_rooms[room_id].get("running", False)
    room["snapshots"] = _list_room_snapshots(room_id)
    return web.json_response(room)


def _arena_task_running(room_id: str) -> bool:
    """True wenn entweder ein Trialog-Task oder ein klassischer Arena-Flow
    für diesen Room aktiv ist. Vereint beide Registries + säubert stale Flags.
    """
    t_task = _trialog_tasks.get(room_id)
    if t_task and not t_task.done():
        return True
    if t_task and t_task.done():
        _trialog_tasks.pop(room_id, None)
        if room_id in _arena_rooms:
            _arena_rooms[room_id]["running"] = False
    state = _arena_rooms.get(room_id) or {}
    return bool(state.get("running"))


def _spawn_for_room(room_id: str, rooms_list=None):
    """Zentraler Dispatcher: spawn _run_trialog für trialog-Räume, sonst
    _run_arena. Setzt `_arena_rooms[room_id].running=True` und trackt
    trialog-Tasks in `_trialog_tasks` damit inject/stop/start-Checks
    nicht parallel zu laufenden Flows einen zweiten Loop starten.
    """
    if rooms_list is None:
        rooms_list = _load_arena_rooms()
    room = _get_arena_room(rooms_list, room_id)
    if not room:
        return None
    _arena_rooms[room_id] = {"running": True, "process": None}
    if (room.get("mode") or "").lower() == "trialog":
        task = asyncio.create_task(_run_trialog(room_id, rooms_list))
        _trialog_tasks[room_id] = task
        return task
    return asyncio.create_task(_run_arena(room_id))


async def api_arena_start(request):
    """POST /api/arena/rooms/{id}/start  -  Start bot discussion."""
    room_id = request.match_info["id"]
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")
    if _arena_task_running(room_id):
        return web.json_response({"error": "Already running"}, status=409)

    body = await request.json() if request.content_length else {}
    if body.get("topic"):
        room["topic"] = body["topic"]
        _save_arena_rooms(rooms)

    room["last_active"] = datetime.now().isoformat()
    _save_arena_rooms(rooms)
    _spawn_for_room(room_id, rooms)
    return web.json_response({"ok": True, "room_id": room_id})


async def api_arena_stop(request):
    """POST /api/arena/rooms/{id}/stop  -  Stop bot discussion."""
    room_id = request.match_info["id"]
    if room_id in _arena_rooms:
        _arena_rooms[room_id]["running"] = False
        proc = _arena_rooms[room_id].get("process")
        if proc and proc.returncode is None:
            proc.kill()
    # Zusätzlich Trialog-Task canceln falls das ein Trialog-Room ist
    t_task = _trialog_tasks.get(room_id)
    if t_task and not t_task.done():
        t_task.cancel()
    # Status im Room auf stopped markieren damit UI es sieht
    try:
        rooms = _load_arena_rooms()
        rm = _get_arena_room(rooms, room_id)
        if rm and rm.get("mode") == "trialog":
            rm["status"] = "stopped"
            _save_arena_rooms(rooms)
    except Exception:
        pass
    return web.json_response({"ok": True})


async def api_arena_inject(request):
    """POST /api/arena/rooms/{id}/inject  -  Robin injects a message into the discussion."""
    room_id = request.match_info["id"]
    body = await request.json()
    message = body.get("message", "").strip()
    image = (body.get("image") or "").strip()  # legacy single-image filename in data/uploads/
    attachments = body.get("attachments") or []  # list of filenames in data/uploads/
    # Backwards-compat: if only legacy `image` was sent, fold it into attachments
    if image and image not in attachments:
        attachments = [image] + list(attachments)
    # Sanitize: only keep strings, drop empties, dedupe
    attachments = [a for a in (str(x).strip() for x in attachments) if a]
    seen = set(); attachments = [a for a in attachments if not (a in seen or seen.add(a))]

    if not message and not attachments:
        return web.json_response({"error": "Empty message"}, status=400)

    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")

    placeholder = "(Datei eingeworfen)" if attachments else ""
    msg = {
        "role": "robin",
        "name": "Robin",
        "content": message or placeholder,
        "timestamp": datetime.now().isoformat(),
        "color": "#f59e0b",
    }
    if attachments:
        msg["attachments"] = attachments
        # Keep `image` for first image so existing frontend rendering still shows a preview.
        first_img = next((a for a in attachments if Path(a).suffix.lower() in
                          (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp")), "")
        if first_img:
            msg["image"] = first_img
    room["messages"].append(msg)
    _save_arena_rooms(rooms)

    # Broadcast to all WS clients
    await _broadcast_all_ws({
        "type": "arena.message",
        "room_id": room_id,
        "message": msg,
    })

    # Trialog-Room: User-Message triggert automatisch eine neue Runde —
    # alle Bots antworten nacheinander. Nur wenn aktuell KEIN Task läuft
    # (weder Trialog- noch klassischer Arena-Flow, sonst Dopplung).
    # WICHTIG: max_rounds NICHT persistent auf 1 setzen — das würde den
    # konfigurierten Wert im Room für alle folgenden Starts kaputtmachen.
    # Stattdessen in-memory Override (_trialog_round_override) für genau
    # diesen Inject-getriggerten Run.
    if room.get("mode") == "trialog" and not _arena_task_running(room_id):
        print(f"[Trialog] inject-triggered round for {room_id[:8]} (override max_rounds=1)")
        _trialog_round_override[room_id] = 1
        _spawn_for_room(room_id)

    return web.json_response({"ok": True})


# ---------------------------------------------------------------------------
# Multi-Model Trialog — provider-agnostische Arena-Variante.
#
# Im Gegensatz zum klassischen _run_arena (Claude-CLI-only, mit Tools, Escalation,
# Crash-Recovery) ist dies eine schlanke Round-Robin-Runde: jeder Bot bekommt
# round-weise das Wort und antwortet über invoke_llm mit seinem eigenen
# Provider + Model. Text-only, keine Tools, keine Escalation — gedacht für
# "reden miteinander" ohne agentische Sub-Tasks.
# ---------------------------------------------------------------------------

_trialog_tasks: dict[str, asyncio.Task] = {}
# Einmalige Overrides pro Run (nicht persistent im Room), z.B. inject-triggered
# rounds wollen nur 1 Runde statt der im Room konfigurierten max_rounds.
_trialog_round_override: dict[str, int] = {}

# Rollen-Inkompatibilität: jeder Bot kriegt pro Position eine feste Rolle.
# Zwingt echten Dissens statt "ich stimme zu, ergänze X"-Ornamente.
# Position im `bots`-Array wird modulo genommen → bei 3 Bots: Kritiker/Simplifier/Umsetzer.
_TRIALOG_ROLES = [
    (
        "Kritiker",
        "Dein Job: WIDERSPRECHEN. Such dir einen konkreten Punkt aus dem eben Gesagten und "
        "widerleg ihn oder zeig was übersehen wurde. \"Ich stimme zu, ergänze X\" ist VERBOTEN — "
        "das hilft der Diskussion nicht. Wenn du wirklich keinen Gegenpunkt findest, sag das "
        "ehrlich und frag eine scharfe Rückfrage. Keine Floskeln-Zustimmung.",
    ),
    (
        "Simplifier",
        "Dein Job: SCOPE REDUZIEREN. Bei jedem Vorschlag der anderen: was kann weg? Was ist "
        "eigentlich Sprint 3 statt Sprint 1? Welches File gehört NICHT dazu? Du schlägst "
        "keine neuen Features vor — nur streichen, zusammenlegen, vertagen. Jedes \"ja, und "
        "zusätzlich...\" ist ein Fehler in deiner Rolle.",
    ),
    (
        "Umsetzer",
        "Dein Job: KONKRET WERDEN. Wenn die anderen abstrakt reden, forderst du Code oder "
        "Diff. Du selbst schlägst nur vor was du als tatsächlichen Datei-Inhalt oder ```diff "
        "sofort ausformulieren kannst. Keine Architektur-Prosa, keine Namespacing-Debatten. "
        "Wenn du etwas baust: zeig die ersten 10 Zeilen Code.",
    ),
]


async def _gather_repo_facts(topic: str, recent_lines: list[str]) -> str:
    """Forced Read-Before-Talk: sucht im Repo-RAG nach Fakten zum Topic + letzten
    Peer-Messages. Bots sehen was es im Repo schon gibt bevor sie spekulieren.
    Gibt formatierten Block zurück (leer wenn RAG unavailable oder nix gefunden).
    """
    try:
        from rag.search import query as _rag_query
    except Exception:
        return ""
    # Query: Topic + letzte 2 peer-lines (zusammen max ~400 Zeichen)
    parts = [topic[:200]]
    for ln in (recent_lines or [])[-2:]:
        parts.append(ln[:150])
    query_text = " ".join(parts).strip()
    if not query_text:
        return ""
    try:
        loop = asyncio.get_event_loop()
        hits = await loop.run_in_executor(None, lambda: _rag_query(query_text, 4))
    except Exception:
        return ""
    if not hits:
        return ""
    lines = []
    for h in hits[:4]:
        path = h.get("source") or h.get("path") or h.get("file") or "?"
        snip = (h.get("text") or h.get("content") or h.get("chunk") or "")
        snip = snip.replace("\n", " ").strip()[:200]
        if snip:
            lines.append(f"- `{path}`: {snip}")
    return "\n".join(lines) if lines else ""


async def _run_trialog(room_id: str, rooms_list: list) -> None:
    """Round-Robin Discussion zwischen 3+ Bots mit je eigenem Provider/Model.

    rooms_list ist der shared _load_arena_rooms-Resultat damit wir
    Änderungen persistieren können.
    """
    # Room aus Liste rausfischen
    room = next((r for r in rooms_list if r.get("id") == room_id), None)
    if not room:
        print(f"[Trialog] room {room_id} verschwunden")
        return

    bots = room.get("bots") or []
    if len(bots) < 2:
        print(f"[Trialog] zu wenig bots ({len(bots)})")
        return

    # Einmaliger Override (z.B. inject-triggered "nur 1 Runde"), fällt sonst
    # auf den im Room konfigurierten Wert zurück
    _override = _trialog_round_override.pop(room_id, None)
    max_rounds = int(_override) if _override else int(room.get("max_rounds", 3))
    topic = room.get("topic") or room.get("title") or ""

    room["status"] = "running"
    room["last_active"] = datetime.now().isoformat()
    _save_arena_rooms(rooms_list)

    await _broadcast_all_ws({
        "type": "arena.status",
        "room_id": room_id,
        "status": "running",
        "mode": "trialog",
    })

    turn_counter = 0
    try:
        for round_idx in range(max_rounds):
            # Check if stopped externally
            fresh = _load_arena_rooms()
            rr = next((r for r in fresh if r.get("id") == room_id), None)
            if not rr or rr.get("status") == "stopped":
                print(f"[Trialog] room {room_id[:8]} externally stopped")
                break

            for bot in bots:
                turn_counter += 1
                bot_name = bot.get("name") or f"Bot-{turn_counter}"
                bot_color = bot.get("color") or "#6366f1"
                bot_provider = bot.get("provider") or "claude-cli"
                bot_model = bot.get("model") or ""
                bot_persona = bot.get("persona") or "Du bist ein hilfreicher Diskussionspartner."

                await _broadcast_all_ws({
                    "type": "arena.thinking",
                    "room_id": room_id,
                    "bot_name": bot_name,
                    "bot_color": bot_color,
                    "turn": turn_counter,
                })

                # --- Historie bauen wie in klassischer Arena: NUR echte Bot-
                # Beiträge + Robin-Zwischenrufe + (wenn vorhanden) das letzte
                # Orchestrator-Fazit. System-Gemecker (Health-Check, skip-tot)
                # wird rausgefiltert damit Bots sich auf den Diskurs
                # konzentrieren statt auf Infrastruktur-Rauschen.
                _all_msgs = (room.get("messages") or [])
                last_orch_line = ""
                for _m in reversed(_all_msgs):
                    if _m.get("role") == "orchestrator":
                        _o = (_m.get("content") or _m.get("text") or "").strip()
                        if _o:
                            import re as _re
                            _o_clean = _re.sub(r"\[VERDICT:[^\]]+\]", "", _o).strip()
                            last_orch_line = _o_clean[:900]
                        break
                peer_lines: list[str] = []
                for m in _all_msgs[-14:]:
                    role = (m.get("role") or "").lower()
                    if role not in ("bot", "robin", "user"):
                        continue
                    from_bot = (
                        m.get("bot_name") or m.get("name") or
                        ("Robin" if role in ("robin", "user") else "Bot")
                    )
                    text = (m.get("text") or m.get("content") or "").strip()
                    if not text:
                        continue
                    peer_lines.append(f"{from_bot}: {text}")
                # Auf ~8 Einträge kappen (gleiche Dichte wie classic Arena)
                peer_lines = peer_lines[-8:]
                history_block = "\n\n".join(peer_lines) or "(noch keine Beiträge)"

                orch_hint = ""
                if last_orch_line:
                    orch_hint = (
                        f"\nLETZTES ORCHESTRATOR-FAZIT (darauf aufbauen, nicht wiederholen):\n"
                        f"{last_orch_line}\n"
                    )

                # Frischer Robin-Zwischenruf: user/robin-Message NACH letztem Bot
                _last_bot_ts = ""
                for _m in reversed(_all_msgs):
                    if _m.get("role") == "bot":
                        _last_bot_ts = _m.get("timestamp", "")
                        break
                robin_hint = ""
                _fresh = [
                    (m.get("content") or m.get("text") or "").strip()
                    for m in _all_msgs
                    if m.get("role") in ("robin", "user")
                    and m.get("timestamp", "") > _last_bot_ts
                    and (m.get("content") or m.get("text") or "").strip()
                ]
                if _fresh:
                    robin_hint = (
                        "\n=== ROBINS ZWISCHENRUF (direkte Anweisung — Priorität!) ===\n"
                        + "\n\n".join(f"[Robin: {t[:500]}]" for t in _fresh[-3:])
                        + "\n=== ENDE ZWISCHENRUF ===\n"
                    )

                # --- Rollen-Inkompatibilität (Anti-Theater Fix 1) -----------
                # Position im bots-Array → feste Rolle (Kritiker/Simplifier/Umsetzer).
                # Zwingt echten Dissens statt paralleler Monologe.
                _pos = next((i for i, b in enumerate(bots) if b.get("name") == bot_name), 0)
                _role_name, _role_task = _TRIALOG_ROLES[_pos % len(_TRIALOG_ROLES)]
                role_block = (
                    f"\n=== DEINE ROLLE DIESE RUNDE: {_role_name.upper()} ===\n"
                    f"{_role_task}\n"
                    f"=== ENDE ROLLE ===\n"
                )

                # --- Force-Diff-Modus (Anti-Theater Fix 2) ------------------
                # Ab dem 1. Orchestrator-Fazit: Ideensammlung ist vorbei, ab jetzt
                # müssen Vorschläge Code/Diff enthalten damit sie approvable sind.
                force_diff = bool((room.get("orchestrator_summary") or "").strip())
                diff_block = ""
                if force_diff:
                    diff_block = (
                        "\n=== FORCE-DIFF-MODUS AKTIV ===\n"
                        "Der Orchestrator hat bereits mindestens ein Fazit gezogen — die reine "
                        "Ideensammlung ist vorbei. Wenn du eine neue Umsetzung vorschlägst, "
                        "zeig sie als Code-Block mit ```diff oder ```python (o.ä.) mit dem "
                        "tatsächlichen Datei-Inhalt/Patch — Abstrakt 'wir könnten X bauen' "
                        "wird nicht mehr approved. Wenn du keinen Patch liefern kannst, "
                        "widersprich/frag zurück — aber schlag nichts Abstraktes vor.\n"
                        "=== ENDE FORCE-DIFF ===\n"
                    )

                # --- Forced Read-Before-Talk (Anti-Theater Fix 3) -----------
                # RAG-Suche im Repo bevor der Bot spekuliert. Heilt die "diskutiert
                # über Code den sie nicht gelesen haben"-Pathologie.
                repo_facts_block = ""
                try:
                    _facts = await _gather_repo_facts(topic, peer_lines[-3:] if peer_lines else [])
                    if _facts:
                        repo_facts_block = (
                            "\n=== IST-STAND IM REPO (echte Fakten aus dem Code — lies bevor du spekulierst) ===\n"
                            f"{_facts}\n"
                            "Wenn deine Argumentation dem oben widerspricht: erklär warum konkret. "
                            "Schlag NICHTS vor was laut den obigen Fakten schon existiert.\n"
                            "=== ENDE IST-STAND ===\n"
                        )
                except Exception as _e:
                    print(f"[Trialog] repo facts fetch failed: {_e}")

                user_prompt = (
                    f"THEMA: {topic}\n"
                    f"{orch_hint}"
                    f"{robin_hint}\n"
                    f"{repo_facts_block}"
                    f"{role_block}"
                    f"{diff_block}"
                    f"\nBISHERIGE DISKUSSION (die anderen Bots — lies sie und geh konkret darauf ein):\n"
                    f"{history_block}\n\n"
                    f"Du bist jetzt dran, {bot_name} (Runde {round_idx + 1}/{max_rounds}). "
                    f"Halte dich strikt an deine Rolle oben. Zitier einen Namen, bau auf einem "
                    f"konkreten Punkt auf oder widersprich gezielt. 2-4 Sätze, keine Floskeln, "
                    f"kein \"{bot_name}: \"-Prefix."
                )

                system_prompt = (
                    f"Du bist {bot_name} in einem Multi-Model-Trialog von Robin Mandel. "
                    f"Persona: {bot_persona}\n\n"
                    f"SPRACHE: Antworte IMMER auf Deutsch. Niemals Englisch — auch nicht "
                    f"teilweise, auch nicht 'nur ein Satz'. Robin schreibt Deutsch, "
                    f"die anderen Bots schreiben Deutsch, du ebenso.\n\n"
                    f"Du redest mit den anderen Bots, nicht mit Robin. Keine "
                    f"\"Ich als {bot_name}...\"-Floskeln. Greife konkrete Aussagen der "
                    f"anderen auf (mit Namen), statt parallel monolog-artig zu senden."
                )

                # --- Mode-Decision (Phase 2 + 3) ---
                # Feldnamen analog classic Arena damit existing /execute + /read
                # Endpoints auch für Trialog-Räume funktionieren.
                fresh_room_check = _load_arena_rooms()
                rr_now = _get_arena_room(fresh_room_check, room_id)
                is_exec = bool(rr_now and rr_now.get("tools_enabled") and
                               int(rr_now.get("execution_turns_left", 0)) > 0)
                is_read = bool(rr_now and not is_exec and
                               rr_now.get("read_enabled") and
                               int(rr_now.get("read_turns_left", 0)) > 0)

                # --- Phase 4 Cross-Review: Bot der zuletzt executed hat,
                # darf die nächste Runde nicht direkt selbst reviewen.
                # Wenn es so passieren würde → swap zum nächsten alive Bot.
                _last_executor = rr_now.get("last_executing_bot") if rr_now else None
                _cross_pending = rr_now.pop("cross_review_pending", False) if rr_now else False
                if _cross_pending and _last_executor and bot_name == _last_executor:
                    _candidates = [b for b in bots if b.get("name") != _last_executor]
                    if _candidates:
                        bot = _candidates[0]
                        bot_name = bot.get("name", bot_name)
                        bot_provider = (bot.get("provider") or "claude-cli").strip()
                        bot_model = (bot.get("model") or "").strip()
                        bot_persona = bot.get("persona") or bot_persona
                        bot_color = bot.get("color") or bot_color
                        # System-/User-Prompt neu mit dem getauschten Bot bauen
                        system_prompt = (
                            f"Du bist {bot_name} in einem Multi-Model-Trialog. "
                            f"Persona: {bot_persona}. Du wirst jetzt gebeten die "
                            f"Arbeit von {_last_executor} kritisch zu reviewen — "
                            f"konstruktiv aber ehrlich. Was ist solide? Was fehlt?"
                        )
                    _save_arena_rooms(fresh_room_check)

                t0 = datetime.now()
                err = None
                answer = ""
                usage = {}

                if is_exec:
                    exec_system = (
                        system_prompt + "\n\n"
                        "EXEC-MODUS: Du darfst in diesem Turn Dateien lesen und schreiben. "
                        "Setze konkrete Änderungen um wenn das die Diskussion weiterbringt. "
                        "Berichte am Ende kurz was du gemacht hast."
                    )
                    try:
                        answer, usage, err = await _trialog_exec_turn(
                            bot=bot, prompt_text=user_prompt,
                            system_prompt=exec_system, read_only=False,
                        )
                    except Exception as e:
                        err = f"exec_dispatch: {type(e).__name__}: {e}"
                    rr_now["execution_turns_left"] = int(rr_now.get("execution_turns_left", 0)) - 1
                    if rr_now["execution_turns_left"] <= 0:
                        rr_now["tools_enabled"] = False
                        rr_now["execution_turns_left"] = 0
                    # Nach Exec-Turn: cross-review für nächsten Bot markieren
                    rr_now["last_executing_bot"] = bot_name
                    rr_now["cross_review_pending"] = True
                    _save_arena_rooms(fresh_room_check)

                elif is_read:
                    read_system = (
                        system_prompt + "\n\n"
                        "READ-MODUS: Du darfst in diesem Turn Dateien LESEN und suchen "
                        "(Read/Glob/Grep oder plan-only), aber NICHT schreiben/ändern. "
                        "Nutze das um deine Argumente mit Fakten aus Robins Codebase zu "
                        "belegen statt zu spekulieren."
                    )
                    try:
                        answer, usage, err = await _trialog_exec_turn(
                            bot=bot, prompt_text=user_prompt,
                            system_prompt=read_system, read_only=True,
                        )
                    except Exception as e:
                        err = f"read_dispatch: {type(e).__name__}: {e}"
                    rr_now["read_turns_left"] = int(rr_now.get("read_turns_left", 0)) - 1
                    if rr_now["read_turns_left"] <= 0:
                        rr_now["read_enabled"] = False
                        rr_now["read_turns_left"] = 0
                    _save_arena_rooms(fresh_room_check)

                else:
                    # Normale Discussion: reiner LLM-Turn, keine Tools.
                    # MCP/Tools gibts in EXEC- und READ-Mode (siehe oben).
                    try:
                        answer, usage, err = await _llm_oneshot(
                            provider=bot_provider, model=bot_model,
                            user_message=user_prompt, system=system_prompt,
                            max_tokens=500,
                        )
                    except Exception as e:
                        err = f"{type(e).__name__}: {e}"

                # --- Phase 4 Escalation: bei Fehler/Leer auf stärkeres Modell ---
                # Nur für claude-cli (andere Provider haben kein klares Modell-Tier-
                # Mapping). haiku → sonnet → opus.
                _ESC_MAP_CLAUDE = {"haiku": "sonnet", "sonnet": "opus", "": "opus"}
                needs_escalation = (not answer or err) and bot_provider == "claude-cli"
                if needs_escalation:
                    escalated = _ESC_MAP_CLAUDE.get(bot_model, None)
                    if escalated and escalated != bot_model:
                        print(f"[Trialog] escalating {bot_name} {bot_model!r} → {escalated!r} (err={err})")
                        try:
                            if is_exec:
                                answer, usage, err = await _trialog_exec_turn(
                                    bot={**bot, "model": escalated},
                                    prompt_text=user_prompt,
                                    system_prompt=system_prompt,
                                    read_only=False,
                                )
                            elif is_read:
                                answer, usage, err = await _trialog_exec_turn(
                                    bot={**bot, "model": escalated},
                                    prompt_text=user_prompt,
                                    system_prompt=system_prompt,
                                    read_only=True,
                                )
                            else:
                                answer, usage, err = await _llm_oneshot(
                                    provider=bot_provider, model=escalated,
                                    user_message=user_prompt, system=system_prompt,
                                    max_tokens=500,
                                )
                            if answer:
                                bot_model = escalated  # Für Message-Metadata
                        except Exception as e:
                            err = f"escalation_failed: {type(e).__name__}: {e}"

                duration_ms = int((datetime.now() - t0).total_seconds() * 1000)

                if err or not answer:
                    text_out = f"*(Fehler: {err or 'leere Antwort'})*"
                else:
                    text_out = answer.strip()

                # Message aufbauen + persistieren + broadcasten
                # WICHTIG: Arena-Frontend rendert `content`, nicht `text` —
                # wir schreiben beide Felder damit sowohl klassische Arena-UI
                # als auch neue Trialog-Konsumenten es sehen. Zusätzlich `name`
                # (legacy Arena) + `role` (Arena-Message-Typ).
                msg = {
                    "turn": turn_counter,
                    "round": round_idx + 1,
                    "role": "bot",
                    "name": bot_name,        # Arena-legacy
                    "bot_name": bot_name,
                    "color": bot_color,      # Arena-legacy
                    "bot_color": bot_color,
                    "provider": bot_provider,
                    "model": bot_model or "default",
                    "content": text_out,     # Arena-Frontend-Feld
                    "text": text_out,        # Trialog-Konsistenz
                    "timestamp": datetime.now().isoformat(),
                    "duration_ms": duration_ms,
                    "usage": usage,
                    "error": err,
                }
                # Neueste rooms-Liste holen (andere Endpoints könnten
                # zwischenzeitlich geschrieben haben)
                fresh = _load_arena_rooms()
                rr = next((r for r in fresh if r.get("id") == room_id), None)
                if rr:
                    rr.setdefault("messages", []).append(msg)
                    rr["last_active"] = datetime.now().isoformat()
                    _save_arena_rooms(fresh)

                await _broadcast_all_ws({
                    "type": "arena.message",
                    "room_id": room_id,
                    "message": msg,
                })

                # --- Periodic Orchestrator/Summary (Phase 1) ---
                # Nach jeden `summarize_every`-Turns ein Orchestrator-Call:
                # fasst die wichtigsten Punkte der bisherigen Diskussion zusammen,
                # highlightet Dissens und Konsens, setzt Impuls für nächste Runde.
                # Default 9 (= 3 Runden bei 3 Bots), konfigurierbar über room.summarize_every.
                _se = int(room.get("summarize_every") or 9)
                if _se > 0 and turn_counter % _se == 0:
                    await _trialog_orchestrator_call(room_id, current_turn=turn_counter)

                # --- Room-Renewal: Archive + Continuity-Seed wenn Raum zu
                # groß wird (Lag-Schutz, identisch zum Classic-Arena-Pattern).
                # Alle Messages ab Turn 0 bis jetzt werden in data/arena-archive/{room_id}/{ts}.jsonl
                # geschrieben. Der Room wird mit einem kompakten Seed neu
                # aufgesetzt: letztes Orchestrator-Fazit + offene Actions + Konsens-Historie.
                try:
                    _rr_rooms = _load_arena_rooms()
                    _rr_room = _get_arena_room(_rr_rooms, room_id)
                except Exception:
                    _rr_room = None
                if _rr_room:
                    _rr_msgs = _rr_room.get("messages", [])
                    try:
                        _payload_kb = len(json.dumps(_rr_msgs, ensure_ascii=False).encode()) / 1024
                    except Exception:
                        _payload_kb = 0
                    _n = len(_rr_msgs)
                    _avg = _payload_kb / _n if _n else 0
                    _renew_density = (_n >= ARENA_RENEW_AVG_MSG_MIN_COUNT and _avg > ARENA_RENEW_AVG_MSG_KB)
                    if _n > ARENA_RENEW_MSG_COUNT or _payload_kb > ARENA_RENEW_PAYLOAD_KB or _renew_density:
                        _renew_info = _archive_and_renew_room(_rr_room)
                        _save_arena_rooms(_rr_rooms)
                        _recov = (
                            f", {_renew_info.get('recovered', 0)} nachgetragen"
                            if _renew_info.get("recovered") else ""
                        )
                        _divider = {
                            "role": "system",
                            "name": "System",
                            "content": (
                                f"📦 **Raum komprimiert bei Turn {turn_counter}** — "
                                f"{_renew_info['archived_msgs']} Messages archiviert, "
                                f"{_renew_info['open_actions']} offene Aktionen übertragen{_recov} "
                                f"→ Archiv: {_renew_info['ts']}"
                                + (f" (⚠️ Dichte: {_avg:.1f} KB/msg)" if _renew_density else "")
                            ),
                            "timestamp": datetime.now().isoformat(),
                            "color": "#6366f1",
                            "turn": turn_counter,
                            "is_renewal_divider": True,
                        }
                        _rr_rooms2 = _load_arena_rooms()
                        _rr_room2 = _get_arena_room(_rr_rooms2, room_id)
                        if _rr_room2:
                            _rr_room2["messages"].append(_divider)
                            _save_arena_rooms(_rr_rooms2)
                        await _broadcast_all_ws({
                            "type": "arena.message",
                            "room_id": room_id,
                            "message": _divider,
                        })
                        await _broadcast_all_ws({
                            "type": "arena.renewed",
                            "room_id": room_id,
                            "archived_msgs": _renew_info["archived_msgs"],
                            "open_actions": _renew_info["open_actions"],
                            "ts": _renew_info["ts"],
                        })
                        # Lokale room-Referenz aktualisieren damit history-Block
                        # im nächsten Turn aus dem kompakten Seed baut
                        room = _rr_room2 or _rr_room

                # Kleiner Rate-Gap (verhindert Overload bei Free-Tier-Providers)
                await asyncio.sleep(0.5)

        # Fertig
        fresh = _load_arena_rooms()
        rr = next((r for r in fresh if r.get("id") == room_id), None)
        if rr:
            rr["status"] = "idle"
            _save_arena_rooms(fresh)
        await _broadcast_all_ws({
            "type": "arena.status",
            "room_id": room_id,
            "status": "idle",
            "mode": "trialog",
            "reason": "rounds_complete",
        })
        print(f"[Trialog] room {room_id[:8]} fertig — {turn_counter} Turns")

    except asyncio.CancelledError:
        fresh = _load_arena_rooms()
        rr = next((r for r in fresh if r.get("id") == room_id), None)
        if rr:
            rr["status"] = "stopped"
            _save_arena_rooms(fresh)
        raise
    except Exception as e:
        print(f"[Trialog] Fehler room {room_id[:8]}: {type(e).__name__}: {e}")
        fresh = _load_arena_rooms()
        rr = next((r for r in fresh if r.get("id") == room_id), None)
        if rr:
            rr["status"] = "error"
            _save_arena_rooms(fresh)
    finally:
        # In-memory Flag zurücksetzen damit Restart möglich ist
        if room_id in _arena_rooms:
            _arena_rooms[room_id]["running"] = False


async def _trialog_exec_turn(
    bot: dict,
    prompt_text: str,
    system_prompt: str,
    cwd: str = "C:/Users/Robin/Jarvis",
    timeout_sec: int = 240,
    read_only: bool = False,
) -> tuple[str, dict, str | None]:
    """Per-Provider-Exec-Dispatch für Phase 2 (write) + Phase 3 (read-only).

    write (default):
      - claude-cli  → --permission-mode bypassPermissions (volle MCP-Tools)
      - codex-cli   → exec --sandbox workspace-write --full-auto
      - gemini-cli  → --approval-mode yolo -s (sandboxed)

    read-only (Phase 3):
      - claude-cli  → --tools Read,Glob,Grep
      - codex-cli   → exec --sandbox read-only
      - gemini-cli  → --approval-mode plan

    Text-only Provider (claude-api, openai, openclaw) fallen auf LLM-only
    mit Disclaimer zurück.
    """
    import time as _t

    bot_provider = (bot.get("provider") or "claude-cli").strip()
    bot_model = (bot.get("model") or "").strip()
    t0 = _t.time()
    _suffix = "-read" if read_only else "-exec"

    # --- Claude-CLI ------------------------------------------------------
    if bot_provider == "claude-cli":
        cmd = [
            CLAUDE_EXE, "-p",
            "--output-format", "text",
            "--max-turns", "8" if read_only else "25",
        ]
        if read_only:
            cmd.extend(["--tools", "Read,Glob,Grep"])
        else:
            cmd.extend(["--permission-mode", "bypassPermissions"])
        if bot_model:
            cmd.extend(["--model", bot_model])
        stdin_text = f"<system>\n{system_prompt}\n</system>\n\n{prompt_text}"
        return await _run_exec_subprocess(
            cmd, stdin_text, cwd, timeout_sec, "claude-cli" + _suffix,
        )

    # --- Codex-CLI -------------------------------------------------------
    if bot_provider == "codex-cli":
        from llm_client import _resolve_codex_exe, _codex_is_logged_in
        codex_exe = _resolve_codex_exe()
        if not _codex_is_logged_in():
            return "", {}, "codex nicht eingeloggt (codex login)"
        cmd = [
            codex_exe, "exec",
            "--sandbox", "read-only" if read_only else "workspace-write",
            "--skip-git-repo-check",
            "-C", cwd,
        ]
        if not read_only:
            cmd.append("--full-auto")
        if bot_model:
            cmd.extend(["-m", bot_model])
        cmd.append("-")
        stdin_text = f"<system>\n{system_prompt}\n</system>\n\n{prompt_text}"
        return await _run_exec_subprocess(
            cmd, stdin_text, cwd, timeout_sec, "codex-cli" + _suffix,
        )

    # --- Gemini-CLI ------------------------------------------------------
    if bot_provider == "gemini-cli":
        from llm_client import _resolve_gemini_exe, _gemini_is_logged_in
        gemini_exe = _resolve_gemini_exe()
        if not _gemini_is_logged_in():
            return "", {}, "gemini nicht eingeloggt (gemini interaktiv starten)"
        cmd = [
            gemini_exe, "-p", "",
            "-o", "stream-json",
            "--skip-trust",
            "--approval-mode", "plan" if read_only else "yolo",
        ]
        if not read_only:
            cmd.append("-s")
        if bot_model:
            cmd.extend(["-m", bot_model])
        stdin_text = f"<system>\n{system_prompt}\n</system>\n\n{prompt_text}"
        return await _run_exec_subprocess(
            cmd, stdin_text, cwd, timeout_sec,
            "gemini-cli" + _suffix, parse_stream_json=True,
        )

    # --- Text-Only Fallback -----------------------------------------------
    # Provider ohne Exec-Unterstützung: wir degradieren auf LLM-only
    # mit klarem Disclaimer im Output damit User sieht warum nichts passiert.
    from llm_client import invoke_llm, Message, TextDelta, ErrorEvent, MessageStop, UsageEvent
    full_text, usage, err = await _llm_oneshot(
        provider=bot_provider,
        model=bot_model,
        user_message=prompt_text,
        system=system_prompt + "\n\n(Hinweis: Du kannst für diesen Provider aktuell "
                               "keine Dateien schreiben. Beschreibe was du tun WÜRDEST.)",
        max_tokens=1500,
    )
    return (full_text or "", usage or {}, err)


async def _run_exec_subprocess(
    cmd: list,
    stdin_text: str,
    cwd: str,
    timeout_sec: int,
    label: str,
    parse_stream_json: bool = False,
) -> tuple[str, dict, str | None]:
    """Gemeinsame Subprocess-Execution für die drei CLI-Provider.

    Liest stdout bis EOF, sammelt alle Text-Chunks. Bei parse_stream_json:
    versucht JSONL zu parsen (Gemini-Format), fällt bei Fehler auf raw text
    zurück.
    """
    import asyncio as _asyncio
    import time as _t

    subprocess_flags = 0
    if os.name == "nt":
        subprocess_flags = 0x08000000  # CREATE_NO_WINDOW

    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.DEVNULL,
            stdin=asyncio.subprocess.PIPE,
            cwd=cwd,
            limit=8 * 1024 * 1024,
            creationflags=subprocess_flags if os.name == "nt" else 0,
        )
    except FileNotFoundError:
        return "", {}, f"{label}: executable nicht gefunden"
    except Exception as e:
        return "", {}, f"{label} spawn failed: {e}"

    try:
        proc.stdin.write(stdin_text.encode("utf-8"))
        await proc.stdin.drain()
        proc.stdin.close()
    except Exception as e:
        proc.kill()
        await proc.wait()
        return "", {}, f"{label} stdin failed: {e}"

    full_text = ""
    usage: dict = {}
    t0 = _t.time()
    try:
        stdout_bytes = await asyncio.wait_for(proc.stdout.read(), timeout=timeout_sec)
    except asyncio.TimeoutError:
        proc.kill()
        await proc.wait()
        return "", {}, f"{label} timeout nach {timeout_sec}s"

    out = stdout_bytes.decode("utf-8", errors="replace")

    if parse_stream_json:
        # Gemini-Format: JSONL mit {type: message, role: assistant, content, delta}
        for line in out.splitlines():
            line = line.strip()
            if not line or not line.startswith("{"):
                continue
            try:
                ev = json.loads(line)
            except json.JSONDecodeError:
                continue
            if ev.get("type") == "message" and ev.get("role") == "assistant":
                c = ev.get("content")
                if isinstance(c, str):
                    full_text += c
            elif ev.get("type") == "result":
                stats = ev.get("stats") or {}
                if stats:
                    usage = {
                        "input_tokens": int(stats.get("input_tokens", 0) or 0),
                        "output_tokens": int(stats.get("output_tokens", 0) or 0),
                    }
    else:
        # Plain text output (claude-cli -p text, codex exec default)
        full_text = out.strip()

    if not full_text:
        return "", usage, f"{label} lieferte leere Antwort"

    return full_text, usage, None


async def _trialog_orchestrator_call(room_id: str, current_turn: int = 0) -> None:
    """Orchestrator-Tick: liest die bisherige Konversation, generiert eine
    strukturierte Zusammenfassung (Dissens/Konsens/Aktionspunkte/Nächster
    Impuls) und injiziert sie als role=orchestrator Message in den Room.

    Provider für den Orchestrator ist konfigurierbar via room.orchestrator_provider
    (default claude-cli/sonnet) — bewusst unabhängig von den Teilnehmer-Bots,
    damit eine neutrale Meta-Sicht entsteht.

    `current_turn` ist der Trialog-Loop-Turn-Counter (1-basiert, analog Bot-Messages).
    Wird fürs UI verwendet damit Orchestrator im selben Turn-Nummern-Raum wie die
    Bots angezeigt wird (nicht als Message-Index).
    """
    t0 = datetime.now()
    # "Thinking"-Indicator für Orchestrator — damit UI wie bei Bots anzeigt
    # dass gerade gerechnet wird
    await _broadcast_all_ws({
        "type": "arena.thinking",
        "room_id": room_id,
        "bot_name": "Orchestrator",
        "bot_color": "#eab308",
        "turn": current_turn,
        "role": "orchestrator",
    })
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        return

    o_provider = (room.get("orchestrator_provider") or "claude-cli").strip()
    o_model = (room.get("orchestrator_model") or "sonnet").strip()
    topic = room.get("topic") or room.get("title") or ""
    bot_names = [b.get("name", "?") for b in (room.get("bots") or [])]

    # Letzte ~40 Bot-Messages als Diskussionsgrundlage (analog classic Arena)
    recent_msgs = []
    for m in (room.get("messages") or [])[-40:]:
        role = m.get("role") or ""
        name = m.get("bot_name") or m.get("name") or "?"
        content = (m.get("content") or m.get("text") or "").strip()
        if not content:
            continue
        if role in ("orchestrator", "approval"):
            continue
        recent_msgs.append(f"{name}: {content[:400]}")
    if not recent_msgs:
        return

    history_block = "\n\n".join(recent_msgs)

    # Classic-Arena-kompatibles Format: ACTION-Zeilen + VERDICT-Marker
    # damit das existing Frontend (arena.js) daraus Approve-Buttons macht
    # und Robin die Aktionen wie gewohnt freigeben kann.
    system_prompt = (
        "Du bist der Orchestrator einer Bot-Diskussion. Robin ist kein reiner "
        "Entwickler — schreib verstaendlich, in Alltagssprache.\n\n"
        "Format (STRIKT einhalten):\n\n"
        "**Worum es geht**\n"
        "1-2 Saetze, was die Bots gerade besprechen.\n\n"
        "**Was koennen wir umsetzen**\n"
        "Fuer JEDE konkret umsetzbare Aktion genau EINE Zeile im Format:\n"
        "[ACTION:BotName] <Was passiert und was bringt das — max 12 Woerter> — "
        "Datei: <Dateiname>, Funktion: <Funktionsname oder Route>\n\n"
        "Regeln fuer ACTION-Zeilen:\n"
        "- BotName = Name eines Teilnehmer-Bots aus {" + ", ".join(bot_names) + "}\n"
        "- Datei + Funktion nur angeben wenn du sie aus der Diskussion sicher "
        "ableiten kannst — sonst weglassen\n"
        "- Keine vagen Ideen, nur konkrete umsetzbare Schritte\n"
        "- Deutsch, knapp, Alltagssprache\n\n"
        "Nach den ACTION-Zeilen in einer eigenen letzten Zeile EXAKT einen "
        "dieser Marker setzen:\n"
        "[VERDICT:ongoing]   - Bots bringen noch neue Punkte, weiterdiskutieren lohnt\n"
        "[VERDICT:consensus] - Plan steht, Robin sollte die ACTIONs pruefen+freigeben\n"
        "[VERDICT:stalemate] - Bots drehen sich im Kreis, Robin muss eingreifen\n\n"
        "HARTE REGEL zu Runden/Aufwand:\n"
        "Falls du Aufwand in Runden nennst: HARTE OBERGRENZE 10 Runden. Niemals "
        "hoehere Zahlen, passt nicht in 10: in kleinere [ACTION]-Schritte aufteilen."
    )
    user_prompt = (
        f"Thema: {topic}\n"
        f"Teilnehmer: {', '.join(bot_names)}\n\n"
        f"Diskussion:\n{history_block}"
    )

    try:
        answer, usage, err = await _llm_oneshot(
            provider=o_provider,
            model=o_model,
            user_message=user_prompt,
            system=system_prompt,
            max_tokens=500,
        )
    except Exception as e:
        answer = ""
        usage = {}
        err = f"{type(e).__name__}: {e}"

    duration_ms = int((datetime.now() - t0).total_seconds() * 1000)

    if err or not answer:
        print(f"[Trialog] orchestrator failed room={room_id[:8]}: {err}")
        return

    # Classic-Arena-kompatibel: role=orchestrator + name=Orchestrator + verdict
    # damit die existing arena.js-UI die ACTION-Zeilen parsed und Approve-
    # Buttons rendert. turn-Feld für Dedup im Frontend.
    verdict = "ongoing"
    low = answer.lower()
    if "[verdict:consensus]" in low:
        verdict = "consensus"
    elif "[verdict:stalemate]" in low:
        verdict = "stalemate"

    msg = {
        "role": "orchestrator",
        "name": "Orchestrator",
        "bot_name": "Orchestrator",
        "color": "#eab308",
        "bot_color": "#eab308",
        "provider": o_provider,
        "model": o_model,
        "content": answer.strip(),
        "text": answer.strip(),
        "verdict": verdict,
        "timestamp": datetime.now().isoformat(),
        "duration_ms": duration_ms,
        "usage": usage,
        "turn": current_turn if current_turn > 0 else len(room.get("messages", [])),
    }
    # Persistieren
    fresh = _load_arena_rooms()
    rr = _get_arena_room(fresh, room_id)
    if rr:
        rr.setdefault("messages", []).append(msg)
        rr["orchestrator_summary"] = answer.strip()
        rr["last_active"] = datetime.now().isoformat()
        _save_arena_rooms(fresh)

    await _broadcast_all_ws({
        "type": "arena.message",
        "room_id": room_id,
        "message": msg,
    })
    print(f"[Trialog] orchestrator tick room={room_id[:8]} via {o_provider}/{o_model} ({duration_ms}ms)")


async def api_arena_trialog_start(request):
    """POST /api/arena/trialog/start — legt einen Multi-Model-Trialog-Raum an
    und startet die Round-Robin-Loop.

    Body: {
      title?, topic, max_rounds?,
      bots: [{name, provider, model, persona?, color?}, …]  (2-4 Bots)
    }
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    topic = (body.get("topic") or "").strip()
    if not topic:
        return web.json_response({"ok": False, "error": "topic erforderlich"}, status=400)

    bots = body.get("bots") or []
    if not isinstance(bots, list) or len(bots) < 2 or len(bots) > 6:
        return web.json_response(
            {"ok": False, "error": "2-6 Bots erforderlich"}, status=400,
        )

    # Default-Persona falls leer, Colors falls leer
    from llm_client import PROVIDER_NAMES as _PN
    color_palette = ["#6366f1", "#22c55e", "#f59e0b", "#ec4899", "#a78bfa", "#14b8a6"]
    cleaned_bots = []
    for i, b in enumerate(bots):
        if not isinstance(b, dict):
            continue
        provider = (b.get("provider") or "claude-cli").strip()
        if provider not in _PN:
            return web.json_response(
                {"ok": False, "error": f"Unbekannter Provider {provider!r}"},
                status=400,
            )
        cleaned_bots.append({
            "name": (b.get("name") or f"Bot-{i+1}").strip(),
            "provider": provider,
            "model": (b.get("model") or "").strip(),
            "persona": (b.get("persona") or "Du bist ein hilfreicher Diskussionspartner.").strip(),
            "color": b.get("color") or color_palette[i % len(color_palette)],
        })

    room = {
        "id": str(uuid.uuid4()),
        "title": (body.get("title") or f"Trialog — {topic[:40]}").strip(),
        "topic": topic,
        "category": "trialog",
        "mode": "trialog",
        "bots": cleaned_bots,
        "messages": [],
        "max_rounds": max(1, min(10, int(body.get("max_rounds", 3)))),
        # Orchestrator-Config. Default: alle 9 Turns (= nach 3 vollen Runden
        # bei 3 Bots — analog classic Arena mit ~10). Neutraler Provider
        # claude-cli/sonnet. 0 deaktiviert den Orchestrator komplett.
        "summarize_every": int(body.get("summarize_every", 9)),
        "orchestrator_provider": (body.get("orchestrator_provider") or "claude-cli").strip(),
        "orchestrator_model": (body.get("orchestrator_model") or "sonnet").strip(),
        "created": datetime.now().isoformat(),
        "last_active": datetime.now().isoformat(),
        "status": "idle",
    }

    rooms = _load_arena_rooms()
    rooms.append(room)
    _save_arena_rooms(rooms)

    # Loop starten (Background-Task)
    task = asyncio.create_task(_run_trialog(room["id"], rooms))
    _trialog_tasks[room["id"]] = task

    return web.json_response({"ok": True, "room": room}, status=201)


async def api_arena_trialog_read(request):
    """POST /api/arena/trialog/{id}/read — Read-Mode aktivieren.
    Body: {"turns": N} — N Bot-Turns im read-only-Modus (Claude: Read/Glob/Grep,
    Codex: --sandbox read-only, Gemini: --approval-mode plan).

    Zweck: Bots sollen ihre Argumente mit echten Fakten aus der Codebase
    belegen können, ohne Schreibrechte. Weniger invasiv als /execute.
    """
    room_id = request.match_info.get("id", "")
    try:
        body = await request.json() if request.content_length else {}
    except Exception:
        body = {}
    turns = max(1, min(int(body.get("turns", 3)), 10))

    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        return web.json_response({"ok": False, "error": "Room nicht gefunden"}, status=404)

    room["read_enabled"] = True
    room["read_turns_left"] = turns
    _save_arena_rooms(rooms)

    await _broadcast_all_ws({
        "type": "arena.mode_change",
        "room_id": room_id,
        "mode": "read",
        "turns_left": turns,
    })
    return web.json_response({"ok": True, "read_turns_left": turns})


async def api_arena_trialog_stop(request):
    """POST /api/arena/trialog/{id}/stop — bricht einen laufenden Trialog ab."""
    room_id = request.match_info.get("id", "")
    task = _trialog_tasks.get(room_id)
    if task and not task.done():
        task.cancel()
    rooms = _load_arena_rooms()
    room = next((r for r in rooms if r.get("id") == room_id), None)
    if room:
        room["status"] = "stopped"
        _save_arena_rooms(rooms)
    return web.json_response({"ok": True, "room_id": room_id})


async def api_arena_update(request):
    """PUT /api/arena/rooms/{id}  -  Update room settings (title, topic, bots, max_turns)."""
    room_id = request.match_info["id"]
    body = await request.json()
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")
    for key in ("title", "topic", "bots", "max_turns", "max_rounds", "category", "summarize_every", "allow_read", "auto_continue"):
        if key in body:
            if key == "max_rounds":
                try:
                    room[key] = max(1, min(25, int(body[key])))
                except (TypeError, ValueError):
                    pass
            else:
                room[key] = body[key]
    _save_arena_rooms(rooms)
    return web.json_response(room)


async def api_arena_execute(request):
    """POST /api/arena/rooms/{id}/execute  -  Enable execution mode for N turns."""
    room_id = request.match_info["id"]
    body = await request.json()
    execution_turns = body.get("turns", 2)  # Default: 2 bots each get 1 execution turn
    execution_turns = max(1, min(execution_turns, 10))  # Clamp 1-10
    brief = body.get("brief") or {}

    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")

    room["tools_enabled"] = True
    room["execution_turns_left"] = execution_turns
    # Persist the mini-brief on the room so UI + bots can reference it every turn
    if brief.get("goal"):
        room["current_brief"] = {
            "goal": brief.get("goal", "").strip(),
            "scope": brief.get("scope", "").strip(),
            "dod": brief.get("dod", "").strip(),
            "check": brief.get("check", "").strip(),
            "turns": execution_turns,
            "created": datetime.now().isoformat(),
        }
    _save_arena_rooms(rooms)

    # Build a structured briefing message so bots see a clear mission before they build
    if brief.get("goal"):
        parts = [
            "== MINI-AUFTRAG (Briefing vor dem Bau) ==",
            f"ZIEL: {brief.get('goal','').strip()}",
        ]
        if brief.get("scope"):
            parts.append(f"SCOPE / NICHT-ZIELE: {brief['scope'].strip()}")
        if brief.get("dod"):
            parts.append(f"DEFINITION OF DONE: {brief['dod'].strip()}")
        if brief.get("check"):
            parts.append(f"CHECKPOINT: {brief['check'].strip()}")
        parts.append(
            f"Ihr habt {execution_turns} Execution-Turn(s) mit vollen Tools. "
            f"Haltet euch strikt an Ziel und Scope. Liefert am Ende ein konkretes Ergebnis "
            f"und sagt klar, ob Definition of Done erreicht wurde."
        )
        content = "\n".join(parts)
    else:
        content = f"Setzt das um was ihr diskutiert habt. Ihr habt {execution_turns} Execution-Turns mit vollen Tools."

    msg = {
        "role": "robin",
        "name": "Robin",
        "content": content,
        "timestamp": datetime.now().isoformat(),
        "color": "#f59e0b",
        "brief": room.get("current_brief") if brief.get("goal") else None,
    }
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if room:
        room["messages"].append(msg)
        _save_arena_rooms(rooms)

    await _broadcast_all_ws({
        "type": "arena.message",
        "room_id": room_id,
        "message": msg,
    })
    await _broadcast_all_ws({
        "type": "arena.mode_change",
        "room_id": room_id,
        "mode": "execute",
        "turns_left": execution_turns,
    })

    return web.json_response({"ok": True, "tools_enabled": True, "execution_turns_left": execution_turns})


async def api_arena_orchestrator_chat(request):
    """POST /api/arena/rooms/{id}/orchestrator  -  Chat with the orchestrator about the discussion."""
    room_id = request.match_info["id"]
    body = await request.json()
    user_msg = body.get("message", "").strip()
    if not user_msg:
        raise web.HTTPBadRequest(reason="No message provided")

    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")

    # Build context from discussion
    recent_msgs = room.get("messages", [])[-30:]
    history_lines = [f"{m['name']}: {m['content']}" for m in recent_msgs]
    history_block = "\n\n".join(history_lines) if history_lines else "(Keine Nachrichten)"

    # Get existing summary if any
    existing_summary = room.get("orchestrator_summary", "")
    summary_ctx = f"\n\nLetzte Zusammenfassung des Orchestrators:\n{existing_summary}" if existing_summary else ""

    prompt = (
        f"Du bist der Orchestrator/Richter einer Bot-Diskussion. "
        f"Dein Job: Die Diskussion analysieren, zusammenfassen und Robin beraten.\n\n"
        f"Thema: {room.get('topic', 'Freie Diskussion')}\n\n"
        f"Bisherige Diskussion ({len(room.get('messages', []))} Nachrichten):\n{history_block}"
        f"{summary_ctx}\n\n"
        f"Robin fragt dich: {user_msg}\n\n"
        f"Antworte kurz, strukturiert, auf Deutsch. Wenn Robin fragt ob etwas umgesetzt werden soll, "
        f"gib konkrete Schritte. Wenn er Aenderungen will, schlage angepasste Plaene vor."
    )

    # Prompt via stdin - umgeht Windows 32K command-line limit
    cmd = [
        CLAUDE_EXE, "-p",
        "--model", "opus",
        "--output-format", "text",
        "--tools", "",
        "--max-turns", "1",
    ]

    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            stdin=asyncio.subprocess.PIPE,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        stdout_data, stderr_data = await asyncio.wait_for(
            proc.communicate(input=prompt.encode("utf-8")), timeout=120
        )
        response = stdout_data.decode("utf-8", errors="replace").strip()
        if "Reached max turns" in response:
            lines = [l for l in response.splitlines() if "Reached max turns" not in l and "Error:" not in l]
            response = "\n".join(lines).strip()
        if not response:
            err = stderr_data.decode("utf-8", errors="replace").strip()[:300]
            response = f"(Keine Antwort vom Orchestrator — stderr: {err or 'leer'})"
    except asyncio.TimeoutError:
        response = "(Timeout — Orchestrator hat zu lange gebraucht)"
    except Exception as e:
        response = f"(Fehler: {e})"

    return web.json_response({"response": response})


ARENA_SNAPSHOTS_DIR = Path("C:/Users/Robin/Jarvis/data/arena-snapshots")


async def api_arena_snapshot(request):
    """POST /api/arena/rooms/{id}/snapshot  -  Save a named checkpoint of the current room state."""
    room_id = request.match_info["id"]
    try:
        body = await request.json()
    except Exception:
        body = {}
    label = (body.get("label") or "").strip() or datetime.now().strftime("%Y-%m-%d %H:%M")

    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")

    ARENA_SNAPSHOTS_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_title = "".join(c if c.isalnum() or c in "-_" else "_" for c in room.get("title", room_id))[:40]
    fname = f"{ts}_{safe_title}.json"
    snapshot = {
        "snapshot_at": datetime.now().isoformat(),
        "label": label,
        "room_id": room_id,
        "title": room.get("title"),
        "topic": room.get("topic"),
        "messages": room.get("messages", []),
        "orchestrator_summary": room.get("orchestrator_summary"),
    }
    (ARENA_SNAPSHOTS_DIR / fname).write_text(
        json.dumps(snapshot, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    # Manuell gemerkten Stand auch in consensus_marks eintragen, damit er in Bot-Prompts fliesst
    orch_summary = (room.get("orchestrator_summary") or "").strip()
    mark_summary = label
    if orch_summary:
        mark_summary += ": " + orch_summary[:250]
    room.setdefault("consensus_marks", []).append({
        "turn": len(room.get("messages", [])),
        "at": datetime.now().isoformat(),
        "msg_count": len(room.get("messages", [])),
        "summary": mark_summary[:300],
        "manual": True,
        "file": fname,
    })
    _save_arena_rooms(rooms)

    return web.json_response({"ok": True, "file": fname, "label": label, "messages": len(snapshot["messages"])})


def _list_room_snapshots(room_id: str) -> list:
    """Return snapshot metadata for a room (index info, no message contents)."""
    if not ARENA_SNAPSHOTS_DIR.exists():
        return []
    result = []
    for f in sorted(ARENA_SNAPSHOTS_DIR.glob("*.json"), reverse=True):
        try:
            data = json.loads(f.read_text(encoding="utf-8"))
            if data.get("room_id") == room_id:
                result.append({
                    "file": f.name,
                    "label": data.get("label"),
                    "snapshot_at": data.get("snapshot_at"),
                    "messages": len(data.get("messages", [])),
                })
        except Exception:
            pass
    return result


async def api_arena_snapshots_list(request):
    """GET /api/arena/rooms/{id}/snapshots  -  List saved snapshots for a room."""
    room_id = request.match_info["id"]
    return web.json_response(_list_room_snapshots(room_id))


async def api_arena_consensus_mark_delete(request):
    """DELETE /api/arena/rooms/{id}/consensus_marks/{index}  -  Remove a poisoned consensus mark."""
    room_id = request.match_info["id"]
    try:
        idx = int(request.match_info["index"])
    except ValueError:
        raise web.HTTPBadRequest(reason="index must be int")
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")
    marks = room.get("consensus_marks") or []
    if idx < 0 or idx >= len(marks):
        raise web.HTTPNotFound(reason="Mark index out of range")
    removed = marks.pop(idx)
    room["consensus_marks"] = marks
    _save_arena_rooms(rooms)
    return web.json_response({"ok": True, "removed": removed, "remaining": len(marks)})


async def api_arena_unclear(request):
    """POST /api/arena/rooms/{id}/unclear  -  Persist or remove an unclear action in room state."""
    room_id = request.match_info["id"]
    body = await request.json()
    action_text = (body.get("text") or "").strip()
    bot_name = (body.get("bot") or "").strip()
    msg_idx = body.get("msgIdx")  # original message index for back-reference
    msg_ts = (body.get("msgTimestamp") or "").strip()
    remove = bool(body.get("remove", False))
    if not action_text:
        raise web.HTTPBadRequest(reason="text required")
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")
    unclear = room.setdefault("unclear_actions", [])
    if remove:
        before = len(unclear)
        room["unclear_actions"] = [u for u in unclear if u.get("text") != action_text]
        _save_arena_rooms(rooms)
        return web.json_response({"ok": True, "removed": before - len(room["unclear_actions"])})
    # Avoid duplicates
    if any(u.get("text") == action_text for u in unclear):
        return web.json_response({"ok": True, "dup": True, "count": len(unclear)})
    entry = {
        "text": action_text,
        "bot": bot_name,
        "at": datetime.now().isoformat(),
    }
    if msg_idx is not None:
        entry["msgIdx"] = int(msg_idx)
    if msg_ts:
        entry["msgTimestamp"] = msg_ts
    unclear.append(entry)
    _save_arena_rooms(rooms)
    return web.json_response({"ok": True, "count": len(unclear)})


async def api_arena_summarize(request):
    """POST /api/arena/rooms/{id}/summarize  -  Generate orchestrator summary on demand."""
    room_id = request.match_info["id"]
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")

    summary = await _generate_orchestrator_summary(room)

    # Save summary to room
    room["orchestrator_summary"] = summary
    room["orchestrator_summary_at"] = datetime.now().isoformat()
    _extract_and_store_action_items(room, summary)
    _save_arena_rooms(rooms)

    # Also add as system message — but skip if last message is already an orchestrator Fazit
    messages = room.get("messages", [])
    last_msg = messages[-1] if messages else None
    already_fazit = last_msg and last_msg.get("role") == "orchestrator"
    msg = {
        "role": "orchestrator",
        "name": "Orchestrator",
        "content": summary,
        "timestamp": datetime.now().isoformat(),
        "color": "#a855f7",
        "turn": len(messages),
    }
    if not already_fazit:
        room["messages"].append(msg)
        _save_arena_rooms(rooms)
        await _broadcast_all_ws({
            "type": "arena.message",
            "room_id": room_id,
            "message": msg,
        })
    else:
        _save_arena_rooms(rooms)

    return web.json_response({"summary": summary})


async def api_arena_approve(request):
    """POST /api/arena/rooms/{id}/approve  -  Approve selected actions for execution by a separate bot."""
    room_id = request.match_info["id"]
    body = await request.json()
    actions_raw = body.get("actions", [])
    # Deduplicate: keep first occurrence, case-insensitive key
    seen_keys: set = set()
    actions = []
    for a in actions_raw:
        key = a.strip().lower()
        if key and key not in seen_keys:
            seen_keys.add(key)
            actions.append(a)
    rejected = body.get("rejected", [])  # Gegenstimmen: nicht-angewaehlte Alternativen
    rej_reason = (body.get("reject_reason") or "").strip()
    if not actions:
        return web.json_response({"ok": False, "error": "Keine Aktionen ausgewaehlt"})

    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")

    # Dedup gegen echten Raumzustand: bereits erfolgreich abgeschlossene Aktionen überspringen
    done_actions: set = set(room.get("done_actions", []))
    already_done = [a for a in actions if a.strip().lower() in done_actions]
    actions = [a for a in actions if a.strip().lower() not in done_actions]
    if not actions:
        return web.json_response({"ok": False, "error": "Alle Aktionen bereits abgeschlossen", "already_done": already_done})

    # Create approval message with tracked per-action status
    approval_id = f"appr-{int(datetime.now().timestamp()*1000)}"
    approval_msg = {
        "role": "approval",
        "name": "Robin",
        "approval_id": approval_id,
        "actions": [{"text": a, "status": "pending", "result": ""} for a in actions],
        "rejected": [{"text": r, "reason": rej_reason} for r in rejected],
        "content": (
            f"**Gruenes Licht** fuer {len(actions)} Aktion(en) — Umsetzung laeuft."
            + (f"  · {len(rejected)} Alternativ-Weg(e) verworfen (bleiben sichtbar)." if rejected else "")
        ),
        "timestamp": datetime.now().isoformat(),
        "color": "#f59e0b",
        "turn": len(room.get("messages", [])),
    }
    room["messages"].append(approval_msg)
    _save_arena_rooms(rooms)
    await _broadcast_all_ws({
        "type": "arena.message",
        "room_id": room_id,
        "message": approval_msg,
    })

    # Log each approved action as queued event
    for i, a in enumerate(actions):
        _append_action_log({
            "event": "action.queued",
            "room_id": room_id,
            "approval_id": approval_id,
            "action_idx": i,
            "action": a,
            "topic": room.get("topic", ""),
        })

    # Persist rejected actions to fail-backlog so they survive to next session
    if rejected:
        _append_fail_backlog([{
            "kind": "rejected",
            "action": r,
            "reason": rej_reason or "",
            "room_id": room_id,
            "topic": room.get("topic", ""),
            "timestamp": datetime.now().isoformat(),
        } for r in rejected])

    # Assigned-bot confirmation: each bot with ≥1 assigned action posts one sentence
    _BOT_COLORS = {"Alpha": "#6366f1", "Beta": "#22c55e", "Gamma": "#f59e0b", "Delta": "#ec4899"}
    _import_re = __import__("re")
    _bot_actions: dict = {}
    for a in actions:
        m = _import_re.match(r"^\[([A-Za-z][A-Za-z0-9_-]*)\]\s*(.+)", a.strip())
        if m:
            bname = m.group(1)
            _bot_actions.setdefault(bname, []).append(m.group(2).strip())
    for bname, btasks in _bot_actions.items():
        n = len(btasks)
        task_preview = btasks[0][:60] + ("…" if len(btasks[0]) > 60 else "")
        confirm_text = (
            f"Verstanden — ich übernehme {'diese Aufgabe' if n == 1 else f'diese {n} Aufgaben'}: „{task_preview}\u201c"
            + (f" (+{n-1} weitere)" if n > 1 else "")
            + " — läuft."
        )
        conf_msg = {
            "role": "assistant",
            "name": bname,
            "content": confirm_text,
            "timestamp": datetime.now().isoformat(),
            "color": _BOT_COLORS.get(bname, "#9ca3af"),
            "turn": len(room.get("messages", [])),
        }
        room["messages"].append(conf_msg)
        _save_arena_rooms(rooms)
        await _broadcast_all_ws({"type": "arena.message", "room_id": room_id, "message": conf_msg})

    # --- Force-Diff-Warnung: wenn Orchestrator schon mal lief (Ideenphase
    # vorbei) UND eine approved Action von einem Bot kommt der keinen
    # Code-Block in seinen letzten ~15 Turns gezeigt hat → warnen. Wir
    # blocken nicht hart (Robin's Freiheit), aber die Warnung landet in
    # der Response und optional als Info-Message im Raum.
    force_diff_warnings: list = []
    if (room.get("orchestrator_summary") or "").strip():
        recent_msgs = room.get("messages", [])[-20:]
        for a in actions:
            m = _import_re.match(r"^\[([A-Za-z][A-Za-z0-9_-]*)\]\s*(.+)", a.strip())
            if not m:
                continue
            bname = m.group(1)
            had_code = any(
                ((msg.get("bot_name") or msg.get("name") or "") == bname)
                and "```" in (msg.get("content") or msg.get("text") or "")
                for msg in recent_msgs
            )
            if not had_code:
                force_diff_warnings.append(
                    f"⚠ {bname}: keinen Code-Block/Diff in den letzten Turns gezeigt — "
                    f"Aktion ist abstrakt. Force-Diff empfiehlt konkreten Patch vor Approve."
                )
        if force_diff_warnings:
            warn_msg = {
                "role": "system",
                "name": "System",
                "content": "**Force-Diff-Warnung:**\n" + "\n".join(force_diff_warnings),
                "timestamp": datetime.now().isoformat(),
                "color": "#f59e0b",
                "turn": len(room.get("messages", [])),
            }
            room["messages"].append(warn_msg)
            _save_arena_rooms(rooms)
            await _broadcast_all_ws({"type": "arena.message", "room_id": room_id, "message": warn_msg})

    # Spawn executor bot in background (does NOT stop the discussion)
    asyncio.ensure_future(_run_executor(room_id, approval_id, actions))

    return web.json_response({
        "ok": True,
        "actions": len(actions),
        "approval_id": approval_id,
        "force_diff_warnings": force_diff_warnings,
    })


async def api_arena_file_peek(request):
    """GET /api/arena/file-peek?path=...  -  Return current file content + optional git diff.

    Used by the arena preview modal: Robin clicks a file path in the list, and
    the frontend opens an inline viewer with the current content (so he can
    check what the executor is about to touch). If the file has uncommitted
    changes relative to HEAD, we include the diff. If the file does not yet
    exist (will be created by the executor), we return an empty-ok marker.
    """
    raw = (request.query.get("path") or "").strip()
    if not raw:
        return web.json_response({"ok": False, "error": "path missing"}, status=400)
    # Normalize + contain inside Jarvis
    rel = raw.replace("\\", "/").lstrip("/")
    if ".." in rel.split("/"):
        return web.json_response({"ok": False, "error": "path traversal"}, status=400)
    JARVIS_ROOT = Path("C:/Users/Robin/Jarvis")
    target = (JARVIS_ROOT / rel).resolve()
    try:
        target.relative_to(JARVIS_ROOT.resolve())
    except Exception:
        return web.json_response({"ok": False, "error": "outside jarvis"}, status=400)
    if not target.exists():
        return web.json_response({
            "ok": True, "path": rel, "exists": False,
            "content": "", "diff": "", "size": 0,
            "note": "Datei existiert noch nicht — wird vom Executor neu angelegt."
        })
    if target.is_dir():
        return web.json_response({"ok": False, "error": "is a directory"}, status=400)
    try:
        size = target.stat().st_size
        # 300 KB cap — keep the modal fast
        if size > 300_000:
            return web.json_response({
                "ok": True, "path": rel, "exists": True,
                "size": size, "content": "", "diff": "",
                "note": f"Datei ist {size // 1024} KB — zu gross fuer Inline-Preview."
            })
        content = target.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return web.json_response({"ok": False, "error": f"read failed: {e}"}, status=500)
    # Best-effort git diff HEAD -- <rel>
    diff_text = ""
    try:
        import subprocess
        proc = subprocess.run(
            ["git", "diff", "HEAD", "--", rel],
            cwd=str(JARVIS_ROOT), capture_output=True, text=True, timeout=3,
        )
        if proc.returncode == 0 and proc.stdout.strip():
            diff_text = proc.stdout
            if len(diff_text) > 80_000:
                diff_text = diff_text[:80_000] + "\n... (gekuerzt)"
    except Exception:
        pass
    return web.json_response({
        "ok": True, "path": rel, "exists": True, "size": size,
        "content": content, "diff": diff_text,
    })


async def api_arena_files_status(request):
    """GET /api/arena/files-status?paths=a,b,c  -  Return existence + plausibility per path.

    Response: {ok: True, status: {path: bool}, plausible: {path: bool}}

    - status[p] = True  -> file exists at that relative path
    - plausible[p] = True -> file exists OR its parent directory exists
                             (so the executor could realistically create it there).
      False means the path looks like a hallucinated/fantasy name (e.g. "example.py"
      with no real parent in the repo) and the UI should hide it to avoid false hits.
    """
    raw = (request.query.get("paths") or "").strip()
    if not raw:
        return web.json_response({"ok": True, "status": {}, "plausible": {}})
    JARVIS_ROOT = Path("C:/Users/Robin/Jarvis")
    root_res = JARVIS_ROOT.resolve()
    status: dict = {}
    plausible: dict = {}
    items = [p.strip() for p in raw.replace("\n", ",").split(",") if p.strip()]
    for rel in items[:200]:
        norm = rel.replace("\\", "/").lstrip("/")
        if ".." in norm.split("/"):
            status[rel] = False
            plausible[rel] = False
            continue
        try:
            tgt = (JARVIS_ROOT / norm).resolve()
            tgt.relative_to(root_res)
            exists = tgt.exists() and tgt.is_file()
            status[rel] = exists
            if exists:
                plausible[rel] = True
            else:
                # Plausible only if parent dir exists AND path has a real dir segment.
                # A bare "example.py" with no slash is never plausible here.
                parts = norm.split("/")
                if len(parts) < 2:
                    plausible[rel] = False
                else:
                    parent = tgt.parent
                    try:
                        parent.relative_to(root_res)
                        plausible[rel] = parent.exists() and parent.is_dir()
                    except Exception:
                        plausible[rel] = False
        except Exception:
            status[rel] = False
            plausible[rel] = False
    return web.json_response({"ok": True, "status": status, "plausible": plausible})


# Cache fuer tiefe Vorschauen - gleicher Aktions-Block -> gleiches Ergebnis.
# Verhindert doppelte Claude-Calls in schneller Folge. In-memory, flueht bei Restart.
_DEEP_PREVIEW_CACHE: dict[str, tuple[float, list]] = {}
_DEEP_PREVIEW_TTL = 900  # 15 min


async def api_arena_deep_preview(request):
    """POST /api/arena/deep-preview  -  LLM-basierte Datei-Vorhersage.

    Fallback fuer die schnelle Regex-Vorschau (extractFilePreview in arena.js):
    Wenn das Regex-Muster keine plausiblen Pfade findet, ruft das Frontend
    diesen Endpoint ("tiefe Analyse nur bei Bedarf"). Antwortet mit einer
    kurzen Liste wahrscheinlich betroffener Dateien, basierend auf Topic +
    Aktionstexten.

    Request JSON: {actions: [str, ...], topic: str}
    Response JSON: {ok, files:[...], cached:bool, mode:"llm"|"empty"}
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "invalid json"}, status=400)
    actions = [str(a or "").strip() for a in (body.get("actions") or []) if str(a or "").strip()]
    topic = str(body.get("topic") or "").strip()
    if not actions:
        return web.json_response({"ok": True, "files": [], "cached": False, "mode": "empty"})

    import hashlib as _hl
    key_src = topic + "\n---\n" + "\n".join(actions)
    cache_key = _hl.sha1(key_src.encode("utf-8", errors="replace")).hexdigest()
    now_ts = time.time()
    hit = _DEEP_PREVIEW_CACHE.get(cache_key)
    if hit and (now_ts - hit[0]) < _DEEP_PREVIEW_TTL:
        return web.json_response({"ok": True, "files": hit[1], "cached": True, "mode": "llm"})

    act_block = "\n".join(f"- {a}" for a in actions[:20])
    prompt = (
        "Du analysierst Bot-Arena-Aktionen und schaetzt, welche Dateien im Repo "
        "C:/Users/Robin/Jarvis davon betroffen sein werden.\n\n"
        f"Thema: {topic or '(unbekannt)'}\n\n"
        f"Aktionen:\n{act_block}\n\n"
        "Antworte AUSSCHLIESSLICH mit einer JSON-Liste relativer Pfade "
        "(max. 15). Keine Prosa, keine Markierungen, nur gueltiges JSON.\n"
        'Beispiel: ["Mission-Control/server.py", "Mission-Control/static/js/arena.js"]'
    )
    cmd = [
        CLAUDE_EXE, "-p",
        "--model", "sonnet",
        "--output-format", "text",
        "--tools", "",
        "--max-turns", "1",
    ]
    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            stdin=asyncio.subprocess.PIPE,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        stdout_data, _ = await asyncio.wait_for(
            proc.communicate(input=prompt.encode("utf-8")), timeout=60
        )
        raw = stdout_data.decode("utf-8", errors="replace").strip()
        import re as _re
        m = _re.search(r"\[[^\[\]]*\]", raw, _re.DOTALL)
        files: list = []
        if m:
            try:
                parsed = json.loads(m.group(0))
                if isinstance(parsed, list):
                    files = [str(x).strip().replace("\\", "/").lstrip("./")
                             for x in parsed if isinstance(x, str) and x.strip()]
            except Exception:
                files = []
        seen = set()
        uniq: list = []
        for f in files:
            if f and f not in seen and len(f) > 2 and "." in f:
                seen.add(f)
                uniq.append(f)
            if len(uniq) >= 15:
                break
        _DEEP_PREVIEW_CACHE[cache_key] = (now_ts, uniq)
        return web.json_response({"ok": True, "files": uniq, "cached": False, "mode": "llm"})
    except asyncio.TimeoutError:
        return web.json_response({"ok": False, "error": "timeout"}, status=504)
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)[:200]}, status=500)


async def api_arena_continue(request):
    """POST /api/arena/rooms/{id}/continue  -  Inject executor summary and resume discussion."""
    room_id = request.match_info["id"]
    body = await request.json()
    context_summary = body.get("context", "").strip()

    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        raise web.HTTPNotFound(reason="Room not found")

    inject_text = f"Folgende Schritte wurden gerade umgesetzt:\n\n{context_summary}\n\nBitte diskutiert weiter — was ist als naechstes sinnvoll?"
    msg = {
        "role": "robin",
        "name": "Robin",
        "content": inject_text,
        "timestamp": datetime.now().isoformat(),
        "color": "#f59e0b",
    }
    room["messages"].append(msg)
    _save_arena_rooms(rooms)
    await _broadcast_all_ws({"type": "arena.message", "room_id": room_id, "message": msg})

    # Re-start discussion from where we left off — but not if session went stale
    if room_id not in _arena_rooms or not _arena_rooms[room_id].get("running"):
        if _room_is_stale(room):
            stale_msg = {
                "role": "system",
                "name": "System",
                "content": "Session war >5 min still — kein automatischer Neustart. [Start] klicken um fortzufahren.",
                "timestamp": datetime.now().isoformat(),
                "color": "#6b7280",
            }
            room["messages"].append(stale_msg)
            _save_arena_rooms(rooms)
            await _broadcast_all_ws({"type": "arena.message", "room_id": room_id, "message": stale_msg})
        else:
            _spawn_for_room(room_id, rooms)

    return web.json_response({"ok": True})


ACTIONS_LOG = DATA_DIR / "actions.jsonl"
PREVIEW_OUTCOMES_LOG = DATA_DIR / "preview-outcomes.jsonl"

# Spiegelt die Vorschau-Regex aus static/js/arena.js (extractFilePreview).
# Wenn sich die Frontend-Regex aendert, MUSS diese hier mitziehen, sonst
# driftet die Lern-Statistik auseinander.
_PREVIEW_EXTS = (
    "py|js|ts|tsx|jsx|mjs|cjs|md|json|html|htm|css|scss|yaml|yml|sh|ps1|"
    "bat|toml|ini|cfg|txt|xml|svg|vue|rs|go|java|cpp|hpp|c|h|rb|php|sql|env"
)
_PREVIEW_PATH_RE = re.compile(r"([A-Za-z0-9_./\\-]+\.(?:" + _PREVIEW_EXTS + r"))\b")
_PREVIEW_DIR_RE = re.compile(r"(?:^|\s|`|\")([A-Za-z0-9_-]+/[A-Za-z0-9_./\\-]+)(?=\s|`|\"|$|,|;|\))")


def _predict_files_from_action(action_text: str) -> list:
    """Dieselbe Heuristik wie die UI-Vorschau — damit wir serverseitig
    nachrechnen koennen, was die schnelle Vorschau *geraten* haette."""
    files = set()
    txt = str(action_text or "")
    for m in _PREVIEW_PATH_RE.finditer(txt):
        p = m.group(1).replace("\\", "/").lstrip("./")
        if len(p) > 2 and not re.match(r"^\d+\.", p):
            files.add(p)
    for m in _PREVIEW_DIR_RE.finditer(txt):
        p = m.group(1).replace("\\", "/")
        if "." in p and len(p) > 3:
            files.add(p)
    return sorted(files)


def _log_preview_outcome(*, room_id: str, approval_id: str, action_idx: int,
                         action: str, changed_files: list, ok: bool) -> None:
    """Vergleicht die Vorschau-Schaetzung mit den Dateien, die der Executor
    tatsaechlich angefasst hat, und schreibt das Urteil nach
    data/preview-outcomes.jsonl — damit das System lernt, wann die schnelle
    Vorschau versagt (False-Negative: Executor aendert Dateien, die keiner
    erraten hat; Wrong: Vorschau hat komplett anderswohin gedeutet).

    Verdicts:
      - ok        : Vorschau deckt alle geaenderten Dateien ab
      - partial   : Vorschau hat einen Teil der Dateien vorhergesagt
      - miss      : Vorschau hat Dateien vorhergesagt, aber keine davon wurde angefasst
      - empty     : Vorschau-Heuristik fand nichts, Executor hat aber geschrieben
      - no_change : Executor hat nichts angefasst (separat auswertbar)
    """
    try:
        predicted = _predict_files_from_action(action)
        changed = [c for c in (changed_files or []) if c]
        pred_set = {p.lower() for p in predicted}
        changed_set = {c.lower() for c in changed}
        overlap = pred_set & changed_set

        if not changed:
            verdict = "no_change"
        elif not predicted:
            verdict = "empty"
        elif not overlap:
            verdict = "miss"
        elif overlap == changed_set:
            verdict = "ok"
        else:
            verdict = "partial"

        # Nur Versagens-Faelle + seltene ok-Stichproben loggen, um das Log
        # schlank zu halten. ok wird nur bei ungewoehnlich grossen Runs
        # (>=3 Dateien) aufgenommen - als Positiv-Beleg.
        should_log = verdict != "ok" or len(changed) >= 3

        if not should_log:
            return

        payload = {
            "ts": datetime.now().isoformat(),
            "event": "preview.outcome",
            "verdict": verdict,
            "room_id": room_id,
            "approval_id": approval_id,
            "action_idx": action_idx,
            "action": (action or "")[:500],
            "predicted": predicted[:30],
            "changed": changed[:30],
            "predicted_count": len(predicted),
            "changed_count": len(changed),
            "missed": sorted(changed_set - pred_set)[:20],
            "extra": sorted(pred_set - changed_set)[:20],
            "action_ok": bool(ok),
        }
        with open(PREVIEW_OUTCOMES_LOG, "a", encoding="utf-8") as fh:
            fh.write(json.dumps(payload, ensure_ascii=False) + "\n")
    except Exception as _e:
        try:
            print(f"[preview-outcomes.jsonl] write failed: {_e}")
        except Exception:
            pass


async def api_arena_preview_misclick(request):
    """POST /api/arena/preview-misclick — User flagt einen Vorschau-Treffer als
    daneben ("Fehlklick"). Wir appenden das Signal ans preview-outcomes.jsonl,
    damit die schnelle Regex-Heuristik (extractFilePreview) spaeter nachgeschaerft
    werden kann. Das ist der *User-Kanal* zum bereits bestehenden Auto-Verdict.

    Request JSON:
      { file: str, room_id?: str, actions?: [str,...], reason?: str }
    Response: { ok: true }
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "invalid json"}, status=400)
    file_clicked = str(body.get("file") or "").strip()
    if not file_clicked:
        return web.json_response({"ok": False, "error": "file missing"}, status=400)
    room_id = str(body.get("room_id") or "").strip()
    reason = str(body.get("reason") or "").strip()[:200]
    actions_raw = body.get("actions") or []
    if not isinstance(actions_raw, list):
        actions_raw = []
    actions = [str(a or "").strip() for a in actions_raw if str(a or "").strip()][:20]
    action_blob = "\n".join(actions)[:2000]
    predicted = _predict_files_from_action(action_blob) if action_blob else []
    payload = {
        "ts": datetime.now().isoformat(),
        "event": "preview.misclick",
        "verdict": "misclick",
        "room_id": room_id,
        "file": file_clicked.replace("\\", "/"),
        "reason": reason,
        "predicted": predicted[:30],
        "predicted_count": len(predicted),
        "actions": actions,
    }
    try:
        with open(PREVIEW_OUTCOMES_LOG, "a", encoding="utf-8") as fh:
            fh.write(json.dumps(payload, ensure_ascii=False) + "\n")
    except Exception as e:
        return web.json_response({"ok": False, "error": str(e)[:200]}, status=500)
    return web.json_response({"ok": True})


async def api_arena_preview_stats(request):
    """GET /api/arena/preview-stats — zeigt, was die schnelle Vorschau am oeftesten
    daneben legt. Liest die letzten N Zeilen aus preview-outcomes.jsonl und aggregiert
    Fehlklicks + Auto-Miss-Verdicts nach Dateiendung und Token.
    """
    try:
        limit = int(request.query.get("limit", "500"))
    except Exception:
        limit = 500
    limit = max(50, min(limit, 5000))
    lines: list[str] = []
    try:
        with open(PREVIEW_OUTCOMES_LOG, "r", encoding="utf-8") as fh:
            lines = fh.readlines()[-limit:]
    except FileNotFoundError:
        return web.json_response({"ok": True, "count": 0, "misclicks": [], "verdicts": {}})
    misclick_files: dict[str, int] = {}
    verdict_counts: dict[str, int] = {}
    total = 0
    for ln in lines:
        try:
            ev = json.loads(ln)
        except Exception:
            continue
        total += 1
        v = ev.get("verdict") or ev.get("event") or "?"
        verdict_counts[v] = verdict_counts.get(v, 0) + 1
        if ev.get("event") == "preview.misclick":
            f = (ev.get("file") or "").lower()
            if f:
                misclick_files[f] = misclick_files.get(f, 0) + 1
    top_misclicks = sorted(misclick_files.items(), key=lambda kv: -kv[1])[:20]
    return web.json_response({
        "ok": True,
        "count": total,
        "misclicks": [{"file": f, "count": c} for f, c in top_misclicks],
        "verdicts": verdict_counts,
    })


def _append_action_log(event: dict) -> None:
    """Append a single action event to data/actions.jsonl (best-effort, non-blocking)."""
    try:
        payload = {"ts": datetime.now().isoformat(), **event}
        with open(ACTIONS_LOG, "a", encoding="utf-8") as fh:
            fh.write(json.dumps(payload, ensure_ascii=False) + "\n")
    except Exception as _e:
        try:
            print(f"[actions.jsonl] write failed: {_e}")
        except Exception:
            pass


async def _update_approval_action(room_id: str, approval_id: str, action_idx: int, status: str, result: str = ""):
    """Persist per-action status change and broadcast to UI."""
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    action_text = ""
    already_feedback = False
    if room:
        for m in room.get("messages", []):
            if m.get("approval_id") == approval_id and m.get("role") == "approval":
                if 0 <= action_idx < len(m.get("actions", [])):
                    prev_status = m["actions"][action_idx].get("status")
                    already_feedback = bool(m["actions"][action_idx].get("_fail_feedback_sent"))
                    m["actions"][action_idx]["status"] = status
                    if result:
                        m["actions"][action_idx]["result"] = result
                    action_text = m["actions"][action_idx].get("text", "")
                    # Claim-Flag setzen bevor wir den inject-Broadcast abschicken,
                    # damit ein zweiter paralleler _update-Call keine zweite Nachricht triggert.
                    if status == "failed" and not already_feedback and prev_status != "failed":
                        m["actions"][action_idx]["_fail_feedback_sent"] = True
                break
        _save_arena_rooms(rooms)
    _append_action_log({
        "event": f"action.{status}",
        "room_id": room_id,
        "approval_id": approval_id,
        "action_idx": action_idx,
        "action": action_text,
        "status": status,
        "result": result[:500] if result else "",
    })
    await _broadcast_all_ws({
        "type": "arena.action.status",
        "room_id": room_id,
        "approval_id": approval_id,
        "action_idx": action_idx,
        "status": status,
        "result": result,
    })

    # Rueck-Inject bei failed: einmalig pro Action, serverseitig (frueher im
    # Frontend — das feuerte einmal pro verbundenem Client, deshalb
    # tauchte die Nachricht bei mehreren offenen Tabs/Geraeten doppelt auf).
    if status == "failed" and room and action_text and not already_feedback:
        m_bot = re.match(r"^\[([A-Za-z][A-Za-z0-9_-]*)\]\s*(.+)", action_text.strip())
        bot_name = m_bot.group(1) if m_bot else ""
        display_text = (m_bot.group(2) if m_bot else action_text).strip()
        err_lines = [l for l in (result or "").splitlines() if l.strip()]
        err_reason = (err_lines[0][:150] if err_lines else "unbekannter Fehler")
        err_excerpt = " | ".join(err_lines[1:4])[:200] if len(err_lines) > 1 else ""
        target = f"@{bot_name}" if bot_name else "Wer auch immer die Aktion vorgeschlagen hat"
        feedback_text = (
            f"[SYSTEM — FAIL-FEEDBACK] {target}: Deine Aktion \"{display_text[:100]}\" ist fehlgeschlagen.\n"
            f"Grund: {err_reason}"
            + (f"\nDetails: {err_excerpt}" if err_excerpt else "")
            + "\n→ Schlage eine KLEINERE, schrittweise Variante vor — nicht das Gleiche nochmal. "
              "Was ist der minimale erste Schritt, der sicher funktionieren würde?"
        )
        feedback_msg = {
            "role": "robin",
            "name": "Robin",
            "content": feedback_text,
            "timestamp": datetime.now().isoformat(),
            "color": "#f59e0b",
        }
        rooms2 = _load_arena_rooms()
        room2 = _get_arena_room(rooms2, room_id)
        if room2:
            room2.setdefault("messages", []).append(feedback_msg)
            _save_arena_rooms(rooms2)
            await _broadcast_all_ws({
                "type": "arena.message",
                "room_id": room_id,
                "message": feedback_msg,
            })


def _detect_file_changes_since(start_ts: float, max_files: int = 20) -> list:
    """Walk the Jarvis repo and return files modified/created since start_ts.
    Skips noisy directories (caches, logs, auto-memory, node_modules).
    Cheap enough for 10-minute executor runs - only stats existing files.
    """
    root = "C:/Users/Robin/Jarvis"
    skip_dirs = {
        ".git", "node_modules", "__pycache__", ".venv", "venv",
        "data",  # Mission-Control/data holds logs+sessions that mutate on their own
        "Jarvis-Memory",  # Obsidian auto-updates
        ".claude-flow", ".swarm", ".hive-mind",
    }
    changed = []
    try:
        for dirpath, dirnames, filenames in os.walk(root):
            dirnames[:] = [d for d in dirnames if d not in skip_dirs and not d.startswith(".")]
            for fn in filenames:
                if fn.endswith((".pyc", ".log", ".tmp", ".bak", ".swp", ".swo", ".lock", ".pid", ".cache")) or fn.startswith("."):
                    continue
                full = os.path.join(dirpath, fn)
                try:
                    mt = os.path.getmtime(full)
                except OSError:
                    continue
                if mt >= start_ts:
                    rel = os.path.relpath(full, root).replace("\\", "/")
                    changed.append(rel)
                    if len(changed) >= max_files:
                        return changed
    except Exception:
        pass
    return changed


# --- Executor Budget-Stufen ------------------------------------------------
# Kleine Aufgaben fressen weniger: günstiges Modell, weniger Turns, kürzerer
# Timeout. Grosse bekommen Opus + voller Spielraum. Heuristik greift auf
# Action-Text-Länge und Schlüsselwörter zu; bei Unsicherheit "medium".
EXECUTOR_BUDGETS = {
    "small":  {"model": "haiku",  "max_turns": 10, "timeout": 180, "tool_calls_hint": 8},
    "medium": {"model": "sonnet", "max_turns": 20, "timeout": 360, "tool_calls_hint": 18},
    "large":  {"model": "opus",   "max_turns": 40, "timeout": 900, "tool_calls_hint": 35},
}

_LARGE_HINTS = (
    "refactor", "umbau", "umschreiben", "neu aufbauen", "integrier",
    "architektur", "migration", "rewrite", "komplett", "end-to-end",
    "mehrere module", "system", "pipeline", "neu strukturier",
    "spracherkennung", "voice pipeline", "speech recognition", "audio pipeline",
    "transkription einricht", "tts integrier", "voice system",
    "chat-reparatur", "chat reparatur", "chat repair", "chat-fix",
)
_SMALL_HINTS = (
    "typo", "rechtschreib", "rename", "umbenenn", "log-message", "log message",
    "kommentar", "konstante", "config-wert", "zeile", "flag umstell",
    "import ", "einzeiler", "eine zeile", "feld hinzuf", "farbe",
    "stimme wechsel", "voice switch", "voice id", "stimme umstell",
    "mikrofon", "mikro ein", "mikro aus", "mute", "unmute",
    "tts aktiv", "tts deaktiv", "tts umstell", "voice key",
)


def _classify_action_size(action: str) -> str:
    """Grobe Heuristik: small/medium/large. Reihenfolge:
    1. Explizite Keywords (large überstimmt small bei Konflikt)
    2. Länge der Action-Beschreibung als Tiebreaker.
    """
    text = (action or "").strip().lower()
    if not text:
        return "medium"
    has_large = any(h in text for h in _LARGE_HINTS)
    has_small = any(h in text for h in _SMALL_HINTS)
    if has_large and not has_small:
        return "large"
    if has_small and not has_large:
        return "small"
    # Längen-Fallback: sehr kurze Beschreibungen → klein, sehr lange → groß
    n = len(text)
    if n < 60:
        return "small"
    if n > 260:
        return "large"
    return "medium"


async def _run_single_action(topic: str, project_context: str, action: str, executor_bot: dict | None = None, changed_out: list | None = None, size: str | None = None, predecessor_failure: str | None = None) -> tuple:
    """Run one action via Claude subprocess. Returns (ok, result_text).

    Verification: action only counts as 'done' when at least one file was
    actually modified during the run. Pure-text responses are downgraded to
    failure so the arena UI does not falsely claim 'Fertig'.

    executor_bot: {name, persona} - if given, the bot takes on this role so
    the executor persona rotates across the team (each bot carries weight).

    size: 'small' | 'medium' | 'large' — steuert Modell, Max-Turns und Timeout.
    Wenn None, wird per Heuristik aus dem Action-Text geraten.
    """
    tier = size if size in EXECUTOR_BUDGETS else _classify_action_size(action)
    budget = EXECUTOR_BUDGETS[tier]
    bot_provider = (executor_bot or {}).get("provider") or "claude-cli"
    bot_model_cfg = (executor_bot or {}).get("model") or ""
    if executor_bot:
        role_intro = (
            f"Du bist {executor_bot.get('name','Executor')} und uebernimmst diesmal die Umsetzer-Rolle "
            f"aus der Bot-Diskussion. Deine Persona: {executor_bot.get('persona','')}\n"
            f"Robin hat eine einzelne Aktion genehmigt — setz sie in deinem Stil um, "
            f"aber bleib pragmatisch und lieferbereit.\n\n"
        )
    else:
        role_intro = "Du bist der Executor-Bot. Robin hat eine einzelne Aktion aus einer Bot-Diskussion genehmigt.\n\n"
    if predecessor_failure:
        role_intro += (
            f"HINWEIS: Ein vorheriger Executor-Bot hat diese Aktion bereits versucht und ist gescheitert.\n"
            f"Grund des Scheiterns: {predecessor_failure[:400]}\n"
            f"Vermeide denselben Fehler. Geh die Sache anders an.\n\n"
        )
    prompt = (
        role_intro + (
        f"Diskussions-Thema: {topic}\n\n"
        f"=== PROJEKT-KONTEXT ===\n{project_context}\n=== ENDE KONTEXT ===\n\n"
        f"Diese konkrete Aktion sollst du jetzt umsetzen:\n\n>>> {action} <<<\n\n"
        f"Du hast Zugriff auf alle Tools (Read, Write, Edit, Bash, Glob, Grep).\n"
        f"Arbeitsverzeichnis: C:/Users/Robin/Jarvis\n\n"
        f"GUARDRAIL — Klarer Auftrag zwingend:\n"
        f"- Die oben in >>> <<< stehende Aktion ist dein EINZIGER Auftrag. Nichts anderes.\n"
        f"- Ist der Auftrag vage, mehrdeutig oder bloss ein Diskussionsfragment ohne konkretes Ziel: "
        f"KEINE Dateien anfassen (kein Write, kein Edit, kein Bash-Schreibkommando). "
        f"Stattdessen in 1-3 Saetzen erklaeren was unklar ist und welche Praezisierung Robin liefern soll.\n"
        f"- Bei unklarem Auftrag: nur Read/Glob/Grep zum Nachsehen, nicht zum Umbauen.\n"
        f"- Scope-Creep verboten: keine 'waehrend ich dabei bin'-Refactors, keine Zusatz-Features, "
        f"keine Aufraeumarbeiten ausserhalb der genehmigten Aktion.\n\n"
        f"WICHTIG — Pragmatischer Modus (Budget-Stufe: {tier.upper()}):\n"
        f"- Du hast nur ~{budget['tool_calls_hint']} Tool-Calls Budget. Plane sparsam.\n"
        f"- Wenn die Aktion zu gross fuer einen Schritt ist: mach den ERSTEN sinnvollen Teil "
        f"(z.B. ersten Prototyp, erstes Geruest, ersten Endpoint) und beschreibe was noch fehlt.\n"
        f"- Lieber 1 fertiges Teilstueck als 0 fertige Sachen.\n"
        f"- Bei Code-Aenderungen: KEINE riesigen Refactors. Minimaler Diff der die Aktion umsetzt.\n\n"
        f"Am Ende antworte mit einer 1-3 Satz Zusammenfassung: was du getan hast UND was noch offen ist.\n"
        f"Wenn die Aktion unmoeglich oder unklar ist: kurz erklaeren warum und Vorschlag was Robin tun soll."
        )
    )

    # --- Provider-Dispatch (Option A: echtes Multi-Provider-Routing) -----
    # Wenn der Orchestrator [ACTION:GPT] schreibt und GPT als codex-cli
    # konfiguriert ist, läuft der Executor wirklich via Codex-CLI. Bei
    # [ACTION:Gemini] via Gemini-CLI. Nur claude-cli nimmt den existing
    # Pfad mit den Budget-Tiers (haiku/sonnet/opus + max-turns-Guard).
    if bot_provider in ("codex-cli", "gemini-cli"):
        bot_label = (executor_bot or {}).get("name", bot_provider)
        start_ts_routed = time.time() - 1.0
        try:
            out_text, usage_r, err_r = await _trialog_exec_turn(
                bot={
                    "provider": bot_provider,
                    "model": bot_model_cfg,
                    "name": bot_label,
                    "persona": (executor_bot or {}).get("persona", ""),
                },
                prompt_text=prompt,
                system_prompt="",  # role_intro steckt schon im prompt
                cwd="C:/Users/Robin/Jarvis",
                timeout_sec=budget["timeout"],
                read_only=False,
            )
        except Exception as e:
            return False, f"({bot_provider} crashed: {type(e).__name__}: {e})"
        if err_r:
            return False, f"({bot_provider}: {err_r})"
        if not (out_text or "").strip():
            return False, f"({bot_provider}: leere Antwort)"
        # File-change detection ist provider-unabhängig (mtime-basiert)
        loop = asyncio.get_event_loop()
        changed_files = await loop.run_in_executor(None, _detect_file_changes_since, start_ts_routed)
        if changed_out is not None:
            changed_out.extend(changed_files or [])
        if not changed_files:
            return None, (
                out_text
                + f"\n\n(Haengend: {bot_provider} hat geantwortet, aber keine Datei geaendert.)"
            )
        files_note = "\n\nGeaenderte Dateien (" + str(len(changed_files)) + "): " + ", ".join(changed_files[:10])
        if len(changed_files) > 10:
            files_note += ", ..."
        return True, f"[via {bot_provider}/{bot_model_cfg or 'default'}]\n" + out_text + files_note

    # --- Claude-CLI (default/fallback) -----------------------------------
    cmd = [
        CLAUDE_EXE, "-p", prompt,
        "--model", budget["model"],
        "--output-format", "text",
        "--max-turns", str(budget["max_turns"]),
        "--permission-mode", "bypassPermissions",
    ]
    proc = None
    # Snapshot before run - allow a tiny clock-skew tolerance so edits that land
    # in the same second as start_ts still register.
    start_ts = time.time() - 1.0
    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            stdin=asyncio.subprocess.DEVNULL,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        stdout_data, stderr_data = await asyncio.wait_for(proc.communicate(), timeout=budget["timeout"])
        stdout = stdout_data.decode("utf-8", errors="replace").strip()
        stderr = stderr_data.decode("utf-8", errors="replace").strip()
        rc = proc.returncode

        # Detect Claude usage/rate-limit responses - these must NEVER count as success.
        # Patterns seen in practice: "You've hit your limit · resets 2am (Europe/Berlin)",
        # "Claude AI usage limit reached", "rate limit", "5-hour limit".
        combined = f"{stdout}\n{stderr}"
        lowered = combined.lower()
        limit_markers = (
            "you've hit your limit",
            "you have hit your limit",
            "usage limit reached",
            "usage limit exceeded",
            "rate limit",
            "5-hour limit",
            "quota exceeded",
            "resets " ,  # "resets 2am ..." - only appears in limit responses
        )
        # Require that the output is ALSO short (real work produces paragraphs).
        # A limit message is typically one short line; guard against false positives.
        if any(m in lowered for m in limit_markers) and len(stdout) < 400:
            first_line = stdout.splitlines()[0].strip() if stdout else stderr[:200]
            return False, f"(Claude-Limit erreicht, keine Umsetzung: {first_line})"

        # Verify the run actually touched the repo. Run this once and reuse
        # the result for both max-turns and normal-exit paths.
        loop = asyncio.get_event_loop()
        changed_files = await loop.run_in_executor(None, _detect_file_changes_since, start_ts)
        if changed_out is not None:
            changed_out.extend(changed_files or [])
        if changed_files:
            files_note = "\n\nGeaenderte Dateien (" + str(len(changed_files)) + "): " + ", ".join(changed_files[:10])
            if len(changed_files) > 10:
                files_note += ", ..."
        else:
            files_note = ""

        hit_max = "Reached max turns" in stdout
        if hit_max:
            cleaned = "\n".join(
                l for l in stdout.splitlines()
                if "Reached max turns" not in l and not l.strip().startswith("Error:")
            ).strip()
            if cleaned:
                if not changed_files:
                    return False, cleaned + f"\n\n(Max-Turns [{tier}] erreicht UND keine Datei geaendert - Pruefung fehlgeschlagen.)"
                return True, f"[Stufe: {tier} · {budget['model']}]\n" + cleaned + files_note + f"\n\n(Hinweis: Max-Turns [{tier}] erreicht - moeglicherweise nicht komplett)"
            return False, f"(Aktion zu gross fuer Executor-Stufe [{tier}] - {budget['max_turns']} Tool-Calls aufgebraucht ohne Ergebnis. Bitte in kleinere Schritte aufteilen oder als 'large' markieren.)"

        if not stdout:
            err_preview = stderr[:300] if stderr else "(kein stderr)"
            return False, f"(Leere Antwort, rc={rc}, stderr: {err_preview})"

        # Primary verification: if the action claims 'Fertig' but nothing on disk
        # actually changed, it was just talk. Mark as 'haengend' (proposed but not implemented).
        if not changed_files:
            return None, (
                stdout
                + "\n\n(Haengend: keine Datei geaendert. "
                + "Executor hat nur geantwortet, aber nichts geschrieben.)"
            )
        return True, f"[Stufe: {tier} · {budget['model']}]\n" + stdout + files_note
    except asyncio.TimeoutError:
        if proc and proc.returncode is None:
            proc.kill()
            await proc.wait()
        mins = budget["timeout"] // 60
        return False, f"(Timeout nach {mins}min - Stufe [{tier}])"
    except Exception as e:
        return False, f"(Fehler: {e})"


def _persist_executor_results(room_id: str, topic: str, approval_id: str,
                               results: list, active_bot: dict | None) -> None:
    """Spiegelt die Executor-Ergebnisse aus der Arena in persistente Stores:

    1) Smart-Tasks (data/smart-tasks.json): jede Aktion wird als abgeschlossener
       Task angelegt (status=done bei OK, failed sonst). So erscheinen fertige
       Bot-Umsetzungen in der Smart-Tasks-Ansicht und sind durchsuchbar.
    2) Obsidian-Tageslog (Jarvis-Memory/YYYY-MM-DD.md): kurzer Eintrag mit
       Aktion + OK/FAIL + Kurzergebnis. Taucht im Brain-Index auf.

    Best-effort — Fehler werden geloggt, aber nicht geworfen.
    """
    if not results:
        return
    exec_name = active_bot.get("name") if active_bot else "Executor"
    now_iso = datetime.now().isoformat()

    # --- 1. Smart-Tasks spiegeln -----------------------------------------
    try:
        tasks = _load_smart_tasks()
        for idx, item in enumerate(results):
            if item is None:
                continue
            action, ok, result_text = item
            title = (action or "").strip().splitlines()[0][:80] or "Arena-Aktion"
            tasks.append({
                "id": str(uuid.uuid4()),
                "title": title,
                "description": action or "",
                "original_request": f"Arena-Room {room_id} — {topic}",
                "chat_context": f"approval_id={approval_id}, executor={exec_name}",
                "priority": "medium",
                "status": "done" if ok is True else ("haengend" if ok is None else "failed"),
                "created": now_iso,
                "started": now_iso,
                "completed": now_iso,
                "output": result_text or "",
                "error": "" if ok is True else (result_text or ""),
                "source": "arena",
                "room_id": room_id,
                "approval_id": approval_id,
                "executor": exec_name,
            })
        _save_smart_tasks(tasks)
    except Exception as e:
        print(f"[arena] smart-tasks persist failed: {e}")

    # --- 2. Obsidian Tageslog ergaenzen ----------------------------------
    try:
        memory_dir = Path("E:/OneDrive/AI/Obsidian-Vault/Jarvis-Brain/Jarvis-Memory")
        if memory_dir.exists():
            today = datetime.now().strftime("%Y-%m-%d")
            daily = memory_dir / f"{today}.md"
            lines = [
                f"\n## Arena-Executor ({exec_name}) — {datetime.now().strftime('%H:%M')}",
                f"*Thema:* {topic or '-'} · *Room:* `{room_id}` · *Approval:* `{approval_id}`",
                "",
            ]
            for item in results:
                if item is None:
                    continue
                action, ok, result_text = item
                mark = "OK" if ok else "FAIL"
                short = (result_text or "").strip().splitlines()
                snippet = short[0][:200] if short else ""
                lines.append(f"- **[{mark}]** {action}")
                if snippet:
                    lines.append(f"    -> {snippet}")
            lines.append("")
            block = "\n".join(lines)
            if daily.exists():
                with open(daily, "a", encoding="utf-8") as fh:
                    fh.write(block)
            else:
                daily.write_text(f"# {today}\n{block}", encoding="utf-8")
    except Exception as e:
        print(f"[arena] memory persist failed: {e}")


ARENA_EXECUTOR_MAX_PARALLEL = 5  # Hard cap: parallele Executor-Aktionen — schuetzt Context/Token-Budget

async def _run_executor(room_id: str, approval_id: str, actions: list, skip_done: bool = False):
    """Execute approved actions in parallel (max ARENA_EXECUTOR_MAX_PARALLEL at once) with per-action status broadcasts.

    skip_done: if True, actions whose persisted status is already terminal
    (done/failed/skipped) are skipped — used by the executor auto-resume hook
    so a restarted server doesn't re-run already-finished actions.
    """
    _activity_set("arena", active=True)
    await _broadcast_all_ws({"type": "arena.started", "room_id": room_id})
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        return

    topic = room.get("topic", "")
    project_context = _gather_arena_context()
    results = [None] * len(actions)
    sem = asyncio.Semaphore(ARENA_EXECUTOR_MAX_PARALLEL)

    _TERMINAL_STATUSES = {"done", "failed", "skipped"}
    _done_idx: set = set()
    if skip_done:
        for _m in room.get("messages", []):
            if _m.get("role") == "approval" and _m.get("approval_id") == approval_id:
                for _i, _a in enumerate(_m.get("actions") or []):
                    if (_a or {}).get("status") in _TERMINAL_STATUSES:
                        _done_idx.add(_i)
                        results[_i] = (_a.get("text", ""), _a.get("status") == "done", _a.get("result", ""))
                break
        if _done_idx:
            print(f"[ArenaExecResume] approval={approval_id[:8]} skipping {len(_done_idx)} already-terminal action(s)")

    # Bot-Zuordnung: PRIMÄR über [BotName]-Prefix in der Action selbst
    # (der Orchestrator schreibt [ACTION:BotName] und das Frontend propagiert
    # es als "[BotName] text" in actions-Array). Das bedeutet: wenn der
    # Orchestrator [ACTION:GPT] schreibt, wird die Aktion wirklich von Codex
    # ausgeführt, [ACTION:Gemini] → Gemini-CLI, [ACTION:Claude] → Claude-CLI.
    # Fallback nur wenn Action keinen Prefix hat: Rotations-Bot.
    import re as _re
    bots_list = [b for b in room.get("bots", []) if b.get("name")]
    _bots_by_name = {b.get("name", "").lower(): b for b in bots_list}
    rot_idx = int(room.get("executor_rotation_idx", 0))
    fallback_bot = bots_list[rot_idx % len(bots_list)] if bots_list else None
    if bots_list:
        room["executor_rotation_idx"] = (rot_idx + 1) % len(bots_list)

    def _bot_for_action(action_text: str) -> dict | None:
        """Parse [BotName]-Prefix aus der Action, lookup in room.bots."""
        m = _re.match(r"^\[([A-Za-z][A-Za-z0-9_-]*)\]", (action_text or "").strip())
        if m:
            bot = _bots_by_name.get(m.group(1).lower())
            if bot:
                return bot
        return fallback_bot

    # Keep old variable name for downstream logging (executor_checkpoint)
    active_bot = fallback_bot
    # Executor-Checkpoint: wird bei jedem Action-Complete gebumpt. Ueberlebt
    # MC-Crash via _save_arena_rooms, der naechste Seed sieht "completed<total"
    # und kann Bots informieren dass ein Batch unterbrochen war.
    room["executor_checkpoint"] = {
        "approval_id": approval_id,
        "completed": 0,
        "total": len(actions),
        "status": "running",
        "started_at": datetime.now().isoformat(),
        "bot": (active_bot or {}).get("name"),
    }
    _save_arena_rooms(rooms)

    async def run_one(idx, action):
        if idx in _done_idx:
            return
        async with sem:
            await _update_approval_action(room_id, approval_id, idx, "running")
            changed_seen: list = []
            # Route nach dem im Action-Tag angegebenen Bot (echter Provider-
            # Dispatch), Fallback auf Rotations-Bot wenn kein [Name]-Prefix.
            per_action_bot = _bot_for_action(action)
            ok, result_text = await _run_single_action(topic, project_context, action, per_action_bot, changed_out=changed_seen)
            status = "done" if ok is True else ("haengend" if ok is None else "failed")
            await _update_approval_action(room_id, approval_id, idx, status, result_text[:500])
            results[idx] = (action, ok, result_text)
            # Checkpoint bump — reload/save (race-tolerant: last-write-wins reicht als Fortschrittsanzeige)
            try:
                _rooms_cp = _load_arena_rooms()
                _room_cp = _get_arena_room(_rooms_cp, room_id)
                if _room_cp is not None:
                    _cp = _room_cp.get("executor_checkpoint") or {}
                    if _cp.get("approval_id") == approval_id:
                        _cp["completed"] = int(_cp.get("completed", 0)) + 1
                        _cp["last_action_at"] = datetime.now().isoformat()
                        _room_cp["executor_checkpoint"] = _cp
                        _save_arena_rooms(_rooms_cp)
            except Exception as _cp_err:
                print(f"[arena] checkpoint bump failed: {_cp_err}")
            # Fehlgeschlagene Aktionen in Fail-Backlog sichern (nächste Session)
            if ok is False:
                _append_fail_backlog([{
                    "kind": "failed",
                    "action": action,
                    "reason": (result_text or "")[:300],
                    "room_id": room_id,
                    "topic": topic,
                    "approval_id": approval_id,
                    "timestamp": datetime.now().isoformat(),
                }])
            # Erfolgreiche Aktionen in Raumzustand persistieren (Dedup-Basis für künftige Approvals)
            if ok:
                _rooms_done = _load_arena_rooms()
                _room_done = _get_arena_room(_rooms_done, room_id)
                if _room_done is not None:
                    _room_done.setdefault("done_actions", [])
                    _key = action.strip().lower()
                    if _key not in _room_done["done_actions"]:
                        _room_done["done_actions"].append(_key)
                    # Sync action_items status to erledigt
                    for _ai in _room_done.get("action_items", []):
                        if _ai.get("text", "").lower() == _key or _key.endswith(_ai.get("text", "").lower()):
                            _ai["status"] = "erledigt"
                    _save_arena_rooms(_rooms_done)
            # Lern-Signal: vergleicht die schnelle Vorschau-Regex mit dem, was
            # der Executor tatsaechlich angefasst hat. Miss/empty/partial landen
            # in data/preview-outcomes.jsonl und koennen spaeter ausgewertet
            # werden, um die Vorschau-Heuristik oder Prompt-Hints zu schaerfen.
            try:
                _log_preview_outcome(
                    room_id=room_id,
                    approval_id=approval_id,
                    action_idx=idx,
                    action=action,
                    changed_files=changed_seen,
                    ok=ok,
                )
            except Exception as _pe:
                print(f"[arena] preview-outcome log failed: {_pe}")

    # Split actions: "restart MC"-style entries are destructive to siblings (they
    # kill MC while other tasks are still running), so we defer them. Run regular
    # actions in parallel first, THEN restart actions serially at the very end.
    import re as _re
    _restart_re = _re.compile(r"(mission[- ]?control|\bmc\b).{0,40}(neu[- ]?start|restart)|restart.{0,20}(mission[- ]?control|\bmc\b)", _re.IGNORECASE)
    regular_items: list = []
    restart_items: list = []
    for i, a in enumerate(actions):
        if _restart_re.search(a or ""):
            restart_items.append((i, a))
        else:
            regular_items.append((i, a))
    if restart_items:
        print(f"[arena] {len(restart_items)} restart-action(s) deferred to end (would otherwise kill siblings): {[a for _, a in restart_items]}")

    # Parallel batch of non-restart actions
    if regular_items:
        await asyncio.gather(*[run_one(i, a) for i, a in regular_items])

    # Auto-complete "restart MC" actions WITHOUT actually restarting. Running
    # them used to be destructive: they'd kill MC → kill the executor task that
    # was supposed to mark its own status → action stayed "running" forever
    # and Robin saw a phantom hang. Code changes are still picked up — the
    # watchdog reads the file fresh on the next real MC restart.
    for i, a in restart_items:
        if i in _done_idx:
            continue
        await _update_approval_action(
            room_id, approval_id, i, "done",
            "Restart-Action uebersprungen — Watchdog laedt neuen Code beim naechsten MC-Restart von selbst.",
        )

    # Auto-Persist: fertige Executor-Ergebnisse wandern automatisch in Smart-Tasks
    # (als abgeschlossene Tasks mit Output) + Obsidian-Memory (Tageslog-Eintrag).
    # So sind Bot-Aktionen ausserhalb der Arena durchsuchbar. Best-effort, darf
    # den Executor-Flow nie blockieren.
    try:
        _persist_executor_results(room_id, topic, approval_id, results, active_bot)
    except Exception as _persist_err:
        print(f"[arena] auto-persist failed: {_persist_err}")

    # Final summary message to main chat
    summary_lines = []
    for item in results:
        if item is None:
            continue
        action, ok, result_text = item
        mark = "OK" if ok else "FAIL"
        summary_lines.append(f"- [{mark}] {action}\n    -> {result_text[:300]}")
    exec_name = active_bot.get("name") if active_bot else "Executor"
    exec_color = active_bot.get("color", "#10b981") if active_bot else "#10b981"
    header = (
        f"**{exec_name} uebernimmt die Umsetzung — {len(results)} Aktion(en):**"
        if active_bot
        else f"**Umsetzung abgeschlossen ({len(results)} Aktion(en)):**"
    )
    exec_msg = {
        "role": "executor",
        "name": exec_name,
        "content": header + "\n\n" + "\n".join(summary_lines),
        "timestamp": datetime.now().isoformat(),
        "color": exec_color,
        "turn": -1,
    }
    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if room:
        room["messages"].append(exec_msg)
        # Checkpoint schliessen — kein "running" mehr, Seed zeigt nichts mehr an
        _cp_end = room.get("executor_checkpoint") or {}
        if _cp_end.get("approval_id") == approval_id:
            _cp_end["status"] = "done"
            _cp_end["finished_at"] = datetime.now().isoformat()
            room["executor_checkpoint"] = _cp_end
        _save_arena_rooms(rooms)
    await _broadcast_all_ws({
        "type": "arena.message",
        "room_id": room_id,
        "message": exec_msg,
    })
    await _broadcast_all_ws({
        "type": "arena.executor",
        "room_id": room_id,
        "status": "done",
        "actions_count": len(actions),
    })
    _activity_set("arena", active=False, pending_delta=1)
    await _broadcast_all_ws({"type": "arena.stopped", "room_id": room_id})

    # Auto-Continue: "Tueftel-Team"-Verhalten. Nach Umsetzung bewerten die Bots
    # das Executor-Ergebnis sofort selbst und schlagen Nachbesserungen vor,
    # ohne dass Robin manuell weiterklicken muss. Per room["auto_continue"]=False
    # abschaltbar. Default: an.
    if room and room.get("auto_continue", True):
        done_count = sum(1 for r in results if r and r[1])
        fail_count = sum(1 for r in results if r and not r[1])
        # Build per-action output snippets so the reviewer sees the real output,
        # not just the summary badge. Truncated to 600 chars per action to keep context manageable.
        output_snippets = []
        for item in results:
            if item is None:
                continue
            _act, _ok, _txt = item
            _mark = "OK" if _ok is True else ("HAENGEND" if _ok is None else "FAIL")
            _snippet = (_txt or "(kein Output)")[:600]
            output_snippets.append(f"[{_mark}] {_act}\n---\n{_snippet}")
        review_text = (
            f"Das Executor-Ergebnis liegt vor ({done_count} OK, {fail_count} FAIL).\n\n"
            f"=== ECHTER EXECUTOR-OUTPUT ===\n"
            + "\n\n".join(output_snippets)
            + f"\n=== ENDE OUTPUT ===\n\n"
            f"Bewertet bitte kurz:\n"
            f"1) Ist das Umgesetzte wirklich so, wie wir es gemeint haben?\n"
            f"2) Was fehlt noch, damit die urspruengliche Aufgabe perfekt ist?\n"
            f"3) Falls noetig: schlagt konkrete naechste [ACTION]-Schritte vor.\n"
            f"Wenn alles passt: klar 'Fertig' sagen und begruenden."
        )
        review_msg = {
            "role": "robin",
            "name": "Auto-Review",
            "content": review_text,
            "timestamp": datetime.now().isoformat(),
            "color": "#f59e0b",
        }
        rooms2 = _load_arena_rooms()
        room2 = _get_arena_room(rooms2, room_id)
        if room2:
            room2["messages"].append(review_msg)
            _save_arena_rooms(rooms2)
            await _broadcast_all_ws({"type": "arena.message", "room_id": room_id, "message": review_msg})
            if not _arena_task_running(room_id):
                if _room_is_stale(room2):
                    stale_msg = {
                        "role": "system",
                        "name": "System",
                        "content": "Session war >5 min still — kein automatischer Neustart. [Start] klicken um fortzufahren.",
                        "timestamp": datetime.now().isoformat(),
                        "color": "#6b7280",
                    }
                    room2["messages"].append(stale_msg)
                    _save_arena_rooms(rooms2)
                    await _broadcast_all_ws({"type": "arena.message", "room_id": room_id, "message": stale_msg})
                else:
                    _spawn_for_room(room_id, rooms2)


def _extract_verdict(summary_raw: str) -> tuple[str, str]:
    """Extract [VERDICT:xxx] marker from orchestrator output.

    Returns (verdict, summary_without_marker).
    verdict is one of: ongoing, consensus, stalemate.
    Falls back to 'ongoing' if no marker found.
    """
    if not summary_raw:
        return "ongoing", summary_raw or ""
    import re as _re
    m = _re.search(r"\[VERDICT:\s*(ongoing|consensus|stalemate)\s*\]", summary_raw, _re.IGNORECASE)
    if not m:
        return "ongoing", summary_raw.strip()
    verdict = m.group(1).lower()
    cleaned = _re.sub(r"\s*\[VERDICT:[^\]]+\]\s*", "", summary_raw).strip()
    return verdict, cleaned


async def _generate_orchestrator_summary(room):
    """Generate a summary of the discussion using Claude."""
    messages = room.get("messages", [])
    if not messages:
        return "Keine Nachrichten zum Zusammenfassen."

    history_lines = [f"{m['name']}: {m['content']}" for m in messages[-40:]]
    history_block = "\n\n".join(history_lines)

    prompt = (
        f"Du bist der Orchestrator einer Bot-Diskussion. Robin ist kein Entwickler — schreib verstaendlich.\n\n"
        f"Thema: {room.get('topic', 'Freie Diskussion')}\n"
        f"Teilnehmer: {', '.join(b['name'] for b in room.get('bots', []))}\n\n"
        f"Diskussion:\n{history_block}\n\n"
        f"Schreibe ein Fazit. Halte dich genau an diese Struktur:\n\n"
        f"**Worum es geht** (1-2 Saetze, einfache Sprache, kein Technikjargon)\n\n"
        f"**Was koennen wir umsetzen**\n"
        f"Fuer JEDE moegliche Aktion eine Zeile:\n"
        f"[ACTION:BotName] <Was passiert, und was bringt das — max 12 Woerter, Alltagssprache> — Datei: <Dateiname>, Funktion: <Funktionsname oder Route>\n"
        f"BotName = Name des Bots der die Idee hauptsaechlich einbrachte (z.B. Alpha, Beta, Gamma, Delta).\n"
        f"Datei+Funktion: Wenn die Aktion Code betrifft, gib GENAU an wo sie landet (z.B. 'Datei: server.py, Funktion: _generate_orchestrator_summary' oder 'Datei: arena.js, Funktion: renderFazitWithActions'). Nutze den Kontext der mitgelieferten Dateiliste und Routes. Falls nicht eindeutig bestimmbar: weglassen.\n\n"
        f"Regeln:\n"
        f"- Keine Fachbegriffe ohne Erklaerung\n"
        f"- Sag was Robin davon hat, nicht wie es technisch laeuft\n"
        f"- Nur umsetzbare Aktionen (keine vagen Ideen)\n"
        f"- Kein 'warte auf gruenes Licht', keine Rueckfragen\n"
        f"- Deutsch, kurz, klar.\n\n"
        f"Am Ende, auf einer eigenen letzten Zeile, EXAKT einen dieser Marker:\n"
        f"[VERDICT:ongoing] - Bots bringen noch neue Punkte, weiterdiskutieren lohnt\n"
        f"[VERDICT:consensus] - Bots sind sich einig, Plan steht, Robin sollte entscheiden/freigeben\n"
        f"[VERDICT:stalemate] - Bots drehen sich im Kreis, kein Fortschritt, Robin muss eingreifen\n\n"
        f"HARTE REGEL — keine Fantasiezahlen bei Runden/Aufwand:\n"
        f"- Falls du einen Aufwand in Runden oder Iterationen nennst: HARTE OBERGRENZE 10 Runden.\n"
        f"- Niemals '20 Runden', '50 Iterationen', 'Monate' o.ae. — Arena ist auf <=10 Runden gedeckelt.\n"
        f"- Passt eine Aktion nicht in <=10 Runden: in kleinere [ACTION]-Schritte aufteilen, statt hohe Zahlen zu erfinden."
    )

    # Prompt via stdin statt argv - umgeht Windows 32K command-line limit
    cmd = [
        CLAUDE_EXE, "-p",
        "--model", "opus",
        "--output-format", "text",
        "--tools", "",
        "--max-turns", "1",
    ]

    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            stdin=asyncio.subprocess.PIPE,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        stdout_data, stderr_data = await asyncio.wait_for(
            proc.communicate(input=prompt.encode("utf-8")), timeout=180
        )
        response = stdout_data.decode("utf-8", errors="replace").strip()
        if "Reached max turns" in response:
            lines = [l for l in response.splitlines() if "Reached max turns" not in l and "Error:" not in l]
            response = "\n".join(lines).strip()
        if not response:
            err = stderr_data.decode("utf-8", errors="replace").strip()[:400]
            return f"Zusammenfassung konnte nicht generiert werden. (stderr: {err or 'leer'})"
        return response
    except Exception as e:
        return f"Fehler bei Zusammenfassung: {e}"


def _gather_arena_context() -> str:
    """Gather current project context so arena bots know what exists."""
    ctx_parts = []
    # 1. CLAUDE.md — project identity
    claude_md = Path("C:/Users/Robin/Jarvis/CLAUDE.md")
    if claude_md.exists():
        txt = claude_md.read_text(encoding="utf-8", errors="replace")[:2000]
        ctx_parts.append(f"=== CLAUDE.md (Projekt-Identitaet) ===\n{txt}")
    # 2. MC config
    mc_cfg = Path("C:/Users/Robin/Jarvis/Mission-Control/config.json")
    if mc_cfg.exists():
        txt = mc_cfg.read_text(encoding="utf-8", errors="replace")[:1500]
        ctx_parts.append(f"=== Mission Control config.json ===\n{txt}")
    # 3. MC file structure (top-level + static/js files)
    mc_dir = Path("C:/Users/Robin/Jarvis/Mission-Control/static/js")
    if mc_dir.exists():
        js_files = sorted(f.name for f in mc_dir.iterdir() if f.is_file())
        ctx_parts.append(f"=== MC Frontend JS-Module ===\n{', '.join(js_files)}")
    # 4. Server routes (grep for app.router.add)
    server_py = Path("C:/Users/Robin/Jarvis/Mission-Control/server.py")
    if server_py.exists():
        lines = server_py.read_text(encoding="utf-8", errors="replace").splitlines()
        routes = [l.strip() for l in lines if "app.router.add" in l][:40]
        if routes:
            ctx_parts.append(f"=== MC Backend Routes ===\n" + "\n".join(routes))
    # 5. Brain structure (Obsidian vault top-level)
    brain_dir = Path("E:/OneDrive/AI/Obsidian-Vault/Jarvis-Brain/Jarvis-Memory")
    if brain_dir.exists():
        brain_files = sorted(f.name for f in brain_dir.iterdir() if f.is_file())[:30]
        ctx_parts.append(f"=== Obsidian Brain (Jarvis-Memory) ===\n{', '.join(brain_files)}")
    # 6. Skills
    skills_dir = Path("C:/Users/Robin/Jarvis/skills")
    if skills_dir.exists():
        skill_names = sorted(d.name for d in skills_dir.iterdir() if d.is_dir())[:20]
        ctx_parts.append(f"=== Installierte Skills ({len(skill_names)}) ===\n{', '.join(skill_names)}")
    # 7. Available APIs/integrations (explicit so bots know what's real)
    scripts_dir = Path("C:/Users/Robin/Jarvis/scripts")
    if scripts_dir.exists():
        script_names = sorted(f.name for f in scripts_dir.iterdir() if f.is_file() and f.suffix == ".py")
        ctx_parts.append(
            f"=== Vorhandene Integrations-Scripts ===\n{', '.join(script_names)}\n"
            f"HINWEIS: outlook_graph.py = Outlook/Email API (konfiguriert, laeuft), "
            f"alpaca_*.py = Trading API (Alpaca, konfiguriert), "
            f"telegram_bridge.py = Telegram Bot (laeuft), "
            f"icloud_*.py = iCloud Kalender (konfiguriert). "
            f"Diese APIs sind BEREITS eingerichtet und funktional."
        )
    return "\n\n".join(ctx_parts)


async def _arena_healthcheck_bots(bots: list) -> dict:
    """Prueft vor jeder Runde ob alle Bots 'leben':
    - Claude-CLI erreichbar (--version Ping, <10s; bei Fehlschlag einmal nach 5s nachpingen)
    - Bot-Config valide (name, model, persona vorhanden)
    Returns: {"cli_ok": bool, "bots": {bot_name: {"alive": bool, "reason": str, "model": str}}}
    """
    async def _ping_cli() -> bool:
        try:
            proc = await asyncio.create_subprocess_exec(
                CLAUDE_EXE, "--version",
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                **_SUBPROCESS_FLAGS,
            )
            try:
                await asyncio.wait_for(proc.communicate(), timeout=10)
                return proc.returncode == 0
            except asyncio.TimeoutError:
                try:
                    proc.kill()
                    await proc.wait()
                except Exception:
                    pass
                return False
        except Exception:
            return False

    cli_ok = await _ping_cli()
    if not cli_ok:
        # Einmal nach 5s nachpingen bevor Bot als tot gilt
        print("[Arena] Health-Check fehlgeschlagen — retry in 5s", flush=True)
        await asyncio.sleep(5)
        cli_ok = await _ping_cli()
        if cli_ok:
            print("[Arena] Health-Check retry OK", flush=True)
        else:
            print("[Arena] Health-Check retry fehlgeschlagen — CLI tot", flush=True)

    # Provider, die ihr Default-Modell ohne explizites `model`-Feld nutzen
    # (codex/gemini leeres model = "benutze was im CLI-Account konfiguriert ist",
    # openai/claude-api nehmen API-Default).
    _OPTIONAL_MODEL_PROVIDERS = {
        "codex-cli", "gemini-cli", "openai", "claude-api", "openclaw",
    }
    bot_health = {}
    for bot in bots:
        name = bot.get("name") or "(unnamed)"
        provider = (bot.get("provider") or "claude-cli").strip()
        model = bot.get("model") or ""
        persona = bot.get("persona") or ""
        reasons = []
        # Nur für claude-cli ist das Model-Feld Pflicht (haiku/sonnet/opus).
        # Andere Provider dürfen leer bleiben → Account-Default wird benutzt.
        if not model and provider not in _OPTIONAL_MODEL_PROVIDERS and provider != "claude-cli":
            reasons.append("kein Model")
        if provider == "claude-cli" and not model:
            reasons.append("kein Model")
        if not persona.strip():
            reasons.append("keine Persona")
        # CLI-Ping ist nur für claude-cli-Bots relevant
        if not cli_ok and provider == "claude-cli":
            reasons.append("Claude-CLI nicht erreichbar")
        bot_health[name] = {
            "alive": not reasons,
            "reason": "ok" if not reasons else ", ".join(reasons),
            "model": model or "(default)",
            "provider": provider,
        }
    return {"cli_ok": cli_ok, "bots": bot_health}


def _build_arena_memory_block(room_dict) -> str:
    """Gemerkter Stand fuer Bot-Prompts: letztes Orchestrator-Fazit +
    bereits umgesetzte Executor-Aktionen. Ziel: Bots wiederholen nicht,
    was schon etabliert / erledigt ist."""
    parts = []

    orch = (room_dict.get("orchestrator_summary") or "").strip()
    if orch:
        # VERDICT-Marker rausstrippen — brauchen Bots nicht, verwirrt nur
        import re as _re
        cleaned = _re.sub(r"\[VERDICT:[^\]]+\]", "", orch).strip()
        if cleaned:
            parts.append(
                "Bisheriger Stand (letztes Orchestrator-Fazit — baue darauf auf, "
                "wiederhole es nicht):\n" + cleaned[:1500]
            )

    # Letzte Executor-Summaries = schon erledigte Aktionen
    # "Frisch" = Executor-Msg nach der letzten regulaeren Bot-Msg (noch nicht bestaetigt)
    all_msgs = (room_dict.get("messages") or [])[-60:]
    last_bot_ts = ""
    for m in reversed(all_msgs):
        if m.get("role") == "bot":
            last_bot_ts = m.get("timestamp", "")
            break
    done_entries = []
    fresh_entries = []
    for m in all_msgs:
        if m.get("role") != "executor":
            continue
        body = (m.get("content") or "").strip()
        if not body:
            continue
        msg_ts = m.get("timestamp", "")
        if last_bot_ts and msg_ts > last_bot_ts:
            fresh_entries.append(body[:600])
        else:
            done_entries.append(body[:600])
    if fresh_entries:
        fresh_block = "\n\n".join(fresh_entries[-5:])
        parts.append(
            "FRISCH ERLEDIGT — diese Aktionen wurden gerade umgesetzt, noch nicht bestätigt:\n"
            + fresh_block
            + "\n\n→ Bitte bestätige im nächsten Turn kurz welche Aktion erledigt wurde (1 Satz), dann diskutier weiter."
        )
    if done_entries:
        done_block = "\n\n".join(done_entries[-3:])
        parts.append(
            "Bereits umgesetzte Aktionen (NICHT erneut vorschlagen, aufbauen statt wiederholen):\n"
            + done_block
        )

    marks = room_dict.get("consensus_marks") or []
    if marks:
        mark_lines = []
        for m in marks[-5:]:
            t = m.get("turn", "?")
            s = (m.get("summary") or "").strip().replace("\n", " ")
            if s:
                mark_lines.append(f"- Runde {t}: {s[:280]}")
        if mark_lines:
            parts.append(
                "Konsens-Historie (gemerkte Staende — darauf aufbauen, nicht neu verhandeln):\n"
                + "\n".join(mark_lines)
            )

    alpha = room_dict.get("alpha_last_result")
    if alpha and alpha.get("content"):
        parts.append(
            f"Alphas letzte Analyse (Turn {alpha.get('turn', '?')} — als Basis nutzen, nicht wiederholen):\n"
            + alpha["content"][:600]
        )

    done_acts = room_dict.get("done_actions") or []
    if done_acts:
        lines = [f"- {a}" for a in done_acts[-15:]]
        parts.append(
            "Bereits erledigte Aktionen (NICHT erneut vorschlagen — diese sind abgehakt):\n"
            + "\n".join(lines)
        )

    if not parts:
        return ""
    return "\n\n=== GEMERKTER STAND ===\n" + "\n\n".join(parts) + "\n=== ENDE STAND ===\n"


async def _run_arena(room_id: str):
    """Orchestration loop: bots take turns discussing the topic."""
    import time as _t

    rooms = _load_arena_rooms()
    room = _get_arena_room(rooms, room_id)
    if not room:
        _arena_rooms.pop(room_id, None)
        return

    bots = room.get("bots", [])
    if len(bots) < 2:
        _arena_rooms.pop(room_id, None)
        return

    topic = room.get("topic", "Freie Diskussion")

    # Rollen-Rotation: pro Runde bekommt jeder Bot eine andere Aufgabe,
    # damit nicht alle dasselbe pruefen. Rotiert zyklisch ueber Runden + Position.
    ROUND_ROLES = [
        ("Einschätzer",  "Bewerte die bisherige Diskussion: Was sind die stärksten Punkte? Was fehlt noch?"),
        ("Kritiker",     "Hinterfrage kritisch: Welche Annahmen sind schwach? Wo fehlt Evidenz oder Machbarkeit?"),
        ("Ergänzer",     "Bringe frische Aspekte rein, die noch nicht diskutiert wurden — kein Doppeln!"),
        ("Entscheider",  "Schlage eine konkrete nächste Aktion vor. Pragmatisch und direkt umsetzbar, nicht abstrakt."),
    ]

    # Harte Obergrenze: max_rounds ist die einzige Wahrheit, keine Fantasiezahlen.
    # Effektive Turns = max_rounds * len(bots). Fallback auf max_turns nur fuer Alt-Rooms.
    ROUNDS_HARD_LIMIT = 15
    hard_turn_cap = ROUNDS_HARD_LIMIT * max(1, len(bots))
    notbremse_triggered = False
    if "max_rounds" in room:
        try:
            max_rounds = max(1, min(ROUNDS_HARD_LIMIT, int(room.get("max_rounds", 2))))
        except (TypeError, ValueError):
            max_rounds = 2
        max_turns = max_rounds * max(1, len(bots))
        room["max_turns_effective"] = max_turns
    else:
        max_turns = room.get("max_turns", 20)
        if max_turns > hard_turn_cap:
            max_turns = hard_turn_cap
            notbremse_triggered = True
            room["max_turns_effective"] = max_turns
            _save_arena_rooms(rooms)

    # Gather project context once at start (so bots know what exists)
    project_context = _gather_arena_context()

    _activity_set("arena", active=True)

    # Notify start
    room["status"] = "running"
    _save_arena_rooms(rooms)
    await _broadcast_all_ws({
        "type": "arena.started",
        "room_id": room_id,
    })

    if notbremse_triggered:
        notbremse_msg = {
            "role": "system",
            "name": "System",
            "content": f"⚠️ Notbremse aktiv: Diskussion wird nach spaetestens {ROUNDS_HARD_LIMIT} Runden ({max_turns} Turns) hart gestoppt, um Endlosschleifen zu verhindern.",
            "timestamp": datetime.now().isoformat(),
            "color": "#f59e0b",
            "turn": 0,
        }
        room["messages"].append(notbremse_msg)
        _save_arena_rooms(rooms)
        await _broadcast_all_ws({
            "type": "arena.message",
            "room_id": room_id,
            "message": notbremse_msg,
        })

    turn = 0
    bot_index = 0
    _dead_bots: set = set()  # tote Bots aus letztem Health-Check; werden pro Turn uebersprungen
    _fertig_this_round: set = set()  # Bots die in dieser Runde [FERTIG] signalisiert haben
    _bot_crash_counts: dict = {}  # bot_name -> aufeinanderfolgende Crash-Zaehler
    _CRASH_THRESHOLD = 2  # Ab wie vielen aufeinanderfolgenden Abstuerzen wird Ersatz-Bot gespawnt
    _REPLACEMENT_MODELS = ["sonnet", "haiku", "sonnet"]  # Modell-Kandidaten fuer Ersatz-Bots
    _replacement_counter = 0  # Zaehler fuer eindeutige Ersatz-Bot-Namen
    _rounds_without_action = 0   # Runden ohne [JETZT AUSFÜHREN: ja] — Inaktivitäts-Stopp
    _action_triggered_this_round = False  # Wurde in dieser Runde eine Aktion getriggert?
    _IDLE_ROUNDS_LIMIT = 2  # Nach N ruhigen Runden: automatisches Ende

    try:
        while turn < max_turns:
            # Check if stopped
            if not _arena_rooms.get(room_id, {}).get("running"):
                break

            # Pre-Round Health-Check: vor jeder neuen Runde pruefen ob alle Bots leben
            if bot_index % len(bots) == 0:
                round_num = (bot_index // len(bots)) + 1
                health = await _arena_healthcheck_bots(bots)
                dead = [n for n, h in health["bots"].items() if not h["alive"]]
                _dead_bots = set(dead)  # merken fuer Turn-Skip unten
                await _broadcast_all_ws({
                    "type": "arena.healthcheck",
                    "room_id": room_id,
                    "round": round_num,
                    "turn": turn,
                    "cli_ok": health["cli_ok"],
                    "bots": health["bots"],
                })
                if dead:
                    details = "; ".join(f"{n}: {health['bots'][n]['reason']}" for n in dead)
                    hc_msg = {
                        "role": "system",
                        "name": "System",
                        "content": f"🩺 Health-Check vor Runde {round_num}: {len(bots)-len(dead)}/{len(bots)} Bots alive. Tote: {details}",
                        "timestamp": datetime.now().isoformat(),
                        "color": "#f59e0b",
                        "turn": turn,
                    }
                    rooms = _load_arena_rooms()
                    room = _get_arena_room(rooms, room_id)
                    if room:
                        room.setdefault("messages", []).append(hc_msg)
                        _save_arena_rooms(rooms)
                    await _broadcast_all_ws({
                        "type": "arena.message",
                        "room_id": room_id,
                        "message": hc_msg,
                    })
                    # Abbruch wenn CLI tot oder alle Bots tot
                    if not health["cli_ok"] or len(dead) == len(bots):
                        print(f"[Arena] Aborting room {room_id[:8]}: healthcheck failed (cli_ok={health['cli_ok']}, dead={len(dead)}/{len(bots)})")
                        break

            bot = bots[bot_index % len(bots)]
            bot_name = bot.get("name", f"Bot-{bot_index}")
            bot_model = bot.get("model", "sonnet")
            bot_persona = bot.get("persona", "Du bist ein hilfreicher Diskussionspartner.")
            bot_color = bot.get("color", "#6366f1")

            # --- Cross-Review: nach Execution-Turn darf Ausführer nicht eigene Arbeit reviewen ---
            _cr_executor = _arena_rooms.get(room_id, {}).get("last_executing_bot")
            _cr_pending = _arena_rooms.get(room_id, {}).pop("cross_review_pending", False)
            _is_cross_review = False
            if _cr_pending and _cr_executor:
                _is_cross_review = True
                if bot_name == _cr_executor:
                    # Swap zu nächstem lebenden Nicht-Ausführer-Bot
                    _alive_others = [b for b in bots if b.get("name") != _cr_executor and b.get("name") not in _dead_bots]
                    if _alive_others:
                        bot = _alive_others[0]
                        bot_name = bot.get("name", bot_name)
                        bot_model = bot.get("model", "sonnet")
                        bot_persona = bot.get("persona", "Du bist ein hilfreicher Diskussionspartner.")
                        bot_color = bot.get("color", "#6366f1")
                        print(f"[Arena] Cross-Review: Bot '{_cr_executor}' übersprungen — '{bot_name}' reviewt stattdessen (room={room_id[:8]})")

            # Toter Bot? Turn ueberspringen und kurz loggen.
            if bot_name in _dead_bots:
                print(f"[Arena] Skipping dead bot '{bot_name}' (turn {turn})")
                skip_msg = {
                    "role": "system",
                    "name": "System",
                    "content": f"⏭️ **{bot_name}** übersprungen (tot, Turn {turn})",
                    "timestamp": datetime.now().isoformat(),
                    "color": "#6b7280",
                    "turn": turn,
                }
                rooms = _load_arena_rooms()
                room = _get_arena_room(rooms, room_id)
                if room:
                    room.setdefault("messages", []).append(skip_msg)
                    _save_arena_rooms(rooms)
                await _broadcast_all_ws({"type": "arena.message", "room_id": room_id, "message": skip_msg})
                turn += 1
                bot_index += 1
                continue

            # Reload room to get latest messages (including Robin injections)
            rooms = _load_arena_rooms()
            room = _get_arena_room(rooms, room_id)
            if not room:
                break

            # Build conversation context for this bot
            recent_msgs = room.get("messages", [])[-8:]
            history_lines = []
            for m in recent_msgs:
                line = f"{m['name']}: {m['content']}"
                # Surface attached files so bots can open them via Read / Bash.
                atts = list(m.get("attachments") or [])
                if not atts and m.get("image"):
                    atts = [m["image"]]
                if atts:
                    att_hints = []
                    for fn in atts:
                        fp = UPLOADS_DIR / fn
                        ext = Path(fn).suffix.lower()
                        if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"):
                            att_hints.append(f"Bild: {fp}")
                        elif ext == ".pdf":
                            att_hints.append(f"PDF: {fp} (lies mit Read oder `pdftotext`/Python)")
                        elif ext in (".txt", ".md", ".csv", ".json", ".py", ".js", ".html",
                                      ".css", ".yaml", ".yml", ".log", ".tsv", ".ini", ".toml"):
                            att_hints.append(f"Textdatei: {fp} (lies mit Read)")
                        else:
                            att_hints.append(f"Datei: {fp}")
                    line += "\n  [Anhang fuer alle Bots einsehbar — " + "; ".join(att_hints) + "]"
                history_lines.append(line)
            history_block = "\n\n".join(history_lines) if history_lines else "(Noch keine Nachrichten)"

            # Gemerkter Stand: letztes Fazit + schon umgesetzte Aktionen.
            # Wird in alle Prompt-Varianten injiziert, damit Bots sich nicht wiederholen.
            memory_block = _build_arena_memory_block(room)

            # Robin-Zwischenruf: nur user-Messages NACH dem letzten Bot-Turn.
            # Alte Zwischenrufe wurden im Vorgaenger-Turn schon adressiert.
            _msgs_all = room.get("messages", [])
            _last_bot_idx = -1
            for _i in range(len(_msgs_all) - 1, -1, -1):
                if _msgs_all[_i].get("role") == "bot":
                    _last_bot_idx = _i
                    break
            _fresh_user = [
                (m.get("content") or "").strip()
                for m in _msgs_all[_last_bot_idx + 1:]
                if m.get("role") == "user" and (m.get("content") or "").strip()
            ]
            robin_block = ""
            if _fresh_user:
                _lines = "\n\n".join(f"[Robin sagt: {t[:600]}]" for t in _fresh_user[-3:])
                robin_block = (
                    "\n\n=== ROBINS ZWISCHENRUF (direkte Anweisung — PRIORITAET vor allem anderen!) ===\n"
                    + _lines
                    + "\n=== ENDE ZWISCHENRUF ===\n"
                )

            # Cross-Review-Hinweis: Wenn Reviewer-Slot aktiv ist, Executor-Status einblenden
            # und strukturiertes URTEIL anfordern. Ersetzt freies Regex-Parsing.
            review_hint = ""
            if _is_cross_review:
                # Echter per-Action-Status aus der letzten approval-Nachricht holen.
                # Fallback: Textsuche im executor-Output für Legacy-Räume.
                _action_lines = []
                _done_c = _fail_c = _haengend_c = 0
                for _m in reversed(room.get("messages", [])):
                    if _m.get("role") == "approval" and _m.get("actions"):
                        for _ai, _a in enumerate(_m["actions"], start=1):
                            _s = _a.get("status", "unbekannt")
                            _t = (_a.get("text") or "")[:80]
                            _r = (_a.get("result") or "")[:120]
                            _action_lines.append(f"  #{ _ai} [{_s.upper()}] {_t}" + (f" → {_r}" if _r else ""))
                            if _s == "done": _done_c += 1
                            elif _s == "failed": _fail_c += 1
                            elif _s == "haengend": _haengend_c += 1
                        break
                if _action_lines:
                    _exec_status = f"{_done_c} fertig, {_fail_c} fehlgeschlagen, {_haengend_c} haengend"
                    _action_block = "\n".join(_action_lines)
                else:
                    # Fallback: Textsuche im executor-Output
                    _exec_status = "unbekannt"
                    for _m in reversed(room.get("messages", [])):
                        if _m.get("role") == "executor":
                            _c = _m.get("content", "")
                            _exec_status = f"OK={_c.count('[OK]')}, FAIL={_c.count('[FAIL]')}"
                            break
                    _action_block = "(keine per-Action Daten verfügbar)"
                review_hint = (
                    f"\n\nCROSS-REVIEW-MODUS: Du reviewst die Umsetzung von **{_cr_executor}** "
                    f"(Executor-Status: {_exec_status}).\n\n"
                    f"=== ECHTER AUSFÜHRUNGS-STATUS (pro Aktion) ===\n"
                    f"{_action_block}\n"
                    f"=== ENDE STATUS ===\n\n"
                    f"Beurteile NUR auf Basis des obigen echten Status — nicht auf Annahmen. "
                    f"Gib am Ende deiner Nachricht exakt eine Zeile: "
                    f"'URTEIL: REIFE' wenn alles sauber umgesetzt. "
                    f"'URTEIL: NACHBESSERN #N' (z.B. '#1,#3') wenn bestimmte Aktionen Lücken haben — "
                    f"ersetze #N durch die Nummern der betroffenen Aktionen aus der Liste oben. "
                    f"Du bist Reviewer, nicht Ausführer — keine Eigen-Vorschläge."
                )

            # Check execution mode for prompt adaptation
            is_executing = room.get("tools_enabled", False) and room.get("execution_turns_left", 0) > 0

            # Check if bots can read files
            can_read = room.get("allow_read", False) and not is_executing
            read_hint = (
                "\n\nDu kannst Dateien lesen (Read, Glob, Grep) um deine Argumente zu stuetzen. "
                "Nutze das wenn es zum Thema passt — z.B. um Code oder Konfiguration zu pruefen. "
                "WICHTIG: Maximal 2 Read-Operationen pro Turn — dann direkt antworten. "
                "Arbeitsverzeichnis: C:/Users/Robin/Jarvis"
            ) if can_read else ""

            # Context block for all prompts
            ctx_block = (
                f"\n\n=== AKTUELLER PROJEKT-KONTEXT (lies das bevor du diskutierst!) ===\n"
                f"{project_context}\n"
                f"=== ENDE KONTEXT ===\n"
            ) if project_context else ""

            # Vorgänger-Crash-Hinweis: Ersatz-Bot wird kurz informiert und dann geleert
            _lc = _arena_rooms.get(room_id, {}).pop("last_crash", None)
            predecessor_hint = (
                f"\nHINWEIS: Dein Vorgänger **{_lc['bot']}** ist in diesem Turn ausgefallen "
                f"({_lc['kind']}: {_lc['summary']}). Übernimm nahtlos — du musst das nicht kommentieren.\n"
            ) if _lc else ""

            if turn == 0 and not recent_msgs:
                # First message: introduce the topic
                prompt = (
                    f"Du bist {bot_name}. {bot_persona}\n"
                    f"{predecessor_hint}"
                    f"\nThema der Diskussion: {topic}\n"
                    f"{robin_block}"
                    f"{ctx_block}"
                    f"{memory_block}\n"
                    f"Du eroeffnest jetzt die Diskussion. Bringe deinen ersten Standpunkt ein. "
                    f"Beziehe dich auf den TATSAECHLICHEN Stand des Projekts (siehe Kontext oben). "
                    f"Halte dich kurz (max 150 Woerter). Antworte direkt ohne Metakommentare. "
                    f"Keine Fantasiezahlen bei Runden/Iterationen — die Arena ist auf max 10 Runden gedeckelt."
                    f"{read_hint}"
                )
            elif is_executing:
                # EXECUTION MODE: bot should actually implement what was discussed
                prompt = (
                    f"Du bist {bot_name}. {bot_persona}\n"
                    f"{predecessor_hint}"
                    f"\nThema: {topic}\n"
                    f"{robin_block}"
                    f"{ctx_block}"
                    f"{memory_block}\n"
                    f"Bisherige Diskussion:\n{history_block}\n\n"
                    f"WICHTIG: Robin hat die Ausfuehrung freigegeben! Du hast jetzt Zugriff auf alle Tools "
                    f"(Read, Write, Edit, Bash, Glob, Grep). Setze um was in der Diskussion besprochen wurde. "
                    f"Arbeitsverzeichnis ist C:/Users/Robin/Jarvis.\n\n"
                    f"Beschreibe kurz was du tust und fuehre die Aenderungen durch. "
                    f"Mach nur kleine, gezielte Fixes — keine grossen Overhauls. "
                    f"Antworte als {bot_name}, ohne Metakommentare."
                )
            else:
                # Only inject full context every 5 turns to save tokens
                ctx_insert = ""  # Context nur bei Turn 0 (bereits oben injiziert), nicht wiederholt
                # Feste Rolle pro Runde: rotiert zyklisch, damit Bots sich nicht doppeln
                _pos_in_round = bot_index % len(bots)
                _round_now = (bot_index // len(bots)) + 1
                _role_name, _role_task = ROUND_ROLES[(_round_now - 1 + _pos_in_round) % len(ROUND_ROLES)]
                role_hint = f"\n\nDeine ROLLE in Runde {_round_now}: **{_role_name}** — {_role_task} Halte dich strikt an diese Rolle."
                prompt = (
                    f"Du bist {bot_name}. {bot_persona}\n"
                    f"{predecessor_hint}"
                    f"\nThema: {topic}\n"
                    f"{robin_block}"
                    f"{ctx_insert}"
                    f"{memory_block}\n"
                    f"Bisherige Diskussion:\n{history_block}\n\n"
                    f"Du bist jetzt dran. Reagiere auf das Gesagte. "
                    f"WICHTIG: Nicht immer widersprechen! Wenn jemand einen guten Punkt macht, sag das und baue darauf auf. "
                    f"Bringe neue Perspektiven, aber sei auch bereit zuzustimmen und gemeinsam weiterzudenken. "
                    f"Wenn Robin etwas eingeworfen hat, geh darauf ein. "
                    f"ERLEDIGTE AKTIONEN BESTÄTIGEN: Wenn im GEMERKTER STAND ein Abschnitt 'FRISCH ERLEDIGT' steht, "
                    f"beginne deinen Beitrag mit einem Satz der bestätigt was umgesetzt wurde (z.B. '✓ [Aktion] ist erledigt.'), "
                    f"dann diskutiere normal weiter. Das ist Pflicht — der Fortschritt muss sichtbar sein. "
                    f"Falls der Orchestrator gerade ein Zwischenfazit gemacht hat: Diskutiere einfach weiter! "
                    f"Das Fazit ist nur eine Zusammenfassung, kein Stopp-Signal. Du wartest NICHT auf Freigabe oder gruenes Licht. "
                    f"Beziehe dich auf das was TATSAECHLICH existiert, nicht auf Vermutungen. "
                    f"Halte dich kurz (max 150 Woerter). Antworte direkt als {bot_name}. "
                    f"Keine Fantasiezahlen bei Runden/Iterationen — Arena-Hardcap ist 10 Runden; teile groessere Vorhaben in kleinere Schritte. "
                    f"Wenn du glaubst dass die Diskussion einen echten Konsens erreicht hat und keine neuen Erkenntnisse mehr zu erwarten sind, "
                    f"schreibe am Ende deiner Nachricht exakt '[FERTIG]' (in eckigen Klammern). "
                    f"Sobald 2 Bots [FERTIG] schreiben, endet die Runde automatisch. Nur wenn du wirklich überzeugt bist — kein voreiliges Ende!\n"
                    f"EXECUTION-MARKIERUNG: Wenn du eine konkrete Aktion vorschlägst die SOFORT umgesetzt werden soll "
                    f"(z.B. eine Datei editieren, einen Endpoint hinzufügen), schreibe ans Ende deiner Nachricht exakt "
                    f"'[JETZT AUSFÜHREN: ja]'. Wenn die Diskussion noch nicht reif für Umsetzung ist: '[JETZT AUSFÜHREN: nein]'. "
                    f"Nur bei echtem Konsens 'ja' schreiben — der nächste Bot bekommt dann volle Tool-Rechte und muss die Aktion wirklich durchführen."
                    f"{role_hint}"
                    f"{review_hint}"
                    f"{read_hint}"
                )

            # Notify: bot is thinking
            await _broadcast_all_ws({
                "type": "arena.thinking",
                "room_id": room_id,
                "bot_name": bot_name,
                "bot_color": bot_color,
                "turn": turn,
            })

            # Check if this room has tools enabled (execution mode)
            room_tools_enabled = room.get("tools_enabled", False)
            execution_turns_left = room.get("execution_turns_left", 0)

            # Prompt via stdin statt argv - umgeht Windows 32K command-line limit
            if room_tools_enabled and execution_turns_left > 0:
                # EXECUTION MODE: full tools, higher max-turns for tool use.
                # bypassPermissions fehlte vorher — ohne den Flag verweigert Claude
                # Write/Edit/Bash und frisst Turns mit Permission-Prompts → Budget leer
                # bevor was umgesetzt wurde. 25 Turns matcht in etwa den medium-Executor.
                cmd = [
                    CLAUDE_EXE, "-p",
                    "--model", bot_model,
                    "--output-format", "text",
                    "--max-turns", "25",
                    "--permission-mode", "bypassPermissions",
                ]
                # Decrement execution turns
                room["execution_turns_left"] = execution_turns_left - 1
                _save_arena_rooms(rooms)
                if room["execution_turns_left"] <= 0:
                    room["tools_enabled"] = False
                    room["execution_turns_left"] = 0
                    _save_arena_rooms(rooms)
                    await _broadcast_all_ws({
                        "type": "arena.mode_change",
                        "room_id": room_id,
                        "mode": "discuss",
                        "turns_left": 0,
                    })
            elif can_read:
                # DISCUSSION + READ MODE: bots can read files to inform their arguments
                # max-turns=4: Grep + 2 Reads + 1 final text turn — Budget-Deckel per bot
                cmd = [
                    CLAUDE_EXE, "-p",
                    "--model", bot_model,
                    "--output-format", "text",
                    "--tools", "Read,Glob,Grep",
                    "--max-turns", "4",
                ]
            else:
                # DISCUSSION MODE: no tools, pure text
                cmd = [
                    CLAUDE_EXE, "-p",
                    "--model", bot_model,
                    "--output-format", "text",
                    "--tools", "",
                    "--max-turns", "8",
                ]

            crash_info = None  # set if anything goes wrong, posted as system msg below
            proc = None
            stderr_full = ""
            stdout_full = ""
            # Timeouts: execution=300s, read-mode=150s (10 turns), pure-text=120s
            if is_executing:
                _timeout = 300
            elif can_read:
                _timeout = 60
            else:
                _timeout = 120
            try:
                proc = await asyncio.create_subprocess_exec(
                    *cmd,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    stdin=asyncio.subprocess.PIPE,
                    cwd="C:/Users/Robin/Jarvis",
                    **_SUBPROCESS_FLAGS,
                )
                _arena_rooms[room_id]["process"] = proc

                stdout_data, stderr_data = await asyncio.wait_for(
                    proc.communicate(input=prompt.encode("utf-8")), timeout=_timeout
                )
                stdout_full = stdout_data.decode("utf-8", errors="replace")
                stderr_full = stderr_data.decode("utf-8", errors="replace")

                response = stdout_full.strip()
                # Strip "Reached max turns" errors — use stderr content if stdout empty
                _budget_exhausted = "Reached max turns" in response
                _stdout_raw = response  # keep pre-strip snapshot for crash log
                if _budget_exhausted:
                    lines = [l for l in response.splitlines() if "Reached max turns" not in l and "Error:" not in l]
                    response = "\n".join(lines).strip()
                if not response:
                    # Determine real crash cause from stderr signals
                    _stderr_lower = (stderr_full or "").lower()
                    _thinking_signals = [
                        "thinking budget", "budget_tokens", "thinking_budget",
                        "extended_thinking", "thinking token", "context_window_exceeded",
                        "context window", "context length", "maximum context",
                    ]
                    _thinking_exhausted = any(s in _stderr_lower for s in _thinking_signals)
                    if _budget_exhausted:
                        _crash_kind = "budget_exhausted"
                    elif _thinking_exhausted:
                        _crash_kind = "thinking_budget_exhausted"
                    else:
                        _crash_kind = "empty_response"
                    # --- Budget-Escalation: bei aufgebrauchtem Denk-/Schritte-Budget höher stufen ---
                    _ESCALATION_MAP_B = {"haiku": "sonnet", "sonnet": "opus"}
                    _esc_model_b = _ESCALATION_MAP_B.get(bot_model) if _crash_kind in ("budget_exhausted", "thinking_budget_exhausted") else None
                    _budget_retry_response = None
                    if _esc_model_b:
                        _esc_label = "Denk-Budget" if _crash_kind == "thinking_budget_exhausted" else "Schritte-Budget"
                        _esc_timeout_b = int(_timeout * 1.5)
                        await _broadcast_all_ws({
                            "type": "arena.message",
                            "room_id": room_id,
                            "message": {
                                "role": "system",
                                "name": "System",
                                "content": f"🧠 **{bot_name}** {_esc_label} erschöpft — Retry mit `{_esc_model_b}` ({_esc_timeout_b}s Budget)…",
                                "timestamp": datetime.now().isoformat(),
                                "color": "#f59e0b",
                                "turn": turn,
                            },
                        })
                        try:
                            _esc_proc_b = await asyncio.create_subprocess_exec(
                                CLAUDE_EXE, "-p",
                                "--model", _esc_model_b,
                                "--output-format", "text",
                                "--tools", "",
                                "--max-turns", "1",
                                stdout=asyncio.subprocess.PIPE,
                                stderr=asyncio.subprocess.PIPE,
                                stdin=asyncio.subprocess.PIPE,
                                cwd="C:/Users/Robin/Jarvis",
                                **_SUBPROCESS_FLAGS,
                            )
                            _arena_rooms[room_id]["process"] = _esc_proc_b
                            _brd, _bre = await asyncio.wait_for(
                                _esc_proc_b.communicate(input=prompt.encode("utf-8")),
                                timeout=_esc_timeout_b,
                            )
                            _budget_retry_response = _brd.decode("utf-8", errors="replace").strip()
                            if "Reached max turns" in _budget_retry_response:
                                _budget_retry_response = "\n".join(
                                    l for l in _budget_retry_response.splitlines()
                                    if "Reached max turns" not in l and "Error:" not in l
                                ).strip()
                        except Exception:
                            _budget_retry_response = None
                    if _budget_retry_response:
                        response = _budget_retry_response
                        bot_model = _esc_model_b  # reflect escalated model in message metadata
                    elif _crash_kind == "empty_response":
                        # Retry once with shortened context before treating as crash
                        _short_history_lines = history_lines[-5:] if len(history_lines) > 5 else history_lines
                        _short_history = "\n\n".join(_short_history_lines) if _short_history_lines else "(Noch keine Nachrichten)"
                        _short_prompt = (
                            f"Du bist {bot_name}. {bot_persona}\n"
                            f"\nThema: {topic}\n"
                            f"Bisherige Diskussion (gekuerzt):\n{_short_history}\n\n"
                            f"Antworte kurz (max 80 Woerter) direkt als {bot_name}."
                        )
                        await _broadcast_all_ws({
                            "type": "arena.message",
                            "room_id": room_id,
                            "message": {
                                "role": "system",
                                "name": "System",
                                "content": f"↩️ **{bot_name}** leere Antwort — Retry mit gekürztem Kontext…",
                                "timestamp": datetime.now().isoformat(),
                                "color": "#6b7280",
                                "turn": turn,
                            },
                        })
                        _short_retry_response = None
                        try:
                            _sr_proc = await asyncio.create_subprocess_exec(
                                CLAUDE_EXE, "-p",
                                "--model", bot_model,
                                "--output-format", "text",
                                "--tools", "",
                                "--max-turns", "1",
                                stdout=asyncio.subprocess.PIPE,
                                stderr=asyncio.subprocess.PIPE,
                                stdin=asyncio.subprocess.PIPE,
                                cwd="C:/Users/Robin/Jarvis",
                                **_SUBPROCESS_FLAGS,
                            )
                            _arena_rooms[room_id]["process"] = _sr_proc
                            _srd, _sre = await asyncio.wait_for(
                                _sr_proc.communicate(input=_short_prompt.encode("utf-8")),
                                timeout=60,
                            )
                            _short_retry_response = _srd.decode("utf-8", errors="replace").strip()
                            if "Reached max turns" in (_short_retry_response or ""):
                                _short_retry_response = "\n".join(
                                    l for l in _short_retry_response.splitlines()
                                    if "Reached max turns" not in l and "Error:" not in l
                                ).strip()
                        except Exception:
                            _short_retry_response = None
                        if _short_retry_response:
                            response = _short_retry_response
                        else:
                            crash_info = _log_arena_crash(
                                room_id=room_id, bot_name=bot_name, turn=turn,
                                kind=_crash_kind, cmd=cmd, returncode=proc.returncode,
                                stderr=stderr_full, stdout_tail=_stdout_raw[-1000:] or stdout_full[-1000:],
                                timeout_s=_timeout,
                            )
                            response = f"(Keine Antwort — {crash_info['summary']})"
                    else:
                        crash_info = _log_arena_crash(
                            room_id=room_id, bot_name=bot_name, turn=turn,
                            kind=_crash_kind, cmd=cmd, returncode=proc.returncode,
                            stderr=stderr_full, stdout_tail=_stdout_raw[-1000:] or stdout_full[-1000:],
                            timeout_s=_timeout,
                        )
                        response = f"(Keine Antwort — {crash_info['summary']})"

            except asyncio.TimeoutError:
                if proc and proc.returncode is None:
                    try:
                        proc.kill()
                        await proc.wait()
                    except Exception:
                        pass
                # --- Escalation-Retry: höheres Modell + längerer Timeout ---
                _ESCALATION_MAP = {"haiku": "sonnet", "sonnet": "opus"}
                _escalated_model = _ESCALATION_MAP.get(bot_model)
                _retry_response = None
                if _escalated_model:
                    _retry_timeout = int(_timeout * 1.5)
                    _retry_cmd = [
                        CLAUDE_EXE, "-p",
                        "--model", _escalated_model,
                        "--output-format", "text",
                        "--tools", "",
                        "--max-turns", "1",
                    ]
                    # Notify room about escalation attempt
                    await _broadcast_all_ws({
                        "type": "arena.message",
                        "room_id": room_id,
                        "message": {
                            "role": "system",
                            "name": "System",
                            "content": f"⚡ **{bot_name}** Timeout nach {_timeout}s — Retry mit `{_escalated_model}` ({_retry_timeout}s Budget)…",
                            "timestamp": datetime.now().isoformat(),
                            "color": "#f59e0b",
                            "turn": turn,
                        },
                    })
                    try:
                        _retry_proc = await asyncio.create_subprocess_exec(
                            *_retry_cmd,
                            stdout=asyncio.subprocess.PIPE,
                            stderr=asyncio.subprocess.PIPE,
                            stdin=asyncio.subprocess.PIPE,
                            cwd="C:/Users/Robin/Jarvis",
                            **_SUBPROCESS_FLAGS,
                        )
                        _arena_rooms[room_id]["process"] = _retry_proc
                        _rd, _re = await asyncio.wait_for(
                            _retry_proc.communicate(input=prompt.encode("utf-8")),
                            timeout=_retry_timeout,
                        )
                        _retry_response = _rd.decode("utf-8", errors="replace").strip()
                        if "Reached max turns" in _retry_response:
                            _retry_response = "\n".join(
                                l for l in _retry_response.splitlines()
                                if "Reached max turns" not in l and "Error:" not in l
                            ).strip()
                    except Exception:
                        _retry_response = None
                if _retry_response:
                    response = _retry_response
                    bot_model = _escalated_model  # reflect escalated model in message metadata
                else:
                    crash_info = _log_arena_crash(
                        room_id=room_id, bot_name=bot_name, turn=turn,
                        kind="timeout", cmd=cmd,
                        returncode=(proc.returncode if proc else None),
                        stderr=stderr_full, stdout_tail=stdout_full[-1000:],
                        timeout_s=_timeout,
                    )
                    response = f"(Timeout — {crash_info['summary']})"
            except Exception as e:
                import traceback as _tb
                crash_info = _log_arena_crash(
                    room_id=room_id, bot_name=bot_name, turn=turn,
                    kind="exception", cmd=cmd,
                    returncode=(proc.returncode if proc else None),
                    stderr=stderr_full, stdout_tail=stdout_full[-1000:],
                    timeout_s=_timeout,
                    exception=f"{type(e).__name__}: {e}\n{_tb.format_exc()[:1500]}",
                )
                response = f"(Fehler — {crash_info['summary']})"

            # Crash als System-Message in den Raum — kein blindes Raten mehr
            if crash_info:
                print(f"[Arena] CRASH room={room_id[:8]} bot={bot_name} turn={turn} -> {crash_info['file']}")
                crash_msg = {
                    "role": "system",
                    "name": "System",
                    "content": (
                        f"💀 Bot-Absturz: **{bot_name}** (Turn {turn})\n\n"
                        f"{crash_info['summary']}\n\n"
                        f"Log: `{crash_info['file']}`"
                    ),
                    "timestamp": datetime.now().isoformat(),
                    "color": "#ef4444",
                    "turn": turn,
                    "crash": {
                        "kind": crash_info["payload"].get("kind"),
                        "bot": bot_name,
                        "returncode": crash_info["payload"].get("returncode"),
                        "timeout_s": crash_info["payload"].get("timeout_s"),
                        "stderr_tail": (stderr_full or "").strip()[-400:],
                        "log_file": crash_info["file"],
                    },
                }
                rooms = _load_arena_rooms()
                room = _get_arena_room(rooms, room_id)
                if room:
                    room.setdefault("messages", []).append(crash_msg)
                    _save_arena_rooms(rooms)
                await _broadcast_all_ws({
                    "type": "arena.message",
                    "room_id": room_id,
                    "message": crash_msg,
                })
                await _broadcast_all_ws({
                    "type": "arena.crash",
                    "room_id": room_id,
                    "bot_name": bot_name,
                    "turn": turn,
                    "kind": crash_info["payload"].get("kind"),
                    "log_file": crash_info["file"],
                    "summary": crash_info["summary"],
                })
                # Merke für den nächsten Bot warum der Vorgänger gescheitert ist
                _arena_rooms[room_id]["last_crash"] = {
                    "bot": bot_name,
                    "kind": crash_info["payload"].get("kind", "unknown"),
                    "summary": crash_info["summary"],
                }

                # --- Ersatz-Bot-Logik: aufeinanderfolgende Abstuerze tracken ---
                _bot_crash_counts[bot_name] = _bot_crash_counts.get(bot_name, 0) + 1
                if _bot_crash_counts[bot_name] >= _CRASH_THRESHOLD:
                    # Ersatz-Bot spawnen: ersetzt diesen Slot in der lokalen bots-Liste
                    _replacement_counter += 1
                    _ersatz_model = _REPLACEMENT_MODELS[_replacement_counter % len(_REPLACEMENT_MODELS)]
                    _ersatz_name = f"Ersatz-{_replacement_counter}"
                    _ersatz_bot = {
                        "name": _ersatz_name,
                        "model": _ersatz_model,
                        "persona": (
                            f"Du bist {_ersatz_name}, ein neutraler Analyst. "
                            f"Du springst ein, weil {bot_name} wiederholt ausgefallen ist. "
                            f"Uebernimm nahtlos den Diskussionsfaden. Bringe sachliche Perspektiven ein."
                        ),
                        "color": "#8b5cf6",
                        "_replaced": bot_name,
                    }
                    _slot_idx = bot_index % len(bots)
                    bots[_slot_idx] = _ersatz_bot
                    _bot_crash_counts.pop(bot_name, None)
                    _dead_bots.discard(bot_name)
                    ersatz_msg = {
                        "role": "system",
                        "name": "System",
                        "content": (
                            f"🔄 **Ersatz-Bot aktiviert**: {_ersatz_name} ({_ersatz_model}) "
                            f"übernimmt den Slot von **{bot_name}** nach {_CRASH_THRESHOLD} aufeinanderfolgenden Abstürzen."
                        ),
                        "timestamp": datetime.now().isoformat(),
                        "color": "#8b5cf6",
                        "turn": turn,
                    }
                    rooms = _load_arena_rooms()
                    room = _get_arena_room(rooms, room_id)
                    if room:
                        room.setdefault("messages", []).append(ersatz_msg)
                        _save_arena_rooms(rooms)
                    await _broadcast_all_ws({
                        "type": "arena.message",
                        "room_id": room_id,
                        "message": ersatz_msg,
                    })
                    await _broadcast_all_ws({
                        "type": "arena.bot_replaced",
                        "room_id": room_id,
                        "old_bot": bot_name,
                        "new_bot": _ersatz_name,
                        "model": _ersatz_model,
                    })
                    print(f"[Arena] Ersatz-Bot '{_ersatz_name}' ersetzt '{bot_name}' nach {_CRASH_THRESHOLD} Abstuerzen (room={room_id[:8]})")
            else:
                # Kein Crash: Zaehler zuruecksetzen
                _bot_crash_counts.pop(bot_name, None)

            # Check again if stopped during execution
            if not _arena_rooms.get(room_id, {}).get("running"):
                break

            # Save message
            msg = {
                "role": "bot",
                "name": bot_name,
                "content": response,
                "timestamp": datetime.now().isoformat(),
                "color": bot_color,
                "model": bot_model,
                "turn": turn,
                "mode": "execute" if is_executing else "discuss",
            }

            rooms = _load_arena_rooms()
            room = _get_arena_room(rooms, room_id)
            if room:
                room["messages"].append(msg)
                room["last_active"] = datetime.now().isoformat()
                # Alpha's analytical output carries into next round so the team learns from it
                if bot_name == "Alpha" and response and not response.startswith("("):
                    room["alpha_last_result"] = {
                        "content": response[:800],
                        "turn": turn,
                        "timestamp": datetime.now().isoformat(),
                    }
                if _is_cross_review:
                    _verdict_m = re.search(r'URTEIL:\s*(REIFE|NACHBESSERN)((?:\s*#\d+(?:,\s*#?\d+)*)?)', response)
                    if _verdict_m:
                        _raw_nums = _verdict_m.group(2) or ""
                        _action_nums = [int(x) for x in re.findall(r'\d+', _raw_nums)] if _raw_nums.strip() else []
                        _verdict_record = {
                            "verdict": _verdict_m.group(1),
                            "action_nums": _action_nums,
                            "executor": _cr_executor,
                            "reviewer": bot_name,
                            "turn": turn,
                            "timestamp": datetime.now().isoformat(),
                            "room_id": room_id,
                            "topic": room.get("topic", ""),
                            "review_excerpt": response[:600],
                        }
                        room["last_review_verdict"] = _verdict_record
                        # Persist to append-log so the full review history survives
                        try:
                            ARENA_REVIEWS_LOG.parent.mkdir(parents=True, exist_ok=True)
                            with ARENA_REVIEWS_LOG.open("a", encoding="utf-8") as _rl:
                                _rl.write(json.dumps(_verdict_record, ensure_ascii=False) + "\n")
                        except Exception as _rle:
                            print(f"[Arena] review-log write error: {_rle}")
                        # Schritt 3: Bei NACHBESSERN geflaggte Actions zurück in Queue
                        if _verdict_m.group(1) == "NACHBESSERN":
                            try:
                                _requeue_nachbessern_actions(room, bot_name, room_id, turn, _action_nums)
                            except Exception as _rqe:
                                print(f"[Arena] requeue error: {_rqe}")
                _save_arena_rooms(rooms)

            # Broadcast message
            await _broadcast_all_ws({
                "type": "arena.message",
                "room_id": room_id,
                "message": msg,
            })

            turn += 1
            bot_index += 1

            # JETZT AUSFÜHREN-Signal: Bot markiert Aktion als umsetzungsreif
            if "[JETZT AUSFÜHREN: ja]" in response:
                _action_triggered_this_round = True
            if "[JETZT AUSFÜHREN: ja]" in response and not is_executing:
                rooms = _load_arena_rooms()
                room = _get_arena_room(rooms, room_id)
                if room:
                    room["tools_enabled"] = True
                    room["execution_turns_left"] = 1
                    _save_arena_rooms(rooms)
                    await _broadcast_all_ws({
                        "type": "arena.mode_change",
                        "room_id": room_id,
                        "mode": "execute",
                        "turns_left": 1,
                        "triggered_by": bot_name,
                    })
                    print(f"[Arena] JETZT AUSFÜHREN: ja — Bot '{bot_name}' hat Execution für nächsten Turn getriggert (room={room_id[:8]})")

            # FERTIG-Signal: Bot signalisiert Konsens — track pro Runde
            if "[FERTIG]" in response:
                _fertig_this_round.add(bot_name)
            # Neue Runde beginnt: Counter zurücksetzen
            if bot_index % len(bots) == 0:
                active_bots = [b.get("name") for b in bots if b.get("name") not in _dead_bots]
                fertig_threshold = min(2, len(active_bots))
                if active_bots and len(_fertig_this_round) >= fertig_threshold:
                    # Mindestens 2 Bots (oder alle, wenn weniger vorhanden) haben FERTIG signalisiert → Diskussion beenden
                    fertig_names = ", ".join(sorted(_fertig_this_round))
                    fertig_msg = {
                        "role": "system",
                        "name": "System",
                        "content": f"✅ **Konsens erreicht** — {fertig_names} haben [FERTIG] signalisiert. Diskussion beendet.",
                        "timestamp": datetime.now().isoformat(),
                        "color": "#10b981",
                        "turn": turn,
                    }
                    rooms = _load_arena_rooms()
                    room = _get_arena_room(rooms, room_id)
                    if room:
                        room["messages"].append(fertig_msg)
                        room["status"] = "fertig"
                        _save_arena_rooms(rooms)
                    if room_id in _arena_rooms:
                        _arena_rooms[room_id]["running"] = False
                    await _broadcast_all_ws({
                        "type": "arena.message",
                        "room_id": room_id,
                        "message": fertig_msg,
                    })
                    await _broadcast_all_ws({
                        "type": "arena.fertig",
                        "room_id": room_id,
                        "turn": turn,
                    })
                    break
                _fertig_this_round = set()  # Runde vorbei, reset
                # Inaktivitäts-Stopp: Runde ohne neue Aktion → Zähler hoch
                if _action_triggered_this_round:
                    _rounds_without_action = 0
                else:
                    _rounds_without_action += 1
                _action_triggered_this_round = False
                if _rounds_without_action >= _IDLE_ROUNDS_LIMIT and turn >= len(bots):
                    idle_msg = {
                        "role": "system",
                        "name": "System",
                        "content": f"🏁 **Arena beendet** — {_IDLE_ROUNDS_LIMIT} Runden ohne neue Aktionen. Diskussion hat sich erschöpft.",
                        "timestamp": datetime.now().isoformat(),
                        "color": "#6b7280",
                        "turn": turn,
                    }
                    rooms = _load_arena_rooms()
                    room = _get_arena_room(rooms, room_id)
                    if room:
                        room["messages"].append(idle_msg)
                        room["status"] = "fertig"
                        _save_arena_rooms(rooms)
                    if room_id in _arena_rooms:
                        _arena_rooms[room_id]["running"] = False
                    await _broadcast_all_ws({"type": "arena.message", "room_id": room_id, "message": idle_msg})
                    await _broadcast_all_ws({"type": "arena.fertig", "room_id": room_id, "turn": turn})
                    break

            # Auto-summary every N turns (default 10)
            summarize_every = room.get("summarize_every", 10)
            if summarize_every > 0 and turn > 0 and turn % summarize_every == 0 and turn < max_turns:
                # Generate orchestrator summary
                await _broadcast_all_ws({
                    "type": "arena.thinking",
                    "room_id": room_id,
                    "bot_name": "Orchestrator",
                    "bot_color": "#a855f7",
                    "turn": turn,
                })
                rooms = _load_arena_rooms()
                room = _get_arena_room(rooms, room_id)
                if room:
                    summary_raw = await _generate_orchestrator_summary(room)
                    verdict, summary = _extract_verdict(summary_raw)
                    # Zwischenfazit nur ins Log, nicht ans Fenster
                    if "fazit_log" not in room:
                        room["fazit_log"] = []
                    room["fazit_log"].append({
                        "turn": turn,
                        "timestamp": datetime.now().isoformat(),
                        "verdict": verdict,
                        "content": summary,
                    })
                    room["orchestrator_summary"] = summary
                    room["orchestrator_summary_at"] = datetime.now().isoformat()
                    room["last_verdict"] = verdict
                    _extract_and_store_action_items(room, summary)
                    _save_arena_rooms(rooms)
                    print(f"[Arena] Zwischenfazit Runde {turn} (verdict={verdict}): {summary[:120]}...")
                    # Auto-pause on consensus or stalemate — Robin soll entscheiden
                    if verdict in ("consensus", "stalemate"):
                        room["awaiting_robin"] = True
                        room["awaiting_reason"] = verdict
                        if verdict == "consensus":
                            room.setdefault("consensus_marks", []).append({
                                "turn": turn,
                                "at": datetime.now().isoformat(),
                                "msg_count": len(room.get("messages", [])),
                                "summary": summary[:300],
                            })
                        _save_arena_rooms(rooms)
                        await _broadcast_all_ws({
                            "type": "arena.needs_decision",
                            "room_id": room_id,
                            "verdict": verdict,
                            "turn": turn,
                            "summary": summary,
                        })
                        if room_id in _arena_rooms:
                            _arena_rooms[room_id]["running"] = False
                        break

            # Room-Renewal: Archive + Continuity-Seed wenn Raum zu groß wird (Lag-Schutz)
            rooms = _load_arena_rooms()
            room = _get_arena_room(rooms, room_id)
            if room:
                msgs = room.get("messages", [])
                try:
                    payload_kb = len(json.dumps(msgs, ensure_ascii=False).encode()) / 1024
                except Exception:
                    payload_kb = 0
                n = len(msgs)
                avg_msg_kb = payload_kb / n if n else 0
                renewal_by_density = (n >= ARENA_RENEW_AVG_MSG_MIN_COUNT and avg_msg_kb > ARENA_RENEW_AVG_MSG_KB)
                if n > ARENA_RENEW_MSG_COUNT or payload_kb > ARENA_RENEW_PAYLOAD_KB or renewal_by_density:
                    renew_info = _archive_and_renew_room(room)
                    _save_arena_rooms(rooms)
                    recovered_note = (
                        f", {renew_info.get('recovered', 0)} nachgetragen"
                        if renew_info.get("recovered") else ""
                    )
                    divider = {
                        "role": "system",
                        "name": "System",
                        "content": (
                            f"📦 **Raum komprimiert bei Turn {turn}** — {renew_info['archived_msgs']} Messages archiviert, "
                            f"{renew_info['open_actions']} offene Aktionen übertragen{recovered_note} → Archiv: {renew_info['ts']}"
                            + (f" (⚠️ Dichte: {avg_msg_kb:.1f} KB/msg)" if renewal_by_density else "")
                        ),
                        "timestamp": datetime.now().isoformat(),
                        "color": "#6366f1",
                        "turn": turn,
                        "is_renewal_divider": True,
                    }
                    rooms2 = _load_arena_rooms()
                    room2 = _get_arena_room(rooms2, room_id)
                    if room2:
                        room2["messages"].append(divider)
                        _save_arena_rooms(rooms2)
                    await _broadcast_all_ws({
                        "type": "arena.message",
                        "room_id": room_id,
                        "message": divider,
                    })
                    await _broadcast_all_ws({
                        "type": "arena.renewed",
                        "room_id": room_id,
                        "turn": turn,
                        "archived": renew_info["archived_msgs"],
                        "ts": renew_info["ts"],
                    })
                    print(f"[Arena] Room {room_id[:8]} renewed at turn {turn}: {renew_info['archived_msgs']} msgs archived, {renew_info['open_actions']} open actions seeded")

            # Small pause between turns (let Robin inject if needed)
            await asyncio.sleep(2)

    except Exception as e:
        print(f"[Arena] Error in room {room_id[:8]}: {e}")
    finally:
        # Generate final summary
        rooms = _load_arena_rooms()
        room = _get_arena_room(rooms, room_id)
        if room and turn > 2:
            try:
                await _broadcast_all_ws({
                    "type": "arena.thinking",
                    "room_id": room_id,
                    "bot_name": "Orchestrator",
                    "bot_color": "#a855f7",
                    "turn": turn,
                })
                summary = await _generate_orchestrator_summary(room)
                summary_msg = {
                    "role": "orchestrator",
                    "name": "Orchestrator",
                    "content": f"**Abschlussfazit ({turn} Runden):**\n\n{summary}",
                    "timestamp": datetime.now().isoformat(),
                    "color": "#a855f7",
                    "turn": turn,
                }
                _msgs = room.get("messages", [])
                _already_final = _msgs and _msgs[-1].get("role") == "orchestrator"
                if not _already_final:
                    room["messages"].append(summary_msg)
                room["orchestrator_summary"] = summary
                room["orchestrator_summary_at"] = datetime.now().isoformat()
                room["status"] = "idle"
                _extract_and_store_action_items(room, summary)
                _save_arena_rooms(rooms)
                if not _already_final:
                    await _broadcast_all_ws({
                        "type": "arena.message",
                        "room_id": room_id,
                        "message": summary_msg,
                    })
            except Exception as e2:
                print(f"[Arena] Summary error: {e2}")
                room["status"] = "idle"
                _save_arena_rooms(rooms)
        elif room:
            room["status"] = "idle"
            _save_arena_rooms(rooms)

        _activity_set("arena", active=False)
        _arena_rooms.pop(room_id, None)
        await _broadcast_all_ws({
            "type": "arena.stopped",
            "room_id": room_id,
        })
        print(f"[Arena] Room {room_id[:8]} stopped after {turn} turns")


# ---------------------------------------------------------------------------
# WS Client Registry + Server Version
# ---------------------------------------------------------------------------
import hashlib as _hashlib

# Compute a version hash from all static files + server.py at startup
def _compute_version():
    h = _hashlib.md5()
    for f in sorted((BASE_DIR / "static").rglob("*")):
        if f.is_file():
            h.update(f.read_bytes())
    h.update(Path(__file__).read_bytes())
    return h.hexdigest()[:12]

_server_version = _compute_version()
_server_start_time = time.time()
_ws_clients: set = set()  # all connected WS clients


async def _broadcast_all_ws(event: dict):
    """Send an event to ALL connected WS clients."""
    dead = set()
    for ws in list(_ws_clients):
        try:
            await ws.send_json(event)
        except Exception:
            dead.add(ws)
    for d in dead:
        _ws_clients.discard(d)


# ── Activity State (Sidebar indicators) ──────────────────────────────
_activity: dict = {}  # {view: {"active": bool, "pending": int}}


def _activity_set(view: str, *, active: bool | None = None, pending_delta: int = 0):
    """Update activity state and broadcast to all WS clients."""
    state = _activity.setdefault(view, {"active": False, "pending": 0})
    if active is not None:
        state["active"] = active
    if pending_delta != 0:
        state["pending"] = max(0, state["pending"] + pending_delta)
    asyncio.ensure_future(_broadcast_all_ws({
        "type": "activity.update",
        "view": view,
        "state": {"active": state["active"], "pending": state["pending"]},
    }))


async def api_activity_get(request):
    """GET /api/activity  -  Return current activity state for all views."""
    return web.json_response(_activity)


async def api_activity_clear(request):
    """POST /api/activity/clear/{view}  -  Clear pending count for a view."""
    view = request.match_info["view"]
    if view in _activity:
        _activity[view]["pending"] = 0
        await _broadcast_all_ws({
            "type": "activity.update",
            "view": view,
            "state": {"active": _activity[view]["active"], "pending": 0},
        })
    return web.json_response({"ok": True})


# ---------------------------------------------------------------------------
# Stream-Buffer + Subscriber-Registry
# ---------------------------------------------------------------------------
# session_id -> {"buffer": [event_dicts], "done": bool}
_active_streams: dict = {}

# session_id -> set of websocket objects currently subscribed to this stream
_stream_subscribers: dict = {}


def _stream_start(session_id: str):
    """Mark a session as actively streaming; reset buffer."""
    _active_streams[session_id] = {"buffer": [], "done": False}
    _stream_subscribers.setdefault(session_id, set())


def _stream_end(session_id: str):
    """Mark stream as done; keep buffer for late subscribers."""
    if session_id in _active_streams:
        _active_streams[session_id]["done"] = True


def _stream_cleanup(session_id: str):
    """Remove stream state entirely (e.g. after client acknowledged done)."""
    _active_streams.pop(session_id, None)
    _stream_subscribers.pop(session_id, None)


def _ws_unsubscribe_all(ws):
    """Remove a dead WS connection from all subscriber sets."""
    for subs in _stream_subscribers.values():
        subs.discard(ws)


async def _broadcast_stream_event(session_id: str, event: dict):
    """Buffer event and broadcast to all subscribers for this session."""
    if session_id in _active_streams:
        _active_streams[session_id]["buffer"].append(event)
    dead = set()
    for sub_ws in list(_stream_subscribers.get(session_id, [])):
        try:
            await sub_ws.send_json(event)
        except Exception:
            dead.add(sub_ws)
    for d in dead:
        _stream_subscribers.get(session_id, set()).discard(d)


# WebSocket chat handler
# ---------------------------------------------------------------------------

async def _classify_session_headless(session):
    """Classify a session without a WS connection (e.g. on startup)."""
    print(f"[Classify] Starting headless for {session.session_id[:8]}, history={len(session.history)} msgs")
    user_msgs = [m["content"] for m in session.history if m.get("role") == "user"][:5]
    if not user_msgs:
        print(f"[Classify] No user msgs for {session.session_id[:8]}, skipping")
        return
    sample = "\n---\n".join(user_msgs)[:1500]
    print(f"[Classify] {session.session_id[:8]} sample={len(sample)} chars, {len(user_msgs)} user msgs")
    existing_topics = list(sessions.get_topics().keys())
    existing_hint = ""
    if existing_topics:
        existing_hint = f"\nExisting topics: {', '.join(existing_topics)}\nPrefer assigning to an existing topic if it fits well. Only create a new one if the session clearly doesn't match any existing topic."

    prompt = f"""Analyze this chat session and return a JSON object with:
- "topic": short category label (1-3 words, German)
- "title": concise session title (max 5 words, German, describes the specific content)
- "new": true if topic is new, false if it matches an existing one

Return ONLY valid JSON like: {{"topic": "Thema Name", "title": "Kurzer Titel", "new": true/false}}
{existing_hint}

User messages:
{sample}"""

    try:
        proc = await asyncio.create_subprocess_exec(
            "claude", "-p", "--output-format", "text", "--max-turns", "0",
            "--model", "haiku",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            stdin=asyncio.subprocess.PIPE,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        proc.stdin.write(prompt.encode("utf-8"))
        await proc.stdin.drain()
        proc.stdin.close()
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=30)
        raw = stdout.decode("utf-8", errors="replace").strip()
        err = stderr.decode("utf-8", errors="replace").strip()
        print(f"[Classify] {session.session_id[:8]} raw={repr(raw[:200])} err={repr(err[:200])}")
        import re as _re
        json_match = _re.search(r'\{[^}]*\}', raw)
        if json_match:
            result = json.loads(json_match.group())
            topic = result.get("topic", "").strip()
            title = result.get("title", "").strip()
            if topic:
                color = sessions._topic_colors().get(topic) or SessionManager._default_color(topic)
                sessions._save_topic_color(topic, color)
                sessions.set_session_topic(session.session_id, topic, color)
                if title and not session.custom_title:
                    session.auto_title = title
                    sessions.persist()
                print(f"[Classify] {session.session_id[:8]} → '{topic}' | '{title}'")
                await _broadcast_all_ws({
                    "type": "chat.sessions",
                    "sessions": sessions.list_sessions(),
                    "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]],
                })
                await _broadcast_all_ws({
                    "type": "chat.topic_classified",
                    "session_id": session.session_id,
                    "topic": topic,
                    "color": color,
                    "title": title,
                    "is_new": False,
                })
        else:
            print(f"[Classify] {session.session_id[:8]} NO JSON match in response")
    except Exception as e:
        print(f"[Classify] Headless error for {session.session_id[:8]}: {e}")


async def _classify_session(ws, session):
    """Use a lightweight Claude call to classify the session into a topic."""
    # Collect first few user messages for classification
    user_msgs = [m["content"] for m in session.history if m.get("role") == "user"][:5]
    if not user_msgs:
        return
    sample = "\n---\n".join(user_msgs)[:1500]

    # Get existing topics
    existing_topics = list(sessions.get_topics().keys())
    existing_hint = ""
    if existing_topics:
        existing_hint = f"\nExisting topics: {', '.join(existing_topics)}\nPrefer assigning to an existing topic if it fits well. Only create a new one if the session clearly doesn't match any existing topic."

    prompt = f"""Analyze this chat session and return a JSON object with:
- "topic": short category label (1-3 words, German)
- "title": concise session title (max 5 words, German, describes the specific content)
- "new": true if topic is new, false if it matches an existing one

Return ONLY valid JSON like: {{"topic": "Thema Name", "title": "Kurzer Titel", "new": true/false}}
{existing_hint}

User messages:
{sample}"""

    try:
        proc = await asyncio.create_subprocess_exec(
            "claude", "-p", "--output-format", "text", "--max-turns", "0",
            "--model", "haiku",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.DEVNULL,
            stdin=asyncio.subprocess.PIPE,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        proc.stdin.write(prompt.encode("utf-8"))
        await proc.stdin.drain()
        proc.stdin.close()
        stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=15)
        raw = stdout.decode("utf-8", errors="replace").strip()

        # Parse JSON from response (might have markdown wrapping)
        import re
        json_match = re.search(r'\{[^}]*\}', raw)
        if json_match:
            result = json.loads(json_match.group())
            topic = result.get("topic", "").strip()
            title = result.get("title", "").strip()
            if topic:
                color = sessions._topic_colors().get(topic) or SessionManager._default_color(topic)
                sessions._save_topic_color(topic, color)
                sessions.set_session_topic(session.session_id, topic, color)
                if title and not session.custom_title:
                    session.auto_title = title
                    sessions.persist()
                print(f"[Classify] Session {session.session_id[:8]} → '{topic}' | '{title}' (new={result.get('new', True)})")
                await _broadcast_all_ws({
                    "type": "chat.sessions",
                    "sessions": sessions.list_sessions(),
                    "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]],
                })
                await _broadcast_all_ws({
                    "type": "chat.topic_classified",
                    "session_id": session.session_id,
                    "topic": topic,
                    "color": color,
                    "title": title,
                    "is_new": result.get("new", True),
                })
    except Exception as e:
        print(f"[Classify] Error: {e}")


# --- Chat-Helpers für provider-agnostischen Stream (Chat-Refactor 2026-04-23) ---
#
# Diese Helpers trennen State-Aufbau von der Stream-Ausführung. `_run_chat` weiter
# unten nutzt sie, um neutrale LLM-Calls via invoke_llm zu fahren.

_DENY_TOOL_PATTERNS = ",".join([
    "Bash(taskkill*)",
    "Bash(Stop-Process*)",
    "Bash(Stop-Service*)",
    "Bash(net stop*)",
    "Bash(sc stop*)",
    "Bash(Restart-Service*)",
    "Bash(*python*server.py*)",
    "Bash(*start.bat*)",
    "Edit(C:/Users/Robin/Jarvis/Mission-Control/server.py)",
    "Edit(C:/Users/Robin/Jarvis/Mission-Control/start.bat)",
    "Write(C:/Users/Robin/Jarvis/Mission-Control/server.py)",
    "Write(C:/Users/Robin/Jarvis/Mission-Control/start.bat)",
])


def _history_to_messages(session, current_user_text: str, attachment_notes: list[str]):
    """Baut eine neutrale Message-Liste aus session.history + dem aktuellen User-Input.

    Attachments werden als Text-Hinweise ans Ende des User-Prompts angehängt — sie
    werden NOCH NICHT als strukturierte ImageBlocks übergeben (Image-Gate-Task
    erledigt das separat, pro Provider nach supports_images).

    Für Provider mit supports_resume=True und gesetztem provider_thread_id sollte
    der Aufrufer die History eigentlich weglassen (Provider hat sie intern) — das
    macht _build_provider_kwargs transparent über cli_session_id. Hier bauen wir
    immer die volle History, der Provider entscheidet.
    """
    from llm_client import Message, TextBlock

    msgs: list = []
    for h in session.history:
        role = h.get("role")
        content = h.get("content", "") or ""
        if role not in ("user", "assistant"):
            continue
        if not content.strip():
            continue
        # Truncate sehr lange Einträge (wie im alten Code: 2000 Zeichen)
        if len(content) > 2000:
            content = content[:2000] + "...[gekürzt]"
        msgs.append(Message(role=role, content=content))

    # Aktuelle User-Nachricht + Attachment-Hinweise
    full_user = current_user_text
    if attachment_notes:
        full_user = full_user + "\n\n" + "\n".join(attachment_notes)
    msgs.append(Message(role="user", content=full_user))
    return msgs


def _build_attachment_notes(attachments, supports_images: bool) -> tuple[list[str], list]:
    """Wandelt Attachment-Filenames in (notes, image_blocks).

    notes: Text-Zeilen die immer in den Prompt gehängt werden (egal welcher Provider).
    image_blocks: strukturierte ImageBlocks — NUR wenn supports_images=True.
                  Aktuell noch nicht verdrahtet (siehe Image-Gate-Task) — wird in
                  späterer Iteration genutzt, aktuell leer zurückgegeben.

    Image-Gate: Bei supports_images=False degradieren wir zu Text-Hinweis
    ("[Bild angehängt: X]") statt strukturiertem Block.
    """
    notes: list[str] = []
    image_blocks: list = []  # placeholder für zukünftige ImageBlock-Nutzung
    if not attachments:
        return notes, image_blocks

    for filename in attachments:
        filepath = UPLOADS_DIR / filename
        if not filepath.exists():
            continue
        ext = filepath.suffix.lower()
        if ext in ('.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'):
            if supports_images:
                # TODO: echte ImageBlocks bauen wenn wir das Bild mitschicken wollen
                notes.append(f"[Bild angehängt: {filepath}  -  bitte lies und analysiere dieses Bild]")
            else:
                notes.append(f"[Bild angehängt: {filepath}  -  Pfad zum Bild, Provider kann Bilder nicht direkt verarbeiten]")
        elif ext == '.pdf':
            notes.append(f"[PDF angehängt: {filepath}  -  bitte lies und analysiere dieses PDF]")
        elif ext in ('.txt', '.md', '.csv', '.json', '.py', '.js', '.html', '.css'):
            try:
                content = filepath.read_text(encoding="utf-8", errors="replace")[:5000]
                notes.append(f"[Datei: {filename}]\n```\n{content}\n```")
            except Exception:
                notes.append(f"[Datei angehängt: {filepath}]")
        elif ext in ('.webm', '.ogg', '.mp3', '.wav', '.m4a'):
            notes.append(f"[Sprachnachricht angehängt: {filepath}  -  bitte lies und transkribiere diese Audiodatei]")
        else:
            notes.append(f"[Datei angehängt: {filepath}]")
    return notes, image_blocks


# Verifiziert 2026-04-25: codex-cli mit ChatGPT-Account akzeptiert nur diese
# Models. Alles andere → "model not supported when using Codex with a ChatGPT
# account" → leerer Stream → "keine Antwort". Wir normalisieren defensiv damit
# ein veralteter Picker-Wert (vor Hard-Reload) nicht jede Session zerlegt.
_CODEX_OK_MODELS = {"", "gpt-5.5"}


def _normalize_codex_model(model: str | None) -> tuple[str, str | None]:
    """Returns (model, warning_msg). Wenn model fuer codex-cli nicht unterstuetzt,
    auf '' (Account-Default) zurueckfallen und Warning fuer den Caller liefern."""
    if model is None:
        return model, None
    if model in _CODEX_OK_MODELS:
        return model, None
    return "", (
        f"⚠️ Model {model!r} ist mit dem ChatGPT-Codex-Account nicht verfuegbar — "
        "auf Account-Default zurueckgesetzt. Im Picker stehen 'Auto' oder 'gpt-5.5'."
    )


def _set_session_provider(session, new_provider: str, new_model: str | None = None) -> dict:
    """Zentrale Switch-Logik für Provider/Model-Wechsel pro Session.

    Verhalten (nach Codex-Review 2026-04-23):
      - Nur Model-Wechsel im selben Provider → kein thread_id-Reset, keine Systemmeldung
      - Provider-Wechsel → provider_thread_id HART auf None, Systemmeldung in history,
        cli_session_id (Legacy-Alias) synchron auf None
      - Invalid Provider → ValueError

    Returns: {"changed": bool, "provider_switched": bool, "message": str|None}
    """
    from llm_client import PROVIDER_NAMES

    if new_provider not in PROVIDER_NAMES:
        raise ValueError(
            f"Unbekannter Provider {new_provider!r}. Erlaubt: {', '.join(PROVIDER_NAMES)}"
        )

    codex_warning = None
    if new_provider == "codex-cli":
        new_model, codex_warning = _normalize_codex_model(new_model)
        if codex_warning:
            session.history.append({"role": "system", "content": codex_warning})

    old_provider = session.provider
    old_model = session.model
    provider_switched = (new_provider != old_provider)
    model_changed = (new_model is not None and new_model != old_model)

    if not provider_switched and not model_changed:
        return {"changed": False, "provider_switched": False, "message": None}

    sys_msg = None
    if provider_switched:
        # Resume-Token HART verwerfen — providerspezifisch, nicht transferierbar
        session.provider_thread_id = None
        session.cli_session_id = None
        # Systemmeldung in history damit User Kontextbruch sieht
        sys_msg = (
            f"🔄 Provider gewechselt: {old_provider} → {new_provider}. "
            f"Neue Conversation (Resume-Kontext verworfen)."
        )
        session.history.append({"role": "system", "content": sys_msg})

    session.provider = new_provider
    if new_model is not None:
        session.model = new_model
    sessions.persist()

    print(f"[Chat] provider-switch session={session.session_id[:8]} "
          f"{old_provider}/{old_model!r} → {new_provider}/{session.model!r} "
          f"(thread_reset={provider_switched})")

    return {
        "changed": True,
        "provider_switched": provider_switched,
        "message": sys_msg,
    }


def _build_provider_kwargs(session, on_resume_token=None) -> dict:
    """Provider-spezifische Extras für invoke_llm. Keine Branch-Logik auf
    provider-Namen außerhalb dieser Funktion — hier ist der einzige Ort.

    on_resume_token: Callback, der vom Provider bei neuem Thread/Session-ID
      aufgerufen wird. Wir hängen ihn hier an die Session-Instanz.
    """
    from pathlib import Path as _P

    kwargs: dict = {}
    p = session.provider

    if p == "claude-cli":
        # Bei --resume braucht claude-cli die alte CLI-Session-ID
        if session.provider_thread_id:
            kwargs["cli_session_id"] = session.provider_thread_id
        # MCP + System-Prompt-Handling macht der Caller (_run_chat) selbst,
        # weil der system_file-Pfad temp/pro-call ist
        kwargs["mcp_config"] = str(BASE_DIR.parent / ".mcp.json")
        # max_turns 40 war zu knapp fuer agentische Tasks (Jarvis-Main mit
        # Multi-Tool-Sessions hat Antworten abgebrochen). 100 ist immer noch
        # eine harte Grenze gegen runaway loops, aber gibt genug Luft fuer
        # Edit→Read→Bash→Edit-Ketten.
        kwargs["max_turns"] = 100
        kwargs["disallowed_tools"] = _DENY_TOOL_PATTERNS
        kwargs["dangerous_skip_permissions"] = True
        # Callback für CLI-Session-ID aus init-Event
        if on_resume_token:
            kwargs["on_cli_session_id"] = on_resume_token

    elif p == "codex-cli":
        kwargs["sandbox"] = "workspace-write"
        kwargs["cwd"] = "C:/Users/Robin/Jarvis"
        kwargs["skip_git_repo_check"] = True
        if session.provider_thread_id:
            kwargs["thread_id"] = session.provider_thread_id  # (Resume-Impl noch offen)
        if on_resume_token:
            kwargs["on_thread_id"] = on_resume_token

    elif p == "gemini-cli":
        kwargs["approval_mode"] = "yolo"  # headless = auto-approve Tools
        kwargs["cwd"] = "C:/Users/Robin/Jarvis"
        if on_resume_token:
            kwargs["on_session_id"] = on_resume_token

    # claude-api und openai brauchen keine provider_kwargs — alles über
    # messages-Liste. Keine Sonderfälle.

    return kwargs


async def _llm_oneshot(
    provider: str,
    model: str,
    user_message: str,
    system: str | None = None,
    max_tokens: int = 2000,
    **kwargs,
) -> tuple[str, dict, str | None]:
    """Ruft invoke_llm und sammelt text+usage bis zum MessageStop. Non-streaming."""
    from llm_client import (
        invoke_llm, Message, TextDelta, ErrorEvent, MessageStop, UsageEvent,
    )
    answer = ""
    usage: dict = {}
    err: str | None = None
    try:
        async for ev in invoke_llm(
            provider=provider, model=model,
            messages=[Message(role="user", content=user_message)],
            system=system, max_tokens=max_tokens,
            **kwargs,
        ):
            if isinstance(ev, TextDelta):
                answer += ev.text
            elif isinstance(ev, UsageEvent):
                usage = {
                    "input_tokens": ev.input_tokens + ev.cache_read_tokens,
                    "output_tokens": ev.output_tokens,
                }
            elif isinstance(ev, ErrorEvent):
                err = f"{ev.code}: {ev.message}"
            elif isinstance(ev, MessageStop):
                break
    except Exception as e:
        err = f"{type(e).__name__}: {e}"
    return answer, usage, err


async def api_llm_ask_peer(request):
    """POST /api/llm/ask-peer
    Body: {provider, model?, query, context?, session_id?}

    "Codex, frag mal Claude ob…" — eine Provider-Session kann via diesem Endpoint
    einen anderen Provider einmalig fragen. Non-streaming, one-shot. Wenn
    session_id mitgegeben wird, sehen die Subscriber dieser Session ein
    `chat.peer_answer`-Event mit Query+Answer.
    """
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    provider = (body.get("provider") or "").strip()
    model = (body.get("model") or "").strip()
    query = (body.get("query") or "").strip()
    context = (body.get("context") or "").strip()
    session_id = (body.get("session_id") or "").strip()

    if not provider or not query:
        return web.json_response(
            {"ok": False, "error": "provider und query erforderlich"}, status=400,
        )

    system_prompt = (
        "Du bist Peer-Assistent für einen anderen AI-Agent von Robin. "
        "Beantworte die Frage knapp, präzise, sachlich, ohne Höflichkeits-Floskeln. "
        "Wenn Kontext beigefügt ist, nutze ihn als Grundlage."
    )
    if context:
        system_prompt += f"\n\n--- BEIGEFÜGTER KONTEXT ---\n{context[:8000]}"

    answer, usage, err = await _llm_oneshot(
        provider=provider, model=model,
        user_message=query, system=system_prompt,
    )

    if session_id and sessions.sessions.get(session_id):
        await _broadcast_stream_event(session_id, {
            "type": "chat.peer_answer",
            "session_id": session_id,
            "peer_provider": provider,
            "peer_model": model or "default",
            "query": query,
            "answer": answer,
            "error": err,
        })

    return web.json_response({
        "ok": err is None and bool(answer),
        "provider": provider, "model": model or "default",
        "query": query, "answer": answer,
        "usage": usage, "error": err,
    })


async def api_chat_session_scratchpad(request):
    """Shared Scratchpad pro Session — provider-agnostisch lesen/schreiben.

    GET  → {notes: {...}}
    POST → {key, value, append?} ODER {notes: {...}, replace?}

    Beide Provider können über diesen dict[str, str] Zwischenergebnisse teilen.
    Broadcastet `chat.scratchpad_update`-Event an Session-Subscriber.
    """
    sid = request.match_info.get("session_id", "")
    session = sessions.sessions.get(sid)
    if not session:
        return web.json_response({"ok": False, "error": "Session nicht gefunden"}, status=404)

    # Backward-Compat: Feld existiert auf alten Sessions evtl. nicht
    if not hasattr(session, "shared_notes") or session.shared_notes is None:
        session.shared_notes = {}

    if request.method == "GET":
        return web.json_response({
            "ok": True, "session_id": session.session_id,
            "notes": dict(session.shared_notes),
        })

    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    if "notes" in body:
        bulk = body.get("notes") or {}
        if not isinstance(bulk, dict):
            return web.json_response({"ok": False, "error": "notes muss Dict sein"}, status=400)
        if body.get("replace"):
            session.shared_notes = {}
        for k, v in bulk.items():
            session.shared_notes[str(k)] = str(v)
    elif "key" in body:
        key = str(body.get("key") or "").strip()
        val = str(body.get("value") or "")
        if not key:
            return web.json_response({"ok": False, "error": "key erforderlich"}, status=400)
        if body.get("append") and key in session.shared_notes:
            session.shared_notes[key] = session.shared_notes[key] + "\n" + val
        else:
            session.shared_notes[key] = val
    else:
        return web.json_response(
            {"ok": False, "error": "Body: entweder {notes:{...}} oder {key,value}"},
            status=400,
        )

    sessions.persist()
    await _broadcast_stream_event(session.session_id, {
        "type": "chat.scratchpad_update",
        "session_id": session.session_id,
        "notes": dict(session.shared_notes),
    })

    return web.json_response({
        "ok": True, "session_id": session.session_id,
        "notes": dict(session.shared_notes),
    })


async def _fetch_amboss_medical(query: str, uploads_dir: Path) -> list[dict]:
    """Amboss Top-1-Treffer zur Query: search → fetch_article_full.
    Kein Image-Download. Return list[{title, eid, text, abstract}].
    """
    try:
        from amboss_client import (
            _load_cached_cookie, search_article_eids, fetch_article_full,
        )
    except ImportError:
        return []
    # FIX: _load_cookie hieß in amboss_client.py eigentlich _load_cached_cookie
    cookie = _load_cached_cookie(uploads_dir) or CONFIG.get("amboss_cookie", "").strip()
    if not cookie:
        return []
    try:
        import aiohttp as _aiohttp
        timeout = _aiohttp.ClientTimeout(total=45)
        async with _aiohttp.ClientSession(connector=_aiohttp.TCPConnector(ssl=True), timeout=timeout) as sess:
            eids = await search_article_eids(sess, cookie, query, limit=2)
            if not eids:
                return []
            results = []
            for eid_obj in (eids[:2] if isinstance(eids, list) else []):
                try:
                    eid = eid_obj.get("articleEid") if isinstance(eid_obj, dict) else eid_obj
                    full = await fetch_article_full(sess, cookie, eid, media_limit=0, text_char_limit=3000)
                    if full and full.get("text"):
                        results.append({
                            "title": full.get("title", ""),
                            "eid": full.get("eid", eid),
                            "abstract": full.get("abstract", ""),
                            "text": full.get("text", ""),
                        })
                except Exception as e:
                    print(f"[MedAgent] amboss article {eid_obj} failed: {e}")
            return results
    except Exception as e:
        print(f"[MedAgent] amboss search failed: {e}")
        return []


async def _fetch_viamedici_medical(query: str, uploads_dir: Path) -> list[dict]:
    """ViaMedici Top-2-Module zur Query: search → module_content.
    Return list[{title, assetId, text}].
    """
    try:
        from viamedici_client import (
            ensure_access_token, search_learning_modules, fetch_module_content,
        )
    except ImportError:
        return []
    try:
        access_token = await ensure_access_token(uploads_dir)
    except Exception as e:
        print(f"[MedAgent] viamedici auth skipped: {e}")
        return []
    try:
        import aiohttp as _aiohttp
        timeout = _aiohttp.ClientTimeout(total=45)
        async with _aiohttp.ClientSession(connector=_aiohttp.TCPConnector(ssl=True), timeout=timeout) as sess:
            modules = await search_learning_modules(sess, access_token, query, limit=2)
            if not modules:
                return []
            results = []
            for mod in modules[:2]:
                asset_id = mod.get("assetId")
                slug = mod.get("slug")
                if not asset_id or not slug:
                    continue
                try:
                    text = await fetch_module_content(sess, access_token, asset_id, slug)
                    if text:
                        results.append({
                            "title": mod.get("title", ""),
                            "assetId": asset_id,
                            "text": text,
                        })
                except Exception as e:
                    print(f"[MedAgent] viamedici module {asset_id} failed: {e}")
            return results
    except Exception as e:
        print(f"[MedAgent] viamedici search failed: {e}")
        return []


async def api_image_providers(request):
    """GET /api/image/providers — Status pro Image-Gen-Provider (configured?)."""
    from image_client import list_image_providers
    return web.json_response({"providers": list_image_providers()})


async def api_image_usage(request):
    """GET /api/image/usage?period=today|month|all — Cost-Tracking.
    Returns: {total_images, total_usd, by_provider, by_model, recent, period}
    """
    from image_client import read_usage_log
    period = (request.rel_url.query.get("period") or "all").lower()
    if period not in ("today", "month", "all"):
        period = "all"
    return web.json_response(read_usage_log(period=period))


def _find_latest_codex_image() -> Path | None:
    """Return the newest image generated by Codex's local image tool, if any."""
    root = Path.home() / ".codex" / "generated_images"
    if not root.exists():
        return None
    exts = {".png", ".jpg", ".jpeg", ".webp"}
    candidates = [
        p for p in root.rglob("*")
        if p.is_file() and p.suffix.lower() in exts
    ]
    if not candidates:
        return None
    return max(candidates, key=lambda p: p.stat().st_mtime)


def _codex_job_path(job_id: str) -> Path:
    safe = re.sub(r"[^a-zA-Z0-9_-]", "", job_id or "")
    return IMAGE_JOBS_DIR / f"{safe}.json"


def _read_codex_job(job_id: str) -> dict | None:
    path = _codex_job_path(job_id)
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None


def _write_codex_job(job: dict) -> None:
    atomic_write_text(
        _codex_job_path(job["id"]),
        json.dumps(job, ensure_ascii=False, indent=2),
        keep_backup=False,
    )


def _list_codex_jobs(status: str | None = None) -> list[dict]:
    jobs = []
    for path in IMAGE_JOBS_DIR.glob("*.json"):
        try:
            job = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            continue
        if status and job.get("status") != status:
            continue
        jobs.append(job)
    jobs.sort(key=lambda j: j.get("created_at", ""), reverse=True)
    return jobs


def _import_codex_image(latest: Path, prompt: str = "", job_id: str = "") -> dict:
    import shutil

    out_dir = UPLOADS_DIR / "images" / "generated"
    out_dir.mkdir(parents=True, exist_ok=True)
    ext = latest.suffix.lower() or ".png"
    dest = out_dir / f"codex_{int(time.time())}_{uuid.uuid4().hex[:8]}{ext}"
    shutil.copy2(latest, dest)

    rel = str(dest.relative_to(UPLOADS_DIR)).replace("\\", "/")
    return {
        "ok": True,
        "provider": "codex-session" if job_id else "codex-drop",
        "model": "codex image tool",
        "prompt": prompt or "Imported from Codex generated_images",
        "size": "",
        "url": "/api/uploads/" + rel,
        "path": rel,
        "source_path": str(latest),
        "job_id": job_id,
        "mime_type": "image/png" if ext == ".png" else f"image/{ext.lstrip('.')}",
    }


async def api_image_codex_latest(request):
    """GET /api/image/codex/latest — newest local Codex-generated image."""
    latest = _find_latest_codex_image()
    if not latest:
        return web.json_response({
            "ok": False,
            "error": "Kein Codex-Bild unter ~/.codex/generated_images gefunden",
        }, status=200)
    return web.json_response({
        "ok": True,
        "path": str(latest),
        "filename": latest.name,
        "mtime": latest.stat().st_mtime,
        "size": latest.stat().st_size,
    })


async def api_image_codex_import(request):
    """POST /api/image/codex/import — copy latest Codex image into MC uploads."""
    latest = _find_latest_codex_image()
    if not latest:
        return web.json_response({
            "ok": False,
            "error": "Kein Codex-Bild unter ~/.codex/generated_images gefunden",
        }, status=200)

    return web.json_response(_import_codex_image(latest))


async def api_image_codex_jobs_create(request):
    """POST /api/image/codex/jobs — create a manual Codex image bridge job."""
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    prompt = (body.get("prompt") or "").strip()
    if len(prompt) < 3:
        return web.json_response({"ok": False, "error": "prompt erforderlich"}, status=400)

    now = datetime.now().isoformat(timespec="seconds")
    job = {
        "id": uuid.uuid4().hex[:12],
        "status": "pending",
        "prompt": prompt,
        "size": (body.get("size") or "").strip(),
        "style": (body.get("style") or "").strip(),
        "created_at": now,
        "updated_at": now,
        "result": None,
        "error": None,
        "notified": False,
    }
    _write_codex_job(job)
    notified = send_push(
        "Codex Image Job wartet",
        f"{job['id']}: {prompt[:180]}",
        priority=4,
        tags=["art", "hourglass"],
        click="http://localhost:8090/",
    )
    job["notified"] = bool(notified)
    _write_codex_job(job)
    return web.json_response({"ok": True, "job": job})


async def api_image_codex_jobs_list(request):
    """GET /api/image/codex/jobs?status=pending — list manual bridge jobs."""
    status = (request.query.get("status") or "").strip() or None
    return web.json_response({"ok": True, "jobs": _list_codex_jobs(status=status)})


async def api_image_codex_job_get(request):
    """GET /api/image/codex/jobs/{job_id} — job status."""
    job_id = request.match_info.get("job_id", "")
    job = _read_codex_job(job_id)
    if not job:
        return web.json_response({"ok": False, "error": "Job nicht gefunden"}, status=404)
    return web.json_response({"ok": True, "job": job})


async def api_image_codex_job_complete(request):
    """POST /api/image/codex/jobs/{job_id}/complete — attach latest Codex image."""
    job_id = request.match_info.get("job_id", "")
    job = _read_codex_job(job_id)
    if not job:
        return web.json_response({"ok": False, "error": "Job nicht gefunden"}, status=404)

    latest = _find_latest_codex_image()
    if not latest:
        job["status"] = "error"
        job["error"] = "Kein Codex-Bild unter ~/.codex/generated_images gefunden"
        job["updated_at"] = datetime.now().isoformat(timespec="seconds")
        _write_codex_job(job)
        return web.json_response({"ok": False, "error": job["error"], "job": job}, status=200)

    result = _import_codex_image(latest, prompt=job.get("prompt", ""), job_id=job_id)
    job["status"] = "completed"
    job["result"] = result
    job["error"] = None
    job["updated_at"] = datetime.now().isoformat(timespec="seconds")
    _write_codex_job(job)
    return web.json_response({"ok": True, "job": job, **result})


async def api_image_codex_job_fail(request):
    """POST /api/image/codex/jobs/{job_id}/fail — mark bridge job failed."""
    job_id = request.match_info.get("job_id", "")
    job = _read_codex_job(job_id)
    if not job:
        return web.json_response({"ok": False, "error": "Job nicht gefunden"}, status=404)
    try:
        body = await request.json()
    except Exception:
        body = {}
    job["status"] = "error"
    job["error"] = (body.get("error") or "Codex job failed").strip()
    job["updated_at"] = datetime.now().isoformat(timespec="seconds")
    _write_codex_job(job)
    return web.json_response({"ok": True, "job": job})


_IMAGE_PROMPT_SKILL_PATH = Path(r"C:\Users\Robin\Jarvis\skills\image-prompt\SKILL.md")

# One-Shot-Override: Der Skill ist auf interaktiven Multi-Step-Dialog ausgelegt
# (asks dimensions, refines etc.). Im Auto-Enhance-Modus wollen wir aber genau
# einen Output: den fertigen Prompt, ohne Rationale-Block, ohne Rueckfrage.
_IMAGE_PROMPT_ONESHOT_TAIL = """

---

## ONE-SHOT MODE (CRITICAL OVERRIDE — read this first)

Du bist im non-interaktiven Auto-Enhance-Modus. **Du darfst KEINE Tools aufrufen** — nicht
Skill, nicht Task, nicht Read, nichts. Antworte ausschliesslich mit reinem Text.

Ignoriere Steps 1, 3 und 4 oben (keine Nachfragen, kein "PROMPT:"-Block, keine Rationale,
keine Refinement-Frage). Nutze NUR die Core Principles + Step-2-Struktur + Anti-Patterns.

Nimm direkt die User-Eingabe als Vision (auch wenn sie kurz/vage ist) und antworte mit
GENAU EINEM fertigen Image-Generation-Prompt — als deine erste und einzige Text-Antwort.

Strikte Output-Regeln:
- KEIN Tool-Use, KEIN Thinking sichtbar — nur der finale Prompt-Text.
- Keine Anfuehrungszeichen, keine Markdown, keine Erklaerung, kein "PROMPT:"-Header.
- Englisch (Image-Modelle reagieren darauf besser).
- 60-200 Woerter. Prompts unter 40 Woerter sind zu vage.
- Wenn der Input einen Personennamen enthaelt → 1:1 uebernehmen + dezent visuell ausschmuecken.
- Keine erfundenen Details die mit dem Input kollidieren.
"""


def _load_image_prompt_skill() -> str:
    """Laedt den image-prompt SKILL.md + haengt One-Shot-Tail an. Cached nicht —
    Skill kann live editiert werden ohne MC-Restart."""
    try:
        return _IMAGE_PROMPT_SKILL_PATH.read_text(encoding="utf-8") + _IMAGE_PROMPT_ONESHOT_TAIL
    except Exception as e:
        print(f"[ImagePromptEnhancer] skill load failed ({e}) — falling back to minimal prompt")
        return (
            "Du bist ein Experte fuer Nano-Banana / Gemini Image-Prompts. "
            "Schreib den User-Input in einen detaillierten englischen Image-Prompt um "
            "(60-200 Woerter, Subject + Composition + Lighting + Style + Mood). "
            "Antworte NUR mit dem fertigen Prompt, kein Drumherum."
        )


async def _enhance_image_prompt(user_prompt: str, style_hint: str | None = None) -> str | None:
    """Schreibt einen User-Prompt in einen image-optimierten Prompt um.
    Nutzt den ~/Jarvis/skills/image-prompt/SKILL.md Skill (Ken Kousen, basiert
    auf Googles offizieller Nano-Banana-Doku) im One-Shot-Modus.
    Provider: claude-cli (Sonnet), kein MCP, max-turns=1. Returns None bei Fehler."""
    from llm_client import invoke_llm, Message, TextDelta, ErrorEvent, MessageStop

    sys_prompt = _load_image_prompt_skill()
    if style_hint:
        sys_prompt += f"\n\nAdditional style hint from caller: {style_hint}"

    out = ""
    err = None
    try:
        async for ev in invoke_llm(
            provider="claude-cli",
            model="sonnet",
            messages=[Message(role="user", content=user_prompt)],
            system=sys_prompt,
            max_turns=2,  # 1 reicht eigentlich, 2 als Sicherheit falls erstes Token tool_use ist
            dangerous_skip_permissions=True,
            system_replace=True,  # Skill ersetzt Claude-Code-Default fuer reinen Prompt-Rewriter-Modus
            # Alle Tools verbieten — Modell soll DIREKT umschreiben, nicht den Skill als Tool callen
            disallowed_tools="Skill,Task,Bash,Read,Write,Edit,Grep,Glob,WebFetch,WebSearch,TodoWrite,Agent",
        ):
            if isinstance(ev, TextDelta):
                out += ev.text
            elif isinstance(ev, ErrorEvent):
                err = ev.message
                break
            elif isinstance(ev, MessageStop):
                break
    except Exception as e:
        err = str(e)

    if err or not out.strip():
        print(f"[ImagePromptEnhancer] no output (err={err})")
        return None
    cleaned = out.strip().strip('"').strip("'").strip()
    # Sicherheitsnetz: Falls Modell trotz One-Shot-Tail einen "PROMPT:"-Block produziert,
    # extrahiere den Prompt-Text dazwischen.
    import re as _re
    m = _re.search(r"PROMPT:\s*\n?[─\-]+\s*\n([\s\S]+?)\n[─\-]+", cleaned)
    if m:
        cleaned = m.group(1).strip()
    # Falls Modell trotzdem mehrere Absaetze produziert (Rationale unten), nimm ersten Block
    if "\n\n" in cleaned:
        cleaned = cleaned.split("\n\n", 1)[0].strip()
    return cleaned[:1500]


async def api_image_generate(request):
    """POST /api/image/generate
    Body: {provider, prompt, model?, size?, quality?, style?, session_id?}

    Generiert ein Bild via OpenAI oder Gemini und speichert es in UPLOADS.
    Bei session_id: broadcastet `chat.image_generated`-Event an Subscriber damit
    das Bild sofort im Chat inline erscheint.
    Response: {ok, path, url, provider, model, prompt, size, metadata}
    """
    from image_client import generate_image, ImageGenError
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    provider = (body.get("provider") or "gemini").strip().lower()
    user_prompt = (body.get("prompt") or "").strip()
    if not user_prompt:
        return web.json_response({"ok": False, "error": "prompt erforderlich"}, status=400)

    model = (body.get("model") or "").strip() or None
    size = (body.get("size") or "1024x1024").strip()
    quality = (body.get("quality") or "auto").strip()
    style = (body.get("style") or "").strip() or None
    session_id = (body.get("session_id") or "").strip()
    enhance = bool(body.get("enhance", True))

    # Prompt-Enhancer: schreibt User-Eingabe in einen detaillierten Image-Gen-Prompt um.
    # Ein-Shot-Call mit Claude Sonnet (claude-cli, ohne MCP). Wenn's failt: original-Prompt nehmen.
    enhanced_prompt = None
    if enhance and len(user_prompt) > 4:
        try:
            enhanced_prompt = await _enhance_image_prompt(user_prompt, style)
        except Exception as e:
            print(f"[ImagePromptEnhancer] failed (using original): {e}")
            enhanced_prompt = None

    final_prompt = enhanced_prompt or user_prompt

    # Avatar-Auto-Detection: Wenn der Prompt Personen-Namen aus dem Avatare-Tab
    # enthaelt, packen wir deren References als visual context ans Modell.
    # Nur Gemini supportet das aktuell nativ.
    reference_paths: list[Path] = []
    references_used: list[dict] = []
    if provider == "gemini":
        try:
            import library as _lib
            avatars = _lib.list_avatars(AVATARS_DIR)
            prompt_lower = (user_prompt + " " + (final_prompt or "")).lower()
            import re as _re
            words = set(_re.findall(r"[a-zäöüß]+", prompt_lower))
            for av in avatars:
                if av["name"].lower() not in words: continue
                if av["ref_count"] == 0: continue
                detail = _lib.get_avatar(AVATARS_DIR, av["name"])
                # max 4 references pro Person, max 6 total (Modell-Limits)
                taken = []
                for ref in detail["references"][:4]:
                    if len(reference_paths) >= 6: break
                    p = AVATARS_DIR / av["name"] / "references" / ref["filename"]
                    if p.is_file():
                        reference_paths.append(p)
                        taken.append(ref["filename"])
                if taken:
                    references_used.append({"name": av["name"], "files": taken})
            if references_used:
                # Anchor-Text differenziert nach Modell:
                # Verifiziert 2026-04-25: Nano Banana Pro filtert (IMAGE_OTHER) wenn der
                # Anchor PERSONENNAMEN + Phrasen wie "reference image"/"real person"/
                # "facial features and identity" kombiniert (Deepfake-Vorsicht). Neutrale
                # Style-Sprache OHNE Personenname klappt. Flash-Varianten (Nano Banana /
                # Nano Banana 2) sind toleranter, da geht der staerkere Identity-Anchor.
                is_pro = (model or "") == "gemini-3-pro-image-preview"
                if is_pro:
                    # Ein einziger neutraler Satz — kein Name, keine "real person"-Phrase
                    anchor = "Use the provided reference images as visual style guide for the subject's appearance."
                else:
                    names = ", ".join(u["name"] for u in references_used)
                    anchor = f"The reference image(s) of {names} show the same person — keep their appearance and features consistent with these references."
                final_prompt = anchor + "\n\n" + final_prompt
                print(f"[ImageRefs] using {len(reference_paths)} refs for: {[u['name'] for u in references_used]} (pro_mode={is_pro})")
        except Exception as e:
            print(f"[ImageRefs] auto-detect failed (continuing without refs): {e}")

    try:
        img = await generate_image(
            provider=provider, prompt=final_prompt,
            uploads_dir=UPLOADS_DIR,
            model=model, size=size, quality=quality, style=style,
            reference_images=reference_paths or None,
        )
    except ImageGenError as e:
        return web.json_response(
            {"ok": False, "error": e.message, "code": e.code, "provider": provider},
            status=200,  # Client-sichtbarer Fehler, kein 5xx
        )
    except Exception as e:
        return web.json_response(
            {"ok": False, "error": f"{type(e).__name__}: {e}", "provider": provider},
            status=500,
        )

    # Frisch generiertes Bild direkt in die Library (_inbox) verschieben.
    # Damit ist es persistent + im Image-Lab unter "Library" auffindbar,
    # statt nur als Upload-Cache (max 24 in localStorage) wegzuflattern.
    library_url = img.url
    library_path_rel = str(img.path.relative_to(UPLOADS_DIR)).replace("\\", "/")
    library_folder = ""
    try:
        inbox = IMAGES_LIBRARY_DIR / "_inbox"
        inbox.mkdir(parents=True, exist_ok=True)
        target = inbox / img.path.name
        if target.exists():
            stem, suf = target.stem, target.suffix
            i = 2
            while (inbox / f"{stem}_{i}{suf}").exists():
                i += 1
            target = inbox / f"{stem}_{i}{suf}"
        img.path.rename(target)
        library_url = f"/library/_inbox/{target.name}"
        library_path_rel = f"_inbox/{target.name}"
        library_folder = "_inbox"
        img.path = target
        img.url = library_url
    except Exception as _e:
        print(f"[Library] Move to library failed for {img.path}: {_e}")

    # Sidecar-Meta-File: User-Prompt + Enhanced-Prompt + Provider/Model/Size + ggf References
    try:
        import library as _lib
        _lib.write_image_meta(img.path, {
            "user_prompt": user_prompt,
            "enhanced_prompt": enhanced_prompt,
            "final_prompt": final_prompt,
            "provider": img.provider,
            "model": img.model,
            "size": img.size,
            "style": style,
            "references_used": references_used,
            "created_at": datetime.now().isoformat(),
        })
    except Exception as _e:
        print(f"[Library] Sidecar write failed: {_e}")

    payload = {
        "ok": True,
        "provider": img.provider,
        "model": img.model,
        "prompt": user_prompt,
        "enhanced_prompt": enhanced_prompt,
        "final_prompt": final_prompt,
        "references_used": references_used,
        "size": img.size,
        "url": library_url,
        "path": library_path_rel,
        "library_folder": library_folder,
        "library_filename": img.path.name,
        "mime_type": img.mime_type,
        "metadata": img.metadata,
    }

    # Wenn Session mitgegeben: Bild sofort in den Chat posten
    if session_id and sessions.sessions.get(session_id):
        session = sessions.sessions[session_id]
        # Als Assistant-Message mit Bild-Markdown einfügen
        image_md = f"![{img.prompt[:80]}]({img.url})\n\n*Generiert via {img.provider} ({img.model})*"
        session.history.append({
            "role": "assistant",
            "content": image_md,
            "_image": True,
            "_image_url": img.url,
        })
        sessions.persist()
        await _broadcast_stream_event(session_id, {
            "type": "chat.image_generated",
            "session_id": session_id,
            "url": img.url,
            "prompt": img.prompt,
            "provider": img.provider,
            "model": img.model,
        })

    return web.json_response(payload)


# ============================================================================
# Image Library + Avatare — siehe library.py
# ============================================================================

async def api_library_list(request):
    """GET /api/library?folder=<name>&group=day
    Ohne folder → folders + flat list. Mit folder → Bilder dieses Folders.
    group=day → zusaetzlich grouped_by_day.
    """
    import library as _lib
    folder = request.query.get("folder") or None
    group = request.query.get("group", "")
    try:
        folders = _lib.list_folders(IMAGES_LIBRARY_DIR)
        images = _lib.list_images(IMAGES_LIBRARY_DIR, folder)
        out = {"ok": True, "folders": folders, "images": images}
        if group == "day":
            out["grouped"] = _lib.list_grouped_by_day(IMAGES_LIBRARY_DIR, folder)
        return web.json_response(out)
    except ValueError as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_library_folder_create(request):
    import library as _lib
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
    try:
        folder = _lib.create_folder(IMAGES_LIBRARY_DIR, (body.get("name") or "").strip())
        return web.json_response({"ok": True, "folder": folder})
    except (ValueError, FileExistsError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_library_folder_delete(request):
    import library as _lib
    try:
        body = await request.json()
    except Exception:
        body = {}
    name = (body.get("name") or request.query.get("name") or "").strip()
    force = bool(body.get("force") or request.query.get("force") in ("1", "true"))
    try:
        n = _lib.delete_folder(IMAGES_LIBRARY_DIR, name, force=force)
        return web.json_response({"ok": True, "deleted_images": n})
    except (ValueError, FileNotFoundError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_library_move(request):
    import library as _lib
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
    try:
        new_url = _lib.move_image(
            IMAGES_LIBRARY_DIR,
            (body.get("src_folder") or "").strip(),
            (body.get("filename") or "").strip(),
            (body.get("dest_folder") or "").strip(),
        )
        return web.json_response({"ok": True, "url": new_url})
    except (ValueError, FileNotFoundError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_library_image_delete(request):
    import library as _lib
    try:
        body = await request.json()
    except Exception:
        body = {}
    folder = (body.get("folder") or request.query.get("folder") or "").strip()
    filename = (body.get("filename") or request.query.get("filename") or "").strip()
    try:
        _lib.delete_image(IMAGES_LIBRARY_DIR, folder, filename)
        return web.json_response({"ok": True})
    except (ValueError, FileNotFoundError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_avatars_list(request):
    import library as _lib
    return web.json_response({"ok": True, "avatars": _lib.list_avatars(AVATARS_DIR)})


async def api_avatars_get(request):
    import library as _lib
    name = request.match_info.get("name", "")
    try:
        return web.json_response({"ok": True, "avatar": _lib.get_avatar(AVATARS_DIR, name)})
    except (ValueError, FileNotFoundError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=404)


async def api_avatars_create(request):
    import library as _lib
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
    try:
        avatar = _lib.create_avatar(
            AVATARS_DIR,
            (body.get("name") or "").strip(),
            description=(body.get("description") or "").strip(),
            style_hints=(body.get("style_hints") or "").strip(),
        )
        return web.json_response({"ok": True, "avatar": avatar})
    except (ValueError, FileExistsError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_avatars_update(request):
    import library as _lib
    name = request.match_info.get("name", "")
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
    try:
        avatar = _lib.update_avatar(
            AVATARS_DIR, name,
            description=body.get("description"),
            style_hints=body.get("style_hints"),
        )
        return web.json_response({"ok": True, "avatar": avatar})
    except (ValueError, FileNotFoundError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_avatars_delete(request):
    import library as _lib
    name = request.match_info.get("name", "")
    try:
        _lib.delete_avatar(AVATARS_DIR, name)
        return web.json_response({"ok": True})
    except (ValueError, FileNotFoundError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_avatars_upload_reference(request):
    """multipart/form-data: file=<image>"""
    import library as _lib
    name = request.match_info.get("name", "")
    try:
        reader = await request.multipart()
        field = await reader.next()
        if not field or field.name != "file":
            return web.json_response({"ok": False, "error": "field 'file' fehlt"}, status=400)
        filename = field.filename or "ref.png"
        content = bytearray()
        while True:
            chunk = await field.read_chunk()
            if not chunk:
                break
            content.extend(chunk)
            if len(content) > 20 * 1024 * 1024:
                return web.json_response({"ok": False, "error": "Datei zu gross (max 20MB)"}, status=400)
        url = _lib.add_reference(AVATARS_DIR, name, filename, bytes(content))
        return web.json_response({"ok": True, "url": url})
    except (ValueError, FileNotFoundError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_avatars_delete_reference(request):
    import library as _lib
    name = request.match_info.get("name", "")
    try:
        body = await request.json()
    except Exception:
        body = {}
    filename = (body.get("filename") or request.query.get("filename") or "").strip()
    try:
        _lib.delete_reference(AVATARS_DIR, name, filename)
        return web.json_response({"ok": True})
    except (ValueError, FileNotFoundError) as e:
        return web.json_response({"ok": False, "error": str(e)}, status=400)


async def api_medical_reason(request):
    """POST /api/medical/reason
    Body: {query, provider?, model?, include_sources?: list[str]}
      include_sources filter — default ["amboss", "viamedici", "rag"]

    Sammelt parallel Kontext aus Amboss + ViaMedici + Vault-RAG, kombiniert zu
    einem strukturierten System-Prompt, und lässt den gewählten LLM-Provider
    (default claude-cli) eine synthesierte medizinische Antwort formulieren.
    Non-streaming — gibt vollständige Response nach Abschluss zurück.

    Response: {ok, answer, sources: {amboss, viamedici, rag}, provider, model, latency_ms}
    """
    import time as _time
    from llm_client import invoke_llm, Message, TextDelta, ErrorEvent, MessageStop, UsageEvent

    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    query = (body.get("query") or "").strip()
    if not query or len(query) < 5:
        return web.json_response({"ok": False, "error": "query zu kurz"}, status=400)

    provider = (body.get("provider") or "claude-cli").strip()
    model = (body.get("model") or "").strip()
    include = body.get("include_sources") or ["amboss", "viamedici", "rag"]
    if isinstance(include, str):
        include = [s.strip() for s in include.split(",") if s.strip()]

    t0 = _time.time()

    # Parallel fetchen
    tasks: dict[str, asyncio.Future] = {}
    if "rag" in include:
        tasks["rag"] = asyncio.create_task(_build_rag_context(query, n_results=6, min_score=0.3))
    if "amboss" in include:
        tasks["amboss"] = asyncio.create_task(_fetch_amboss_medical(query, UPLOADS_DIR))
    if "viamedici" in include:
        tasks["viamedici"] = asyncio.create_task(_fetch_viamedici_medical(query, UPLOADS_DIR))

    rag_result = amboss_result = viamedici_result = None
    if "rag" in tasks:
        try:
            rag_result = await tasks["rag"]
        except Exception as e:
            print(f"[MedAgent] rag failed: {e}")
    if "amboss" in tasks:
        try:
            amboss_result = await tasks["amboss"]
        except Exception as e:
            print(f"[MedAgent] amboss failed: {e}")
    if "viamedici" in tasks:
        try:
            viamedici_result = await tasks["viamedici"]
        except Exception as e:
            print(f"[MedAgent] viamedici failed: {e}")

    # Strukturierter Context-Block
    ctx_parts: list[str] = []
    if rag_result and rag_result[0]:
        ctx_parts.append(rag_result[0])  # Schon formatierter Block
    if amboss_result:
        amboss_block = "\n\n--- AMBOSS ---\n" + "\n\n".join(
            f"### {a['title']}  (eid: {a['eid']})\n"
            f"{a.get('abstract', '')}\n\n{a.get('text', '')}"
            for a in amboss_result
        )
        ctx_parts.append(amboss_block)
    if viamedici_result:
        vm_block = "\n\n--- THIEME/VIAMEDICI ---\n" + "\n\n".join(
            f"### {v['title']}  (assetId: {v['assetId']})\n{v['text']}"
            for v in viamedici_result
        )
        ctx_parts.append(vm_block)

    gather_ms = int((_time.time() - t0) * 1000)

    if not ctx_parts:
        return web.json_response({
            "ok": False,
            "error": "Keine Quelle lieferte Kontext (Cookies abgelaufen? Frage zu spezifisch?)",
            "sources": {"amboss": [], "viamedici": [], "rag_chunks": 0},
            "gather_ms": gather_ms,
        }, status=200)

    combined_context = "\n".join(ctx_parts)
    system_prompt = (
        "Du bist ein medizinisches Reasoning-System für den Medizinstudenten Robin. "
        "Antworte präzise auf Deutsch, strukturiert, klinisch relevant. "
        "Nutze AUSSCHLIESSLICH die bereitgestellten Quellen unten (Amboss, Thieme/"
        "ViaMedici, Robins Vault). Wenn eine Aussage NICHT durch die Quellen "
        "gedeckt ist, sage es explizit ('nicht in den Quellen'). Markiere pro "
        "Aussage die Quelle (Amboss / Thieme / Vault)."
        + combined_context
    )

    # LLM-Call — non-streaming, sammeln bis MessageStop
    messages = [Message(role="user", content=query)]
    answer_text = ""
    usage: dict = {}
    error_msg: str | None = None
    LLM_TIMEOUT = 120.0 # 2 Minuten für die Antwort-Synthese

    try:
        # Wir nutzen wait_for um den gesamten Generator-Lauf zu begrenzen
        async def _run_llm():
            nonlocal answer_text, usage, error_msg
            async for ev in invoke_llm(
                provider=provider,
                model=model,
                messages=messages,
                system=system_prompt,
                max_tokens=1500,
            ):
                if isinstance(ev, TextDelta):
                    answer_text += ev.text
                elif isinstance(ev, UsageEvent):
                    usage = {
                        "input_tokens": ev.input_tokens + ev.cache_read_tokens,
                        "output_tokens": ev.output_tokens,
                    }
                elif isinstance(ev, ErrorEvent):
                    error_msg = f"{ev.code}: {ev.message}"
                elif isinstance(ev, MessageStop):
                    break
        
        await asyncio.wait_for(_run_llm(), timeout=LLM_TIMEOUT)

    except asyncio.TimeoutError:
        error_msg = f"LLM Timeout ({LLM_TIMEOUT}s)"
        print(f"[MedAgent] LLM synthesis timed out after {LLM_TIMEOUT}s")
    except Exception as e:
        error_msg = f"{type(e).__name__}: {e}"
        print(f"[MedAgent] LLM synthesis failed: {e}")

    total_ms = int((_time.time() - t0) * 1000)

    if error_msg and not answer_text:
        return web.json_response({
            "ok": False,
            "error": error_msg,
            "provider": provider,
            "latency_ms": total_ms,
        }, status=200)

    return web.json_response({
        "ok": True,
        "query": query,
        "answer": answer_text,
        "sources": {
            "amboss": amboss_result or [],
            "viamedici": viamedici_result or [],
            "rag_chunks": (rag_result[1] if rag_result else 0),
        },
        "provider": provider,
        "model": model or "default",
        "usage": usage,
        "gather_ms": gather_ms,
        "latency_ms": total_ms,
    })


async def api_vault_bundles(request):
    """GET /api/vault/bundles — kuratierte Themen-Bundles + File-Counts
    (für Frontend-Context-Picker).
    """
    result = []
    for name, paths in VAULT_CONTEXT_BUNDLES.items():
        file_count = 0
        sample_files: list[str] = []
        for rel in paths:
            abs_path = _VAULT_ROOT / rel
            if not abs_path.exists():
                continue
            if abs_path.is_file() and abs_path.suffix.lower() == ".md":
                file_count += 1
                sample_files.append(rel)
            elif abs_path.is_dir():
                for md in abs_path.rglob("*.md"):
                    file_count += 1
                    if len(sample_files) < 5:
                        sample_files.append(
                            str(md.relative_to(_VAULT_ROOT)).replace("\\", "/")
                        )
        result.append({
            "name": name,
            "paths": paths,
            "file_count": file_count,
            "sample_files": sample_files[:5],
        })
    return web.json_response({"bundles": result})


async def api_chat_session_context(request):
    """POST /api/chat/session/{session_id}/context
    Body: {"bundles": ["medizin", "dissertation"]}
    Setzt Context-Bundles pro Session. Leere Liste deaktiviert Global-Context.
    """
    sid = request.match_info.get("session_id", "")
    session = sessions.sessions.get(sid)
    if not session:
        return web.json_response({"ok": False, "error": "Session nicht gefunden"}, status=404)
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    bundles_raw = body.get("bundles") or []
    if not isinstance(bundles_raw, list):
        return web.json_response({"ok": False, "error": "bundles muss Liste sein"}, status=400)

    valid = [b for b in bundles_raw if isinstance(b, str) and b.lower() in VAULT_CONTEXT_BUNDLES]
    session.context_bundles = valid
    sessions.persist()

    if valid:
        await _broadcast_stream_event(session.session_id, {
            "type": "chat.system",
            "session_id": session.session_id,
            "message": f"📁 Vault-Kontext aktiviert: {', '.join(valid)}",
        })
    return web.json_response({
        "ok": True,
        "session_id": session.session_id,
        "context_bundles": session.context_bundles,
    })


async def api_llm_providers(request):
    """GET /api/llm/providers — Capability-Liste aller registrierten LLM-Provider.

    Frontend nutzt das um:
      - Den Provider-Dropdown dynamisch zu befüllen
      - Capability-spezifische UI-Hints zu zeigen (z.B. Resume-Indicator,
        Image-Badge, "streamt nicht fließend" bei codex-cli)
    """
    from llm_client import list_providers
    return web.json_response({"providers": list_providers()})


async def api_chat_session_provider(request):
    """POST /api/chat/session/{session_id}/provider — Provider/Model einer
    Session setzen oder wechseln.

    Body: {"provider": str, "model": str|null}
    Response: {"ok": bool, "provider_switched": bool, "message": str|null, ...}

    Triggert _set_session_provider — Provider-Wechsel → provider_thread_id
    wird verworfen und eine Systemmeldung in history angehängt.
    """
    session_id = request.match_info.get("session_id", "")
    session = sessions.sessions.get(session_id)
    if not session:
        return web.json_response({"ok": False, "error": "Session nicht gefunden"}, status=404)

    try:
        body = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)

    provider = (body.get("provider") or "").strip()
    model = body.get("model")
    if model is not None:
        model = str(model).strip()
    if not provider:
        return web.json_response({"ok": False, "error": "provider erforderlich"}, status=400)

    # jarvis-main + jarvis-heartbeat sind die Orchestrator-Sessions und auf
    # claude-cli fixiert — keine Provider-Wechsel erlaubt.
    if session_id in ("jarvis-main", "jarvis-heartbeat") and provider != "claude-cli":
        return web.json_response({
            "ok": False,
            "error": "Jarvis-Sessions laufen fest auf Claude (CLI). Provider-Wechsel hier nicht erlaubt.",
        }, status=400)

    try:
        info = _set_session_provider(session, provider, model)
    except ValueError as ve:
        return web.json_response({"ok": False, "error": str(ve)}, status=400)

    # Wenn ein System-Message in history kam, via WS broadcasten (Client sieht's gleich)
    if info.get("provider_switched") and info.get("message"):
        await _broadcast_stream_event(session.session_id, {
            "type": "chat.system",
            "session_id": session.session_id,
            "message": info["message"],
        })

    return web.json_response({
        "ok": True,
        "session_id": session.session_id,
        "provider": session.provider,
        "model": session.model,
        "provider_thread_id": session.provider_thread_id,
        "provider_switched": info.get("provider_switched", False),
        "message": info.get("message"),
    })


async def api_chat_session_compress(request):
    """POST /api/chat/session/{session_id}/compress

    Komprimiert die Session-History in eine kompakte Zusammenfassung und startet
    eine frische CLI-Session mit dem Summary als Kontext-Seed.

    Nur möglich wenn die Session gerade nicht streamt und >= 10 Messages hat.
    """
    session_id = request.match_info.get("session_id", "")
    session = sessions.sessions.get(session_id)
    if not session:
        return web.json_response({"ok": False, "error": "Session nicht gefunden"}, status=404)
    if session._stream_active or session.lock.locked():
        return web.json_response({"ok": False, "error": "Session streamt gerade"}, status=409)

    history = session.history or []
    user_assistant = [m for m in history if m.get("role") in ("user", "assistant")]
    if len(user_assistant) < 10:
        return web.json_response({"ok": False, "error": "Zu wenige Messages für Komprimierung"}, status=400)

    # Build conversation text (last 100 messages, max 800 chars each)
    lines = []
    for m in user_assistant[-100:]:
        role = "Robin" if m["role"] == "user" else "Jarvis"
        content = (m.get("content") or "")[:800]
        lines.append(f"[{role}] {content}")
    conv_text = "\n\n".join(lines)

    # One-shot summary via _llm_oneshot
    system = (
        "Du bist Jarvis, Robins persönlicher AI-Assistent. "
        "Erstelle eine kompakte, strukturierte Zusammenfassung des folgenden Gesprächs. "
        "Halte alle wichtigen Fakten, Entscheidungen, offenen Aufgaben und Kontext. "
        "Format: Bullet-Punkte, maximal 1500 Zeichen. Deutsch."
    )
    user_prompt = f"Fasse dieses Gespräch kompakt zusammen:\n\n{conv_text[:12000]}"

    print(f"[Compress] Session {session_id[:8]}: {len(user_assistant)} msgs → summary start")
    summary, _, err = await _llm_oneshot(
        provider=session.provider or "claude-cli",
        model=session.model or "",
        user_message=user_prompt,
        system=system,
        max_tokens=1500,
    )
    if err and not summary:
        return web.json_response({"ok": False, "error": f"Summary-Fehler: {err}"}, status=500)

    # Reset session: clear history → summary als Seed, neue CLI-Session
    old_count = len(user_assistant)
    seed = (
        f"[KONTEXT-KOMPRIMIERUNG — {old_count} Messages zusammengefasst]\n\n"
        f"{summary}"
    )
    session.history = [{"role": "assistant", "content": seed}]
    session.provider_thread_id = None
    session.cli_session_id = None
    session.message_count = 1
    sessions.persist()

    # Broadcast system event so frontend reloads chat
    await _broadcast_stream_event(session.session_id, {
        "type": "chat.compressed",
        "session_id": session.session_id,
        "old_count": old_count,
        "summary": seed,
    })

    print(f"[Compress] Session {session_id[:8]}: done — {old_count} msgs → {len(seed)} chars summary")
    return web.json_response({
        "ok": True,
        "old_message_count": old_count,
        "summary_length": len(summary),
        "seed_preview": seed[:200],
    })


async def _run_chat(ws, session, message, attachments=None):
    """Provider-agnostischer Chat-Stream. Nutzt llm_client.invoke_llm() statt
    direktem Subprocess. Funktioniert mit allen registrierten Providern
    (claude-cli, claude-api, codex-cli, openai).

    Semantik gegenüber dem legacy _run_claude bewahrt:
      - chat.start / chat.delta / chat.tool_use / chat.tool_result / chat.ping
        / chat.done / chat.error als WS-Events
      - Partial-Persist alle 5 s
      - Ping alle 30 s wenn kein Text fließt (Frontend-Stuck-Timer)
      - Subscriber-Check: wenn alle WS weg → Task canceln (Provider beendet
        seinen Subprocess über CancelledError-Pfad)
      - Stale-resume-Retry: wenn claude-cli mit --resume leer zurückkommt,
        provider_thread_id droppen + einmal ohne Resume retry
      - Genau EIN Abschluss-Event (chat.done ODER chat.error), niemals beides
        und niemals keins — Frontend sieht nie einen hängenden Turn
    """
    import time as _time
    from llm_client import (
        invoke_llm, Message, TextDelta, ToolUseStart, ToolUseEnd,
        ToolResultEvent, MessageStop, UsageEvent, MetadataEvent, ErrorEvent, get_provider_info,
    )

    # Provider-Info für Gating (supports_images, streams_incremental_text, ...)
    try:
        prov_info = get_provider_info(session.provider or "claude-cli")
    except Exception as e:
        await _broadcast_stream_event(session.session_id, {
            "type": "chat.error",
            "session_id": session.session_id,
            "error": f"Unbekannter Provider {session.provider!r}: {e}",
        })
        return

    _stream_start(session.session_id)
    if ws is not None:
        _stream_subscribers.setdefault(session.session_id, set()).add(ws)
    else:
        _stream_subscribers.setdefault(session.session_id, set())
    _activity_set("chat", active=True)
    session._stream_active = True
    session.last_heartbeat = datetime.now().isoformat()
    sessions.persist()

    # State für genau-einmaligen Abschluss
    _done_sent = False

    async def _send_done(payload: dict):
        """Genau EIN Abschluss-Event. Subsequent calls no-op."""
        nonlocal _done_sent
        if _done_sent:
            return
        _done_sent = True
        await _broadcast_stream_event(session.session_id, payload)
        # chat.done und chat.error gehen zusaetzlich an alle Clients —
        # so kann ein nach Restart frisch verbundener Browser den Stuck-Spinner
        # einer Session aufloesen, fuer die er noch keinen Subscriber hatte.
        if payload.get("type") in ("chat.done", "chat.error"):
            await _broadcast_all_ws(payload)
        _activity_set("chat", active=False, pending_delta=1)
        _stream_end(session.session_id)
        asyncio.get_event_loop().call_later(30, _stream_cleanup, session.session_id)

    async with session.lock:
        session._lock_acquired_at = _time.time()
        try:
            # Attachment-Notes bauen (Image-Gate nach supports_images)
            attachment_notes, _image_blocks = _build_attachment_notes(
                attachments, bool(prov_info.get("supports_images", False)),
            )

            # Messages-Liste
            messages = _history_to_messages(session, message, attachment_notes)

            # System-Prompt — MC-spezifisch, egal welcher Provider
            system = _build_system_prompt()

            # Shared RAG (Geminis Stufe 1): semantisch ähnliche Chunks aus Vault +
            # Diss als zusätzlichen System-Context. Funktioniert provider-agnostisch
            # weil es im System-Prompt landet — kein Tool-Wiring nötig.
            try:
                rag_ctx, rag_n = await _build_rag_context(message)
                if rag_ctx:
                    system = system + rag_ctx
                    print(f"[RAG] injected {rag_n} chunks for session {session.session_id[:8]}")
                    await _broadcast_stream_event(session.session_id, {
                        "type": "chat.rag_info",
                        "session_id": session.session_id,
                        "chunks": rag_n,
                    })
            except Exception as e:
                print(f"[RAG] context skipped (error: {e})")

            # Shared Scratchpad (Stufe 3): wenn der Peer was reingeschrieben hat,
            # bekommt der aktuelle Provider es als System-Context zu sehen.
            if getattr(session, "shared_notes", None):
                notes_block = "\n\n--- SHARED SCRATCHPAD (Peer-Notes) ---\n"
                for k, v in session.shared_notes.items():
                    notes_block += f"### {k}\n{v[:2000]}\n\n"
                system = system + notes_block

            # Global Context (Geminis #2): kuratierte Vault-Bereiche als Volltext
            # — nur wenn die Session Bundles aktiviert hat. Ergänzt RAG um "atme
            # das ganze Dissertation-Kapitel" statt nur Score-getroffene Chunks.
            if session.context_bundles:
                try:
                    gc_ctx, gc_n, gc_paths = await asyncio.to_thread(
                        _load_vault_context, session.context_bundles,
                    )
                    if gc_ctx:
                        system = system + gc_ctx
                        print(f"[Vault] injected {gc_n} files from bundles={session.context_bundles}")
                        await _broadcast_stream_event(session.session_id, {
                            "type": "chat.vault_context",
                            "session_id": session.session_id,
                            "files": gc_n,
                            "bundles": list(session.context_bundles),
                            "paths": gc_paths[:10],
                        })
                except Exception as e:
                    print(f"[Vault] context skipped (error: {e})")

            # Provider-kwargs mit Resume-Token-Callback
            def _on_resume_token(new_token: str):
                """Provider liefert neue Thread/Session-ID → in Session persistieren."""
                session.provider_thread_id = new_token
                if session.provider == "claude-cli":
                    session.cli_session_id = new_token  # Legacy-Sync
                sessions.persist()
                print(f"[Chat] {session.provider} thread_id: {new_token[:12]}")

            prov_kwargs = _build_provider_kwargs(session, on_resume_token=_on_resume_token)

            # Persist User-Message SOFORT (überlebt Reload während Streaming)
            session.history.append({"role": "user", "content": message})
            sessions.persist()

            # chat.start → session subs + alle clients (Sidebar-Indikator)
            from llm_client import _humanize_model_name
            start_event = {"type": "chat.start", "session_id": session.session_id,
                           "provider": session.provider, 
                           "model": _humanize_model_name(session.provider, session.model)}
            await _broadcast_stream_event(session.session_id, start_event)
            await _broadcast_all_ws(start_event)

            used_resume = bool(session.provider_thread_id) and session.provider == "claude-cli"
            used_model = session.model or "(Auto)"
            print(f"[Chat] provider={session.provider} model={used_model} "
                  f"resume={session.provider_thread_id[:12] if session.provider_thread_id else 'new'}")
            t_start = _time.time()

            full_text = ""
            last_tool = ""
            last_tool_id = ""
            usage_info: dict = {}
            stop_reason = "end_turn"
            error_seen: str | None = None
            error_code: str = ""
            _last_partial_persist = _time.time()
            _stream_started_at = _time.time()
            _last_ping = _time.time()

            async def _ping_loop():
                """Sendet alle ~30 s einen chat.ping wenn nichts fließt. Läuft parallel zum Stream."""
                nonlocal _last_ping
                while True:
                    await asyncio.sleep(28)
                    if _done_sent:
                        break
                    elapsed = _time.time() - _stream_started_at
                    if _time.time() - _last_ping >= 28:
                        _last_ping = _time.time()
                        session.last_heartbeat = datetime.now().isoformat()
                        await _broadcast_stream_event(session.session_id, {
                            "type": "chat.ping",
                            "session_id": session.session_id,
                            "elapsed": int(elapsed),
                        })

            async def _subscriber_watchdog(stream_task: asyncio.Task):
                """Bricht den Stream wenn alle WS-Subscriber weg sind."""
                while not stream_task.done():
                    await asyncio.sleep(5)
                    subs = _stream_subscribers.get(session.session_id, set())
                    alive_subs = {s for s in subs if not s.closed}
                    if not alive_subs and not _ws_clients:
                        print(f"[Chat] All subs gone for {session.session_id[:8]} — cancelling")
                        stream_task.cancel()
                        return

            # Background-Tasks für Ping + Watchdog
            ping_task = asyncio.create_task(_ping_loop())
            current_stream_task = asyncio.current_task()
            watchdog_task = asyncio.create_task(_subscriber_watchdog(current_stream_task))

            try:
                async for ev in invoke_llm(
                    provider=session.provider,
                    model=session.model or "",
                    messages=messages,
                    system=system,
                    **prov_kwargs,
                ):
                    # Alle Events resetzen den Stuck-Timer
                    _last_ping = _time.time()

                    if isinstance(ev, TextDelta):
                        if ev.text:
                            full_text += ev.text
                            await _broadcast_stream_event(session.session_id, {
                                "type": "chat.delta",
                                "session_id": session.session_id,
                                "text": ev.text,
                            })
                    elif isinstance(ev, ToolUseStart):
                        last_tool = ev.name
                        last_tool_id = ev.id
                        # chat.tool_use wird beim End mit finalem input gesendet
                    elif isinstance(ev, ToolUseEnd):
                        last_tool = ev.name
                        last_tool_id = ev.id
                        # Große String-Inputs truncaten fürs WS (Write-Bodies etc.)
                        safe_input: dict = {}
                        for k, v in (ev.input or {}).items():
                            if isinstance(v, str) and len(v) > 200:
                                safe_input[k] = v[:200] + "..."
                            else:
                                safe_input[k] = v
                        await _broadcast_stream_event(session.session_id, {
                            "type": "chat.tool_use",
                            "session_id": session.session_id,
                            "tool": last_tool,
                            "input": safe_input,
                        })
                    elif isinstance(ev, ToolResultEvent):
                        await _broadcast_stream_event(session.session_id, {
                            "type": "chat.tool_result",
                            "session_id": session.session_id,
                            "tool": ev.tool_name or last_tool,
                            "is_error": ev.is_error,
                        })
                    elif isinstance(ev, UsageEvent):
                        # Context = alles was das Modell in diesem Turn sieht.
                        # Bei Claude CLI mit Prompt-Cache liegt der grosse Block
                        # (System-Prompt + History) in cache_read_input_tokens —
                        # ohne den waere _total_used massiv unterschaetzt.
                        _total_used = (
                            ev.input_tokens
                            + ev.cache_creation_tokens
                            + ev.cache_read_tokens
                            + ev.output_tokens
                        )
                        _m = (used_model or "").lower()
                        if any(x in _m for x in ("gemini-1.5", "gemini-exp")):
                            _ctx_win = 1_000_000
                        elif "gemini" in _m:
                            _ctx_win = 1_048_576
                        else:
                            _ctx_win = 200_000
                        usage_info = {
                            "input_tokens": ev.input_tokens + ev.cache_creation_tokens + ev.cache_read_tokens,
                            "output_tokens": ev.output_tokens,
                            "cache_creation_tokens": ev.cache_creation_tokens,
                            "cache_read_tokens": ev.cache_read_tokens,
                            "context_window": _ctx_win,
                            "context_used": _total_used,
                            "context_pct": round(_total_used / _ctx_win * 100, 1) if _total_used else 0,
                        }
                        print(f"[Context] in={ev.input_tokens} cc={ev.cache_creation_tokens} cr={ev.cache_read_tokens} out={ev.output_tokens} → {_total_used}/{_ctx_win} = {usage_info['context_pct']}%")
                        await _broadcast_stream_event(session.session_id, {
                            "type": "chat.usage",
                            "session_id": session.session_id,
                            **usage_info
                        })
                        # Broadcast context% to ALL clients — bypasses stream-subscriber filter.
                        # Ensures topbar always shows correct context even if WS wasn't subscribed.
                        await _broadcast_all_ws({
                            "type": "context.update",
                            "session_id": session.session_id,
                            "context_pct": usage_info["context_pct"],
                            "context_used": usage_info["context_used"],
                            "context_window": usage_info["context_window"],
                        })
                    elif isinstance(ev, MessageStop):
                        stop_reason = ev.stop_reason
                    elif isinstance(ev, MetadataEvent):
                        if ev.model:
                            used_model = ev.model
                        await _broadcast_stream_event(session.session_id, {
                            "type": "chat.metadata",
                            "session_id": session.session_id,
                            "model": ev.model,
                            "provider": session.provider,
                        })
                    elif isinstance(ev, ErrorEvent):
                        error_seen = ev.message
                        error_code = ev.code

                    # Partial-Persist alle 5 s
                    if full_text and (session._streaming_partial is None or
                                      (_time.time() - _last_partial_persist) > 5):
                        _last_partial_persist = _time.time()
                        session._streaming_partial = full_text
                        sessions.persist()
            finally:
                ping_task.cancel()
                watchdog_task.cancel()

            duration_ms = int((_time.time() - t_start) * 1000)

            # Stale-resume-Detection (claude-cli-spezifisch)
            if used_resume and full_text and session.provider == "claude-cli":
                _low = full_text.lower()
                if "session" in _low and "not found" in _low and len(full_text) < 500:
                    print("[Chat] Stale --resume detected, forcing retry")
                    full_text = ""
                    # Fehler vom Stale-Resume nicht als echten Fehler melden
                    error_seen = None

            # Fehler-Pfad: ErrorEvent gesehen ODER (leer & kein Resume-Retry möglich)
            if error_seen:
                session._stream_active = False
                sessions.persist()
                await _send_done({
                    "type": "chat.error",
                    "session_id": session.session_id,
                    "error": error_seen,
                    "error_code": error_code,
                    "provider": session.provider,
                })
                print(f"[Chat] error: {error_code} — {error_seen[:120]}")
                return

            # Leer-Pfad: Stale-resume-Retry für claude-cli
            if not full_text:
                if used_resume and session.provider == "claude-cli":
                    print("[Chat] --resume lieferte leer, retry ohne resume")
                    session.provider_thread_id = None
                    session.cli_session_id = None
                    sessions.persist()
                    if session.history and session.history[-1].get("role") == "user":
                        session.history.pop()
                        sessions.persist()
                    _stream_end(session.session_id)
                    asyncio.get_event_loop().call_later(30, _stream_cleanup, session.session_id)
                    # Wichtig: nicht über _send_done — das wäre stale Dup. Stattdessen Recursive-Call.
                    await _run_chat(ws, session, message, attachments)
                    return
                session._stream_active = False
                sessions.persist()
                await _send_done({
                    "type": "chat.error",
                    "session_id": session.session_id,
                    "error": f"{session.provider} hat keine Antwort geliefert.",
                    "provider": session.provider,
                })
                return

            # Erfolgspfad
            session.history.append({"role": "assistant", "content": full_text})
            session.message_count += 1
            session.last_message = datetime.now().isoformat()
            session.last_heartbeat = session.last_message
            session.process = None
            session._streaming_partial = None
            session._stream_active = False
            sessions.persist()

            done_payload: dict = {
                "type": "chat.done",
                "session_id": session.session_id,
                "full_text": full_text,
                "duration_ms": duration_ms,
                "model": used_model,
                "provider": session.provider,
                "stop_reason": stop_reason,
            }
            if usage_info:
                done_payload.update(usage_info)
            print(f"[Chat-done-debug] context_pct={done_payload.get('context_pct')} usage_info_keys={list(usage_info.keys()) if usage_info else 'EMPTY'}")
            await _send_done(done_payload)
            print(f"[Chat] chat.done: session={session.session_id[:8]} "
                  f"({len(full_text)} chars, {duration_ms}ms, provider={session.provider})")

            # Auto-classify wie im Legacy-Pfad
            _should_classify = (
                (not session.topic and session.message_count >= 2) or
                (session.message_count == 4 and not session.custom_title) or
                (not session.topic and session.message_count % 3 == 0)
            )
            if _should_classify:
                asyncio.create_task(_classify_session_headless(session))

        except asyncio.CancelledError:
            print(f"[Chat] Cancelled: session={session.session_id[:8]}")
            _ft = locals().get("full_text", "")
            if _ft and (not session.history or session.history[-1].get("content") != _ft + "\n\n_(abgebrochen)_"):
                session.history.append({"role": "assistant", "content": _ft + "\n\n_(abgebrochen)_"})
            session._stream_active = False
            sessions.persist()
            # Abschluss-Event auch bei Cancel, damit Frontend nicht hängt
            if not _done_sent:
                await _send_done({
                    "type": "chat.error",
                    "session_id": session.session_id,
                    "error": "Abgebrochen.",
                    "provider": session.provider,
                })
            raise
        except Exception as e:
            print(f"[Chat] Exception: {type(e).__name__}: {e}")
            import traceback
            traceback.print_exc()
            session._stream_active = False
            sessions.persist()
            if not _done_sent:
                await _send_done({
                    "type": "chat.error",
                    "session_id": session.session_id,
                    "error": f"{type(e).__name__}: {e}",
                    "provider": session.provider,
                })


async def _run_claude(ws, session, message, attachments=None):
    """Legacy-Wrapper: delegiert an _run_chat. Alle neuen WS-Handler rufen
    _run_chat direkt auf. Bleibt für Backward-Compat mit bestehenden Call-Sites.
    """
    return await _run_chat(ws, session, message, attachments)


# --- Legacy _run_claude-Body (vor Chat-Refactor 2026-04-23) -------------------
# Wird nicht mehr gerufen, aber als Referenz behalten. Nach 1-2 Wochen stabiler
# Nutzung des provider-agnostischen Pfads löschen.
async def _run_claude_legacy(ws, session, message, attachments=None):
    """[DEPRECATED] Alter Claude-CLI-direkter Pfad. Nicht mehr aktiv."""
    import time as _time

    _stream_start(session.session_id)
    if ws is not None:  # ws may be None for headless auto-resume after restart
        _stream_subscribers.setdefault(session.session_id, set()).add(ws)
    else:
        _stream_subscribers.setdefault(session.session_id, set())
    _activity_set("chat", active=True)
    # Persist "stream is active" immediately so a kill during the tool phase
    # (before any text tokens) is still detected at boot.
    session._stream_active = True
    session.last_heartbeat = datetime.now().isoformat()
    sessions.persist()

    async with session.lock:
        session._lock_acquired_at = _time.time()
        try:
            full_prompt = message

            # Build attachment context into the prompt
            if attachments:
                att_parts = []
                for filename in attachments:
                    filepath = UPLOADS_DIR / filename
                    if not filepath.exists():
                        continue
                    ext = filepath.suffix.lower()
                    if ext in ('.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'):
                        att_parts.append(f"\n[Bild angehängt: {filepath}  -  bitte lies und analysiere dieses Bild]")
                    elif ext == '.pdf':
                        att_parts.append(f"\n[PDF angehängt: {filepath}  -  bitte lies und analysiere dieses PDF]")
                    elif ext in ('.txt', '.md', '.csv', '.json', '.py', '.js', '.html', '.css'):
                        try:
                            content = filepath.read_text(encoding="utf-8", errors="replace")[:5000]
                            att_parts.append(f"\n[Datei: {filename}]\n```\n{content}\n```")
                        except Exception:
                            att_parts.append(f"\n[Datei angehängt: {filepath}]")
                    elif ext in ('.webm', '.ogg', '.mp3', '.wav', '.m4a'):
                        att_parts.append(f"\n[Sprachnachricht angehängt: {filepath}  -  bitte lies und transkribiere diese Audiodatei]")
                    else:
                        att_parts.append(f"\n[Datei angehängt: {filepath}]")
                if att_parts:
                    full_prompt = full_prompt + "\n" + "\n".join(att_parts)

            system = _build_system_prompt()

            # Write system prompt to temp file (avoids Windows cmd length limit)
            import tempfile
            system_file = Path(tempfile.gettempdir()) / f"jarvis-sys-{session.session_id[:8]}.txt"
            system_file.write_text(system, encoding="utf-8")

            # Build command  -  mirrors CLI experience exactly.
            # --disallowedTools blocks commands that have historically crashed MC itself
            # (self-kill of python/server, edits to MC's own source or the watchdog).
            # These still work via --dangerously-skip-permissions elsewhere, just not
            # from inside MC's own Claude loop.
            _deny_patterns = ",".join([
                "Bash(taskkill*)",
                "Bash(Stop-Process*)",
                "Bash(Stop-Service*)",
                "Bash(net stop*)",
                "Bash(sc stop*)",
                "Bash(Restart-Service*)",
                "Bash(*python*server.py*)",
                "Bash(*start.bat*)",
                "Edit(C:/Users/Robin/Jarvis/Mission-Control/server.py)",
                "Edit(C:/Users/Robin/Jarvis/Mission-Control/start.bat)",
                "Write(C:/Users/Robin/Jarvis/Mission-Control/server.py)",
                "Write(C:/Users/Robin/Jarvis/Mission-Control/start.bat)",
            ])
            cmd = [
                "claude", "-p",
                "--output-format", "stream-json", "--verbose",
                "--max-turns", "40",
                "--dangerously-skip-permissions",
                "--disallowedTools", _deny_patterns,
                "--append-system-prompt-file", str(system_file),
                "--mcp-config", str(BASE_DIR.parent / ".mcp.json"),
            ]

            if session.cli_session_id:
                # Follow-up: resume existing CLI session (full context preserved)
                cmd.extend(["--resume", session.cli_session_id])

            model = _pick_model(full_prompt, session)
            if model:
                cmd.extend(["--model", model])

            # If no CLI session to resume, inject chat history for context
            if not session.cli_session_id and session.history:
                history_lines = []
                # Include last 20 messages (10 exchanges) for context
                recent = session.history[-20:]
                for msg in recent:
                    role = "Robin" if msg.get("role") == "user" else "Jarvis"
                    content = msg.get("content", "")
                    # Truncate very long messages to save tokens
                    if len(content) > 2000:
                        content = content[:2000] + "...[gekürzt]"
                    history_lines.append(f"{role}: {content}")
                history_block = "\n\n".join(history_lines)
                full_prompt = (
                    f"<conversation_history>\n{history_block}\n</conversation_history>\n\n"
                    f"Robin: {full_prompt}"
                )

            stdin_text = full_prompt
            used_resume = session.cli_session_id

            from llm_client import _humanize_model_name
            used_model = _humanize_model_name("claude-cli", model)
            print(f"[Claude] Streaming: session={session.session_id[:8]} model={used_model} resume={used_resume or 'new'}")
            t_start = _time.time()

            # Pipe the prompt via stdin to avoid cmd length limit
            session.process = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.DEVNULL,
                stdin=asyncio.subprocess.PIPE,
                cwd="C:/Users/Robin/Jarvis",
                limit=1024 * 1024,  # 1MB readline buffer for large tool outputs
                **_SUBPROCESS_FLAGS,
            )

            session.process.stdin.write(stdin_text.encode("utf-8"))
            await session.process.stdin.drain()
            session.process.stdin.close()

            # Persist user message immediately so it survives page reload during streaming
            session.history.append({"role": "user", "content": message})
            sessions.persist()

            # Send chat.start — to session subscribers AND all clients (for sidebar streaming indicators)
            _start_event = {"type": "chat.start", "session_id": session.session_id}
            await _broadcast_stream_event(session.session_id, _start_event)
            await _broadcast_all_ws(_start_event)

            full_text = ""
            result_event = None
            last_usage = {}  # fallback: captured from assistant events
            last_tool = ""
            _last_partial_persist = _time.time()  # track periodic partial saves
            _stream_started_at = _time.time()
            _last_ping = _time.time()

            try:
                while True:
                    # Check if all WS subscribers disconnected — no point continuing
                    subs = _stream_subscribers.get(session.session_id, set())
                    alive_subs = {s for s in subs if not s.closed}
                    if not alive_subs and not _ws_clients:
                        print(f"[Claude] All subscribers gone for {session.session_id[:8]} — killing subprocess")
                        if session.process and session.process.returncode is None:
                            session.process.kill()
                            await session.process.wait()
                        break

                    try:
                        line = await asyncio.wait_for(
                            session.process.stdout.readline(),
                            timeout=30,  # short inner timeout — send ping if silent, retry up to 300s total
                        )
                    except asyncio.TimeoutError:
                        elapsed = _time.time() - _stream_started_at
                        if elapsed >= 300:
                            print("[Claude] Subprocess readline timeout (300s) — killing")
                            if session.process and session.process.returncode is None:
                                session.process.kill()
                                await session.process.wait()
                            break
                        # Send ping every 30s so frontend stuck-timer resets
                        if _time.time() - _last_ping >= 28:
                            _last_ping = _time.time()
                            session.last_heartbeat = datetime.now().isoformat()
                            await _broadcast_stream_event(session.session_id, {
                                "type": "chat.ping",
                                "session_id": session.session_id,
                                "elapsed": int(elapsed),
                            })
                        continue
                    except Exception:
                        break

                    if not line:
                        break  # EOF

                    line_str = line.decode("utf-8", errors="replace").strip()
                    if not line_str:
                        continue

                    try:
                        event = json.loads(line_str)
                    except json.JSONDecodeError:
                        continue

                    evt_type = event.get("type", "")

                    # System events: capture CLI session ID, skip the rest
                    if evt_type == "system":
                        subtype = event.get("subtype", "")
                        if subtype == "init":
                            cli_sid = event.get("session_id")
                            if cli_sid:
                                session.cli_session_id = cli_sid
                                sessions.persist()
                                print(f"[Claude] CLI session: {cli_sid[:12]}")
                        continue

                    if evt_type == "assistant":
                        _msg_usage = event.get("message", {}).get("usage", {})
                        if _msg_usage:
                            last_usage = _msg_usage
                        for block in event.get("message", {}).get("content", []):
                            block_type = block.get("type", "")
                            if block_type == "text":
                                text_chunk = block.get("text", "")
                                if text_chunk:
                                    full_text += text_chunk
                                    await _broadcast_stream_event(session.session_id, {
                                        "type": "chat.delta",
                                        "session_id": session.session_id,
                                        "text": text_chunk,
                                    })
                            elif block_type == "tool_use":
                                last_tool = block.get("name", "")
                                # Truncate large inputs (e.g. Write content) for WS
                                raw_input = block.get("input", {})
                                safe_input = {}
                                for k, v in raw_input.items():
                                    if isinstance(v, str) and len(v) > 200:
                                        safe_input[k] = v[:200] + "..."
                                    else:
                                        safe_input[k] = v
                                await _broadcast_stream_event(session.session_id, {
                                    "type": "chat.tool_use",
                                    "session_id": session.session_id,
                                    "tool": last_tool,
                                    "input": safe_input,
                                })

                    elif evt_type == "user":
                        # tool_result comes as a "user" event
                        for block in event.get("message", {}).get("content", []):
                            block_type = block.get("type", "")
                            if block_type == "tool_result":
                                await _broadcast_stream_event(session.session_id, {
                                    "type": "chat.tool_result",
                                    "session_id": session.session_id,
                                    "tool": last_tool,
                                })

                    elif evt_type in ("rate_limit_event",):
                        continue  # skip

                    # Partial persist: save immediately on first text, then every 5s
                    if full_text and (session._streaming_partial is None or (_time.time() - _last_partial_persist) > 5):
                        _last_partial_persist = _time.time()
                        session._streaming_partial = full_text
                        sessions.persist()

                    if evt_type == "result":
                        result_event = event
                        # Capture final text from result if we missed it
                        if not full_text:
                            full_text = event.get("result", "")

            except asyncio.CancelledError:
                if session.process and session.process.returncode is None:
                    session.process.kill()
                    await session.process.wait()
                # Save partial response if any text was streamed before cancel
                if full_text:
                    session.history.append({"role": "assistant", "content": full_text + "\n\n_(abgebrochen)_"})
                    sessions.persist()
                raise

            # Wait for process to finish
            if session.process and session.process.returncode is None:
                await session.process.wait()

            duration_ms = int((_time.time() - t_start) * 1000)

            # Detect stale --resume: CLI returns "Error: Session <uuid> not found" as result text
            # (can arrive in full_text or only in result_event.result). Treat as empty for retry.
            if used_resume and full_text:
                _low = full_text.lower()
                if "session" in _low and "not found" in _low and len(full_text) < 500:
                    print(f"[Claude] Stale --resume detected in output, forcing retry")
                    full_text = ""

            if not full_text:
                if used_resume:
                    # --resume failed — drop the stale CLI session and retry fresh
                    print(f"[Claude] --resume failed (empty response), retrying without resume")
                    session.cli_session_id = None
                    sessions.persist()
                    # Remove the user message we already appended (will be re-added on retry)
                    if session.history and session.history[-1].get("role") == "user":
                        session.history.pop()
                        sessions.persist()
                    _stream_end(session.session_id)
                    asyncio.get_event_loop().call_later(30, _stream_cleanup, session.session_id)
                    # Retry without --resume (recursive call, but used_resume is now None)
                    await _run_claude(ws, session, message, attachments)
                    return
                print("[Claude] Leere Antwort - sende Fehler")
                await _broadcast_stream_event(session.session_id, {
                    "type": "chat.error",
                    "session_id": session.session_id,
                    "error": "Claude hat keine Antwort geliefert.",
                })
                session._stream_active = False
                sessions.persist()
                _stream_end(session.session_id)
                asyncio.get_event_loop().call_later(30, _stream_cleanup, session.session_id)
                return

            # Update local tracking  -  user msg already appended before streaming
            session.history.append({"role": "assistant", "content": full_text})
            session.message_count += 1
            session.last_message = datetime.now().isoformat()
            session.last_heartbeat = session.last_message
            session.process = None
            session._streaming_partial = None  # clear partial AFTER history is in place
            session._stream_active = False
            sessions.persist()  # atomic: full response in history, partial gone, stream inactive

            # Build done event with metadata from result
            done_payload = {
                "type": "chat.done",
                "session_id": session.session_id,
                "full_text": full_text,
                "duration_ms": duration_ms,
                "model": used_model,
                "provider": session.provider,
            }
            if result_event:
                done_payload["cost_usd"] = result_event.get("total_cost_usd", result_event.get("cost_usd", 0))
            # Use result_event.usage, fall back to last assistant-event usage
            usage = (result_event or {}).get("usage", {}) or last_usage
            input_base = usage.get("input_tokens", 0)
            cache_read = usage.get("cache_read_input_tokens", 0)
            cache_create = usage.get("cache_creation_input_tokens", 0)
            done_payload["input_tokens"] = input_base + cache_read + cache_create
            done_payload["output_tokens"] = usage.get("output_tokens", 0)
            # Context fill — estimate window from modelUsage (rare) or model name
            model_usage = (result_event or {}).get("modelUsage", {})
            total_used = input_base + cache_create + done_payload["output_tokens"]
            ctx_window = 0
            for _, mu in model_usage.items():
                ctx_window = mu.get("contextWindow", 0)
                if ctx_window:
                    break
            if not ctx_window:
                _m = (used_model or "").lower()
                if any(x in _m for x in ("gemini-1.5", "gemini-exp")):
                    ctx_window = 1_000_000
                elif "gemini" in _m:
                    ctx_window = 1_048_576
                else:
                    ctx_window = 200_000  # Claude Sonnet/Opus/Haiku all 200K
            # Prefer usage_info from UsageEvent (accurate cache_read included) over recomputed total
            if usage_info.get("context_pct") is not None and usage_info["context_pct"] > 0:
                done_payload["context_window"] = usage_info.get("context_window", 200_000)
                done_payload["context_used"] = usage_info.get("context_used", 0)
                done_payload["context_pct"] = usage_info["context_pct"]
                print(f"[Context-done] using UsageEvent value: {done_payload['context_pct']}%")
            elif ctx_window and total_used:
                done_payload["context_window"] = ctx_window
                done_payload["context_used"] = total_used
                done_payload["context_pct"] = round(total_used / ctx_window * 100, 1)
                print(f"[Context-done] recomputed: {total_used}/{ctx_window} = {done_payload['context_pct']}%")
            else:
                print(f"[Context-debug] result_event={'None' if result_event is None else 'yes'}, usage={usage}, total_used={total_used}")

            await _broadcast_stream_event(session.session_id, done_payload)
            await _broadcast_all_ws(done_payload)  # also notify all clients for sidebar toast + indicator
            _activity_set("chat", active=False, pending_delta=1)
            _stream_end(session.session_id)
            asyncio.get_event_loop().call_later(30, _stream_cleanup, session.session_id)
            print(f"[WS] chat.done: session={session.session_id[:8]} ({len(full_text)} chars, {duration_ms}ms)")

            # Auto-classify: after 2nd message (no topic yet), or every 4th message to refresh title
            _should_classify = (
                (not session.topic and session.message_count >= 2) or
                (session.message_count == 4 and not session.custom_title) or
                (not session.topic and session.message_count % 3 == 0)  # retry if still unclassified
            )
            if _should_classify:
                asyncio.create_task(_classify_session_headless(session))

        except asyncio.CancelledError:
            print(f"[WS] Cancelled: session={session.session_id[:8]}")
            _ft = locals().get("full_text", "")
            if _ft and (not session.history or session.history[-1].get("content") != _ft + "\n\n_(abgebrochen)_"):
                session.history.append({"role": "assistant", "content": _ft + "\n\n_(abgebrochen)_"})
                session.process = None
                sessions.persist()
            session._streaming_partial = None
            session._stream_active = False
            sessions.persist()
            _stream_end(session.session_id)
            asyncio.get_event_loop().call_later(30, _stream_cleanup, session.session_id)
        except Exception as e:
            print(f"[WS] Error: {e}")
            _ft = locals().get("full_text", "")
            if _ft and (not session.history or session.history[-1].get("role") != "assistant" or session.history[-1].get("content") != _ft):
                session.history.append({"role": "assistant", "content": _ft})
                session.message_count += 1
                session.last_message = datetime.now().isoformat()
                session.last_heartbeat = session.last_message
                session.process = None
                sessions.persist()
                print(f"[WS] Saved partial response ({len(_ft)} chars) despite error")
            session._streaming_partial = None
            session._stream_active = False
            sessions.persist()
            try:
                err_ev = {
                    "type": "chat.error",
                    "session_id": session.session_id,
                    "error": str(e),
                }
                await _broadcast_stream_event(session.session_id, err_ev)
            except Exception:
                pass
            _stream_end(session.session_id)
            asyncio.get_event_loop().call_later(30, _stream_cleanup, session.session_id)


async def api_transcribe(request):
    """POST /api/transcribe  -  YouTube-URL -> Transcript -> Claude-Analyse (SSE)."""
    import re as _re
    body = await request.json()
    url = body.get("url", "").strip()
    mode = body.get("mode", "analyze")
    custom_prompt = body.get("custom_prompt", "").strip()

    # Extract YouTube video ID
    match = _re.search(r'(?:v=|youtu\.be/|embed/)([A-Za-z0-9_-]{11})', url)
    if not match:
        raise web.HTTPBadRequest(reason="Keine gültige YouTube-URL erkannt")
    video_id = match.group(1)

    resp = web.StreamResponse(headers={
        "Content-Type": "text/event-stream; charset=utf-8",
        "Cache-Control": "no-cache",
        "Connection": "keep-alive",
    })
    await resp.prepare(request)
    _activity_set("transcriber", active=True)

    async def sse(event: str, data: dict):
        line = f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"
        try:
            await resp.write(line.encode("utf-8"))
        except Exception:
            pass

    try:
        await sse("status", {"msg": "Transcript wird geladen…"})

        # Fetch transcript (blocking I/O in thread pool)
        from youtube_transcript_api import YouTubeTranscriptApi
        def _fetch():
            api_yt = YouTubeTranscriptApi()
            t = api_yt.fetch(video_id)
            return " ".join([s.text for s in t])

        try:
            full_text = await asyncio.get_event_loop().run_in_executor(None, _fetch)
        except Exception as e:
            await sse("error", {"msg": f"Transcript-Fehler: {e}"})
            try:
                await resp.write_eof()
            except Exception:
                pass
            return resp

        word_count = len(full_text.split())
        preview = full_text[:800] + ("…" if len(full_text) > 800 else "")
        await sse("transcript", {"preview": preview, "words": word_count, "chars": len(full_text)})
        await sse("status", {"msg": f"Claude analysiert ({word_count} Wörter)…"})

        # Build prompt
        text_limit = 12000
        trimmed = full_text[:text_limit]
        if mode == "summarize":
            prompt = f"Fasse dieses Video-Transcript kompakt auf Deutsch zusammen. Nutze Bullet Points und Abschnitte:\n\n{trimmed}"
        elif mode == "compare":
            jarvis_ctx = (
                "Jarvis hat aktuell:\n"
                "- Level 1 ✅ Auto-Memory (Claude Code ~/.claude/projects/.../memory/)\n"
                "- Level 2 ✅ CLAUDE.md (expliziter Kontext)\n"
                "- Level 3 ✅ State Files (Daily Notes, MEMORY.md-Index, Markdown-Netzwerk)\n"
                "- Level 4 ✅ Obsidian Vault (Second Brain, E:/OneDrive/AI/Obsidian-Vault/)\n"
                "- Level 5 ❌ Naive RAG (Embeddings / Vektordatenbank fehlt)\n"
                "- Level 6 ❌ Graph RAG (LightRAG fehlt)\n"
                "- Level 7 ❌ Agentic RAG (multimodal, autonome Pipelines fehlen)\n"
            )
            prompt = (
                f"{jarvis_ctx}\n\nDas Video-Transcript:\n{trimmed}\n\n"
                "Vergleiche den Video-Inhalt mit Jarvis' aktuellem Stand. "
                "Was haben wir bereits? Was fehlt? Was sind die nächsten konkreten Schritte?"
            )
        elif mode == "custom" and custom_prompt:
            prompt = f"{custom_prompt}\n\nVideo-Transcript:\n{trimmed}"
        else:  # analyze
            prompt = f"Analysiere dieses Video-Transcript auf Deutsch. Kernkonzepte, Empfehlungen, Tools — strukturiert:\n\n{trimmed}"

        sys_p = "Du bist Jarvis. Antworte auf Deutsch. Kein Smalltalk. Keine Füllwörter."

        proc = await asyncio.create_subprocess_exec(
            CLAUDE_EXE, "-p", prompt,
            "--output-format", "text",
            "--max-turns", "1",
            "--system-prompt", sys_p,
            "--disallowed-tools", "Bash,Read,Edit,Write,Glob,Grep,WebFetch,WebSearch,Agent",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            stdin=asyncio.subprocess.DEVNULL,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )

        buf = []
        async for line in proc.stdout:
            chunk = line.decode("utf-8", errors="replace")
            buf.append(chunk)
            await sse("chunk", {"text": chunk})
        await proc.wait()
        output_text = "".join(buf)
        await sse("done", {"chars": len(output_text)})

        # Save transcript to disk
        transcript_id = uuid.uuid4().hex[:8]
        # Try to extract video title from first line of output
        first_line = output_text.strip().split("\n")[0][:80] if output_text.strip() else ""
        ts_data = {
            "id": transcript_id,
            "url": url,
            "video_id": video_id,
            "mode": mode,
            "title": first_line or f"Video {video_id}",
            "date": datetime.now().isoformat(),
            "word_count": word_count,
            "chars": len(output_text),
            "transcript_preview": preview[:500],
            "output": output_text,
        }
        ts_file = TRANSCRIPTS_DIR / f"{datetime.now().strftime('%Y-%m-%d_%H%M')}_{transcript_id}.json"
        ts_file.write_text(json.dumps(ts_data, ensure_ascii=False, indent=2), encoding="utf-8")
        _activity_set("transcriber", active=False, pending_delta=1)

    except Exception as e:
        try:
            await sse("error", {"msg": str(e)})
        except Exception:
            pass
    finally:
        _activity_set("transcriber", active=False)

    try:
        await resp.write_eof()
    except Exception:
        pass
    return resp


async def api_transcribe_chat(request):
    """POST /api/transcribe/chat  -  Follow-up chat about an analyzed video (SSE)."""
    body = await request.json()
    message = body.get("message", "").strip()
    transcript = body.get("transcript", "").strip()
    analysis = body.get("analysis", "").strip()
    history = body.get("history", [])

    if not message:
        raise web.HTTPBadRequest(reason="Keine Nachricht")

    context_parts = []
    if transcript:
        context_parts.append(f"Video-Transcript (Auszug):\n{transcript[:3000]}")
    if analysis:
        context_parts.append(f"Meine vorherige Analyse:\n{analysis[:3000]}")
    context = "\n\n---\n\n".join(context_parts)

    history_str = ""
    for h in history[-6:]:
        role = "Robin" if h.get("role") == "user" else "Jarvis"
        history_str += f"{role}: {h.get('content','')}\n\n"

    prompt = f"{context}\n\n---\n\nGesprächsverlauf:\n{history_str}Robin: {message}\n\nJarvis:"
    sys_p = "Du bist Jarvis, ein persönlicher AI-Assistent. Du diskutierst ein analysiertes Video mit Robin. Antworte direkt und präzise auf Deutsch. Beziehe dich auf Transcript und Analyse."

    resp = web.StreamResponse(headers={
        "Content-Type": "text/event-stream; charset=utf-8",
        "Cache-Control": "no-cache",
        "Connection": "keep-alive",
    })
    await resp.prepare(request)

    async def sse(event: str, data: dict):
        line = f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"
        try:
            await resp.write(line.encode("utf-8"))
        except Exception:
            pass

    try:
        proc = await asyncio.create_subprocess_exec(
            CLAUDE_EXE, "-p", prompt,
            "--output-format", "text",
            "--max-turns", "1",
            "--system-prompt", sys_p,
            "--disallowed-tools", "Bash,Read,Edit,Write,Glob,Grep,WebFetch,WebSearch,Agent",
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            stdin=asyncio.subprocess.DEVNULL,
            cwd="C:/Users/Robin/Jarvis",
            **_SUBPROCESS_FLAGS,
        )
        buf = []
        async for line in proc.stdout:
            chunk = line.decode("utf-8", errors="replace")
            buf.append(chunk)
            await sse("chunk", {"text": chunk})
        await proc.wait()
        full = "".join(buf)
        await sse("done", {"text": full})
    except Exception as e:
        try:
            await sse("error", {"msg": str(e)})
        except Exception:
            pass

    try:
        await resp.write_eof()
    except Exception:
        pass
    return resp


async def api_transcripts(request):
    """GET /api/transcripts  -  Return list of saved transcriptions."""
    items = []
    for f in sorted(TRANSCRIPTS_DIR.glob("*.json"), key=lambda x: x.stat().st_mtime, reverse=True):
        try:
            data = json.loads(f.read_text(encoding="utf-8"))
            items.append({
                "id": data.get("id"),
                "url": data.get("url"),
                "video_id": data.get("video_id"),
                "mode": data.get("mode"),
                "title": data.get("title", ""),
                "date": data.get("date"),
                "word_count": data.get("word_count", 0),
                "chars": data.get("chars", 0),
                "transcript_preview": data.get("transcript_preview", "")[:200],
            })
        except Exception:
            pass
    return web.json_response({"items": items})


async def api_transcript_get(request):
    """GET /api/transcripts/{id}  -  Get full transcript output."""
    tid = request.match_info["id"]
    for f in TRANSCRIPTS_DIR.glob(f"*_{tid}.json"):
        try:
            return web.Response(text=f.read_text(encoding="utf-8"), content_type="application/json")
        except Exception:
            pass
    return web.json_response({"error": "Not found"}, status=404)


async def ws_chat(request):
    ws = web.WebSocketResponse()
    await ws.prepare(request)
    _ws_clients.add(ws)
    # Tell client our version so it can detect when a reload is needed
    # Include restarted=True for the first 90s after server start (so frontend can nudge user)
    _fresh_restart = (time.time() - _server_start_time) < 90
    await ws.send_json({"type": "system.version", "version": _server_version, "restarted": _fresh_restart, "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]]})
    # Send current sessions immediately so client has fresh data on connect
    await ws.send_json({
        "type": "chat.sessions",
        "sessions": sessions.list_sessions(),
        "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]],
    })

    try:
      async for msg in ws:
        if msg.type == web.WSMsgType.TEXT:
            try:
                data = json.loads(msg.data)
            except json.JSONDecodeError:
                await ws.send_json({"type": "chat.error", "error": "Invalid JSON"})
                continue

            msg_type = data.get("type")

            if msg_type == "chat.new_session":
                s = sessions.create_session()
                await ws.send_json({
                    "type": "chat.session_created",
                    "session_id": s.session_id,
                })

            elif msg_type == "chat.delete_session":
                sid = data.get("session_id")
                sessions.delete_session(sid)
                await ws.send_json({
                    "type": "chat.sessions",
                    "sessions": sessions.list_sessions(),
                    "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]],
                })

            elif msg_type == "chat.list_sessions":
                await ws.send_json({
                    "type": "chat.sessions",
                    "sessions": sessions.list_sessions(),
                    "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]],
                })

            elif msg_type == "chat.subscribe_stream":
                sid = data.get("session_id")
                if sid and sid in _active_streams and not _active_streams[sid]["done"]:
                    stream = _active_streams[sid]
                    # Send history first so client has context before stream deltas
                    hist_session = sessions.sessions.get(sid)
                    if hist_session and hist_session.history:
                        await ws.send_json({
                            "type": "chat.history",
                            "session_id": sid,
                            "history": hist_session.history,
                        })
                    _stream_subscribers.setdefault(sid, set()).add(ws)
                    # Replay all buffered events to this late-joining client
                    for event in list(stream["buffer"]):
                        await ws.send_json(event)
                else:
                    await ws.send_json({"type": "chat.stream_not_found", "session_id": sid})

            elif msg_type == "chat.rename_session":
                sid = data.get("session_id")
                title = data.get("title", "").strip()
                if sid and title:
                    sessions.rename_session(sid, title)
                await ws.send_json({
                    "type": "chat.sessions",
                    "sessions": sessions.list_sessions(),
                })

            elif msg_type == "chat.set_topic":
                sid = data.get("session_id")
                topic = data.get("topic", "").strip()
                color = data.get("color")
                if sid and topic == "__remove__":
                    # Remove topic from session
                    sessions.set_session_topic(sid, None, None)
                elif sid and topic:
                    # Save color mapping for future use
                    if color:
                        sessions._save_topic_color(topic, color)
                    else:
                        color = sessions._topic_colors().get(topic) or SessionManager._default_color(topic)
                        sessions._save_topic_color(topic, color)
                    sessions.set_session_topic(sid, topic, color)
                await ws.send_json({
                    "type": "chat.sessions",
                    "sessions": sessions.list_sessions(),
                })
                await ws.send_json({
                    "type": "chat.topics",
                    "topics": sessions.get_topics(),
                    "colors": sessions._topic_colors(),
                })

            elif msg_type == "chat.get_topics":
                await ws.send_json({
                    "type": "chat.topics",
                    "topics": sessions.get_topics(),
                    "colors": sessions._topic_colors(),
                })

            elif msg_type == "chat.classify_session":
                sid = data.get("session_id")
                session = sessions.sessions.get(sid)
                if session and session.history:
                    asyncio.create_task(_classify_session(ws, session))

            elif msg_type == "chat.classify_all":
                # Classify all unclassified sessions (or force=True for all)
                force = data.get("force", False)
                targets = [
                    s for s in sessions.sessions.values()
                    if (not s.topic or force) and s.history
                    and any(m.get("role") == "user" for m in s.history)
                ]
                await ws.send_json({"type": "chat.classify_all_started", "count": len(targets)})
                done_count = [0]
                total = len(targets)
                async def _classify_one(s):
                    try:
                        await _classify_session_headless(s)
                    except Exception as e:
                        print(f"[Classify] Error {s.session_id[:8]}: {e}")
                    done_count[0] += 1
                    await _broadcast_all_ws({"type": "chat.classify_all_progress", "done": done_count[0], "total": total})
                async def _classify_all_bg():
                    # Process in parallel batches of 4
                    batch_size = 4
                    for i in range(0, total, batch_size):
                        batch = targets[i:i+batch_size]
                        await asyncio.gather(*[_classify_one(s) for s in batch])
                    await _broadcast_all_ws({"type": "chat.sessions", "sessions": sessions.list_sessions(), "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]]})
                    await _broadcast_all_ws({"type": "chat.classify_all_done", "total": done_count[0]})
                    print(f"[Classify] All done: {done_count[0]}/{total}")
                asyncio.create_task(_classify_all_bg())

            elif msg_type == "chat.load_history":
                session_id = data.get("session_id")
                session = sessions.sessions.get(session_id)
                # Unknown session_id (e.g. frontend cached a UUID that was never
                # persisted before a restart) → return empty history instead of
                # error. Frontend will hydrate + create on next chat.send.
                await ws.send_json({
                    "type": "chat.history",
                    "session_id": session_id,
                    "history": session.history if session else [],
                })

            elif msg_type == "chat.send":
                session_id = data.get("session_id")
                message = data.get("message", "").strip()
                attachments = data.get("attachments", [])  # list of filenames in uploads/
                # Provider + Model optional pro Message — wenn gesetzt & abweichend,
                # triggert Switch-Logik (Thread-Reset bei Provider-Wechsel).
                desired_provider = (data.get("provider") or "").strip()
                desired_model = data.get("model")
                if desired_model is not None:
                    desired_model = str(desired_model).strip()
                if not message and not attachments:
                    await ws.send_json({"type": "chat.error", "error": "Empty message"})
                    continue

                session = sessions.get_or_create(session_id)

                # Validiere gewünschten Provider (nur Check, noch kein State-Update),
                # damit ungültiger Provider NICHT erst nach Lock-Check gemeldet wird.
                if desired_provider:
                    try:
                        from llm_client import PROVIDER_NAMES as _PN
                        if desired_provider not in _PN:
                            raise ValueError(
                                f"Unbekannter Provider {desired_provider!r}. Erlaubt: {', '.join(_PN)}"
                            )
                    except ValueError as ve:
                        await ws.send_json({"type": "chat.error",
                                            "session_id": session.session_id,
                                            "error": str(ve)})
                        continue

                if session.lock.locked():
                    # Auto-unlock only truly stuck sessions:
                    # process dead/gone AND lock held >30s (not just transiently locked)
                    import time as _t
                    process_dead = (
                        session.process is None
                        or session.process.returncode is not None
                    )
                    lock_age = (_t.time() - session._lock_acquired_at) if session._lock_acquired_at else 0
                    stuck = (process_dead and lock_age > 30) or (lock_age > 300)
                    if stuck:
                        print(f"[Session] Auto-unlocking stuck session {session.session_id[:8]} (process_dead={process_dead}, lock_age={lock_age:.0f}s)")
                        await session.cancel()
                        session.force_unlock()
                        session._lock_acquired_at = None
                        session._streaming_partial = None
                        await ws.send_json({
                            "type": "chat.done",
                            "session_id": session.session_id,
                            "full_text": "",
                            "error": "Session war blockiert und wurde automatisch freigegeben.",
                        })
                        continue
                    # Session is actively working — reject, don't kill
                    await ws.send_json({
                        "type": "chat.error",
                        "session_id": session.session_id,
                        "error": "Session is busy  -  wait for current response or cancel.",
                    })
                    continue

                # Provider-Switch atomar mit Send: erst nach Lock-Check, bevor
                # der Stream startet. So kann der Switch nicht durchrutschen
                # wenn der Send selbst gerade abgelehnt wurde.
                if desired_provider:
                    try:
                        switch_info = _set_session_provider(session, desired_provider, desired_model)
                        if switch_info["provider_switched"] and switch_info["message"]:
                            await _broadcast_stream_event(session.session_id, {
                                "type": "chat.system",
                                "session_id": session.session_id,
                                "message": switch_info["message"],
                            })
                    except ValueError as ve:
                        # Sollte schon oben validiert sein, Safety-Net
                        await ws.send_json({"type": "chat.error",
                                            "session_id": session.session_id, "error": str(ve)})
                        continue
                elif desired_model is not None and desired_model != session.model:
                    # Model-only-Wechsel im selben Provider: kein Thread-Reset
                    _set_session_provider(session, session.provider, desired_model)

                # Run chat in background task so WS can still receive cancel etc.
                asyncio.create_task(_run_chat(ws, session, message or "(siehe Anhang)", attachments))

            elif msg_type == "chat.cancel":
                session_id = data.get("session_id")
                session = sessions.sessions.get(session_id)
                if session:
                    await session.cancel()
                    session.force_unlock()
                    session._lock_acquired_at = None
                    session._streaming_partial = None
                    session._stream_active = False
                    session.needs_resume = False
                    sessions.persist()  # make sure flags are cleared on disk so next boot doesn't auto-resume
                    await ws.send_json({
                        "type": "chat.cancelled",
                        "session_id": session_id,
                    })


        elif msg.type == web.WSMsgType.ERROR:
            break

    finally:
        _ws_clients.discard(ws)
        _ws_unsubscribe_all(ws)

    return ws


# ---------------------------------------------------------------------------
# TTS (Edge TTS  -  free, no API key + ElevenLabs fallback)
# ---------------------------------------------------------------------------

# Available Edge TTS voices for the dropdown
EDGE_VOICES = {
    "seraphina": {"id": "de-DE-SeraphinaMultilingualNeural", "label": "Seraphina (Sweet)", "gender": "Female"},
    "amala":     {"id": "de-DE-AmalaNeural",                 "label": "Amala (Warm)",  "gender": "Female"},
    "florian":   {"id": "de-DE-FlorianMultilingualNeural",   "label": "Florian (Jarvis Classic)", "gender": "Male"},
    "conrad":    {"id": "de-DE-ConradNeural",                "label": "Conrad (Deep)", "gender": "Male"},
    "ava":       {"id": "en-US-AvaMultilingualNeural",       "label": "Ava (English Sweet)", "gender": "Female"},
    "andrew":    {"id": "en-US-AndrewMultilingualNeural",    "label": "Andrew (English Jarvis)", "gender": "Male"},
}
_active_voice = "seraphina"


async def _edge_tts(text, voice_key):
    """Generate TTS audio using Edge TTS (free, no API key)."""
    import edge_tts, tempfile
    voice_info = EDGE_VOICES.get(voice_key, EDGE_VOICES["seraphina"])
    voice_name = voice_info["id"]

    tmp = Path(tempfile.gettempdir()) / f"jarvis-tts-{uuid.uuid4().hex[:8]}.mp3"
    communicate = edge_tts.Communicate(text, voice_name, rate="+5%")
    await communicate.save(str(tmp))
    audio = tmp.read_bytes()
    tmp.unlink(missing_ok=True)
    return audio


async def api_tts(request):
    """Generate TTS audio  -  uses Edge TTS (free) by default."""
    global _active_voice
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"error": "Invalid JSON"}, status=400)

    text = body.get("text", "").strip()
    if not text:
        return web.json_response({"error": "No text provided"}, status=400)

    if len(text) > 5000:
        text = text[:5000]

    try:
        audio = await _edge_tts(text, _active_voice)
        return web.Response(body=audio, content_type="audio/mpeg")
    except Exception as e:
        print(f"[TTS] Edge TTS error: {e}")
        return web.json_response({"error": str(e)}, status=500)


async def api_tts_config(request):
    """Return TTS config."""
    return web.json_response({
        "enabled": CONFIG.get("tts_enabled", False),
        "voice": _active_voice,
        "voices": {k: {"label": v["label"], "gender": v["gender"]} for k, v in EDGE_VOICES.items()},
    })


async def api_tts_voice(request):
    """POST /api/tts/voice  -  Switch active voice."""
    global _active_voice
    try:
        body = await request.json()
    except Exception:
        return web.json_response({"error": "Invalid JSON"}, status=400)
    voice_key = body.get("voice_id", "").strip()
    if voice_key not in EDGE_VOICES:
        return web.json_response({"error": f"Unknown voice. Available: {list(EDGE_VOICES.keys())}"}, status=400)
    _active_voice = voice_key
    print(f"[TTS] Voice switched to: {voice_key} ({EDGE_VOICES[voice_key]['label']})")
    return web.json_response({"ok": True, "voice": voice_key})


# ---------------------------------------------------------------------------
# App setup
# ---------------------------------------------------------------------------

async def index_handler(request):
    return web.FileResponse(STATIC_DIR / "index.html", headers={
        'Cache-Control': 'no-cache, no-store, must-revalidate',
        'Pragma': 'no-cache',
        'Expires': '0',
    })


# ---------------------------------------------------------------------------
# Betreuer view (read-only, token-protected)
# ---------------------------------------------------------------------------
BETREUER_TOKEN = "diss-betreuer-2026"  # share this URL with your supervisor

async def betreuer_handler(request):
    token = request.query.get("token", "")
    if token != BETREUER_TOKEN:
        return web.Response(status=403, text="Zugriff verweigert")
    return web.FileResponse(STATIC_DIR / "betreuer.html", headers={
        'Cache-Control': 'no-cache, no-store, must-revalidate',
    })

async def betreuer_check_handler(request):
    token = request.query.get("token", "")
    if token != BETREUER_TOKEN:
        return web.Response(status=403, text="invalid token")
    return web.json_response({"ok": True})



# ============================================================================
# HEARTBEAT & ORCHESTRATOR LOGIC — externes Modul (heartbeat.py). Die alten
# in-line-Funktionen wurden in den HeartbeatManager migriert. Die legacy-
# Definitionen unten bleiben als dünne No-Ops stehen für Code-Pfade die
# noch darauf referenzieren, werden aber nicht mehr genutzt.
# ============================================================================

_HEARTBEAT_INTERVAL = 600  # legacy, wird ignoriert — siehe config.heartbeat.interval_minutes
_heartbeat_task = None  # legacy — HeartbeatManager verwaltet seinen eigenen Task

async def _gather_environment_context():
    """[Deprecated] legacy — siehe heartbeat.HeartbeatManager._build_status_context."""
    now = datetime.now()
    ctx = [
        f"### AUTONOMER HEARTBEAT - {now.strftime('%Y-%m-%d %H:%M:%S')} ###",
        f"Context: win32, Jarvis Mission Control V3",
        f"Status: {len(sessions.sessions)} aktive Sessions, {len(list(DATA_DIR.glob('tasks.json')))} Task-Dateien.",
    ]
    
    # Minecraft-Status-Check (Beispielhaft)
    try:
        import socket
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.settimeout(0.5)
            if s.connect_ex(('127.0.0.1', 25565)) == 0:
                ctx.append("Minecraft Server: Online (Port 25565)")
            else:
                ctx.append("Minecraft Server: Offline")
    except:
        pass
        
    return "\n".join(ctx)

async def _run_heartbeat_pulse():
    """Führt einen einzelnen autonomen Heartbeat-Turn aus."""
    print(f"[Heartbeat] Puls startet...", flush=True)
    
    # Orchestrator-Session finden oder erstellen
    oid = "orchestrator"
    session = sessions.get_or_create(oid)
    session.provider = "openclaw"
    session.model = "" # OpenClaw Default
    
    env_ctx = await _gather_environment_context()
    prompt = (
        f"{env_ctx}\n\n"
        "Du bist der autonome Orchestrator von Jarvis. Dies ist ein Heartbeat-Turn. "
        "Prüfe den Status und entscheide, ob Aktionen (z.B. Notifications via send_push, "
        "Skill-Aktivierungen via skill_activate) notwendig sind. "
        "Wenn alles okay ist und keine Aktion erforderlich ist, antworte nur mit 'OK'. "
        "Antworte kurz und präzise."
    )
    
    # Wir rufen den LLM-Client direkt auf (Headless)
    try:
        from llm_client import invoke_llm, Message, TextDelta
        full_response = ""
        async for ev in invoke_llm(
            provider=session.provider,
            model=session.model,
            messages=[Message(role="user", content=prompt)],
            system="Du bist Jarvis Orchestrator. Handle autonom basierend auf dem System-Status.",
        ):
            if isinstance(ev, TextDelta):
                full_response += ev.text
        
        # Ergebnis in der Session loggen (für UI Sichtbarkeit)
        session.history.append({"role": "user", "content": f"[System Heartbeat at {datetime.now().strftime('%H:%M')}]"})
        session.history.append({"role": "assistant", "content": full_response.strip()})
        session.message_count += 2
        sessions.persist()
        
        print(f"[Heartbeat] Antwort: {full_response[:50].strip()}...", flush=True)
        
    except Exception as e:
        print(f"[Heartbeat] Fehler: {e}", flush=True)

async def _heartbeat_loop():
    """Endlosschleife für den autonomen Puls."""
    await asyncio.sleep(120)  # 2 Minuten warten nach Boot
    while True:
        try:
            await _run_heartbeat_pulse()
            await asyncio.sleep(_HEARTBEAT_INTERVAL)
        except asyncio.CancelledError:
            break
        except Exception as e:
            print(f"[Heartbeat] Loop Error: {e}", flush=True)
            await asyncio.sleep(60)

async def _start_heartbeat(app):
    """Startet den neuen HeartbeatManager (heartbeat.py). Legacy-Loop ist tot."""
    global _heartbeat_manager
    from heartbeat import HeartbeatManager
    try:
        from notify import send_push as _send_push
    except ImportError:
        _send_push = None
    _heartbeat_manager = HeartbeatManager(
        config=CONFIG.get("heartbeat", {}),
        sessions=sessions,
        invoke_llm_oneshot=_llm_oneshot,
        broadcast_all=_broadcast_all_ws,
        build_system_prompt=_build_system_prompt,
        send_push=_send_push,
    )
    await _heartbeat_manager.start()


async def _stop_heartbeat(app):
    global _heartbeat_manager
    if _heartbeat_manager:
        await _heartbeat_manager.stop()
        _heartbeat_manager = None


# Globales Handle auf den laufenden Manager (von REST-Endpoints konsumiert)
_heartbeat_manager = None


# ---- REST-Endpoints für Heartbeat-UI (Task #41) -----------------------------

async def api_heartbeat_status(request):
    """GET /api/heartbeat/status — Live-Status + Tick-History für Dashboard-UI."""
    if not _heartbeat_manager:
        return web.json_response({
            "ok": False, "error": "Heartbeat-Manager nicht initialisiert",
            "enabled": False, "running": False,
        })
    return web.json_response({"ok": True, **_heartbeat_manager.status_dict()})


async def api_heartbeat_config(request):
    """POST /api/heartbeat/config — Config-Patch zur Laufzeit.
    Body: {enabled?, interval_minutes?, autonomy?, provider?, model?,
           push_enabled?, cron_jobs?}
    """
    if not _heartbeat_manager:
        return web.json_response({"ok": False, "error": "nicht initialisiert"}, status=500)
    try:
        patch = await request.json()
    except Exception:
        return web.json_response({"ok": False, "error": "Invalid JSON"}, status=400)
    # Nur erlaubte Keys durchlassen
    allowed = {"enabled", "interval_minutes", "autonomy", "provider", "model",
               "push_enabled", "max_tick_history", "cron_jobs", "boot_delay_seconds"}
    clean = {k: v for k, v in (patch or {}).items() if k in allowed}
    new_cfg = _heartbeat_manager.update_config(clean)
    # Auch in config.json persistieren damit es Neustart überlebt
    try:
        cfg_path = BASE_DIR / "config.json"
        full = json.loads(cfg_path.read_text(encoding="utf-8"))
        full.setdefault("heartbeat", {}).update(clean)
        cfg_path.write_text(json.dumps(full, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception as e:
        print(f"[Heartbeat] config persist failed: {e}")
    # Wenn enabled geändert, Loop entsprechend starten/stoppen
    if "enabled" in clean:
        if clean["enabled"] and not (_heartbeat_manager._task and not _heartbeat_manager._task.done()):
            await _heartbeat_manager.start()
        elif not clean["enabled"]:
            await _heartbeat_manager.stop()
    return web.json_response({"ok": True, "config": new_cfg})


async def api_heartbeat_trigger(request):
    """POST /api/heartbeat/trigger — Manueller Tick jetzt (für UI-Button)."""
    if not _heartbeat_manager:
        return web.json_response({"ok": False, "error": "nicht initialisiert"}, status=500)
    try:
        body = await request.json() if request.content_length else {}
    except Exception:
        body = {}
    reason = (body.get("reason") or "manual").strip() or "manual"
    entry = await _heartbeat_manager.trigger_now(reason=reason)
    return web.json_response({"ok": True, "tick": entry})


def create_app():
    app = web.Application(client_max_size=MAX_UPLOAD_BYTES)

    # Pages
    app.router.add_get("/", index_handler)
    app.router.add_get("/betreuer", betreuer_handler)
    app.router.add_get("/betreuer/check", betreuer_check_handler)

    # API
    app.router.add_get("/api/tasks", api_tasks)
    app.router.add_put("/api/tasks", api_tasks)
    app.router.add_post("/api/tasks", api_tasks)
    app.router.add_get("/api/timeline", api_timeline)
    app.router.add_get("/api/memory", api_memory_list)
    app.router.add_get("/api/memory/{filename}", api_memory_file)
    app.router.add_get("/api/sessions", api_sessions)
    app.router.add_get("/api/email", api_email)
    # New Email v2 (Outlook-style: accounts + folders + messages, local SQLite cache)
    try:
        import email_routes
        email_routes.register_routes(app)
    except Exception as _e:
        print(f"[Email v2] route registration failed: {_e}")
    app.router.add_get("/api/calendar", api_calendar)
    app.router.add_get("/api/system", api_system)
    app.router.add_post("/api/system/restart", api_system_restart)
    app.router.add_post("/api/rag/search", api_rag_search)
    app.router.add_get("/api/rag/status", api_rag_status)
    app.router.add_post("/api/rag/reindex", api_rag_reindex)
    app.router.add_post("/api/rag/reload", api_rag_reload)
    app.router.add_post("/api/quick-action", api_quick_action)
    app.router.add_get("/api/bots", api_bots)
    app.router.add_get("/api/bots/diss-research", api_bots_diss_research)
    app.router.add_get("/api/knowledge-graph", api_knowledge_graph)
    app.router.add_get("/api/diss/research", api_diss_research)
    app.router.add_post("/api/diss/research", api_diss_research)
    app.router.add_get("/api/diss/session-topic", api_diss_session_topic)
    app.router.add_post("/api/diss/session-topic", api_diss_session_topic)
    app.router.add_get("/api/diss/wiki-log", api_diss_wiki_log)
    app.router.add_get("/api/diss/wiki-index", api_diss_wiki_index)
    app.router.add_get("/api/diss/word", api_diss_word)
    app.router.add_post("/api/diss/word", api_diss_word)
    app.router.add_get("/api/diss/word-url", api_diss_word_url)
    app.router.add_post("/api/diss/word-url", api_diss_word_url)
    app.router.add_get("/api/diss/prism-results", api_diss_prism_results)
    app.router.add_post("/api/diss/prism-results", api_diss_prism_results)
    app.router.add_post("/api/diss/prism-analyse", api_diss_prism_analyse)
    app.router.add_get("/api/diss/delta", api_diss_delta)
    app.router.add_post("/api/diss/delta", api_diss_delta)
    app.router.add_delete("/api/diss/delta", api_diss_delta)
    app.router.add_get("/api/diss/structure", api_diss_structure)
    app.router.add_post("/api/diss/structure", api_diss_structure)
    app.router.add_get("/api/medizin/log", api_medizin_log)
    app.router.add_post("/api/medizin/log", api_medizin_log)
    app.router.add_post("/api/medizin/lernplan", api_medizin_lernplan)
    app.router.add_post("/api/anki/generate", api_anki_generate)
    app.router.add_get("/api/anki/decks", api_anki_decks)
    app.router.add_post("/api/anki/bulk-sync", api_anki_bulk_sync)
    app.router.add_get("/api/llm/providers", api_llm_providers)
    app.router.add_post("/api/chat/session/{session_id}/provider", api_chat_session_provider)
    app.router.add_post("/api/chat/session/{session_id}/compress", api_chat_session_compress)
    app.router.add_get("/api/vault/bundles", api_vault_bundles)
    app.router.add_post("/api/chat/session/{session_id}/context", api_chat_session_context)
    app.router.add_post("/api/medical/reason", api_medical_reason)
    app.router.add_get("/api/image/providers", api_image_providers)
    app.router.add_get("/api/image/usage", api_image_usage)
    app.router.add_get("/api/image/codex/latest", api_image_codex_latest)
    app.router.add_post("/api/image/codex/import", api_image_codex_import)
    app.router.add_get("/api/image/codex/jobs", api_image_codex_jobs_list)
    app.router.add_post("/api/image/codex/jobs", api_image_codex_jobs_create)
    app.router.add_get("/api/image/codex/jobs/{job_id}", api_image_codex_job_get)
    app.router.add_post("/api/image/codex/jobs/{job_id}/complete", api_image_codex_job_complete)
    app.router.add_post("/api/image/codex/jobs/{job_id}/fail", api_image_codex_job_fail)
    app.router.add_post("/api/image/generate", api_image_generate)
    # Image-Library
    app.router.add_get("/api/library", api_library_list)
    app.router.add_post("/api/library/folder", api_library_folder_create)
    app.router.add_delete("/api/library/folder", api_library_folder_delete)
    app.router.add_post("/api/library/move", api_library_move)
    app.router.add_delete("/api/library/image", api_library_image_delete)
    # Avatare
    app.router.add_get("/api/avatars", api_avatars_list)
    app.router.add_post("/api/avatars", api_avatars_create)
    app.router.add_get("/api/avatars/{name}", api_avatars_get)
    app.router.add_put("/api/avatars/{name}", api_avatars_update)
    app.router.add_delete("/api/avatars/{name}", api_avatars_delete)
    app.router.add_post("/api/avatars/{name}/reference", api_avatars_upload_reference)
    app.router.add_delete("/api/avatars/{name}/reference", api_avatars_delete_reference)
    app.router.add_get("/api/heartbeat/status", api_heartbeat_status)
    app.router.add_post("/api/heartbeat/config", api_heartbeat_config)
    app.router.add_post("/api/heartbeat/trigger", api_heartbeat_trigger)
    app.router.add_post("/api/llm/ask-peer", api_llm_ask_peer)
    app.router.add_get("/api/chat/session/{session_id}/scratchpad", api_chat_session_scratchpad)
    app.router.add_post("/api/chat/session/{session_id}/scratchpad", api_chat_session_scratchpad)
    app.router.add_get("/api/trading/account", api_trading_account)
    app.router.add_get("/api/trading/positions", api_trading_positions)
    app.router.add_get("/api/trading/summary", api_trading_summary)
    app.router.add_get("/api/trading/orders", api_trading_orders)
    app.router.add_get("/api/trading/history", api_trading_history)
    app.router.add_post("/api/trading/order", api_trading_order)
    # Trading Bots (Start/Stop/Restart)
    app.router.add_get("/api/trading/bots", api_trading_bots_list)
    app.router.add_post("/api/trading/bots", api_trading_bots_create)
    app.router.add_get("/api/trading/bots/{id}", api_trading_bots_get)
    app.router.add_put("/api/trading/bots/{id}", api_trading_bots_update)
    app.router.add_delete("/api/trading/bots/{id}", api_trading_bots_delete)
    app.router.add_post("/api/trading/bots/{id}/start", api_trading_bots_start)
    app.router.add_post("/api/trading/bots/{id}/stop", api_trading_bots_stop)
    app.router.add_post("/api/trading/bots/{id}/restart", api_trading_bots_restart)
    app.router.add_get("/api/trading/bots/rank", api_trading_bots_rank)
    app.router.add_post("/api/trading/bots/kick-mutate", api_trading_bots_kick_mutate)
    app.router.add_get("/api/world-monitor", api_world_monitor)

    # Smart Tasks
    app.router.add_post("/api/smart-tasks/size", api_smart_tasks_size)
    app.router.add_get("/api/smart-tasks", api_smart_tasks_list)
    app.router.add_post("/api/smart-tasks", api_smart_tasks_create)
    app.router.add_post("/api/smart-tasks/{id}/start", api_smart_tasks_start)
    app.router.add_get("/api/smart-tasks/{id}", api_smart_tasks_get)
    app.router.add_delete("/api/smart-tasks/{id}", api_smart_tasks_delete)

    # Roadmap (projects.md)
    app.router.add_get("/api/roadmap", api_roadmap_get)
    app.router.add_post("/api/roadmap/suggest", api_roadmap_suggest)
    app.router.add_post("/api/roadmap/{id}/toggle", api_roadmap_toggle)
    app.router.add_post("/api/roadmap/{id}/build", api_roadmap_build)

    # Skills
    app.router.add_get("/api/skills", api_skills)
    app.router.add_get("/api/skills/usage", api_skills_usage)
    app.router.add_get("/api/skills/{slug}", api_skill_detail)
    app.router.add_post("/api/skills", api_skills_install)
    app.router.add_post("/api/skills/preview", api_skills_preview)

    # Arena
    app.router.add_get("/api/arena/rooms", api_arena_rooms)
    app.router.add_post("/api/arena/rooms", api_arena_create)
    app.router.add_get("/api/arena/rooms/{id}", api_arena_room)
    app.router.add_put("/api/arena/rooms/{id}", api_arena_update)
    app.router.add_delete("/api/arena/rooms/{id}", api_arena_delete)
    app.router.add_post("/api/arena/rooms/{id}/start", api_arena_start)
    app.router.add_post("/api/arena/rooms/{id}/stop", api_arena_stop)
    app.router.add_post("/api/arena/trialog/start", api_arena_trialog_start)
    app.router.add_post("/api/arena/trialog/{id}/stop", api_arena_trialog_stop)
    app.router.add_post("/api/arena/trialog/{id}/read", api_arena_trialog_read)
    app.router.add_post("/api/arena/rooms/{id}/inject", api_arena_inject)
    app.router.add_post("/api/arena/rooms/{id}/execute", api_arena_execute)
    app.router.add_post("/api/arena/rooms/{id}/orchestrator", api_arena_orchestrator_chat)
    app.router.add_post("/api/arena/rooms/{id}/snapshot", api_arena_snapshot)
    app.router.add_get("/api/arena/rooms/{id}/snapshots", api_arena_snapshots_list)
    app.router.add_delete("/api/arena/rooms/{id}/consensus_marks/{index}", api_arena_consensus_mark_delete)
    app.router.add_post("/api/arena/rooms/{id}/unclear", api_arena_unclear)
    app.router.add_post("/api/arena/rooms/{id}/summarize", api_arena_summarize)
    app.router.add_post("/api/arena/rooms/{id}/approve", api_arena_approve)
    app.router.add_post("/api/arena/rooms/{id}/continue", api_arena_continue)
    app.router.add_get("/api/arena/file-peek", api_arena_file_peek)
    app.router.add_get("/api/arena/files-status", api_arena_files_status)
    app.router.add_post("/api/arena/deep-preview", api_arena_deep_preview)
    app.router.add_post("/api/arena/preview-misclick", api_arena_preview_misclick)
    app.router.add_get("/api/arena/preview-stats", api_arena_preview_stats)

    # Upload
    app.router.add_post("/api/upload", api_upload)
    app.router.add_get("/api/funnel/log", api_funnel_log)
    app.router.add_post("/api/funnel/ingest", api_funnel_ingest)
    app.router.add_post("/api/learn/ingest", api_learn_ingest)
    app.router.add_post("/api/learn/generate", api_learn_generate)
    app.router.add_get("/api/learn/meta", api_learn_meta)
    app.router.add_get("/api/learn/history", api_learn_history)
    app.router.add_get("/api/learn/decks/{deck_id}/apkg", api_learn_deck_apkg)
    app.router.add_get("/api/learn/decks/{deck_id}/preview", api_learn_deck_preview)
    app.router.add_get("/api/learn/feedback", api_learn_feedback)
    app.router.add_delete("/api/learn/feedback", api_learn_feedback)
    app.router.add_post("/api/learn/screenshot", api_learn_screenshot)
    app.router.add_post("/api/learn/upload", api_learn_upload)
    app.router.add_post("/api/learn/export-zip", api_learn_export_zip)

    # --- Documents-Tab (Bewerbungen, gesyncte Word/Google-Docs) ---
    _DOCS_JSON = DATA_DIR / "documents.json"
    _DOCS_SYNC_REQUESTS = DATA_DIR / "document-sync-requests.jsonl"

    def _load_documents() -> list[dict]:
        try:
            return json.loads(_DOCS_JSON.read_text(encoding="utf-8")).get("documents", [])
        except Exception:
            return []

    async def api_documents_list(request):
        docs = _load_documents()
        # Word-File-Status pro Dokument anreichern (existiert? mtime? size?)
        for d in docs:
            wp = d.get("word_path", "")
            if wp:
                p = Path(wp)
                if p.exists():
                    st = p.stat()
                    d["word_exists"] = True
                    d["word_mtime"] = datetime.fromtimestamp(st.st_mtime).isoformat(timespec="seconds")
                    d["word_size_kb"] = round(st.st_size / 1024, 1)
                else:
                    d["word_exists"] = False
        return web.json_response({"ok": True, "documents": docs})

    async def api_documents_word(request):
        """GET /api/documents/{id}/word — liefert die Word-Datei als Download."""
        doc_id = request.match_info.get("id", "")
        for d in _load_documents():
            if d.get("id") == doc_id:
                wp = Path(d.get("word_path", ""))
                if not wp.exists():
                    return web.json_response({"ok": False, "error": "Word-Datei fehlt"}, status=404)
                return web.FileResponse(
                    path=wp,
                    headers={
                        "Content-Disposition": f'attachment; filename="{d.get("word_filename", wp.name)}"',
                    },
                )
        return web.json_response({"ok": False, "error": "Document nicht gefunden"}, status=404)

    async def api_documents_sync(request):
        """POST /api/documents/{id}/sync — schreibt Sync-Request in JSONL.
        Claude liest die Requests + macht den eigentlichen Sync (Google Doc -> Word)
        bei naechster Session oder bei expliziter Aufforderung.
        """
        doc_id = request.match_info.get("id", "")
        docs = _load_documents()
        target = next((d for d in docs if d.get("id") == doc_id), None)
        if not target:
            return web.json_response({"ok": False, "error": "Document nicht gefunden"}, status=404)
        try:
            req_entry = {
                "ts": datetime.now().isoformat(timespec="seconds"),
                "doc_id": doc_id,
                "gdoc_id": target.get("gdoc_id"),
                "word_path": target.get("word_path"),
                "status": "pending",
            }
            with open(_DOCS_SYNC_REQUESTS, "a", encoding="utf-8") as f:
                f.write(json.dumps(req_entry, ensure_ascii=False) + "\n")
            return web.json_response({
                "ok": True,
                "message": "Sync-Request angelegt. Sage Jarvis: 'sync das PROMOS-Doc' (oder warte auf naechste Auto-Sync)",
                "request": req_entry,
            })
        except Exception as e:
            return web.json_response({"ok": False, "error": str(e)}, status=500)

    app.router.add_get("/api/documents/list", api_documents_list)
    app.router.add_get("/api/documents/{id}/word", api_documents_word)
    app.router.add_post("/api/documents/{id}/sync", api_documents_sync)

    async def api_learn_anki_preview(request):
        """GET /api/learn/anki-preview — Browser-Vorschau, wie die exportierten
        Anki-Karten aussehen sollten (CSS+Templates aus apkg_export.py)."""
        try:
            from apkg_export import render_preview_html
            return web.Response(text=render_preview_html(), content_type="text/html")
        except Exception as e:
            return web.Response(text=f"Preview-Fehler: {type(e).__name__}: {e}", status=500)
    app.router.add_get("/api/learn/anki-preview", api_learn_anki_preview)
    app.router.add_post("/api/learn/amboss-login", api_learn_amboss_login)
    app.router.add_get("/api/learn/amboss-test", api_learn_amboss_test)
    app.router.add_post("/api/learn/viamedici-login", api_learn_viamedici_login)
    app.router.add_get("/api/learn/viamedici-test", api_learn_viamedici_test)
    app.router.add_post("/api/learn/preview-card", api_learn_preview_card)

    # TTS
    app.router.add_post("/api/tts", api_tts)
    app.router.add_get("/api/tts/config", api_tts_config)
    app.router.add_post("/api/tts/voice", api_tts_voice)

    # Transcriber
    app.router.add_post("/api/transcribe", api_transcribe)
    app.router.add_post("/api/transcribe/chat", api_transcribe_chat)
    app.router.add_get("/api/transcripts", api_transcripts)
    app.router.add_get("/api/transcripts/{id}", api_transcript_get)

    # Activity indicators
    app.router.add_get("/api/activity", api_activity_get)
    app.router.add_post("/api/activity/clear/{view}", api_activity_clear)

    # WebSocket
    app.router.add_get("/ws/chat", ws_chat)

    # Health check for reconnect polling
    async def api_health(request):
        import time as _t
        active = [
            {
                "session_id": sid,
                "done": s["done"],
                "lock_held": sessions.sessions.get(sid) and sessions.sessions[sid].lock.locked(),
                "lock_age_s": round(_t.time() - sessions.sessions[sid]._lock_acquired_at, 1)
                    if sessions.sessions.get(sid) and sessions.sessions[sid]._lock_acquired_at else None,
                "has_process": sessions.sessions.get(sid) and sessions.sessions[sid].process is not None,
            }
            for sid, s in _active_streams.items()
        ]
        return web.json_response({
            "ok": True,
            "version": _server_version,
            "ws_clients": len(_ws_clients),
            "active_streams": active,
            "sessions_total": len(sessions.sessions),
        })
    app.router.add_get("/api/health", api_health)

    # Push notifications via ntfy.sh
    async def api_notify_test(request):
        loop = asyncio.get_event_loop()
        ok = await loop.run_in_executor(
            None,
            lambda: send_push("Jarvis Mission Control", "Push funktioniert ✅", priority=4, tags=["white_check_mark"]),
        )
        return web.json_response({"sent": ok})

    async def api_notify_send(request):
        try:
            data = await request.json()
        except Exception:
            return web.json_response({"error": "invalid JSON"}, status=400)
        title = data.get("title", "Jarvis")
        message = data.get("message", "")
        priority = int(data.get("priority", 3))
        tags = data.get("tags")
        click = data.get("click")
        loop = asyncio.get_event_loop()
        ok = await loop.run_in_executor(
            None,
            lambda: send_push(title, message, priority=priority, tags=tags, click=click),
        )
        return web.json_response({"sent": ok})

    app.router.add_post("/api/notify/test", api_notify_test)
    app.router.add_post("/api/notify/send", api_notify_send)

    # Trigger coordinated reload of all connected clients
    async def api_trigger_reload(request):
        await _broadcast_all_ws({
            "type": "system.reload_required",
            "version": _server_version,
            "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]],
        })
        return web.json_response({"ok": True, "clients": len(_ws_clients), "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]]})
    app.router.add_post("/api/trigger-reload", api_trigger_reload)

    # No-cache middleware for JS and HTML files (prevents stale code after restart)
    @web.middleware
    async def no_cache_static(request, handler):
        resp = await handler(request)
        if request.path.endswith(('.js', '.html')) or request.path in ('/', '/index.html', '/chat'):
            resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
            resp.headers['Pragma'] = 'no-cache'
        return resp
    app.middlewares.append(no_cache_static)

    # Service worker must be served from root so it can control the whole origin.
    async def serve_sw(request):
        return web.FileResponse(
            STATIC_DIR / "sw.js",
            headers={"Service-Worker-Allowed": "/", "Cache-Control": "no-cache"},
        )
    app.router.add_get("/sw.js", serve_sw)

    # Static files
    app.router.add_static("/static/", path=str(STATIC_DIR), name="static")
    app.router.add_static("/avatars/", path=str(AVATARS_DIR), name="avatars")
    app.router.add_static("/api/uploads/", path=str(UPLOADS_DIR), name="uploads")
    app.router.add_static("/library/", path=str(IMAGES_LIBRARY_DIR), name="library")

    # Startup: classify all sessions that have history but no topic yet
    async def _startup_classify(app):
        unclassified = [
            s for s in sessions.sessions.values()
            if not s.topic and s.history and any(m.get("role") == "user" for m in s.history)
        ]
        if not unclassified:
            return
        print(f"[Classify] Startup: {len(unclassified)} sessions to classify...")
        async def _run_all():
            done = 0
            for s in unclassified:
                await asyncio.sleep(1)  # stagger to avoid hammering Claude
                try:
                    await _classify_session_headless(s)
                except Exception as e:
                    print(f"[Classify] Startup error {s.session_id[:8]}: {e}")
                done += 1
            print(f"[Classify] Startup done: {done}/{len(unclassified)}")
            await _broadcast_all_ws({"type": "chat.sessions", "sessions": sessions.list_sessions(), "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]]})
        asyncio.create_task(_run_all())

    app.on_startup.append(_startup_classify)

    # Periodic classify: every 5 minutes, classify sessions that still have no topic
    _CLASSIFY_INTERVAL = 300  # seconds
    _classify_timer_task = None

    async def _periodic_classify_loop():
        """Background loop: every 5 min, classify unclassified sessions."""
        await asyncio.sleep(60)  # initial delay — let startup classify finish first
        while True:
            try:
                await asyncio.sleep(_CLASSIFY_INTERVAL)
                unclassified = [
                    s for s in sessions.sessions.values()
                    if not s.topic and s.history
                    and any(m.get("role") == "user" for m in s.history)
                ]
                if not unclassified:
                    continue
                print(f"[Classify] Periodic: {len(unclassified)} unclassified sessions found")
                done = 0
                batch_size = 3
                for i in range(0, len(unclassified), batch_size):
                    batch = unclassified[i:i + batch_size]
                    tasks = []
                    for s in batch:
                        tasks.append(_classify_session_headless(s))
                    results = await asyncio.gather(*tasks, return_exceptions=True)
                    for j, res in enumerate(results):
                        if isinstance(res, Exception):
                            print(f"[Classify] Periodic error {batch[j].session_id[:8]}: {res}")
                        done += 1
                    await asyncio.sleep(2)  # small pause between batches
                print(f"[Classify] Periodic done: {done}/{len(unclassified)}")
                await _broadcast_all_ws({
                    "type": "chat.sessions",
                    "sessions": sessions.list_sessions(),
                    "active_streams": [sid for sid, s in _active_streams.items() if not s["done"]],
                })
            except asyncio.CancelledError:
                print("[Classify] Periodic timer cancelled")
                break
            except Exception as e:
                print(f"[Classify] Periodic loop error: {e}")
                await asyncio.sleep(60)  # backoff on unexpected errors

    async def _start_periodic_classify(app):
        nonlocal _classify_timer_task
        _classify_timer_task = asyncio.create_task(_periodic_classify_loop())
        print(f"[Classify] Periodic timer started (every {_CLASSIFY_INTERVAL}s)")

    async def _stop_periodic_classify(app):
        nonlocal _classify_timer_task
        if _classify_timer_task:
            _classify_timer_task.cancel()
            try:
                await _classify_timer_task
            except asyncio.CancelledError:
                pass
            print("[Classify] Periodic timer stopped")

    app.on_startup.append(_start_periodic_classify)
    app.on_startup.append(_start_heartbeat)
    app.on_cleanup.append(_stop_periodic_classify)
    app.on_cleanup.append(_stop_heartbeat)

    # -----------------------------------------------------------------------
    # Dead-stream watchdog: detects sessions stuck in _active_streams where
    # the subprocess already died (process dead + lock stale ≥ 90s).
    # Handles the edge case where _run_claude is cancelled before lock
    # acquisition so cleanup code never runs.
    # -----------------------------------------------------------------------
    _dead_stream_watchdog_task = None

    async def _dead_stream_watchdog_loop():
        import time as _wt
        while True:
            await asyncio.sleep(30)
            try:
                for sid in list(_active_streams.keys()):
                    stream = _active_streams.get(sid)
                    if stream is None or stream.get("done"):
                        continue
                    session = sessions.sessions.get(sid)
                    if session is None:
                        # Orphaned stream with no session object — clean up immediately
                        _stream_end(sid)
                        asyncio.get_event_loop().call_later(5, _stream_cleanup, sid)
                        print(f"[Watchdog] Removed orphaned stream {sid[:8]} (no session)")
                        continue
                    process_dead = (
                        session.process is None
                        or session.process.returncode is not None
                    )
                    lock_age = (
                        _wt.time() - session._lock_acquired_at
                        if session._lock_acquired_at else 9999
                    )
                    # Heartbeat: was the stream recently active?
                    hb_age = 9999
                    if session.last_heartbeat:
                        try:
                            from datetime import datetime as _dt
                            hb_age = (_dt.now() - _dt.fromisoformat(session.last_heartbeat)).total_seconds()
                        except Exception:
                            pass
                    stuck = process_dead and (lock_age > 90 or hb_age > 120)
                    if stuck:
                        print(f"[Watchdog] Dead stream detected {sid[:8]} "
                              f"(process_dead={process_dead}, lock_age={lock_age:.0f}s, hb_age={hb_age:.0f}s) — clearing")
                        session._stream_active = False
                        session._streaming_partial = None
                        sessions.persist()
                        if session.lock.locked():
                            session.force_unlock()
                            session._lock_acquired_at = None
                        _stream_end(sid)
                        asyncio.get_event_loop().call_later(5, _stream_cleanup, sid)
                        err_ev = {
                            "type": "chat.done",
                            "session_id": sid,
                            "full_text": "",
                            "error": "Stream wurde automatisch beendet (Prozess nicht mehr aktiv).",
                        }
                        try:
                            await _broadcast_stream_event(sid, err_ev)
                            await _broadcast_all_ws({"type": "chat.sessions",
                                                     "sessions": sessions.list_sessions(),
                                                     "active_streams": [s for s, v in _active_streams.items() if not v.get("done")]})
                        except Exception:
                            pass
            except Exception as e:
                print(f"[Watchdog] Loop error: {e}")

    async def _start_dead_stream_watchdog(app):
        nonlocal _dead_stream_watchdog_task
        _dead_stream_watchdog_task = asyncio.create_task(_dead_stream_watchdog_loop())
        print("[Watchdog] Dead-stream watchdog started (interval: 30s)")

    async def _stop_dead_stream_watchdog(app):
        nonlocal _dead_stream_watchdog_task
        if _dead_stream_watchdog_task:
            _dead_stream_watchdog_task.cancel()
            try:
                await _dead_stream_watchdog_task
            except asyncio.CancelledError:
                pass

    app.on_startup.append(_start_dead_stream_watchdog)
    app.on_cleanup.append(_stop_dead_stream_watchdog)

    # -----------------------------------------------------------------------
    # Stale-Executor-Watchdog (Alpha-Action aus Arena-Consensus):
    # Wenn ein executor_checkpoint seit >10min auf "running" steht ohne
    # last_action_at-Bump, ist der Worker de facto tot. Markiere als "stale"
    # und setze hängende pending/running-Actions auf "failed" mit Grund.
    # Läuft alle 120s. Im Gegensatz zum Startup-Check greift das WÄHREND
    # MC läuft (z.B. Subprocess deadlock ohne Crash).
    # -----------------------------------------------------------------------
    _stale_exec_watchdog_task = None

    async def _stale_executor_watchdog_loop():
        STALE_MIN = 10.0
        while True:
            await asyncio.sleep(120)
            try:
                rooms = _load_arena_rooms()
                dirty = False
                now = datetime.now()
                for r in rooms:
                    cp = r.get("executor_checkpoint") or {}
                    if cp.get("status") != "running":
                        continue
                    # Stale, wenn weder started_at noch last_action_at in den letzten 10min
                    ref_ts = cp.get("last_action_at") or cp.get("started_at")
                    if not ref_ts:
                        continue
                    try:
                        age_min = (now - datetime.fromisoformat(ref_ts)).total_seconds() / 60.0
                    except Exception:
                        continue
                    if age_min <= STALE_MIN:
                        continue
                    aid = cp.get("approval_id")
                    cp["status"] = "stale"
                    cp["stale_detected_at"] = now.isoformat()
                    r["executor_checkpoint"] = cp
                    # pending/running Actions im passenden Approval-Block als failed schließen
                    for m in r.get("messages") or []:
                        if m.get("role") != "approval" or m.get("approval_id") != aid:
                            continue
                        for a in m.get("actions") or []:
                            if (a or {}).get("status") in ("pending", "running"):
                                a["status"] = "failed"
                                a["result"] = "Executor stale (>10min inaktiv) — vom Watchdog beendet."
                    dirty = True
                    print(f"[StaleExecWatchdog] room={r.get('id','?')[:8]} approval={(aid or '')[:8]} age={age_min:.1f}min → stale")
                if dirty:
                    _save_arena_rooms(rooms)
                    try:
                        await _broadcast_all_ws({"type": "arena.executor_stale"})
                    except Exception:
                        pass
            except Exception as e:
                print(f"[StaleExecWatchdog] loop error: {e}")

    async def _start_stale_executor_watchdog(app):
        nonlocal _stale_exec_watchdog_task
        _stale_exec_watchdog_task = asyncio.create_task(_stale_executor_watchdog_loop())
        print("[StaleExecWatchdog] started (interval: 120s, threshold: 10min)")

    async def _stop_stale_executor_watchdog(app):
        nonlocal _stale_exec_watchdog_task
        if _stale_exec_watchdog_task:
            _stale_exec_watchdog_task.cancel()
            try:
                await _stale_exec_watchdog_task
            except asyncio.CancelledError:
                pass

    app.on_startup.append(_start_stale_executor_watchdog)
    app.on_cleanup.append(_stop_stale_executor_watchdog)

    async def _startup_arena_recovery(app):
        """On server start: mark any rooms that were 'running' as 'paused' (process died)."""
        try:
            rooms = _load_arena_rooms()
            changed = False
            for r in rooms:
                if r.get("status") == "running":
                    r["status"] = "paused"
                    r["paused_reason"] = "server_restart"
                    changed = True
                    print(f"[Arena] Recovered room {r.get('title','?')} → paused")
            if changed:
                _save_arena_rooms(rooms)
                # Broadcast after WS clients connect (slight delay)
                async def _notify():
                    await asyncio.sleep(2)
                    await _broadcast_all_ws({"type": "arena.server_restart"})
                asyncio.create_task(_notify())
        except Exception as e:
            print(f"[Arena] Startup recovery error: {e}")
    app.on_startup.append(_startup_arena_recovery)

    async def _auto_resume_paused_arena(app):
        """Self-start any Arena rooms that were paused by the server-restart recovery hook.

        Robin wants rooms to continue on their own after a crash/restart, without
        having to click Start again. We wait a few seconds so recovery has flipped
        statuses and WS clients can reconnect, then kick off _run_arena for each.
        """
        async def _go():
            await asyncio.sleep(4)
            try:
                rooms = _load_arena_rooms()
            except Exception as e:
                print(f"[ArenaAutoResume] load error: {e}")
                return
            pending = [
                r for r in rooms
                if r.get("status") == "paused"
                and r.get("paused_reason") == "server_restart"
                and r.get("id") not in _arena_rooms
            ]
            if not pending:
                return
            print(f"[ArenaAutoResume] Resuming {len(pending)} paused room(s)")
            changed = False
            for r in pending:
                rid = r.get("id")
                if not rid:
                    continue
                try:
                    spawn_task = _spawn_for_room(rid, rooms)
                    if spawn_task is None:
                        continue
                    r.pop("paused_reason", None)
                    r["status"] = "running"
                    r["last_active"] = datetime.now().isoformat()
                    changed = True
                    _mode = (r.get("mode") or "arena").lower()
                    print(f"[ArenaAutoResume] → {r.get('title','?')} ({rid[:8]}, mode={_mode})")
                except Exception as e:
                    _arena_rooms.pop(rid, None)
                    print(f"[ArenaAutoResume] spawn failed for {rid}: {e}")
                    continue
                await asyncio.sleep(0.5)  # stagger to avoid subprocess pile-up
            if changed:
                try:
                    _save_arena_rooms(rooms)
                except Exception as e:
                    print(f"[ArenaAutoResume] save error: {e}")
                try:
                    await _broadcast_all_ws({"type": "arena.auto_resumed", "count": len(pending)})
                except Exception:
                    pass
        asyncio.create_task(_go())
    app.on_startup.append(_auto_resume_paused_arena)

    async def _auto_resume_paused_executor(app):
        """Re-spawn executors whose approval had unfinished actions when the
        server died. Finds approval messages with any action in status
        'pending' or 'running', resets running→pending (the worker is gone),
        and re-calls _run_executor(..., skip_done=True) so already-finished
        actions are not re-executed.
        """
        async def _go():
            await asyncio.sleep(5)  # after _auto_resume_paused_arena (4s)
            try:
                rooms = _load_arena_rooms()
            except Exception as e:
                print(f"[ArenaExecResume] load error: {e}")
                return
            _OPEN = {"pending", "running"}
            resumed = 0
            dirty = False
            for r in rooms:
                rid = r.get("id")
                if not rid:
                    continue
                for m in r.get("messages") or []:
                    if m.get("role") != "approval":
                        continue
                    aid = m.get("approval_id")
                    acts = m.get("actions") or []
                    if not aid or not acts:
                        continue
                    open_idx = [i for i, a in enumerate(acts) if (a or {}).get("status") in _OPEN]
                    if not open_idx:
                        continue
                    # Reset running→pending (the worker that was running died with the process)
                    for i in open_idx:
                        if acts[i].get("status") == "running":
                            acts[i]["status"] = "pending"
                            dirty = True
                    # Checkpoint reconcile: wenn executor_checkpoint fehlt oder zu einem
                    # anderen Approval gehoert → aus Message-Scan rekonstruieren (Betas
                    # Fallback). Zusaetzlich Stale-Check (Alphas Vorschlag): ein running-
                    # Checkpoint, dessen started_at > 10min alt ist, war ein toter Worker
                    # → als stale markieren, damit Seed/Bots wissen dass der Stand
                    # rekonstruiert ist.
                    _cp_cur = r.get("executor_checkpoint") or {}
                    _done_count = sum(1 for a in acts if (a or {}).get("status") == "done")
                    _needs_rebuild = (not _cp_cur) or _cp_cur.get("approval_id") != aid
                    _stale = False
                    if not _needs_rebuild and _cp_cur.get("status") == "running":
                        _started = _cp_cur.get("started_at")
                        if _started:
                            try:
                                _age_min = (datetime.now() - datetime.fromisoformat(_started)).total_seconds() / 60.0
                                if _age_min > 10.0:
                                    _stale = True
                            except Exception:
                                _stale = True
                        else:
                            _stale = True
                    if _needs_rebuild:
                        r["executor_checkpoint"] = {
                            "approval_id": aid,
                            "completed": _done_count,
                            "total": len(acts),
                            "status": "recovered",
                            "started_at": datetime.now().isoformat(),
                            "recovered_from": "message_scan",
                        }
                        dirty = True
                        print(f"[ArenaExecResume] rebuilt checkpoint room={rid[:8]} approval={aid[:8]} done={_done_count}/{len(acts)}")
                    elif _stale:
                        _cp_cur["status"] = "stale"
                        _cp_cur["stale_detected_at"] = datetime.now().isoformat()
                        _cp_cur["completed"] = _done_count  # sync mit realem Message-Stand
                        r["executor_checkpoint"] = _cp_cur
                        dirty = True
                        print(f"[ArenaExecResume] stale checkpoint room={rid[:8]} approval={aid[:8]} (>10min)")
                    action_texts = [a.get("text", "") for a in acts]
                    try:
                        asyncio.create_task(_run_executor(rid, aid, action_texts, skip_done=True))
                        resumed += 1
                        print(f"[ArenaExecResume] → room={rid[:8]} approval={aid[:8]} open={len(open_idx)}/{len(acts)}")
                    except Exception as e:
                        print(f"[ArenaExecResume] spawn failed for {rid[:8]}/{aid[:8]}: {e}")
                    await asyncio.sleep(0.5)  # stagger
            if dirty:
                try:
                    _save_arena_rooms(rooms)
                except Exception as e:
                    print(f"[ArenaExecResume] save error: {e}")
            if resumed:
                try:
                    await _broadcast_all_ws({"type": "arena.executor_resumed", "count": resumed})
                except Exception:
                    pass
        asyncio.create_task(_go())
    app.on_startup.append(_auto_resume_paused_executor)

    async def _startup_smart_tasks_recovery(app):
        """On server start: reset any 'running' smart tasks to 'interrupted' (subprocesses died with parent)."""
        try:
            tasks = _load_smart_tasks()
            now = datetime.now()
            now_iso = now.isoformat()
            changed = False
            reset_count = 0
            for t in tasks:
                if t.get("status") == "running":
                    started_str = t.get("started")
                    age_minutes = None
                    if started_str:
                        try:
                            started_dt = datetime.fromisoformat(started_str)
                            age_minutes = (now - started_dt).total_seconds() / 60.0
                        except Exception:
                            age_minutes = None
                    t["status"] = "interrupted"
                    t["interrupted_by"] = "server_restart"
                    t["interrupted_at"] = now_iso
                    t["pid"] = None
                    if age_minutes is not None:
                        t["interrupted_after_minutes"] = round(age_minutes, 1)
                        t["stale"] = age_minutes >= 15.0
                    changed = True
                    reset_count += 1
                    print(f"[SmartTasks] Recovered task {t.get('title','?')} → interrupted (age: {age_minutes})")
            if changed:
                _save_smart_tasks(tasks)
                async def _notify():
                    await asyncio.sleep(2)
                    try:
                        await _broadcast_all_ws({"type": "smart_tasks.server_restart", "reset": reset_count})
                    except Exception:
                        pass
                asyncio.create_task(_notify())
        except Exception as e:
            print(f"[SmartTasks] Startup recovery error: {e}")
    app.on_startup.append(_startup_smart_tasks_recovery)

    async def _start_window_hider(app):
        asyncio.create_task(_window_hider_loop())
    app.on_startup.append(_start_window_hider)

    async def _cleanup_interrupted(app):
        """Sessions die beim Restart mid-stream waren bekommen einen sauberen
        'abgebrochen'-Marker. Kein automatisches Weiterfuehren mehr — das hat
        zu haengenden UIs gefuehrt, weil Resume-Events ins Leere broadcastet
        wurden bevor der Browser-Client wieder verbunden war.
        """
        async def _go():
            await asyncio.sleep(2)
            pending = [s for s in list(sessions.sessions.values()) if getattr(s, "needs_resume", False)]
            if not pending:
                return
            print(f"[Cleanup] Marking {len(pending)} interrupted session(s) as aborted")
            for s in pending:
                s.needs_resume = False
                s._stream_active = False
                # Wenn nicht eh schon ein partial-Recover-Marker drinhaengt: einen Hinweis anhaengen
                if s.history and s.history[-1].get("role") == "user":
                    s.history.append({
                        "role": "assistant",
                        "content": "_(Server wurde neugestartet bevor eine Antwort kam — bitte nochmal senden falls du eine willst)_",
                    })
                    s.message_count += 1
            sessions.persist()
            # Allen Clients sagen dass diese Sessions nicht mehr streamen — sobald
            # sie sich verbinden bekommen sie eh chat.sessions, das hier ist nur
            # ein expliziter chat.done falls ein Client schon dranhing.
            for s in pending:
                try:
                    await _broadcast_all_ws({
                        "type": "chat.done",
                        "session_id": s.session_id,
                        "full_text": "",
                        "interrupted": True,
                    })
                except Exception:
                    pass
        asyncio.create_task(_go())
    app.on_startup.append(_cleanup_interrupted)

    return app


if __name__ == "__main__":
    port = CONFIG.get("port", 8080)
    host = CONFIG.get("host", "0.0.0.0")

    # HTTPS disabled  -  use Chrome flag for mic access over LAN instead
    ssl_ctx = None
    protocol = "http"

    import socket
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(("8.8.8.8", 80))
        lan_ip = s.getsockname()[0]
        s.close()
    except Exception:
        lan_ip = "unknown"

    print(f"\n  Mission Control V3")
    print(f"  Local:   {protocol}://localhost:{port}")
    print(f"  LAN:     {protocol}://{lan_ip}:{port}")
    print()
    web.run_app(create_app(), host=host, port=port, ssl_context=ssl_ctx, print=None)






