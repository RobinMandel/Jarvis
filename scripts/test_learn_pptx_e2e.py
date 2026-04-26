#!/usr/bin/env python3
"""
E2E-Test: PPTX-Upload → Image-Extraktion → Karten-Generierung
Prüft ob Bilder aus einer PPTX in den generierten Anki-Karten ankommen.
"""
import sys, io, json, struct, zlib, tempfile, requests
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")
from pathlib import Path

MC_URL = "http://localhost:8090"
SUBJECT = "Anatomie_E2ETest"

# ── 1. Minimales PNG im Speicher erzeugen (1x1 roter Pixel) ──────────────────

def make_minimal_png() -> bytes:
    def chunk(name: bytes, data: bytes) -> bytes:
        c = struct.pack(">I", len(data)) + name + data
        crc = zlib.crc32(name + data) & 0xFFFFFFFF
        return c + struct.pack(">I", crc)

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr = chunk(b"IHDR", ihdr_data)
    raw = b"\x00\xff\x00\x00"  # filter byte + R G B
    idat = chunk(b"IDAT", zlib.compress(raw))
    iend = chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


# ── 2. PPTX mit eingebettetem PNG erzeugen ───────────────────────────────────

def make_test_pptx(png_bytes: bytes) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    import io as _io

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # Titel
    title = slide.shapes.title
    if title:
        title.text = "Herz – Anatomie Grundlagen"

    # Bild einbetten
    img_stream = _io.BytesIO(png_bytes)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(2), Inches(2), Inches(2))

    # Text-Box mit Lerninhalt
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Inches(3.5), Inches(2), Inches(4), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("Das Herz liegt im Mediastinum. Es hat 4 Kammern: "
              "2 Vorhöfe (Atria) und 2 Ventrikel. "
              "Der linke Ventrikel pumpt sauerstoffreiches Blut in die Aorta.")

    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── 3. Upload ─────────────────────────────────────────────────────────────────

def upload_pptx(pptx_bytes: bytes) -> dict:
    files = {"files[]": ("test_herz.pptx", pptx_bytes,
                         "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    data  = {"subject": SUBJECT, "type": "vorlesung"}
    r = requests.post(f"{MC_URL}/api/learn/ingest", files=files, data=data, timeout=30)
    r.raise_for_status()
    return r.json()


# ── 4. Karten generieren ──────────────────────────────────────────────────────

def generate_cards() -> dict:
    payload = {"subject": SUBJECT, "topics": ["Herz", "Ventrikel", "Atrium"]}
    r = requests.post(f"{MC_URL}/api/learn/generate",
                      json=payload, timeout=180)
    r.raise_for_status()
    return r.json()


# ── 5. Uploads-Verzeichnis auf extrahierte Bilder prüfen ─────────────────────

def check_uploads_dir() -> list:
    img_root = Path("C:/Users/Robin/Jarvis/Mission-Control/data/uploads/images")
    if not img_root.exists():
        return []
    return [str(p.relative_to(img_root.parent.parent)) for p in img_root.rglob("*.png")]


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("E2E-Test: PPTX → Image-Extraktion → Karten-Generierung")
    print("=" * 60)

    print("\n[1] PNG erstellen …")
    png = make_minimal_png()
    print(f"    PNG: {len(png)} bytes ✓")

    print("\n[2] PPTX erstellen …")
    pptx = make_test_pptx(png)
    print(f"    PPTX: {len(pptx)} bytes ✓")

    print(f"\n[3] Upload → {MC_URL}/api/learn/ingest …")
    try:
        upload_result = upload_pptx(pptx)
        images_extracted = upload_result.get("images_extracted", "?")
        print(f"    Antwort: {json.dumps(upload_result, indent=2)}")
        print(f"\n    ► Extrahierte Bilder: {images_extracted}")
        if images_extracted == 0:
            print("    ⚠️  WARNUNG: Keine Bilder extrahiert — shape_type=13 nicht gefunden.")
            print("       Ursache oft: Bild als Placeholder eingebettet, nicht als Picture-Shape.")
    except Exception as e:
        print(f"    ✗ Upload-Fehler: {e}")
        sys.exit(1)

    print("\n[4] Extrahierte Bilder im Dateisystem:")
    imgs = check_uploads_dir()
    for p in imgs:
        print(f"    {p}")
    if not imgs:
        print("    (keine gefunden)")

    print(f"\n[5] Karten generieren → {MC_URL}/api/learn/generate …")
    try:
        gen_result = generate_cards()
        ok    = gen_result.get("ok")
        cards = gen_result.get("cards", [])
        print(f"    ok={ok}, Karten={len(cards)}")

        cards_with_img = [c for c in cards if c.get("image")]
        cards_no_img   = [c for c in cards if not c.get("image")]

        print(f"\n    ► Karten MIT Bild:  {len(cards_with_img)}")
        print(f"    ► Karten OHNE Bild: {len(cards_no_img)}")

        if cards_with_img:
            print("\n    Erste Karte mit Bild:")
            c = cards_with_img[0]
            print(f"      Front:  {c.get('front', '')[:80]}")
            print(f"      Image:  {c.get('image', '')}")
        else:
            print("\n    ⚠️  KEINE Karte hat ein Bild-Feld gesetzt.")
            if not gen_result.get("ok"):
                print(f"      Fehler: {gen_result.get('error')}")
            if cards:
                print(f"      Erste Karte: {json.dumps(cards[0], ensure_ascii=False, indent=2)}")

    except Exception as e:
        print(f"    ✗ Generate-Fehler: {e}")
        sys.exit(1)

    print("\n" + "=" * 60)
    print("Ergebnis-Zusammenfassung:")
    print(f"  Upload:           {'✓ ' + str(images_extracted) + ' Bild(er) extrahiert' if images_extracted else '⚠️  0 Bilder'}")
    print(f"  Karten gesamt:    {len(cards)}")
    print(f"  Karten mit Bild:  {len(cards_with_img)}")
    status = "PASS ✓" if cards_with_img else ("FAIL ✗" if cards else "PARTIAL ⚠️")
    print(f"  Test-Status:      {status}")
    print("=" * 60)


if __name__ == "__main__":
    main()
